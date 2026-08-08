"""assemble.py OOXML 不变量回归测试（纯 Python，不需要浏览器）。

每个测试对应一个已修复的 bug / lessons-learned checkpoint：
- generic-only font-family 不得把关键字写进 typeface
- 负坐标 shape 的 <a:off> 必须保留负值
- 满页半透明遮罩必须保留；与背景同色的满页 shape 才跳过
- 圆形 带背景+四边同色边框 必须有 outline
"""
import json
import re
import zipfile

import pytest
from pptx import Presentation

from assemble import assemble, first_font, DEFAULT_LATIN_FALLBACK

WHITE = "rgb(255, 255, 255)"


def _slide(records, bg=WHITE):
    return {"slide": {"background": bg, "theme": "test"}, "records": records}


def _shape(x, y, w, h, deco):
    return {"kind": "shape", "rect": {"x": x, "y": y, "w": w, "h": h}, "deco": deco}


def _render_xml(tmp_path, records, bg=WHITE):
    out = tmp_path / "t.pptx"
    assemble({"slides": [_slide(records, bg)]}, out)
    return zipfile.ZipFile(out).read("ppt/slides/slide1.xml").decode("utf-8")


def test_first_font_generic_only_falls_back():
    # 整条 stack 全是 generic 关键字 → 不能把 "system-ui" 写进 OOXML typeface
    assert first_font("system-ui, sans-serif") == DEFAULT_LATIN_FALLBACK
    assert first_font("sans-serif") == DEFAULT_LATIN_FALLBACK
    # 正常路径不受影响：第一个非 generic 项透传
    assert first_font('"My Custom Font", sans-serif') == "My Custom Font"


def test_negative_offset_preserved(tmp_path):
    # position:absolute; left:-200px 的装饰 shape：OOXML <a:off x> 支持负值，必须原样保留
    xml = _render_xml(tmp_path, [
        _shape(-200, -150, 500, 500,
               {"hasBg": True, "bg": "rgb(233, 69, 96)", "borderRadius": "50%"}),
    ])
    assert '<a:off x="-1270000" y="-952500"/>' in xml  # -200/-150 px × 6350 EMU


def test_fullpage_translucent_overlay_kept(tmp_path):
    # inset:0 + rgba 压暗层：不是"与背景等价"，必须发 shape（带 alpha）
    xml = _render_xml(tmp_path, [
        _shape(0, 0, 1920, 1080,
               {"hasBg": True, "bg": "rgba(10, 10, 40, 0.45)"}),
    ])
    assert xml.count("<p:sp>") == 1
    assert re.search(r'<a:alpha val="45000"', xml)


def test_fullpage_same_as_bg_skipped(tmp_path):
    # 与 slide 背景同色不透明、无边框的满页 shape：视觉冗余且在 WPS 里可选可拖 → 跳过
    xml = _render_xml(tmp_path, [
        _shape(0, 0, 1920, 1080, {"hasBg": True, "bg": WHITE}),
    ])
    assert xml.count("<p:sp>") == 0


def _text(runs, x=100, y=100, w=800, h=100):
    return {"kind": "text", "rect": {"x": x, "y": y, "w": w, "h": h},
            "style": {"fontSize": 32, "lineHeight": "48px", "textAlign": "left"},
            "runs": runs, "deco": {},
            "text": "".join(r.get("text", "") for r in runs)}


def test_run_color_alpha_emitted(tmp_path):
    # rgba 文字色（muted 弱化文字）的 alpha 必须进 <a:alpha>，否则变不透明主色
    xml = _render_xml(tmp_path, [
        _text([{"text": "muted label", "fontSize": 18, "fontWeight": "400",
                "fontFamily": "Arial", "color": "rgba(240, 232, 210, 0.58)"}]),
    ])
    assert '<a:alpha val="58000"/>' in xml


def test_ea_ambiguous_symbol_gets_latin_ea(tmp_path):
    # × ÷ ° 等码点会被 PowerPoint 路由到 ea 槽；非 CJK run 的 ea 应指回 latin 同字体
    xml = _render_xml(tmp_path, [
        _text([{"text": "3.2×", "fontSize": 90, "fontWeight": "800",
                "fontFamily": "Arial", "color": "rgb(240, 232, 210)"}]),
    ])
    assert '<a:ea typeface="Arial"/>' in xml
    # 纯 ASCII run 不写 ea
    xml2 = _render_xml(tmp_path, [
        _text([{"text": "plain", "fontSize": 20, "fontWeight": "400",
                "fontFamily": "Arial", "color": "rgb(0, 0, 0)"}]),
    ])
    assert "<a:ea" not in xml2


def test_circle_with_bg_and_border_has_outline(tmp_path):
    # border-radius:50% + background + border 的头像圆：描边并入 shape outline，不得丢失
    sides = {f"border{s}": True for s in ("Top", "Bottom", "Left", "Right")}
    widths = {f"border{s}Width": 6 for s in ("Top", "Bottom", "Left", "Right")}
    colors = {f"border{s}Color": "rgb(51, 51, 51)" for s in ("Top", "Bottom", "Left", "Right")}
    xml = _render_xml(tmp_path, [
        _shape(800, 400, 200, 200,
               {"hasBg": True, "bg": WHITE, "borderRadius": "50%",
                **sides, **widths, **colors}),
    ])
    assert 'prst="ellipse"' in xml
    m = re.search(r'<a:ln[^>]*w="(\d+)"[^>]*>.*?srgbClr val="333333"', xml, re.S)
    assert m, "圆形描边丢失"
    assert int(m.group(1)) == 6 * 6350  # 6px 线宽


# ---------------------------------------------------------------- 单行 CJK 宽度冗余


def _text_rec(w, h, runs, display="block", align_items=None):
    """runs: list[str] 或 list[dict]（dict 需带 text 键）。"""
    norm = [{"text": r} if isinstance(r, str) else dict(r) for r in runs]
    rec = {"kind": "text", "rect": {"x": 0, "y": 0, "w": w, "h": h}, "runs": norm}
    st = {"display": display}
    if align_items:
        st["alignItems"] = align_items
    rec["style"] = st
    return rec


def test_single_line_cjk_width_padded():
    from assemble import _text_box_size_px
    rec = _text_rec(100, 20, ["数据标签"])
    w, h = _text_box_size_px(rec, 14.0, True)
    assert w == pytest.approx(100 * 1.025)
    assert h == pytest.approx(max(20 * 1.4, 14.0 * 1.6))


def test_single_line_ascii_width_padded():
    # 单行拉丁/粗体标签同样加 2.5% 冗余（PPT DirectWrite 度量比浏览器略宽）
    from assemble import _text_box_size_px
    rec = _text_rec(100, 20, ["SELECT"])
    w, _ = _text_box_size_px(rec, 14.0, True)
    assert w == pytest.approx(100 * 1.025)


def test_multi_line_cjk_not_padded():
    from assemble import _text_box_size_px
    rec = _text_rec(100, 20, ["数据标签"])
    w, _ = _text_box_size_px(rec, 14.0, False)
    assert w == pytest.approx(100.0)


def test_explicit_break_cjk_not_padded():
    from assemble import _text_box_size_px
    rec = _text_rec(100, 20, [{"text": "数据", "linebreak": True}])
    w, _ = _text_box_size_px(rec, 14.0, True)
    assert w == pytest.approx(100.0)


def test_explicit_break_height_covers_line_count():
    # 显式 <br> 多行：BCR h 亚像素取整比「行数×行高」短 → 取 max，避免盒高系统偏短
    from assemble import _text_box_size_px
    rec = _text_rec(200, 40, [{"text": "line1"}, {"text": "line2", "linebreak": True}])
    rec["style"]["fontSize"] = 32
    rec["style"]["lineHeight"] = "48px"
    w, h = _text_box_size_px(rec, 32.0, True)
    assert h == pytest.approx(96.0)  # 2 行 × 48px
    assert w == pytest.approx(200.0)  # 显式分行不撑宽


def test_flex_anchor_cjk_not_padded():
    from assemble import _text_box_size_px
    rec = _text_rec(100, 20, ["数据"], display="flex", align_items="center")
    w, _ = _text_box_size_px(rec, 14.0, True)
    assert w == pytest.approx(100.0)


# ---------------------------------------------------------------- asset 透明占位符


def _asset(asset_id, x, y, w, h):
    return {"kind": "asset", "assetId": asset_id, "rect": {"x": x, "y": y, "w": w, "h": h}}


def _open_pptx(tmp_path, records, bg=WHITE):
    out = tmp_path / "t.pptx"
    assemble({"slides": [_slide(records, bg)]}, out)
    return Presentation(out)


def test_asset_placeholder_is_invisible_rectangle(tmp_path):
    # 占位符是普通 sp（非 picture）、无填充、无描边 → 渲染不可见；名字精确 OFFIPY_ASSET::<id>
    xml = _render_xml(tmp_path, [_asset("asset-s01-001", 20, 30, 120, 60)])
    assert "<p:pic" not in xml
    m = re.search(r"<p:spPr[^>]*>(.*?)</p:spPr>", xml, re.S)
    assert m and "<a:noFill/>" in m.group(1)  # 无填充
    assert re.search(r"<a:ln[^>]*>\s*<a:noFill/>", xml)  # 无描边
    assert "OFFIPY_ASSET::asset-s01-001" in xml


def test_asset_placeholder_exact_emu_rect(tmp_path):
    xml = _render_xml(tmp_path, [_asset("asset-s01-001", 20, 30, 120, 60)])
    assert '<a:off x="127000" y="190500"/>' in xml  # 20/30 px × 6350
    assert '<a:ext cx="762000" cy="381000"/>' in xml  # 120/60 px × 6350


def test_asset_placeholder_present_in_shape_collection(tmp_path):
    prs = _open_pptx(tmp_path, [_asset("asset-s01-001", 20, 30, 120, 60)])
    shapes = list(prs.slides[0].shapes)
    assert len(shapes) == 1
    sp = shapes[0]
    assert sp.name == "OFFIPY_ASSET::asset-s01-001"
    assert (sp.left, sp.top, sp.width, sp.height) == (127000, 190500, 762000, 381000)


def test_asset_placeholder_order_follows_record_order(tmp_path):
    records = [
        _shape(10, 10, 50, 50, {"hasBg": True, "bg": "rgb(1, 2, 3)", "borderRadius": "0px"}),
        _asset("asset-s01-001", 20, 30, 120, 60),
        _text([{"text": "hello", "fontSize": 12, "fontWeight": "400",
                "fontFamily": "Arial", "color": "rgb(0, 0, 0)"}]),
    ]
    xml = _render_xml(tmp_path, records)
    # 记录迭代顺序：shape 在 asset 前、text 在 asset 后（text 统一后置）
    assert xml.index("<a:solidFill>") < xml.index("OFFIPY_ASSET::asset-s01-001") < xml.index("hello")


def test_duplicate_asset_id_rejected(tmp_path):
    out = tmp_path / "t.pptx"
    with pytest.raises(ValueError, match="asset-s01-001"):
        assemble({"slides": [
            _slide([_asset("asset-s01-001", 0, 0, 10, 10)]),
            _slide([_asset("asset-s01-001", 0, 0, 10, 10)]),
        ]}, out)
    assert not out.exists()


def test_empty_asset_id_rejected(tmp_path):
    out = tmp_path / "t.pptx"
    with pytest.raises(ValueError, match="assetId"):
        assemble({"slides": [_slide([_asset("", 0, 0, 10, 10)])]}, out)
    assert not out.exists()


def test_no_asset_records_no_validation(tmp_path):
    # 无 asset record → 校验不介入、输出与既有行为一致
    xml = _render_xml(tmp_path, [
        _shape(10, 10, 50, 50, {"hasBg": True, "bg": "rgb(1, 2, 3)", "borderRadius": "0px"}),
    ])
    assert "OFFIPY_ASSET" not in xml
    assert xml.count("<p:sp>") == 1


def test_sld_sz_declares_screen16x9(tmp_path):
    # #61：python-pptx 只改 cx/cy 不改 sldSz@type（新建默认 screen4x3），Office 按
    # type 识别宽高比 → type 与 16:9 宽高矛盾。assemble 必须显式标 screen16x9。
    out = tmp_path / "t.pptx"
    assemble({"slides": [_slide([])]}, out)
    xml = zipfile.ZipFile(out).read("ppt/presentation.xml").decode("utf-8")
    assert 'type="screen16x9"' in xml
    assert 'cx="12192000"' in xml
    assert 'cy="6858000"' in xml


# === 不可信测量边界硬化（F1）===
# DOM 覆盖攻击（恶意 HTML 覆写 getBoundingClientRect/getComputedStyle）或手写测量 JSON
# 可注入字符串/NaN/Inf 到数值字段、XML 非法控制字符到文本字段——装配必须降级不崩溃。

def test_safe_float_guards_bad_values():
    from assemble import _safe_float
    assert _safe_float("abc", 0.0) == 0.0
    assert _safe_float(float("nan"), 0.0) == 0.0
    assert _safe_float(float("inf"), 5.0) == 5.0
    assert _safe_float(None, 0.0) == 0.0
    assert _safe_float("3.5", 0.0) == 3.5


def test_sanitize_text_strips_xml_invalid_controls():
    from assemble import _sanitize_text
    assert _sanitize_text("a\x00b\x1f\x0b\x0e") == "a�b���"
    # \t \n \r 是 XML 合法控制字符，必须保留
    assert _sanitize_text("ok\t\n\r") == "ok\t\n\r"
    assert _sanitize_text(None) == ""
    assert _sanitize_text(123) == "123"


def test_assemble_hostile_measurement_no_crash(tmp_path):
    # rect 数值全是字符串/None、fontSize NaN、inkBottom 字符串、padding 垃圾、
    # text 带 \x00 \x1f、slide 缺 background——装配降级为兜底值，不抛异常，
    # 输出 PPTX 合法且控制字符被替换成 U+FFFD。
    hostile = {
        "slides": [{
            "slide": {"theme": "t"},
            "records": [
                {"kind": "shape", "rect": {"x": "oops", "y": None, "w": "abc", "h": 100},
                 "deco": {"hasBg": True, "bg": "rgb(1, 2, 3)", "borderTopWidth": "zzz"}},
                {"kind": "text", "rect": {"x": 0, "y": 0, "w": "nan", "h": "inf"},
                 "runs": [{"text": "bad\x00char\x1f", "fontSize": "NaN", "inkBottom": "x",
                           "color": "rgb(0, 0, 0)"}],
                 "style": {"fontSize": "abc", "paddingTop": "xx"}},
                {"kind": "text", "rect": {"w": 100, "h": 30},
                 "runs": [{"text": "normal", "fontSize": 16, "color": "rgb(0, 0, 0)"}],
                 "style": {}},
            ],
        }]
    }
    out = tmp_path / "t.pptx"
    assemble(hostile, out)  # must not raise
    assert out.exists()
    xml = zipfile.ZipFile(out).read("ppt/slides/slide1.xml").decode("utf-8")
    assert "\x00" not in xml and "\x1f" not in xml
    assert "�" in xml


def test_assemble_missing_rect_and_slide_meta(tmp_path):
    # rect 缺失 / 非 dict → 全 0；slide 整体缺失 → 仍装配成功
    out = tmp_path / "t.pptx"
    assemble({"slides": [
        {"records": [{"kind": "text",
                      "runs": [{"text": "hello", "fontSize": 16, "color": "rgb(0,0,0)"}],
                      "style": {}}]},
        {"slide": None, "records": [{"kind": "text", "rect": "not-a-dict",
                                     "runs": [{"text": "x", "fontSize": 16}], "style": {}}]},
    ]}, out)
    assert out.exists()
    xml = zipfile.ZipFile(out).read("ppt/slides/slide1.xml").decode("utf-8")
    assert "hello" in xml


def test_assemble_malformed_record_structure_skips(tmp_path):
    # #68：047dddf 只规范化 record 值、不校验结构——缺 kind 的 dict record 在
    # assemble_slide 直取 rec["kind"] 崩 KeyError，非 dict record 在下游 rec.get 崩
    # AttributeError（含非 dict sdata 撞 _validate_asset_placeholders）。
    # 装配边界须过滤结构畸形 record（非 dict / 无 str kind），降级为跳过不崩。
    out = tmp_path / "t.pptx"
    data = {
        "slides": [
            {
                "slide": {"background": WHITE},
                "records": [
                    42,
                    None,
                    {"rect": {"x": 0, "y": 0, "w": 100, "h": 30}},  # 缺 kind
                    {"kind": 7, "rect": {"x": 0, "y": 0, "w": 100, "h": 30}},  # kind 非 str
                    {"kind": "text", "rect": {"x": 0, "y": 0, "w": 100, "h": 30},
                     "runs": [{"text": "ok", "fontSize": 16, "color": "rgb(0,0,0)"}],
                     "style": {}},
                ],
            },
            None,  # 非 dict sdata → _validate_asset_placeholders 也不崩
        ]
    }
    assemble(data, out)  # must not raise
    assert out.exists()
    xml = zipfile.ZipFile(out).read("ppt/slides/slide1.xml").decode("utf-8")
    assert "ok" in xml


def test_self_check_load_measurement_texts_guards_bad_numbers(tmp_path):
    # F4：self_check 读不可信 measurement JSON，slide.width/height 与 rect 坐标
    # 注入字符串/NaN 时降级为兜底，不崩（否则自检门禁被静默吞掉）
    from self_check import _load_measurement_texts
    m = tmp_path / "meas.json"
    m.write_text(
        json.dumps({"slides": [{
            "slide": {"width": "abc", "height": "nan"},
            "records": [
                {"kind": "text", "runs": [{"text": "hi"}],
                 "rect": {"x": "xx", "y": None, "w": "inf", "h": 10}},
                {"kind": "shape", "runs": [{"text": "skip-me"}], "rect": {"x": 0, "y": 0}},
            ],
        }]}),
        encoding="utf-8",
    )
    out = _load_measurement_texts(m, (1920, 1080))
    assert 1 in out
    refs = out[1]
    assert len(refs) == 1  # 只有 text kind；shape 跳过
    assert refs[0]["text"] == "hi"
    assert all(v == 0.0 for v in (refs[0]["x"], refs[0]["y"], refs[0]["w"]))
    assert refs[0]["h"] == 10.0


def test_measure_parse_single_index_graceful():
    # F3：single_index 非整数给 SystemExit 清晰信息，而非 ValueError traceback
    from measure import _parse_single_index
    assert _parse_single_index(None) is None
    assert _parse_single_index("3") == 3
    with pytest.raises(SystemExit, match="single_index"):
        _parse_single_index("abc")
