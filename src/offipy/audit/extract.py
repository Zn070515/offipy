"""PPTX Shape 提取：递归遍历 + 文本/XML 属性读取。

硬约束：只有读文件的函数内部才 `from pptx import Presentation`；
`import offipy.audit.extract` 本身不触发 python-pptx。

`_ShapeRecord` 是内部结构（不作稳定公共 API）。坐标语义：
- 顶层 shape：left/top/width/height 已是幻灯片绝对英寸；
- group 内子 shape：读取时为**局部坐标**（相对 group 坐标空间），
  绝对化由本模块 absolutize_records 完成（Affine2D 见 audit/geometry.py）。

形状类型判定（python-pptx 1.0.2 实证）：
- 隐藏：`p:cNvPr[@hidden='1'|'true']` 必须用 `.//` 后代搜索（cNvPr 嵌套在
  nvSpPr/nvPicPr/nvGrpSpPr 下，不是 shape 元素的直接子元素）；
- 连接线：`shape_type == MSO_SHAPE_TYPE.LINE` 优先，`}cxnSp` 标签兜底；
  不能靠 width/height==0（斜线/肘形连接线可能宽高非零）；
- group：`shape_type == MSO_SHAPE_TYPE.GROUP`（python-pptx 的 GroupShape
  没有 `.is_group` 属性）。
"""

from __future__ import annotations

import hashlib
import zipfile
from dataclasses import dataclass, field
from typing import TYPE_CHECKING, Any

from offipy.exceptions import ConversionError

from .geometry import Affine2D, Rect
from .models import AuditWarning

if TYPE_CHECKING:
    from pathlib import Path

_EMU_PER_INCH = 914400.0


def _to_inches(value: Any) -> float | None:
    """Length(EMU int) 或 None → 英寸；非 None 时必定返回 float。"""
    if value is None:
        return None
    return float(value) / _EMU_PER_INCH


def _safe_float(value: Any, default: float | None = None) -> float | None:
    """解析原始 XML 属性浮点值；损坏值（非数字）→ default，不抛 ValueError。"""
    if value is None:
        return default
    try:
        return float(value)
    except (ValueError, TypeError):
        return default


# ---------------------------------------------------------------- 内部记录


@dataclass
class _TextRun:
    text: str
    font_size: float | None  # pt；None = 继承默认
    bold: bool | None
    font_name: str | None
    color: tuple[int, int, int, float] | None = None  # (r,g,b,a)；a=0-1；None=无颜色证据


@dataclass
class _Paragraph:
    text: str  # 段落全部文本（runs + 软换行语义，可能含 \v）
    runs: list[_TextRun]
    # a:br 软换行拆出的视觉行分组（每行一组 run）；无 a:br 时 = [runs]
    segments: list[list[_TextRun]] = field(default_factory=list)
    # a:lnSpc 行距：spcPts 绝对点值（pt）/ spcPct 百分比（%）；None = 未显式设置
    line_spacing_pts: float | None = None
    line_spacing_pct: float | None = None


@dataclass
class _TextFrameData:
    text: str
    paragraphs: list[_Paragraph]
    word_wrap: bool | None
    autofit_mode: str  # NONE / SHAPE_TO_FIT_TEXT / TEXT_TO_FIT_SHAPE / UNKNOWN
    margin_left: float | None  # in
    margin_right: float | None
    margin_top: float | None
    margin_bottom: float | None
    autofit_font_scale: float | None  # 0-1 分数；None = 无 normAutofit 或未写 fontScale
    autofit_norm_auto_fit: bool  # a:normAutofit 存在（缩小字体适应 Shape）
    autofit_sp_auto_fit: bool  # a:spAutoFit 存在（扩大 Shape 适应文字）


@dataclass
class _GroupTransform:
    """p:grpSpPr/a:xfrm 的原始几何（单位英寸 / 度 / 布尔）。"""

    off_x: float
    off_y: float
    ext_cx: float
    ext_cy: float
    ch_off_x: float
    ch_off_y: float
    ch_ext_cx: float
    ch_ext_cy: float
    rotation_deg: float
    flip_h: bool
    flip_v: bool


@dataclass
class _ShapeRecord:
    slide_index: int  # 1-based
    shape_id: int
    name: str
    shape_type: str  # MSO_SHAPE_TYPE.name 或 "UNKNOWN"
    left: float | None  # 见模块 docstring：局部或绝对（inches）
    top: float | None
    width: float | None
    height: float | None
    rotation: float  # 度
    z_order: int  # 所在容器（slide/group）内从底到顶的序号
    text: str
    has_text_frame: bool
    word_wrap: bool | None
    autofit_mode: str
    is_group: bool
    is_connector: bool
    is_hidden: bool
    has_table: bool
    placeholder_type: str | None
    parent_shape_id: int | None  # None = 顶层
    group_path: tuple[int, ...]  # 祖先 group 的 shape_id 链
    transform: _GroupTransform | None = None  # 仅 group 有
    paragraphs: list[_Paragraph] = field(default_factory=list)
    tf_margin_left: float | None = None
    tf_margin_right: float | None = None
    tf_margin_top: float | None = None
    tf_margin_bottom: float | None = None
    role: str = "unknown"
    # 绝对化阶段填充：累计旋转非轴对齐 → AABB 近似；祖先 group 无 xfrm → 无法精确定位
    is_rotated: bool = False
    geometry_unknown: bool = False
    autofit_font_scale: float | None = None
    autofit_norm_auto_fit: bool = False
    autofit_sp_auto_fit: bool = False
    image_sha256: str | None = None  # PICTURE 图片内容 hash（compare 匹配用）
    # 显式填充类型：none(noFill 透明)/solid/gradient/blip/pattern/unknown（无显式填充，
    # 继承 style，按不透明处理）。overlap 遮挡判定用——透明上层不遮挡下方内容。
    fill_kind: str = "unknown"
    fill_color: tuple[int, int, int, float] | None = None  # spPr solidFill 的 (r,g,b,a)


@dataclass
class _SlideExtract:
    slide_index: int  # 1-based
    shapes: list[_ShapeRecord]


@dataclass
class _PresentationExtract:
    slide_size: tuple[float, float]  # (宽, 高) 英寸
    slides: list[_SlideExtract]
    warnings: list[AuditWarning] = field(default_factory=list)


# ---------------------------------------------------------------- 提取


def extract_presentation(path: str | Path) -> _PresentationExtract:
    """打开 PPTX，逐页递归提取 shape 记录并换算为幻灯片绝对坐标。

    惰性加载 python-pptx。group 子元素局部坐标在此完成绝对化
    （Affine2D），几何解析警告进 result.warnings。
    """
    from pptx import Presentation

    try:
        prs = Presentation(str(path))
    except Exception as e:
        # 只把 python-pptx/lxml/zip/ValueError 解析错误转 ConversionError；
        # 未预期异常原样重抛（不伪装成解析错误）。
        from lxml.etree import XMLSyntaxError
        from pptx.exc import PythonPptxError

        if not isinstance(e, (zipfile.BadZipFile, ValueError, XMLSyntaxError, PythonPptxError)):
            raise
        raise ConversionError(f"PPTX 文件无法解析（ZIP/XML 损坏）: {path} ({e})") from e
    warnings: list[AuditWarning] = []
    try:
        slide_w = _to_inches(prs.slide_width) or 0.0
        slide_h = _to_inches(prs.slide_height) or 0.0
    except (ValueError, TypeError) as e:
        # 演示文稿级 sldSz@cx/@cy 非数字 → prs.slide_width 抛 ValueError（per-shape
        # 级已兜底，演示文稿级漏网，#70）。顶层尺寸降级 0.0 + 告警，不连带整文件崩。
        slide_w = slide_h = 0.0
        warnings.append(
            AuditWarning(
                slide_index=None,
                shape_id=None,
                code="audit.extract.slidesize_corrupt",
                message=f"演示文稿尺寸 sldSz@cx/@cy 原始 XML 损坏，降级 0.0: {e}",
            )
        )
    slides = []
    for index, slide in enumerate(prs.slides, start=1):
        records: list[_ShapeRecord] = []
        for z_order, shape in enumerate(slide.shapes):
            records.extend(
                _flatten(shape, index, z_order, parent=None, group_path=(), warnings=warnings)
            )
        warnings.extend(absolutize_records(records))
        slides.append(_SlideExtract(slide_index=index, shapes=records))
    return _PresentationExtract(slide_size=(slide_w, slide_h), slides=slides, warnings=warnings)


def _flatten(
    shape: object,
    slide_index: int,
    z_order: int,
    parent: int | None,
    group_path: tuple[int, ...],
    warnings: list[AuditWarning],
) -> list[_ShapeRecord]:
    try:
        rec = _build_record(shape, slide_index, z_order, parent, group_path)
    except (ValueError, TypeError, AttributeError) as e:
        # python-pptx 属性解析（rotation/left/width 等）对损坏 XML 抛 ValueError；
        # 单个损坏 shape 跳过并告警，不连带整页/整文件崩。
        try:
            sid = shape.shape_id  # type: ignore[attr-defined]
        except Exception:
            sid = None
        warnings.append(
            AuditWarning(
                slide_index=slide_index,
                shape_id=sid,
                code="audit.extract.shape_skip_corrupt",
                message=f"shape {sid if sid is not None else '?'} 原始 XML 损坏，跳过: {e}",
            )
        )
        return []
    if rec.is_group:
        child_path = (*group_path, rec.shape_id)
        children: list[_ShapeRecord] = []
        for c_z, child in enumerate(shape.shapes):  # type: ignore[attr-defined]
            children.extend(
                _flatten(
                    child,
                    slide_index,
                    c_z,
                    parent=rec.shape_id,
                    group_path=child_path,
                    warnings=warnings,
                )
            )
        return [rec, *children]
    return [rec]


def _build_record(
    shape: object,
    slide_index: int,
    z_order: int,
    parent: int | None,
    group_path: tuple[int, ...],
) -> _ShapeRecord:
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    shape_type = getattr(shape.shape_type, "name", None) or "UNKNOWN"  # type: ignore[attr-defined]
    is_group = shape.shape_type == MSO_SHAPE_TYPE.GROUP or shape._element.tag.endswith("}grpSp")  # type: ignore[attr-defined]
    is_connector = shape.shape_type == MSO_SHAPE_TYPE.LINE or shape._element.tag.endswith("}cxnSp")  # type: ignore[attr-defined]
    is_hidden = bool(shape._element.xpath('.//p:cNvPr[@hidden="1" or @hidden="true"]'))  # type: ignore[attr-defined]

    image_sha256 = None
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:  # type: ignore[attr-defined]
        try:
            image_sha256 = hashlib.sha256(shape.image.blob).hexdigest()  # type: ignore[attr-defined]
        except Exception:
            image_sha256 = None

    has_tf = bool(getattr(shape, "has_text_frame", False))
    tf = _read_text_frame(shape) if has_tf else None

    placeholder_type = None
    if getattr(shape, "is_placeholder", False):
        ph_type = getattr(shape.placeholder_format.type, "name", None)  # type: ignore[attr-defined]
        placeholder_type = ph_type or "PLACEHOLDER"

    transform = _read_group_transform(shape) if is_group else None

    return _ShapeRecord(
        slide_index=slide_index,
        shape_id=int(shape.shape_id),  # type: ignore[attr-defined]
        name=str(shape.name),  # type: ignore[attr-defined]
        shape_type=shape_type,
        fill_kind=_read_fill(shape),
        fill_color=_read_fill_color(shape),
        left=_to_inches(shape.left),  # type: ignore[attr-defined]
        top=_to_inches(shape.top),  # type: ignore[attr-defined]
        width=_to_inches(shape.width),  # type: ignore[attr-defined]
        height=_to_inches(shape.height),  # type: ignore[attr-defined]
        rotation=float(shape.rotation or 0.0),  # type: ignore[attr-defined]
        z_order=z_order,
        text=tf.text if tf else "",
        has_text_frame=has_tf,
        word_wrap=tf.word_wrap if tf else None,
        autofit_mode=tf.autofit_mode if tf else "NONE",
        is_group=is_group,
        is_connector=is_connector,
        is_hidden=is_hidden,
        has_table=bool(getattr(shape, "has_table", False)),
        placeholder_type=placeholder_type,
        parent_shape_id=parent,
        group_path=group_path,
        transform=transform,
        paragraphs=tf.paragraphs if tf else [],
        tf_margin_left=tf.margin_left if tf else None,
        tf_margin_right=tf.margin_right if tf else None,
        tf_margin_top=tf.margin_top if tf else None,
        tf_margin_bottom=tf.margin_bottom if tf else None,
        autofit_font_scale=tf.autofit_font_scale if tf else None,
        autofit_norm_auto_fit=tf.autofit_norm_auto_fit if tf else False,
        autofit_sp_auto_fit=tf.autofit_sp_auto_fit if tf else False,
        image_sha256=image_sha256,
    )


def _read_text_frame(shape: object) -> _TextFrameData:
    from pptx.oxml.ns import qn

    tf = shape.text_frame  # type: ignore[attr-defined]
    paragraphs: list[_Paragraph] = []
    text_parts: list[str] = []
    for para in tf.paragraphs:
        runs = []
        segments: list[list[_TextRun]] = [[]]
        pp_runs = para.runs
        ri = 0
        # 按文档序遍历 a:p 子元素，a:br 处切开视觉行分组（runs 不含 a:br，无法仅靠 runs 还原顺序）
        for child in para._p:
            if child.tag == qn("a:r"):
                run = pp_runs[ri]
                ri += 1
                font = run.font
                size = font.size.pt if font.size is not None else None
                tr = _TextRun(run.text, size, font.bold, font.name, _run_color(run))
                runs.append(tr)
                segments[-1].append(tr)
            elif child.tag == qn("a:br"):
                segments.append([])
        p_text = para.text
        spc_pts, spc_pct = _read_line_spacing(para)
        paragraphs.append(
            _Paragraph(
                text=p_text,
                runs=runs,
                segments=segments,
                line_spacing_pts=spc_pts,
                line_spacing_pct=spc_pct,
            )
        )
        text_parts.append(p_text)
    auto = tf.auto_size
    autofit = getattr(auto, "name", None) or "UNKNOWN"
    norm = tf._txBody.xpath("./a:bodyPr/a:normAutofit")
    sp = tf._txBody.xpath("./a:bodyPr/a:spAutoFit")
    font_scale = None
    if norm:
        raw = norm[0].get("fontScale")
        if raw:
            fs = _safe_float(raw)
            font_scale = fs / 100000.0 if fs is not None else None
    return _TextFrameData(
        text="\n".join(text_parts),
        paragraphs=paragraphs,
        word_wrap=tf.word_wrap,
        autofit_mode=autofit,
        margin_left=_to_inches(tf.margin_left),
        margin_right=_to_inches(tf.margin_right),
        margin_top=_to_inches(tf.margin_top),
        margin_bottom=_to_inches(tf.margin_bottom),
        autofit_font_scale=font_scale,
        autofit_norm_auto_fit=bool(norm),
        autofit_sp_auto_fit=bool(sp),
    )


def _read_line_spacing(para: Any) -> tuple[float | None, float | None]:
    """段落行距 a:lnSpc → (spcPts 点值, spcPct 百分比)。

    spcPts val 单位 1/100pt（绝对值）；spcPct val 单位 1/1000%（相对单行高）。
    两者互斥（OOXML 同一 lnSpc 内只有一个子元素）；None = 未显式设置。
    """
    from pptx.oxml.ns import qn

    pPr = para._p.find(qn("a:pPr"))
    if pPr is None:
        return None, None
    lnSpc = pPr.find(qn("a:lnSpc"))
    if lnSpc is None:
        return None, None
    spcPts = lnSpc.find(qn("a:spcPts"))
    if spcPts is not None and spcPts.get("val") is not None:
        val = _safe_float(spcPts.get("val"))
        if val is not None:
            return val / 100.0, None
    spcPct = lnSpc.find(qn("a:spcPct"))
    if spcPct is not None and spcPct.get("val") is not None:
        val = _safe_float(spcPct.get("val"))
        if val is not None:
            return None, val / 1000.0
    return None, None


def _read_fill(shape: object) -> str:
    """读取 p:spPr 显式填充类型，区分是否 noFill 透明。

    a:noFill → "none"（透明，不遮挡下方内容）；solid/grad/blip/patt → 不透明；
    无显式填充元素 → "unknown"（继承 style，按不透明处理——避免把继承主题填充的
    AutoShape 误判成透明而漏报真实遮挡）。
    """
    from pptx.oxml.ns import qn

    spPr = shape._element.xpath("./p:spPr")  # type: ignore[attr-defined]
    if not spPr:
        return "unknown"
    el = spPr[0]
    for tag, kind in (
        ("a:noFill", "none"),
        ("a:solidFill", "solid"),
        ("a:gradFill", "gradient"),
        ("a:blipFill", "blip"),
        ("a:pattFill", "pattern"),
    ):
        if el.find(qn(tag)) is not None:
            return kind
    return "unknown"


def _srgb_color_alpha(el: object, default_alpha: float = 1.0) -> tuple[int, int, int, float] | None:
    """解析 a:srgbClr + 嵌套 a:alpha → (r,g,b,a)；无 val / 非 6 位 hex → None。"""
    from pptx.oxml.ns import qn

    if el is None or el.get("val") is None or len(el.get("val")) != 6:  # type: ignore[attr-defined]
        return None
    v = el.get("val")  # type: ignore[attr-defined]
    try:
        r, g, b = (int(v[i : i + 2], 16) for i in (0, 2, 4))
    except ValueError:
        return None
    a = default_alpha
    alpha = el.find(qn("a:alpha"))  # type: ignore[attr-defined]
    if alpha is not None and alpha.get("val") is not None:
        try:
            a = max(0.0, min(1.0, int(alpha.get("val")) / 100000.0))
        except ValueError:
            a = default_alpha
    return (r, g, b, a)


def _run_color(run: Any) -> tuple[int, int, int, float] | None:
    """run 前景色：a:rPr/a:solidFill/a:srgbClr（含 alpha）。schemeClr/sysClr 不解析。"""
    from pptx.oxml.ns import qn

    rPr = run._r.find(qn("a:rPr"))
    if rPr is None:
        return None
    solid = rPr.find(qn("a:solidFill"))
    if solid is None:
        return None
    return _srgb_color_alpha(solid.find(qn("a:srgbClr")))


def _read_fill_color(shape: object) -> tuple[int, int, int, float] | None:
    """shape 填充色：p:spPr/a:solidFill/a:srgbClr（含 alpha）。"""
    from pptx.oxml.ns import qn

    spPr = shape._element.xpath("./p:spPr")  # type: ignore[attr-defined]
    if not spPr:
        return None
    solid = spPr[0].find(qn("a:solidFill"))
    if solid is None:
        return None
    return _srgb_color_alpha(solid.find(qn("a:srgbClr")))


def _read_group_transform(shape: object) -> _GroupTransform | None:
    from pptx.oxml.ns import qn

    xfrm = shape._element.xpath("./p:grpSpPr/a:xfrm")  # type: ignore[attr-defined]
    if not xfrm:
        return None
    el = xfrm[0]

    def _attr(tag: str, attr: str, default: float = 0.0) -> float:
        node = el.find(qn(tag))
        if node is None or node.get(attr) is None:
            return default
        v = _safe_float(node.get(attr))
        return v / _EMU_PER_INCH if v is not None else default

    rot_raw = el.get("rot")
    rotation_deg = (_safe_float(rot_raw) or 0.0) / 60000.0 if rot_raw else 0.0
    return _GroupTransform(
        off_x=_attr("a:off", "x"),
        off_y=_attr("a:off", "y"),
        ext_cx=_attr("a:ext", "cx"),
        ext_cy=_attr("a:ext", "cy"),
        ch_off_x=_attr("a:chOff", "x"),
        ch_off_y=_attr("a:chOff", "y"),
        ch_ext_cx=_attr("a:chExt", "cx"),
        ch_ext_cy=_attr("a:chExt", "cy"),
        rotation_deg=rotation_deg,
        flip_h=el.get("flipH") in ("1", "true"),
        flip_v=el.get("flipV") in ("1", "true"),
    )


# ---------------------------------------------------------------- 绝对化


def absolutize_records(records: list[_ShapeRecord]) -> list[AuditWarning]:
    """把 group 子元素局部坐标换算为幻灯片绝对坐标（就地修改）。

    每层 Group 按序变换：子坐标原点归一化（-chOff）→ ext/chExt 缩放 →
    水平/垂直翻转 → 绕 Group 中心旋转 → Group off 平移 → 乘祖先矩阵。

    无法精确解析的变换：记 AuditWarning（group.no_transform）并把受影响
    子记录标记 geometry_unknown=True，rules 会跳过需精确位置的规则
    （不允许悄悄按零旋转处理）。
    """
    warnings: list[AuditWarning] = []
    group_mats: dict[int, Affine2D] = {}
    unknown_groups: set[int] = set()
    group_axis_aligned: dict[int, bool] = {}

    for rec in records:
        if not rec.is_group:
            continue
        if rec.transform is None:
            unknown_groups.add(rec.shape_id)
            warnings.append(
                AuditWarning(
                    slide_index=rec.slide_index,
                    shape_id=rec.shape_id,
                    code="group.no_transform",
                    message=f"group #{rec.shape_id} 缺少 a:xfrm，子元素无法精确定位",
                )
            )
            continue
        gm = _group_matrix(rec.transform)
        if rec.parent_shape_id is not None and rec.parent_shape_id in group_mats:
            gm = group_mats[rec.parent_shape_id].compose(gm)
        group_mats[rec.shape_id] = gm
        parent_aligned = (
            group_axis_aligned.get(rec.parent_shape_id, True)
            if rec.parent_shape_id is not None
            else True
        )
        group_axis_aligned[rec.shape_id] = parent_aligned and _is_axis_aligned(
            rec.transform.rotation_deg
        )

    for rec in records:
        if rec.left is None or rec.top is None or rec.width is None or rec.height is None:
            continue
        if any(aid in unknown_groups for aid in rec.group_path):
            rec.geometry_unknown = True
            continue
        parent_mat = (
            group_mats.get(rec.parent_shape_id) if rec.parent_shape_id is not None else None
        )
        own = _own_rotate_about_center(rec)
        mat = parent_mat.compose(own) if parent_mat is not None else own
        abs_rect = mat.transform_rect(Rect(rec.left, rec.top, rec.width, rec.height))
        rec.left, rec.top, rec.width, rec.height = (
            abs_rect.x,
            abs_rect.y,
            abs_rect.width,
            abs_rect.height,
        )
        ancestor_aligned = (
            group_axis_aligned.get(rec.parent_shape_id, True)
            if rec.parent_shape_id is not None
            else True
        )
        rec.is_rotated = not (_is_axis_aligned(rec.rotation) and ancestor_aligned)
    return warnings


def _group_matrix(t: _GroupTransform) -> Affine2D:
    """Group 变换矩阵：局部坐标 → 幻灯片坐标（含缩放/翻转/绕中心旋转/平移）。"""
    sx = t.ext_cx / t.ch_ext_cx if t.ch_ext_cx else 1.0
    sy = t.ext_cy / t.ch_ext_cy if t.ch_ext_cy else 1.0
    # compose(other) = self∘other，后调用的先应用；从「最后应用的」最外层变换
    # 按逆序 compose，使结果 = Translate(C)·Rotate·Flip·T(-ext/2)·Scale·T(-chOff)。
    # Flip 须在 T(-ext/2) 之后应用（flip 作用于中心相对坐标），且旋转/翻转都绕
    # 组中心 C=off+ext/2（ECMA-376 对 group xfrm 的语义）。
    m = Affine2D.identity()
    m = m.compose(Affine2D.translate(t.off_x + t.ext_cx / 2.0, t.off_y + t.ext_cy / 2.0))
    m = m.compose(Affine2D.rotate(t.rotation_deg))
    if t.flip_v:
        m = m.compose(Affine2D.flip_v())
    if t.flip_h:
        m = m.compose(Affine2D.flip_h())
    m = m.compose(Affine2D.translate(-t.ext_cx / 2.0, -t.ext_cy / 2.0))
    m = m.compose(Affine2D.scale(sx, sy))
    return m.compose(Affine2D.translate(-t.ch_off_x, -t.ch_off_y))


def _own_rotate_about_center(rec: _ShapeRecord) -> Affine2D:
    """自身 rotation 绕局部 bbox 中心旋转（无旋转返回恒等）。"""
    if (
        not rec.rotation
        or rec.left is None
        or rec.top is None
        or rec.width is None
        or rec.height is None
    ):
        return Affine2D.identity()
    cx = rec.left + rec.width / 2.0
    cy = rec.top + rec.height / 2.0
    m = Affine2D.translate(cx, cy)
    m = m.compose(Affine2D.rotate(rec.rotation))
    return m.compose(Affine2D.translate(-cx, -cy))


def _is_axis_aligned(deg: float) -> bool:
    """旋转是否为 90° 整数倍（轴对齐，AABB 即真实包围盒）。"""
    rem = deg % 90.0
    return rem < 1e-6 or abs(rem - 90.0) < 1e-6
