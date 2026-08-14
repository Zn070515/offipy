"""draw.io 图 → PPTX 原生可编辑形状（保留作者版式与配色）。

输入 .drawio 文件 → vendored drawio_extract.py 提取 IR（节点自带绝对坐标与样式）
→ offipy 薄层：选页 + 坐标归一化 fit（layout_drawio）→ 复用 diagrams.render_to_slide
渲染成可编辑形状（autoshape + connector）。可独立生成整页 PPTX（drawio_to_pptx），
也可作为 deck 后处理注入（postprocess_drawio：HTML <div class="drawio"
data-drawio="..."> → 替换为可编辑形状）。

vendored 提取器用 importlib 加载，保持上游文件原样；安全边界（DTD/ENTITY 拒绝、
压缩上限）继承自 vendor 的 parse_file，offipy 不自行解析 draw.io XML。
"""

from __future__ import annotations

import importlib.util
import json
import os
import sys
from dataclasses import dataclass
from html.parser import HTMLParser
from pathlib import Path
from urllib.parse import urlparse
from urllib.request import url2pathname

from offipy.diagrams import (
    PX_TO_EMU,
    DiagramLayout,
    PlacedEdge,
    PlacedNode,
    render_to_slide,
)

_EXTRACT_REL = Path("_vendor/diagram-design/skills/diagram-design/scripts/drawio_extract.py")

_extractor = None


def _load_extractor():
    """惰性加载 vendored drawio_extract（importlib，保持上游原样）。"""
    global _extractor
    if _extractor is None:
        script = Path(__file__).resolve().parent / _EXTRACT_REL
        if not script.exists():
            raise RuntimeError(f"vendored drawio_extract 缺失: {script}")
        name = "offipy_vendored_drawio_extract"
        spec = importlib.util.spec_from_file_location(name, script)
        mod = importlib.util.module_from_spec(spec)
        sys.modules[name] = mod
        assert spec.loader is not None
        spec.loader.exec_module(mod)
        _extractor = mod
    return _extractor


@dataclass
class DrawioNode:
    id: str
    label: str = ""
    shape: str = "rect"
    parent: str | None = None
    x: float = 0.0
    y: float = 0.0
    w: float = 0.0
    h: float = 0.0
    fill: str = ""
    stroke: str = ""
    font_color: str = ""
    font_size: float | None = None
    dashed: bool = False
    rounded: bool = False
    container: bool = False


@dataclass
class DrawioEdge:
    source: str
    target: str
    label: str = ""
    dashed: bool = False
    bidirectional: bool = False
    undirected: bool = False
    stroke: str = ""


@dataclass
class DrawioDiagram:
    nodes: list[DrawioNode]
    edges: list[DrawioEdge]


def _parse_font_size(raw: str) -> float | None:
    """style 里的 fontSize（字符串）→ float；空/非数值 → None（走 12pt 默认）。"""
    if not raw or not raw.strip():
        return None
    try:
        return float(raw)
    except ValueError:
        return None


def parse_drawio(source, *, page=None) -> DrawioDiagram:
    """.drawio 文件 → DrawioDiagram（vendored parse_file + select_pages 选页）。

    page: int → 页码（0 起，对齐 draw.io 索引）；str → 数字串按 index、否则按页名
    原样传给 vendored select_pages；None → 第一页。单页输出，不暴露 "all"。

    安全边界：只调 vendored parse_file（继承 DTD/ENTITY 拒绝与压缩上限），
    绝不自行 ET.fromstring 未经验证的原始输入。
    """
    path = Path(source)
    if not path.is_file():
        raise FileNotFoundError(f"源文件不存在: {path}")
    ex = _load_extractor()
    selector = None if page is None else (str(page) if isinstance(page, int) else page)
    try:
        pages = ex.parse_file(path)
        selected = ex.select_pages(pages, selector)
    except SystemExit as e:
        code = e.code if isinstance(e.code, str) else str(e)
        raise ValueError(f"无法解析 draw.io 源码: {code}") from None
    if not selected:
        raise ValueError("无法解析 draw.io 源码: 未找到所选页")
    p = selected[0]
    nodes = [
        DrawioNode(
            id=n.id,
            label=n.label,
            shape=n.shape,
            parent=n.parent,
            x=n.x,
            y=n.y,
            w=n.w,
            h=n.h,
            fill=n.fill,
            stroke=n.stroke,
            font_color=n.font_color,
            font_size=_parse_font_size(n.font_size),
            dashed=n.dashed,
            rounded=n.rounded,
            container=n.container,
        )
        for n in p.nodes
    ]
    # 悬空边（vendored 已把指向不存在节点的 source/target 置 None）直接过滤
    edges = [
        DrawioEdge(
            source=e.source,
            target=e.target,
            label=e.label,
            dashed=e.dashed,
            bidirectional=e.bidirectional,
            undirected=e.undirected,
            stroke=e.stroke,
        )
        for e in p.edges
        if e.source and e.target
    ]
    if not nodes and not edges:
        raise ValueError("无法解析 draw.io 源码: 未提取到任何节点或边")
    return DrawioDiagram(nodes, edges)


def _render_shape(shape: str, rounded: bool) -> str:
    """canonical shape → 渲染层 shape 名（rounded 决定 rect 圆角与否）。"""
    if shape == "rect":
        return "round" if rounded else "rectangle"
    if shape in (
        "ellipse",
        "triangle",
        "rhombus",
        "hexagon",
        "parallelogram",
        "cylinder",
    ):
        return shape
    return "round" if rounded else "rectangle"


def layout_drawio(
    diagram: DrawioDiagram, *, max_w: float = 12.0, max_h: float = 6.75
) -> DiagramLayout:
    """把 draw.io IR 布局成坐标（inches）：归一化 + 等比 fit + 非绑定轴居中 + shape 名 + 配色。

    vendored parse_page 已解析绝对坐标（父链累加）；这里只做：等比缩放 fit 到
    max_w×max_h（不改变相对位置与配色）、非绑定轴居中留白对称、canonical shape →
    渲染层 shape 名、字号按 scale 换算（font_pt，12pt 默认）。容器节点保留自身几何
    （draw.io 已摆好），render 层画背景框。
    """
    nodes = diagram.nodes
    if not nodes:
        raise ValueError("无法布局 draw.io 图: 无节点")
    xs = [n.x for n in nodes] + [n.x + n.w for n in nodes]
    ys = [n.y for n in nodes] + [n.y + n.h for n in nodes]
    x0, y0 = min(xs), min(ys)
    x1, y1 = max(xs), max(ys)
    raw_w, raw_h = x1 - x0, y1 - y0
    scale = min(
        1.0,
        max_w / raw_w if raw_w > 0 else 1.0,
        max_h / raw_h if raw_h > 0 else 1.0,
    )
    content_w, content_h = raw_w * scale, raw_h * scale
    # 非绑定轴居中（绑定轴 content 正好填满画布 → off=0）；留白对称，不再贴左上角
    off_x = (max_w - content_w) / 2
    off_y = (max_h - content_h) / 2
    placed = [
        PlacedNode(
            id=n.id,
            label=n.label,
            shape=_render_shape(n.shape, n.rounded),
            x=(n.x - x0) * scale + off_x,
            y=(n.y - y0) * scale + off_y,
            w=n.w * scale,
            h=n.h * scale,
            is_container=n.container,
            parent=n.parent,
            fill=n.fill,
            stroke=n.stroke,
            font_color=n.font_color,
            font_pt=(n.font_size or 12.0) * scale * 72,
        )
        for n in nodes
    ]
    edges = _layout_edges(diagram, placed)
    return DiagramLayout(placed, edges, content_w, content_h)


def _edge_anchors(s: PlacedNode, t: PlacedNode) -> tuple[tuple[float, float], tuple[float, float]]:
    """按源/目标相对位置选边缘中点（drawio 无方向概念，固定规则，不做 BT/RL 反转）。"""
    scx, scy = s.x + s.w / 2, s.y + s.h / 2
    tcx, tcy = t.x + t.w / 2, t.y + t.h / 2
    if abs(tcx - scx) >= abs(tcy - scy):
        if tcx > scx:
            return (s.x + s.w, s.y + s.h / 2), (t.x, t.y + t.h / 2)
        return (s.x, s.y + s.h / 2), (t.x + t.w, t.y + t.h / 2)
    if tcy > scy:
        return (s.x + s.w / 2, s.y + s.h), (t.x + t.w / 2, t.y)
    return (s.x + s.w / 2, s.y), (t.x + t.w / 2, t.y + t.h)


def _layout_edges(diagram: DrawioDiagram, placed: list[PlacedNode]) -> list[PlacedEdge]:
    by_id = {n.id: n for n in placed}
    out: list[PlacedEdge] = []
    for e in diagram.edges:
        s, t = by_id.get(e.source), by_id.get(e.target)
        if s is None or t is None:
            continue  # 悬空边（parse 已过滤，此处防御）
        a1, a2 = _edge_anchors(s, t)
        style = "dashed" if e.dashed else "solid"
        out.append(
            PlacedEdge(
                e.source,
                e.target,
                e.label,
                style,
                "arrow",
                e.undirected,
                a1[0],
                a1[1],
                a2[0],
                a2[1],
                stroke=e.stroke,
            )
        )
    return out


def drawio_to_pptx(source, out_path: str, *, page=None) -> str:
    """.drawio 文件 → 可编辑 PPTX（16:9 整页）。返回 out_path。

    python-pptx 惰性 import（deck extra），未安装时给可操作错误。
    """
    diagram = parse_drawio(source, page=page)
    lay = layout_drawio(diagram)
    try:
        from pptx import Presentation
        from pptx.util import Emu
    except ImportError as e:  # pragma: no cover
        raise RuntimeError("drawio_to_pptx 需要 python-pptx，请安装 offipy[deck]") from e
    prs = Presentation()
    prs.slide_width = Emu(12192000)
    prs.slide_height = Emu(6858000)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    render_to_slide(slide, lay)
    prs.save(out_path)
    return str(out_path)


class _DrawioHTMLParser(HTMLParser):
    """扫 <div class="drawio" data-drawio="..."> → [{slide, path}]。slide 1-based。"""

    def __init__(self) -> None:
        super().__init__()
        self.slide_index = 0
        self.decls: list[dict] = []

    def handle_starttag(self, tag: str, attrs) -> None:
        d = {k: (v or "") for k, v in attrs}
        if tag == "section" and "data-pptx-slide" in d:
            self.slide_index += 1
            return
        if tag == "div" and "drawio" in d.get("class", "").split() and d.get("data-drawio"):
            self.decls.append({"slide": self.slide_index, "path": d["data-drawio"]})


def parse_drawio_declarations(html_text: str) -> list[dict]:
    """HTML → [{slide, path}]（slide 1-based）。div.drawio 在 data-pptx-slide
    <section> 之外 → ValueError；同页多个 drawio 声明 → ValueError（对齐 charts
    每页仅支持一个容器）。"""
    p = _DrawioHTMLParser()
    p.feed(html_text)
    p.close()
    for d in p.decls:
        if d["slide"] <= 0:
            raise ValueError(
                'drawio 声明出现在 slide 之外——<div class="drawio"> 必须在 '
                "data-pptx-slide 的 <section> 内"
            )
    seen = [d["slide"] for d in p.decls]
    dupes = sorted({s for s in seen if seen.count(s) > 1})
    if dupes:
        raise ValueError(f"每页仅支持一个 drawio 图（第 {dupes} 页声明了多个）")
    return p.decls


def load_drawio_boxes(measurements_path: str) -> dict[int, dict]:
    """measurements.json → {slide_index(1-based): {"x","y","w","h"}}（px）。

    匹配 record className 分词含 "drawio"（div.drawio 的 rect）。
    """
    with open(measurements_path, encoding="utf-8") as f:
        data = json.load(f)
    boxes: dict[int, dict] = {}
    for i, slide in enumerate(data.get("slides", []), start=1):
        for rec in slide.get("records", []):
            cls = (rec.get("className") or "").split()
            if "drawio" in cls:
                boxes[i] = dict(rec["rect"])
                break
    return boxes


def _measurements_path(pptx_path: str) -> str:
    p = Path(pptx_path)
    return str(p.with_name(f"{p.stem}_audit") / "_cache" / "measurements.json")


def _resolve_source_path(path: str) -> Path:
    """data-drawio 源路径归一化：file:// URI → 真实路径（deck 管线改写相对→file://）。"""
    if path.startswith("file://"):
        return Path(url2pathname(urlparse(path).path))
    return Path(path)


def _remove_bbox_shapes(slide, box_emu: dict) -> None:
    """删除与注入矩形几何一致（同位置同尺寸）的占位形状（div.drawio 占位块等）。

    占位矩形是确定尺寸的（box_emu = px→EMU），用户任意形状几乎不可能精确重合；
    旧「中心点落在 bbox 内」判定会误删用户内容。容差 0.01in 吸收 px→EMU 舍入。
    """
    tol = int(0.01 * 914400)  # 0.01in，EMU
    bx, by, bw, bh = box_emu["x"], box_emu["y"], box_emu["w"], box_emu["h"]
    for shape in list(slide.shapes):
        if (
            abs(shape.left - bx) <= tol
            and abs(shape.top - by) <= tol
            and abs(shape.width - bw) <= tol
            and abs(shape.height - bh) <= tol
        ):
            slide.shapes._spTree.remove(shape._element)


def inject_drawio(pptx_path: str, decls: list[dict], boxes: dict[int, dict]) -> None:
    """把每块 drawio 渲染成可编辑形状，替换 slide 内对应 bbox 占位。"""
    missing = [d["path"] for d in decls if not _resolve_source_path(d["path"]).is_file()]
    if missing:
        raise RuntimeError("drawio 源文件缺失: " + "、".join(missing))
    from pptx import Presentation

    prs = Presentation(pptx_path)
    for decl in decls:
        box = boxes.get(decl["slide"])
        if box is None:
            continue
        if box["w"] <= 0 or box["h"] <= 0:
            raise RuntimeError(
                f"第 {decl['slide']} 页 drawio 容器测量尺寸为 0（w={box['w']}, "
                f"h={box['h']}）——请给 div.drawio 设固定宽/高"
            )
        box_emu = {k: int(v * PX_TO_EMU) for k, v in box.items()}
        slide = prs.slides[decl["slide"] - 1]
        try:
            diagram = parse_drawio(decl["path"])
        except ValueError as e:
            raise ValueError(f"第 {decl['slide']} 页 drawio 源解析失败: {e}") from None
        # bbox 是 px；转换器画布 1920px = 13.333in = 12192000 EMU → 1px = 1/144 in。
        # layout 的 max_w/max_h 单位是 inches，px→inch 除以 144（不能除以 PX_TO_EMU）。
        lay = layout_drawio(diagram, max_w=box["w"] / 144, max_h=box["h"] / 144)
        _remove_bbox_shapes(slide, box_emu)
        render_to_slide(slide, lay, offset_x=box_emu["x"], offset_y=box_emu["y"])
    prs.save(pptx_path)


def postprocess_drawio(html_path: str, pptx_path: str) -> None:
    """转换后调用（对齐 charts.postprocess_charts 签名）：HTML 含 <div class="drawio"
    data-drawio="..."> → 读 measurements → 注入可编辑形状。无声明 → 原样返回。
    声明/数据非法 → ValueError（由 deck._postprocess 归一化）。
    """
    with open(html_path, encoding="utf-8") as f:
        html_text = f.read()
    decls = parse_drawio_declarations(html_text)
    if not decls:
        return
    meas_path = _measurements_path(pptx_path)
    if not os.path.exists(meas_path):
        raise RuntimeError(
            f"找不到 convert 审计产物 {meas_path}——drawio 注入需要 measurements.json，"
            "请勿用 --no-visual-audit"
        )
    boxes = load_drawio_boxes(meas_path)
    missing = [d["slide"] for d in decls if d["slide"] not in boxes]
    if missing:
        raise RuntimeError(
            f'第 {missing} 页没测到 <div class="drawio"> 容器——请确认 div 有 '
            'class="drawio" 且被渲染进 visual audit'
        )
    # data-drawio 路径归一化：file:// URI → 真实绝对路径（deck 管线改写相对→file://）；
    # 相对路径基于 HTML 所在目录解析。一律改写为绝对路径再传 inject（parse_drawio 不接受 file://）。
    base = Path(html_path).parent
    for d in decls:
        p = _resolve_source_path(d["path"])
        if not p.is_absolute():
            p = (base / p).resolve()
        d["path"] = str(p)
    inject_drawio(pptx_path, decls, boxes)
