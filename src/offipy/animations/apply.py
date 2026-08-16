"""动画注入执行层：独立 API（按形状名）+ 管线入口（OFFIPY_ELEM 锚）。

共用注入核心：按 slide 分组 → 定位 spid → 构建 timing/transition → 按
p:sld schema 顺序（cSld → clrMapOvr → transition → timing）插入 → 保存 → 报告。
"""

from __future__ import annotations

import json
from pathlib import Path
from typing import Any

from lxml import etree
from pptx import Presentation

from offipy.exceptions import InvalidArgumentError

from .spec import AnimationSpec, TransitionSpec, parse_declaration
from .timing import AnimationUnit, build_timing
from .transition import build_transition


def _coerce_animations(items: list[Any] | None) -> list[AnimationSpec]:
    out: list[AnimationSpec] = []
    for item in items or []:
        if isinstance(item, AnimationSpec):
            out.append(item)
        elif isinstance(item, dict):
            try:
                out.append(AnimationSpec(**item))  # __post_init__ 校验 → InvalidArgumentError
            except TypeError as exc:
                raise InvalidArgumentError(f"动画声明 dict 含非法字段：{exc}") from exc
        else:
            raise InvalidArgumentError(
                f"动画声明必须是 dict 或 AnimationSpec（实际 {type(item).__name__}）"
            )
    return out


def _coerce_transitions(items: list[Any] | None) -> list[TransitionSpec]:
    out: list[TransitionSpec] = []
    for item in items or []:
        if isinstance(item, TransitionSpec):
            out.append(item)
        elif isinstance(item, dict):
            try:
                out.append(TransitionSpec(**item))
            except TypeError as exc:
                raise InvalidArgumentError(f"过渡声明 dict 含非法字段：{exc}") from exc
        else:
            raise InvalidArgumentError(
                f"过渡声明必须是 dict 或 TransitionSpec（实际 {type(item).__name__}）"
            )
    return out


def _resolve_spids(slide: Any, target: str) -> list[int]:
    """按形状名精确匹配 → spid 列表（同名多形状视为同一单元）。"""
    return [sh.shape_id for sh in slide.shapes if sh.name == target]


def _insert_sld_elements(
    sld_el: etree._Element,
    transition_el: etree._Element | None,
    timing_el: etree._Element | None,
) -> None:
    """按 p:sld schema 顺序插入：cSld → clrMapOvr → transition → timing。"""
    idx = len(list(sld_el))
    for i, child in enumerate(sld_el):
        tag = etree.QName(child).localname
        if tag == "clrMapOvr" or tag == "transition":
            idx = i + 1
        elif tag == "timing":
            # 幂等性已在上游拦截，此处仅防御性覆盖
            idx = i
            break
    if transition_el is not None:
        sld_el.insert(idx, transition_el)
        idx += 1
    if timing_el is not None:
        sld_el.insert(idx, timing_el)


def _inject_slide(
    slide: Any,
    slide_idx: int,
    animations: list[AnimationSpec],
    transitions: list[TransitionSpec],
) -> tuple[int, bool, list[dict[str, Any]], list[dict[str, Any]]]:
    """对单页注入。返回 (anim_applied, trans_injected, unmatched, skipped)。

    trans_injected 表示本页实际写入了 <p:transition>（供 apply_animations 汇总
    transitions_applied）。slide 已有 <p:transition>（无 timing）时跳过全部过渡
    请求，不产生第二个 transition（CT_Slide transition maxOccurs=1），并把它们
    记入 skipped。slide_idx 为 1-based。
    """
    units: list[AnimationUnit] = []
    unmatched: list[dict[str, Any]] = []
    for a in animations:
        if a.slide != slide_idx:
            continue
        spids = _resolve_spids(slide, a.target)
        if not spids:
            unmatched.append({"slide": a.slide, "target": a.target})
            continue
        units.append(
            AnimationUnit(
                spids=spids,
                effect=a.effect,
                direction=a.direction,
                trigger=a.trigger,
                duration_ms=round(a.duration * 1000),
                delay_ms=round(a.delay * 1000),
            )
        )

    transitions_for_slide = [t for t in transitions if t.slide == slide_idx]
    if not units and not transitions_for_slide:
        return 0, False, unmatched, []

    # 幂等性：slide 已有 <p:timing> → 拒绝（spec §错误处理）
    existing_timing = slide._element.findall(
        ".//{http://schemas.openxmlformats.org/presentationml/2006/main}timing"
    )
    if existing_timing:
        raise InvalidArgumentError(
            f"slide {slide_idx} 已有动画（<p:timing> 已存在），请先移除再注入"
        )

    timing_el = build_timing(units)
    # 每页一个过渡（spec）：同页多余声明跳过并记入报告，不静默覆盖。
    has_existing_transition = bool(
        slide._element.findall(
            ".//{http://schemas.openxmlformats.org/presentationml/2006/main}transition"
        )
    )
    transition_el = None
    trans_injected = False
    skipped: list[dict[str, Any]] = []
    if transitions_for_slide:
        if has_existing_transition:
            # slide 已有过渡 → 跳过全部请求，不产生第二个 transition（CT_Slide maxOccurs=1）
            skipped = [
                {"slide": slide_idx, "kind": t.kind, "reason": "slide 已有过渡"}
                for t in transitions_for_slide
            ]
        else:
            transition_el = build_transition(
                transitions_for_slide[0].kind, transitions_for_slide[0].speed
            )
            trans_injected = True
            skipped = [
                {"slide": slide_idx, "kind": extra.kind, "reason": "每页仅一个过渡"}
                for extra in transitions_for_slide[1:]
            ]
    _insert_sld_elements(slide._element, transition_el, timing_el)
    return len(units), trans_injected, unmatched, skipped


def apply_animations(
    pptx: str,
    animations: list[AnimationSpec] | list[dict[str, Any]] | None = None,
    transitions: list[TransitionSpec] | list[dict[str, Any]] | None = None,
    raise_on_all_unmatched: bool = True,
) -> dict[str, Any]:
    """对现成 .pptx 注入动画/过渡。返回注入报告。

    每个 target 按形状名精确匹配；找不到 → 报告 unmatched 不硬失败；
    全部动画 target 未命中且 raise_on_all_unmatched=True → InvalidArgumentError
    （管线入口传 False，把全 miss 降级为告警，不硬失败）。已带 <p:timing> 的
    slide → InvalidArgumentError（幂等性）。同页多个过渡 → 只注入第一个，多余记 skipped。
    """
    anim_specs = _coerce_animations(animations)
    trans_specs = _coerce_transitions(transitions)
    if not anim_specs and not trans_specs:
        return {"animations_applied": 0, "transitions_applied": 0, "unmatched": [], "skipped": []}

    prs = Presentation(str(pptx))
    total_anim = 0
    total_trans = 0
    all_unmatched: list[dict[str, Any]] = []
    all_skipped: list[dict[str, Any]] = []
    for slide_idx, slide in enumerate(prs.slides, 1):
        applied, trans_injected, unmatched, skipped = _inject_slide(
            slide, slide_idx, anim_specs, trans_specs
        )
        total_anim += applied
        if trans_injected:
            total_trans += 1  # 每页至多注入一个过渡（实际写入才计）
        all_unmatched.extend(unmatched)
        all_skipped.extend(skipped)

    # 全 miss 是纯错误路径：不写回文件，保持原 pptx 不变
    if (
        raise_on_all_unmatched
        and anim_specs
        and total_anim == 0
        and len(all_unmatched) == len(anim_specs)
    ):
        raise InvalidArgumentError(
            f"全部 {len(anim_specs)} 个动画 target 均未匹配到形状："
            + "、".join(f"slide {u['slide']} {u['target']}" for u in all_unmatched)
        )

    prs.save(str(pptx))

    return {
        "animations_applied": total_anim,
        "transitions_applied": total_trans,
        "unmatched": all_unmatched,
        "skipped": all_skipped,
    }


def apply_transitions(
    pptx: str,
    transitions: list[TransitionSpec] | list[dict[str, Any]],
) -> dict[str, Any]:
    """只注入页面过渡（薄封装）。"""
    return apply_animations(pptx, animations=None, transitions=transitions)


def _measurements_path(pptx_path: str) -> Path:
    """与 charts.py 一致：<pptx stem>_audit/_cache/measurements.json。"""
    p = Path(pptx_path)
    return p.with_name(f"{p.stem}_audit") / "_cache" / "measurements.json"


def _warning(message: str) -> dict[str, Any]:
    """anim 告警条目（对齐 deck._measure_warnings 消费的 {kind, message} 结构）。"""
    return {"kind": "anim", "message": message}


def postprocess_animations(html_path: str, pptx_path: str) -> dict[str, Any]:  # noqa: ARG001
    """管线入口：读 measurements.json → 按 OFFIPY_ELEM::<elem_id> 定位形状 → 注入。

    spec §管线接入：record 有 anim_decl 但 assemble 没产出形状 → unmatched/告警，
    不硬失败（转换不因动画缺失而失败）。无声明 → no-op。measurements 缺失或
    不可读 → no-op。告警（含 spec 解析告警）追加进 measurements.json 的 _warnings
    （deck._measure_warnings 的 "anim" 码透出）。html_path 保持管线签名一致（未用）。
    """
    meas_path = _measurements_path(pptx_path)
    if not meas_path.is_file():
        return {"animations_applied": 0, "transitions_applied": 0, "unmatched": [], "skipped": []}
    try:
        data = json.loads(meas_path.read_text(encoding="utf-8"))
    except (json.JSONDecodeError, UnicodeError, OSError):
        return {"animations_applied": 0, "transitions_applied": 0, "unmatched": [], "skipped": []}

    animations: list[AnimationSpec] = []
    transitions: list[TransitionSpec] = []
    unmatched: list[dict[str, Any]] = []
    warnings: list[dict[str, Any]] = []

    slides = data.get("slides") or []
    for i, sdata in enumerate(slides, start=1):
        if not isinstance(sdata, dict):
            continue
        slide_meta = sdata.get("slide") or {}
        tdecl = slide_meta.get("transition_decl")
        if isinstance(tdecl, dict):
            kind = tdecl.get("kind")
            speed = tdecl.get("speed", "medium")
            if kind in {"fade", "wipe", "push", "cover"}:
                transitions.append(TransitionSpec(slide=i, kind=kind, speed=speed))
            else:
                warnings.append(_warning(f"slide {i} 未知过渡类型 {kind!r}，跳过"))
        for rec in sdata.get("records") or []:
            if not isinstance(rec, dict):
                continue
            raw = rec.get("anim_decl")
            if not raw:
                continue
            spec = parse_declaration(raw, rec)
            if spec is None:
                # spec.py 已把解析告警写成 {kind, message} dict 落进 raw['_warnings']
                warnings.extend(
                    w
                    for w in raw.get("_warnings") or []
                    if isinstance(w, dict) and isinstance(w.get("message"), str)
                )
                continue
            elem_id = rec.get("elem_id")
            if not elem_id:
                warnings.append(_warning(f"slide {i} 有动画声明但缺 elem_id，跳过"))
                continue
            target = f"OFFIPY_ELEM::{elem_id}"
            animations.append(
                AnimationSpec(
                    slide=i,
                    target=target,
                    effect=spec.effect,
                    direction=spec.direction,
                    trigger=spec.trigger,
                    duration=spec.duration,
                    delay=spec.delay,
                )
            )

    if not animations and not transitions:
        report: dict[str, Any] = {
            "animations_applied": 0,
            "transitions_applied": 0,
            "unmatched": unmatched,
            "skipped": [],
        }
    else:
        # 管线语义：全 miss 不抛（raise_on_all_unmatched=False），降级为告警。
        report = apply_animations(
            pptx_path, animations=animations, transitions=transitions, raise_on_all_unmatched=False
        )
        for u in report.get("unmatched", []):
            target = u.get("target") or ""
            if target.startswith("OFFIPY_ELEM::"):
                unmatched.append(
                    {"elem_id": target[len("OFFIPY_ELEM::") :], "slide": u.get("slide")}
                )
            else:
                unmatched.append(u)
        warnings.extend(
            _warning(
                f"slide {u.get('slide')} 元素 {u.get('elem_id')} 未产出可动画形状（unmatched）"
            )
            for u in unmatched
        )
        report["unmatched"] = unmatched

    # 有告警才写回 measurements.json（内部产物；无告警不动缓存）
    if warnings:
        _append_warnings(meas_path, warnings)
        report["warnings"] = [w["message"] for w in warnings]
    return report


def _append_warnings(meas_path: Path, warnings: list[dict[str, Any]]) -> None:
    try:
        data = json.loads(meas_path.read_text(encoding="utf-8"))
    except (json.JSONDecodeError, UnicodeError, OSError):
        return
    if not isinstance(data, dict):
        return
    existing = data.get("_warnings")
    if not isinstance(existing, list):
        existing = []
    existing.extend(warnings)
    data["_warnings"] = existing
    meas_path.write_text(json.dumps(data, ensure_ascii=False, indent=1), encoding="utf-8")
