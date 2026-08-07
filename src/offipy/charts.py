# src/offipy/charts.py
"""原生图表：HTML 声明 → 替换为 PowerPoint 原生可编辑图表。

HTML 给图表容器打 `data-chart="<type>"`（type ∈ bar/line/pie），数据来自：
  1. 容器自身 `data-chart-data` JSON 属性（categories + series），或
  2. 页面 `<script type="application/json" data-chart-target="<css选择器>">` 块。
convert.py 转换照常跑（图表区渲染成占位形状）；转换后读 convert 审计产物
`<out>_audit/_cache/measurements.json`（含每元素 className + rect 坐标），
用 python-pptx 把图表区占位形状替换成原生图表。1 px = 6350 EMU（转换器画布
1920×1080 px = 12192000×6858000 EMU）。
"""

from __future__ import annotations

import json
import re
from collections.abc import Set as AbstractSet
from dataclasses import dataclass
from html.parser import HTMLParser

from .design import THEMES, Theme

PX_TO_EMU = 6350

CHART_TYPES = ("bar", "line", "pie")


@dataclass
class ChartSeries:
    name: str
    values: list[float]


@dataclass
class ChartData:
    categories: list[str]
    series: list[ChartSeries]


@dataclass
class ChartDecl:
    slide_index: int  # 1-based，对齐 HTML section 顺序 / convert only-slides
    chart_type: str
    data: ChartData
    slide_classes: frozenset[str] = frozenset()  # 包含它的 <section class="..."> 的 class 集合
    colors_override: tuple[str, ...] | None = None  # 解析后的 data-chart-colors（大写 #RRGGBB）


@dataclass
class _Container:
    slide_index: int
    chart_type: str
    data_json: str | None = None
    el_id: str | None = None
    slide_classes: frozenset[str] = frozenset()
    colors_override: tuple[str, ...] | None = None


def _parse_data(json_text: str) -> ChartData:
    raw = json.loads(json_text)
    if not isinstance(raw, dict):
        raise ValueError("图表数据必须是对象 {categories, series}")
    categories = raw.get("categories")
    series_raw = raw.get("series")
    if (
        not isinstance(categories, list)
        or not categories
        or not all(isinstance(c, str) for c in categories)
    ):
        raise ValueError("图表数据 categories 必须是非空字符串列表")
    if not isinstance(series_raw, list) or not series_raw:
        raise ValueError("图表数据 series 必须是非空列表")
    series: list[ChartSeries] = []
    for sr in series_raw:
        if not isinstance(sr, dict):
            raise ValueError("图表数据 series 每项必须是对象 {name, values}")
        name = sr.get("name", "")
        values = sr.get("values")
        if (
            not isinstance(values, list)
            or not values
            or not all(isinstance(v, (int, float)) and not isinstance(v, bool) for v in values)
        ):
            raise ValueError("图表数据 series.values 必须是非空数值列表")
        series.append(ChartSeries(name=str(name), values=[float(v) for v in values]))
    return ChartData(categories=categories, series=series)


# 确定性调色板的固定 H/S/L 偏移（Task 0 冻结）。每个元素 = (hue 偏移°, sat 偏移, light 偏移)。
# 首色恒为 accent 本身；其余五项围绕 accent 做固定小偏移，纯函数、无随机。
_PALETTE_HSL_OFFSETS: tuple[tuple[float, float, float], ...] = (
    (0.0, 0.0, 0.0),  # 0: accent 本身
    (20.0, 0.0, 0.08),  # 1: 更亮的靛蓝
    (-25.0, 0.0, 0.06),  # 2: 偏青的亮蓝
    (45.0, 0.0, -0.02),  # 3: 紫罗兰
    (-60.0, -0.10, 0.05),  # 4: 亮青绿
    (150.0, -0.25, 0.10),  # 5: 暖珊瑚（蓝的补色方向，柔和化）
)

_HEX_RE = re.compile(r"^#?[0-9A-Fa-f]{6}$")  # accent：可带可不带 #
_HEX_STRICT_RE = re.compile(r"^#[0-9A-Fa-f]{6}$")  # data-chart-colors 项：必须带 #
_THEME_STYLE_RE = re.compile(  # <style data-theme="NAME"> 块（injected 全量或空占位）
    r"<style\s+data-theme=['\"]([A-Za-z0-9-]+)['\"]", re.IGNORECASE
)
_VARIANT_SELECTOR_RE = re.compile(r"^\.slide\.([A-Za-z0-9_-]+)$")  # 冻结内置反色选择器


def parse_chart_colors_override(raw: str) -> tuple[str, ...]:
    """解析 data-chart-colors：严格 JSON 的 #RRGGBB 字符串列表（>=1 项），非法即 ValueError。

    归一化为大写 #RRGGBB。颜色数少于 series/point 数时由调用方按 _cycle_colors 循环补足。
    """
    try:
        value = json.loads(raw)
    except json.JSONDecodeError:
        raise ValueError(f"data-chart-colors 必须是 #RRGGBB 列表: {raw!r}") from None
    if not isinstance(value, list) or not value:
        raise ValueError(f"data-chart-colors 必须是 #RRGGBB 列表: {raw!r}")
    colors: list[str] = []
    for entry in value:
        if not isinstance(entry, str) or not _HEX_STRICT_RE.fullmatch(entry):
            raise ValueError(f"data-chart-colors 必须是 #RRGGBB 列表: {raw!r}")
        colors.append(entry.upper())
    return tuple(colors)


def _rgb_to_hsl(r: int, g: int, b: int) -> tuple[float, float, float]:
    """RGB(0-255) → HSL(H∈[0,360), S∈[0,1], L∈[0,1])。纯 stdlib math，确定性。"""
    rf, gf, bf = r / 255.0, g / 255.0, b / 255.0
    mx, mn = max(rf, gf, bf), min(rf, gf, bf)
    light = (mx + mn) / 2.0
    if mx == mn:
        return 0.0, 0.0, light
    delta = mx - mn
    s = delta / (1.0 - abs(2.0 * light - 1.0))
    if mx == rf:
        h = 60.0 * (((gf - bf) / delta) % 6.0)
    elif mx == gf:
        h = 60.0 * (((bf - rf) / delta) + 2.0)
    else:
        h = 60.0 * (((rf - gf) / delta) + 4.0)
    return h, s, light


def _hsl_to_rgb(h: float, s: float, light: float) -> tuple[int, int, int]:
    """HSL → RGB(0-255)。h 先 mod 360；s/light 假定调用方已 clamp。纯 stdlib math，确定性。"""
    h = h % 360.0
    c = (1.0 - abs(2.0 * light - 1.0)) * s
    hp = h / 60.0
    x = c * (1.0 - abs(hp % 2.0 - 1.0))
    if hp < 1.0:
        r, g, b = c, x, 0.0
    elif hp < 2.0:
        r, g, b = x, c, 0.0
    elif hp < 3.0:
        r, g, b = 0.0, c, x
    elif hp < 4.0:
        r, g, b = 0.0, x, c
    elif hp < 5.0:
        r, g, b = x, 0.0, c
    else:
        r, g, b = c, 0.0, x
    m = light - c / 2.0
    return round((r + m) * 255.0), round((g + m) * 255.0), round((b + m) * 255.0)


def derive_chart_palette(accent: str) -> tuple[str, str, str, str, str, str]:
    """由单一 accent 确定性推导 6 色调色板（首色 = accent 本身）。

    RGB → HSL 后对 _PALETTE_HSL_OFFSETS 做固定 H/S/L 偏移、确定性 clamp，再转回大写
    #RRGGBB。纯函数无随机：同输入 → 字节级相同（Python 3.10-3.13）。
    """
    if not _HEX_RE.fullmatch(accent):
        raise ValueError(f"accent 必须是 #RRGGBB（或 RRGGBB）: {accent!r}")
    hex_part = accent.lstrip("#")
    r, g, b = int(hex_part[0:2], 16), int(hex_part[2:4], 16), int(hex_part[4:6], 16)
    h, s, light = _rgb_to_hsl(r, g, b)
    colors: list[str] = []
    for dh, ds, dl in _PALETTE_HSL_OFFSETS:
        nh = (h + dh) % 360.0
        ns = min(1.0, max(0.0, s + ds))
        nl = min(1.0, max(0.0, light + dl))
        nr, ng, nb = _hsl_to_rgb(nh, ns, nl)
        colors.append(f"#{nr:02X}{ng:02X}{nb:02X}")
    c0, c1, c2, c3, c4, c5 = colors
    return c0, c1, c2, c3, c4, c5


def _cycle_colors(colors: tuple[str, ...], n: int) -> tuple[str, ...]:
    """把调色板循环扩展/截断到 n 个（series 或 point 数）。colors 必须非空。"""
    if not colors:
        raise ValueError("调色板不能为空")
    return tuple(colors[i % len(colors)] for i in range(n))


def _theme_name_from_html(html_text: str) -> str | None:
    """从 HTML 的 `<style data-theme="NAME">` 块提取主题名；无则 None。

    只做裸匹配、不校验 THEMES——injected 全量块和空占位块都能识别。
    """
    m = _THEME_STYLE_RE.search(html_text)
    return m.group(1) if m else None


def _effective_accent(theme: Theme, slide_classes: AbstractSet[str]) -> str | None:
    """页面实际生效的 accent：命中 `.slide.<class>` → variant，否则 base。

    只认冻结内置形式 `.slide.<class>`；其它选择器形态静默回落 base，不实现通用 CSS 语义。
    """
    if theme.variant_selector:
        m = _VARIANT_SELECTOR_RE.match(theme.variant_selector)
        if m and m.group(1) in slide_classes:
            variant_accent = theme.variant_vars.get("--accent")
            if variant_accent is not None:
                return variant_accent
    return theme.base_vars.get("--accent")


def _decl_hex_palette(decl: ChartDecl, theme: Theme | None) -> tuple[str, ...] | None:
    """按优先级解析一张图表的调色板（大写 #RRGGBB 元组）：

    1. data-chart-colors 显式覆盖 → 原样返回；
    2. 主题 accent（variant 命中优先）→ derive_chart_palette；
    3. 都没有 → None（调用方回退到固定 _palette()）。
    """
    if decl.colors_override is not None:
        return decl.colors_override
    if theme is not None:
        accent = _effective_accent(theme, decl.slide_classes)
        if accent is not None:
            return derive_chart_palette(accent)
    return None


class _ChartHTMLParser(HTMLParser):
    """提取：每页的 .chart 容器 + 页内 <script data-chart-target> 数据块。"""

    def __init__(self) -> None:
        super().__init__()
        self.slide_index = 0
        self.cur_containers: list[_Container] = []
        self.all_containers: list[_Container] = []
        self.scripts: list[tuple[str, str]] = []  # (target_selector, json_text)
        self._script_target: str | None = None
        self._script_buf: str | None = None
        self._current_slide_classes: frozenset[str] = frozenset()

    def handle_starttag(self, tag: str, attrs: list[tuple[str, str | None]]) -> None:
        d = {k: (v or "") for k, v in attrs}
        if tag == "section" and "data-pptx-slide" in d:
            if self.cur_containers:
                self.all_containers.extend(self.cur_containers)
            self.slide_index += 1
            self._current_slide_classes = frozenset(d.get("class", "").split())
            self.cur_containers = []
            return
        if tag == "div" and "data-chart" in d:
            cls = d.get("class", "").split()
            if "chart" in cls:
                colors_override = None
                if "data-chart-colors" in d:
                    colors_override = parse_chart_colors_override(d["data-chart-colors"])
                self.cur_containers.append(
                    _Container(
                        slide_index=self.slide_index,
                        chart_type=d["data-chart"],
                        data_json=d.get("data-chart-data"),
                        el_id=d.get("id") or None,
                        slide_classes=self._current_slide_classes,
                        colors_override=colors_override,
                    )
                )
            return
        if tag == "script" and d.get("type") == "application/json" and d.get("data-chart-target"):
            self._script_target = d["data-chart-target"]
            self._script_buf = ""
            return

    def handle_endtag(self, tag: str) -> None:
        if tag == "script" and self._script_buf is not None:
            self.scripts.append((self._script_target or "", self._script_buf.strip()))
            self._script_target = None
            self._script_buf = None
            return

    def handle_data(self, data: str) -> None:
        if self._script_buf is not None:
            self._script_buf += data


def _resolve_data(container: _Container, scripts: list[tuple[str, str]]) -> ChartData:
    if container.data_json:
        try:
            return _parse_data(container.data_json)
        except (json.JSONDecodeError, ValueError) as e:
            raise ValueError(
                f"第 {container.slide_index} 页 data-chart-data 解析失败: {e}"
            ) from None
    # 从 script 块找：data-chart-target 指向该容器
    match = None
    for sel, json_text in scripts:
        if sel.startswith("#") and container.el_id == sel[1:]:
            match = json_text
            break
    if match is None:
        raise ValueError(
            f"第 {container.slide_index} 页图表 (data-chart={container.chart_type}) "
            "缺数据：请在容器加 data-chart-data，或用 <script data-chart-target> 提供"
        )
    try:
        return _parse_data(match)
    except (json.JSONDecodeError, ValueError) as e:
        raise ValueError(
            f"第 {container.slide_index} 页 <script data-chart-target> 数据解析失败: {e}"
        ) from None


def parse_chart_declarations(html_text: str) -> list[ChartDecl]:
    """解析 HTML → 图表声明列表（slide_index 1-based）。图表类型非法或数据缺失即 ValueError。"""
    p = _ChartHTMLParser()
    p.feed(html_text)
    p.close()
    # 关键：flush 最后一段 section 的容器（feed 结束时最后一页的 cur_containers 还在缓存里）
    if p.cur_containers:
        p.all_containers.extend(p.cur_containers)
    decls: list[ChartDecl] = []
    seen: set[int] = set()  # v1: 每页最多一个图表容器
    for container in p.all_containers:
        if container.chart_type not in CHART_TYPES:
            raise ValueError(
                f"第 {container.slide_index} 页图表类型非法: "
                f"{container.chart_type!r}（可选: bar/line/pie）"
            )
        if container.slide_index <= 0:
            raise ValueError(
                f"第 {container.slide_index} 页图表出现在 slide 之外——图表容器必须在 "
                "data-pptx-slide 的 <section> 内"
            )
        if container.slide_index in seen:
            raise ValueError(f"第 {container.slide_index} 页有多个图表容器——v1 每页仅支持一个图表")
        seen.add(container.slide_index)
        data = _resolve_data(container, p.scripts)
        decls.append(
            ChartDecl(
                slide_index=container.slide_index,
                chart_type=container.chart_type,
                data=data,
                slide_classes=container.slide_classes,
                colors_override=container.colors_override,
            )
        )
    return decls


def load_chart_boxes(measurements_path: str) -> dict[int, dict]:
    """读 measurements.json → {slide_index(1-based): {"x","y","w","h"}}（px，图表容器 bbox）。

    匹配规则：记录 className 分词后恰含 "chart"（chart-note 不算）。每页取第一个匹配
    （v1 契约：每页最多一个图表容器，与 parse_chart_declarations 一致）。
    """
    import json as _json

    with open(measurements_path, encoding="utf-8") as f:
        data = _json.load(f)
    boxes: dict[int, dict] = {}
    for i, slide in enumerate(data.get("slides", []), start=1):
        for rec in slide.get("records", []):
            cls = (rec.get("className") or "").split()
            if "chart" in cls:
                boxes[i] = dict(rec["rect"])
                break
    return boxes


def _measurements_path(pptx_path: str) -> str:
    from pathlib import Path

    p = Path(pptx_path)
    return str(p.with_name(f"{p.stem}_audit") / "_cache" / "measurements.json")


def inject_native_charts(
    pptx_path: str,
    decls: list[ChartDecl],
    boxes: dict[int, dict],
    theme: Theme | None = None,
) -> None:
    """把图表区占位形状替换成原生图表。boxes 缺某页 → 该页跳过（调用方保证完整）。"""
    from pptx import Presentation
    from pptx.enum.chart import XL_CHART_TYPE

    XL_TYPE = {
        "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "line": XL_CHART_TYPE.LINE_MARKERS,
        "pie": XL_CHART_TYPE.PIE,
    }
    prs = Presentation(pptx_path)
    for decl in decls:
        box = boxes.get(decl.slide_index)
        if box is None:
            continue
        slide = prs.slides[decl.slide_index - 1]
        _replace_with_chart(slide, decl, box, XL_TYPE, theme)
    prs.save(pptx_path)


def _palette() -> list:
    """惰性构造配色（RGBColor 需 import python-pptx，顶层 import 会拖慢无图表路径）。"""
    from pptx.dml.color import RGBColor

    return [
        RGBColor(0x22, 0x51, 0xFF),  # 主蓝（对齐 mckinsey --accent）
        RGBColor(0x0E, 0x93, 0x87),  # 青
        RGBColor(0xF5, 0x9E, 0x0B),  # 琥珀
        RGBColor(0xE0, 0x5D, 0x5D),  # 珊瑚
        RGBColor(0x7C, 0x3A, 0xED),  # 紫
        RGBColor(0x66, 0x70, 0x85),  # 灰（对齐 --muted）
    ]


def _hex_to_rgb(colors: tuple[str, ...]) -> list:
    """大写 #RRGGBB 元组 → RGBColor 列表（惰性 import python-pptx，与 _palette 一致）。"""
    from pptx.dml.color import RGBColor

    return [RGBColor(int(c[1:3], 16), int(c[3:5], 16), int(c[5:7], 16)) for c in colors]


def _replace_with_chart(
    slide, decl: ChartDecl, box: dict, xl_type, theme: Theme | None = None
) -> None:
    from pptx.chart.data import CategoryChartData

    x = int(box["x"] * PX_TO_EMU)
    y = int(box["y"] * PX_TO_EMU)
    cx = int(box["w"] * PX_TO_EMU)
    cy = int(box["h"] * PX_TO_EMU)
    cxm, cym = x + cx // 2, y + cy // 2
    # 移除中心落在图表 bbox 内的占位形状（容器 surface 矩形 + bar 矩形）
    for shape in list(slide.shapes):
        if (
            shape.left <= cxm <= shape.left + shape.width
            and shape.top <= cym <= shape.top + shape.height
        ):
            slide.shapes._spTree.remove(shape._element)
    cd = CategoryChartData()
    cd.categories = decl.data.categories
    for s in decl.data.series:
        cd.add_series(s.name, s.values)
    gframe = slide.shapes.add_chart(xl_type[decl.chart_type], x, y, cx, cy, cd)
    chart = gframe.chart
    chart.has_title = False
    chart.has_legend = len(decl.data.series) > 1
    hex_palette = _decl_hex_palette(decl, theme)
    colors = _hex_to_rgb(hex_palette) if hex_palette else _palette()
    if decl.chart_type == "bar":
        for i, series in enumerate(chart.plots[0].series):
            series.format.fill.solid()
            series.format.fill.fore_color.rgb = colors[i % len(colors)]
    elif decl.chart_type == "line":
        for i, series in enumerate(chart.plots[0].series):
            series.format.line.color.rgb = colors[i % len(colors)]
    elif decl.chart_type == "pie":
        plot = chart.plots[0]
        for i, point in enumerate(plot.series[0].points):
            point.format.fill.solid()
            point.format.fill.fore_color.rgb = colors[i % len(colors)]


def postprocess_charts(html_path: str, pptx_path: str) -> None:
    """转换后调用：HTML 含图表声明 → 读 measurements → 注入原生图表。

    无图表声明 → 原样返回。measurements.json 缺失（--no-visual-audit）→ RuntimeError，
    让调用方知道图表不会变成原生。图表声明/数据非法 → ValueError（从 parse 上抛）。
    """
    import os

    with open(html_path, encoding="utf-8") as f:
        html_text = f.read()
    if "data-chart" not in html_text:
        return
    decls = parse_chart_declarations(html_text)
    if not decls:
        return
    meas_path = _measurements_path(pptx_path)
    if not os.path.exists(meas_path):
        raise RuntimeError(
            f"找不到 convert 审计产物 {meas_path}——图表注入需要 measurements.json，"
            "请勿用 --no-visual-audit"
        )
    boxes = load_chart_boxes(meas_path)
    missing = [d.slide_index for d in decls if d.slide_index not in boxes]
    if missing:
        raise RuntimeError(
            f'第 {missing} 页没测到图表容器（class="chart"）——请确认 deck 用了 '
            "chart-dominant 布局且 deck make 带了 --layouts"
        )
    theme_name = _theme_name_from_html(html_text)
    theme = THEMES.get(theme_name) if theme_name else None
    inject_native_charts(pptx_path, decls, boxes, theme)
