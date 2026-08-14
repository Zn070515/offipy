# src/offipy/icons.py
"""原生图标：内联 <svg data-icon> → PowerPoint freeform 矢量图标。

HTML 给图标容器打空 `<svg class="icon" data-icon="<set>:<name>" viewBox="..."
width=".." height=".."></svg>`（path 由注入生成）。convert 照常跑（SVG 截图成
PNG 占位）；转换后读 `<out>_audit/_cache/measurements.json`（kind='svg' record
含 outerHTML/rect/color），用 python-pptx freeform 把占位替换成矢量图标并删
占位。1 px = 6350 EMU（转换器画布 1920×1080 px = 12192000×6858000 EMU）。

图标集（vendored 到 src/offipy/assets/icons/）：
  ph: Phosphor fill 权重（256 viewBox，单 path fill="currentColor"，fill 模式）
  lu: Lucide（24 viewBox，fill="none" stroke="currentColor"，stroke 模式）
图标颜色取 measurements 里 svg 的计算 color（HTML 设 color: var(--accent)）。

顶层不 import python-pptx（镜像 charts.py）：from pptx... 全在函数内惰性 import。
"""

from __future__ import annotations

import math
import re
from dataclasses import dataclass
from html.parser import HTMLParser
from pathlib import Path

PX_TO_EMU = 6350
ASSETS_DIR = Path(__file__).resolve().parent / "assets" / "icons"

# data-icon 前缀 → assets 子目录 / 渲染模式
SET_DIRS = {"ph": "phosphor", "lu": "lucide"}
SET_MODE = {"ph": "fill", "lu": "stroke"}

_CURVE_SAMPLES = 12  # 贝塞尔压平采样段数
_CIRCLE_SAMPLES = 32  # circle/ellipse 多边形近似段数
_ARC_SAMPLES = 24  # 圆弧压平采样段数

_TOKEN_RE = re.compile(r"[A-Za-z]|-?\d*\.?\d+(?:[eE][-+]?\d+)?")
_NUM_RE = re.compile(r"-?\d*\.?\d+(?:[eE][-+]?\d+)?")


@dataclass
class _SubPath:
    points: list[tuple[float, float]]
    close: bool
    filled: bool = False  # 元素级 fill="currentColor"（Lucide 混 fill 图标的内实心区域）


@dataclass
class IconDecl:
    slide_index: int  # 1-based，对齐 HTML section 顺序
    data_icon: str  # "ph:name"
    view_box: tuple[float, float, float, float]  # (x, y, w, h)


def _tokenize(d: str) -> list[str]:
    return _TOKEN_RE.findall(d)


def _cubic_points(p0, p1, p2, p3, n: int) -> list[tuple[float, float]]:
    pts = [p0]
    for i in range(1, n + 1):
        t = i / n
        mt = 1 - t
        x = mt**3 * p0[0] + 3 * mt**2 * t * p1[0] + 3 * mt * t**2 * p2[0] + t**3 * p3[0]
        y = mt**3 * p0[1] + 3 * mt**2 * t * p1[1] + 3 * mt * t**2 * p2[1] + t**3 * p3[1]
        pts.append((x, y))
    return pts


def _quad_points(p0, p1, p2, n: int) -> list[tuple[float, float]]:
    pts = [p0]
    for i in range(1, n + 1):
        t = i / n
        mt = 1 - t
        x = mt**2 * p0[0] + 2 * mt * t * p1[0] + t**2 * p2[0]
        y = mt**2 * p0[1] + 2 * mt * t * p1[1] + t**2 * p2[1]
        pts.append((x, y))
    return pts


def _arc_points(p0, radii, rot, large, sweep, p1, n: int) -> list[tuple[float, float]]:
    """SVG 椭圆弧 → 折线（endpoint→center 参数化，F.6.5）。"""
    rx, ry = abs(radii[0]), abs(radii[1])
    x0, y0 = p0
    x1, y1 = p1
    if (x0, y0) == (x1, y1):
        return [p0]
    if rx == 0 or ry == 0:
        return [p0, p1]  # 退化弧按 SVG F.6.2 等价于 lineto 到终点
    phi = math.radians(rot)
    cos_p, sin_p = math.cos(phi), math.sin(phi)
    dx, dy = (x0 - x1) / 2, (y0 - y1) / 2
    xp = cos_p * dx + sin_p * dy
    yp = -sin_p * dx + cos_p * dy
    lam = xp * xp / (rx * rx) + yp * yp / (ry * ry)
    if lam > 1:
        s = math.sqrt(lam)
        rx *= s
        ry *= s
    sign = -1.0 if large == sweep else 1.0
    num = rx * rx * ry * ry - rx * rx * yp * yp - ry * ry * xp * xp
    den = rx * rx * yp * yp + ry * ry * xp * xp
    coef = sign * math.sqrt(max(0.0, num / den))
    cxp = coef * (rx * yp / ry)
    cyp = coef * (-ry * xp / rx)
    cx = cos_p * cxp - sin_p * cyp + (x0 + x1) / 2
    cy = sin_p * cxp + cos_p * cyp + (y0 + y1) / 2

    def angle(ux, uy, vx, vy):
        dot = ux * vx + uy * vy
        m = math.hypot(ux, uy) * math.hypot(vx, vy)
        a = math.acos(max(-1.0, min(1.0, dot / m)))
        if ux * vy - uy * vx < 0:
            a = -a
        return a

    theta1 = angle(1, 0, (xp - cxp) / rx, (yp - cyp) / ry)
    dtheta = angle((xp - cxp) / rx, (yp - cyp) / ry, (-xp - cxp) / rx, (-yp - cyp) / ry)
    if sweep == 0 and dtheta > 0:
        dtheta -= 2 * math.pi
    elif sweep == 1 and dtheta < 0:
        dtheta += 2 * math.pi
    pts = [p0]
    for i in range(1, n + 1):
        a = theta1 + dtheta * (i / n)
        ca, sa = math.cos(a), math.sin(a)
        pts.append((cos_p * rx * ca - sin_p * ry * sa + cx, sin_p * rx * ca + cos_p * ry * sa + cy))
    pts[-1] = p1  # 终点由 SVG 规范保证精确等于 p1；采样末点存在浮点漂移，收拢到 p1
    return pts


def _parse_path(d: str) -> list[_SubPath]:
    """SVG path d 属性 → 压平后的子路径列表。曲线命令已折线化。"""
    toks = _tokenize(d)
    subpaths: list[_SubPath] = []
    pts: list[tuple[float, float]] = []
    close = False
    cx = cy = 0.0
    start = (0.0, 0.0)
    cmd: str | None = None
    prev_cubic: tuple[float, float] | None = None
    prev_quad: tuple[float, float] | None = None
    i = 0
    n = len(toks)
    first_cmd = True  # SVG path 首个命令必须是 moveto

    def read_num() -> float:
        nonlocal i
        if i >= n:
            raise ValueError(f"SVG path 参数不完整: {d!r}")
        v = float(toks[i])
        i += 1
        return v

    def read_point(rel: bool) -> tuple[float, float]:
        nonlocal cx, cy
        x, y = read_num(), read_num()
        if rel:
            x += cx
            y += cy
        return x, y

    while i < n:
        t = toks[i]
        if _NUM_RE.fullmatch(t):
            if cmd is None:
                raise ValueError(f"SVG path 首 token 必须是命令字母: {d!r}")
        else:
            cmd = t
            i += 1
        if cmd is None:
            raise ValueError(f"SVG path 无命令: {d!r}")
        c = cmd.upper()
        rel = cmd.islower()
        if first_cmd:
            if c != "M":
                raise ValueError(f"SVG path 首命令必须是 M/m: {d!r}")
            first_cmd = False
        if c == "M":
            if pts:
                subpaths.append(_SubPath(pts, close))
            p = read_point(rel)
            pts = [p]
            start = p
            cx, cy = p
            cmd = "l" if rel else "L"  # M 后隐式点 → lineto
            prev_cubic = prev_quad = None
            close = False
        elif c == "L":
            p = read_point(rel)
            pts.append(p)
            cx, cy = p
            prev_cubic = prev_quad = None
        elif c == "H":
            x = read_num()
            cx = cx + x if rel else x
            pts.append((cx, cy))
            prev_cubic = prev_quad = None
        elif c == "V":
            y = read_num()
            cy = cy + y if rel else y
            pts.append((cx, cy))
            prev_cubic = prev_quad = None
        elif c == "C":
            p1 = read_point(rel)
            p2 = read_point(rel)
            p3 = read_point(rel)
            for p in _cubic_points((cx, cy), p1, p2, p3, _CURVE_SAMPLES)[1:]:
                pts.append(p)
            cx, cy = p3
            prev_cubic = p2
            prev_quad = None
        elif c == "S":
            p1 = (2 * cx - prev_cubic[0], 2 * cy - prev_cubic[1]) if prev_cubic else (cx, cy)
            p2 = read_point(rel)
            p3 = read_point(rel)
            for p in _cubic_points((cx, cy), p1, p2, p3, _CURVE_SAMPLES)[1:]:
                pts.append(p)
            cx, cy = p3
            prev_cubic = p2
            prev_quad = None
        elif c == "Q":
            p1 = read_point(rel)
            p2 = read_point(rel)
            for p in _quad_points((cx, cy), p1, p2, _CURVE_SAMPLES)[1:]:
                pts.append(p)
            cx, cy = p2
            prev_quad = p1
            prev_cubic = None
        elif c == "T":
            p1 = (2 * cx - prev_quad[0], 2 * cy - prev_quad[1]) if prev_quad else (cx, cy)
            p2 = read_point(rel)
            for p in _quad_points((cx, cy), p1, p2, _CURVE_SAMPLES)[1:]:
                pts.append(p)
            cx, cy = p2
            prev_quad = p1
            prev_cubic = None
        elif c == "A":
            rx = read_num()
            ry = read_num()
            rot = read_num()
            large = read_num()
            sweep = read_num()
            p2 = read_point(rel)
            for p in _arc_points((cx, cy), (rx, ry), rot, large, sweep, p2, _ARC_SAMPLES)[1:]:
                pts.append(p)
            cx, cy = p2
            prev_cubic = prev_quad = None
        elif c == "Z":
            close = True
            subpaths.append(_SubPath(pts, close))
            close = False
            pts = []
            cx, cy = start
            cmd = None
            prev_cubic = prev_quad = None
        else:
            raise ValueError(f"SVG path 非法命令 {cmd!r}: {d!r}")
    if pts:
        subpaths.append(_SubPath(pts, close))
    if not subpaths:
        raise ValueError(f"SVG path 无有效子路径: {d!r}")
    return subpaths


def _parse_points_list(s: str) -> list[tuple[float, float]]:
    nums = [float(v) for v in re.findall(_NUM_RE, s)]
    try:
        return list(zip(nums[0::2], nums[1::2], strict=True))
    except ValueError as exc:
        raise ValueError(f"points 坐标数必须为偶数（got {len(nums)}）: {s!r}") from exc


def _geom_points(tag: str, attrs: dict) -> tuple[list[tuple[float, float]], bool]:
    """基本几何元素 → (顶点列表, 是否闭合)。"""
    if tag == "line":
        x1 = float(attrs.get("x1", 0))
        y1 = float(attrs.get("y1", 0))
        x2 = float(attrs.get("x2", 0))
        y2 = float(attrs.get("y2", 0))
        return [(x1, y1), (x2, y2)], False
    if tag == "polyline":
        return _parse_points_list(attrs.get("points", "")), False
    if tag == "polygon":
        return _parse_points_list(attrs.get("points", "")), True
    if tag == "rect":
        # rx/ry 圆角忽略（MVP 直角矩形）
        x = float(attrs.get("x", 0))
        y = float(attrs.get("y", 0))
        w = float(attrs.get("width", 0))
        h = float(attrs.get("height", 0))
        return [(x, y), (x + w, y), (x + w, y + h), (x, y + h)], True
    if tag == "circle":
        cx = float(attrs.get("cx", 0))
        cy = float(attrs.get("cy", 0))
        r = float(attrs.get("r", 0))
        return [
            (
                cx + r * math.cos(2 * math.pi * i / _CIRCLE_SAMPLES),
                cy + r * math.sin(2 * math.pi * i / _CIRCLE_SAMPLES),
            )
            for i in range(_CIRCLE_SAMPLES)
        ], True
    if tag == "ellipse":
        cx = float(attrs.get("cx", 0))
        cy = float(attrs.get("cy", 0))
        rx = float(attrs.get("rx", 0))
        ry = float(attrs.get("ry", 0))
        return [
            (
                cx + rx * math.cos(2 * math.pi * i / _CIRCLE_SAMPLES),
                cy + ry * math.sin(2 * math.pi * i / _CIRCLE_SAMPLES),
            )
            for i in range(_CIRCLE_SAMPLES)
        ], True
    raise ValueError(f"不支持的 SVG 元素: {tag!r}")


def load_icon_svg(data_icon: str) -> str:
    """从 vendored assets 读图标 SVG 源码。"ph:name" → assets/phosphor/<name>-fill.svg。

    ph 资产按 <name>-fill.svg 命名（fill 权重），用户写 ph:check / ph:check-fill
    都解析到 check-fill.svg；已带 -fill 后缀不重复加。lu 资产是纯名，原样解析。
    """
    if ":" not in data_icon:
        raise ValueError(f"图标必须带集前缀 ph:/lu:，如 'ph:check': {data_icon!r}")
    set_, name = data_icon.split(":", 1)
    if set_ not in SET_DIRS:
        raise ValueError(f"未知图标集 {set_!r}（可选: {', '.join(SET_DIRS)}）")
    if not re.fullmatch(r"[a-z0-9][a-z0-9-]*", name):
        raise ValueError(f"非法图标名 {name!r}（限小写字母/数字/连字符）")
    if set_ == "ph" and not name.endswith("-fill"):
        name = f"{name}-fill"
    path = ASSETS_DIR / SET_DIRS[set_] / f"{name}.svg"
    if not path.exists():
        raise ValueError(
            f"图标 {data_icon} 不在资产库（{path}）——先跑 uv run python "
            "scripts/fetch_icons.py 抓取全量图标"
        )
    return path.read_text(encoding="utf-8")


def _svg_to_subpaths(svg_text: str) -> tuple[list[_SubPath], float]:
    """资产 SVG 源码 → (子路径列表, stroke-width)。子路径已压平。"""
    # vendored 图标无 transform/<g> 覆盖，忽略；解析走统一 parse_svg 守卫，
    # 拒绝 DOCTYPE/ENTITY（billion-laughs）并包装畸形 XML 为 InvalidArgumentError
    from offipy.assets._xml import parse_svg

    root = parse_svg(svg_text)
    stroke_width = float(root.get("stroke-width", "2"))
    subpaths: list[_SubPath] = []
    for el in root.iter():
        tag = el.tag.split("}")[-1].lower()
        filled = el.get("fill") == "currentColor"
        if tag == "path":
            d = el.get("d")
            if d:
                for sp in _parse_path(d):
                    sp.filled = filled
                    subpaths.append(sp)
        elif tag in ("line", "polyline", "polygon", "rect", "circle", "ellipse"):
            pts, close = _geom_points(tag, el.attrib)
            if pts:
                subpaths.append(_SubPath(pts, close, filled))
    return subpaths, stroke_width


# ---------- HTML 声明解析 ----------


class _IconHTMLParser(HTMLParser):
    """提取每个 <section data-pptx-slide> 内的 <svg data-icon>。"""

    def __init__(self) -> None:
        super().__init__()
        self.slide_index = 0
        self.cur: list[IconDecl] = []
        self.all: list[IconDecl] = []

    def handle_starttag(self, tag: str, attrs: list[tuple[str, str | None]]) -> None:
        d = {k: (v or "") for k, v in attrs}
        if tag == "section" and "data-pptx-slide" in d:
            if self.cur:
                self.all.extend(self.cur)
            self.slide_index += 1
            self.cur = []
            return
        if tag == "svg" and "data-icon" in d:
            # HTMLParser 会把属性名小写化：viewBox → viewbox，两种都兜底。
            # 缺省 viewBox 按集区分：ph 是 256 坐标系，lu 是 24 坐标系，给错会静默缩小。
            set_ = d["data-icon"].split(":", 1)[0] if ":" in d["data-icon"] else "ph"
            default_vb = "0 0 24 24" if set_ == "lu" else "0 0 256 256"
            vb = d.get("viewbox") or d.get("viewBox") or default_vb
            parts = vb.split()
            if len(parts) != 4:
                raise ValueError(
                    f"第 {self.slide_index} 页图标 viewBox 非法: {vb!r}（须为 'x y w h' 四值）"
                )
            self.cur.append(
                IconDecl(
                    slide_index=self.slide_index,
                    data_icon=d["data-icon"],
                    view_box=(
                        float(parts[0]),
                        float(parts[1]),
                        float(parts[2]),
                        float(parts[3]),
                    ),
                )
            )


def parse_icon_declarations(html_text: str) -> list[IconDecl]:
    """解析 HTML → 图标声明列表（slide_index 1-based）。图标集非法即 ValueError。"""
    p = _IconHTMLParser()
    p.feed(html_text)
    p.close()
    if p.cur:  # flush 最后一页
        p.all.extend(p.cur)
    for decl in p.all:
        if decl.slide_index <= 0:
            raise ValueError(
                f"第 {decl.slide_index} 页图标出现在 slide 之外——<svg data-icon> 必须在 "
                "data-pptx-slide 的 <section> 内"
            )
        if ":" not in decl.data_icon or decl.data_icon.split(":", 1)[0] not in SET_DIRS:
            raise ValueError(
                f"第 {decl.slide_index} 页图标集非法: {decl.data_icon!r}（须带 ph:/lu: 前缀）"
            )
    return p.all


# ---------- measurements 定位 ----------


def _measurements_path(pptx_path: str) -> str:
    """与 charts._measurements_path 一致：<stem>_audit/_cache/measurements.json。"""
    from pathlib import Path

    p = Path(pptx_path)
    return str(p.with_name(f"{p.stem}_audit") / "_cache" / "measurements.json")


def load_icon_boxes(measurements_path: str) -> dict[int, list[dict]]:
    """读 measurements.json → {slide_index: [svg record, ...]}（保序，仅 kind='svg'）。"""
    import json as _json

    with Path(measurements_path).open(encoding="utf-8") as f:
        data = _json.load(f)
    boxes: dict[int, list[dict]] = {}
    for i, slide in enumerate(data.get("slides", []), start=1):
        svgs = [
            {"rect": rec["rect"], "color": rec.get("color"), "outerHTML": rec.get("outerHTML", "")}
            for rec in slide.get("records", [])
            if rec.get("kind") == "svg" and "rect" in rec
        ]
        if svgs:
            boxes[i] = svgs
    return boxes


def _match_svg(decl: IconDecl, svgs: list[dict]) -> dict | None:
    """在 svg records 里找 outerHTML 含对应 data-icon 的那条。"""
    pat = re.compile(r'data-icon=["\']' + re.escape(decl.data_icon) + r'["\']')
    for svg in svgs:
        if pat.search(svg.get("outerHTML", "")):
            return svg
    return None


# ---------- 注入 ----------


def _parse_color(color: str | None):
    """measure 的 color（'rgb(r, g, b)' / '#rrggbb'）→ RGBColor；无法解析返回 None。"""
    if not color:
        return None
    m = re.search(r"rgba?\(\s*(\d+)\s*,\s*(\d+)\s*,\s*(\d+)", color)
    if m:
        from pptx.dml.color import RGBColor

        return RGBColor(int(m.group(1)), int(m.group(2)), int(m.group(3)))
    m = re.fullmatch(r"#([0-9a-fA-F]{6})", color.strip())
    if m:
        from pptx.dml.color import RGBColor

        return RGBColor.from_string(m.group(1))
    return None


def _set_round_stroke(shape) -> None:
    """Lucide 源 SVG 用 stroke-linecap/linejoin=round；freeform 折线默认 flat，显式设回。

    只对 stroke 模式生效（line 可见）。fill 模式的线已被 noFill 掉，无需设置。
    """
    from pptx.oxml.ns import qn

    ln = shape.line._get_or_add_ln()
    ln.set("cap", "rnd")
    if ln.find(qn("a:round")) is None:
        ln.append(ln.makeelement(qn("a:round"), {}))


def _theme_accent_fallback(html_text: str):
    """HTML <style data-theme="..."> 声明的主题 → 该主题 --accent（RGBColor）。

    容器没测到 color 时兜底用主题强调色；未知/缺失主题返回 None（再落缺省主蓝）。
    """
    m = re.search(r'<style\s+data-theme="([a-z0-9-]+)"', html_text)
    if not m:
        return None
    from .design import THEMES

    theme = THEMES.get(m.group(1))
    if theme is None:
        return None
    return _parse_color(theme.base_vars.get("--accent"))


def _style_shape(
    shape,
    mode: str,
    color: str | None,
    line_width_emu: int,
    filled: bool = False,
    fallback=None,
) -> None:
    rgb = _parse_color(color)
    if rgb is None:
        from pptx.dml.color import RGBColor

        rgb = fallback if fallback is not None else RGBColor(0x22, 0x51, 0xFF)  # 缺省主蓝
    fill_mode = (mode == "fill") or filled
    if fill_mode:
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb
        shape.line.fill.background()
    else:
        # 纯 stroke 子路径：显式 noFill（background()），否则 PowerPoint 会用
        # 主题 fillRef（accent1）按隐式闭合把开放路径涂成实心蓝块。
        shape.fill.background()
        shape.line.color.rgb = rgb
        shape.line.width = line_width_emu
        _set_round_stroke(shape)


def _build_icon_shapes(
    slide,
    mode: str,
    subpaths: list[_SubPath],
    stroke_width: float,
    view_box: tuple[float, float, float, float],
    rect: dict,
    color: str | None,
    fallback=None,
) -> list:
    """在 rect（px）位置画图标：每个子路径一个 freeform shape。返回创建的 shape 列表。"""
    vbx, vby, vbw, _ = view_box
    scale = rect["w"] / vbw if vbw else 1.0
    ox = rect["x"] - vbx * scale
    oy = rect["y"] - vby * scale
    line_width_emu = int(stroke_width * scale * PX_TO_EMU)

    def to_emu(p: tuple[float, float]) -> tuple[int, int]:
        return int((p[0] * scale + ox) * PX_TO_EMU), int((p[1] * scale + oy) * PX_TO_EMU)

    created: list = []
    for sp in subpaths:
        if len(sp.points) < 2:
            continue
        first = to_emu(sp.points[0])
        fb = slide.shapes.build_freeform(first[0], first[1])
        rest = [to_emu(p) for p in sp.points[1:]]
        if rest:
            fb.add_line_segments(rest, close=sp.close)
        shape = fb.convert_to_shape()
        _style_shape(shape, mode, color, line_width_emu, sp.filled, fallback)
        created.append(shape)
    return created


def render_icon_payload(
    slide,
    svg_text: str,
    *,
    mode: str,
    view_box: tuple[float, float, float, float],
    rect: dict,
    color: str | None = None,
    fallback=None,
) -> list:
    """把 asset 图标的 SVG payload 渲染成 freeform 形状，返回创建的 shape 列表。

    mode（fill/stroke）由 provider 决定（ph → fill，lu → stroke），是 freeform
    引擎的内部元数据，不进公共 Asset 模型。rect 是 px dict {x, y, w, h}。
    """
    subpaths, stroke_width = _svg_to_subpaths(svg_text)
    return _build_icon_shapes(slide, mode, subpaths, stroke_width, view_box, rect, color, fallback)


def _remove_placeholder(slide, rect: dict) -> None:
    """删除中心落在图标 bbox 内的 PNG picture 占位（convert 的 SVG 光栅图）。"""
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    cxm = (rect["x"] + rect["w"] / 2) * PX_TO_EMU
    cym = (rect["y"] + rect["h"] / 2) * PX_TO_EMU
    for shape in list(slide.shapes):
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE and (
            shape.left <= cxm <= shape.left + shape.width
            and shape.top <= cym <= shape.top + shape.height
        ):
            slide.shapes._spTree.remove(shape._element)


def inject_icons(
    pptx_path: str, matched: dict[int, list[tuple[IconDecl, dict]]], fallback=None
) -> None:
    """把图标占位替换成 freeform 矢量图标。matched: {slide_index: [(decl, svg_record), ...]}。

    fallback: 容器没测到 color 时的兜底色（主题 --accent 的 RGBColor），None 落缺省主蓝。
    """
    from pptx import Presentation

    prs = Presentation(pptx_path)
    for slide_index, pairs in matched.items():
        slide = prs.slides[slide_index - 1]
        for decl, svg in pairs:
            _remove_placeholder(slide, svg["rect"])
            mode = SET_MODE[decl.data_icon.split(":", 1)[0]]
            subpaths, stroke_width = _svg_to_subpaths(load_icon_svg(decl.data_icon))
            _build_icon_shapes(
                slide,
                mode,
                subpaths,
                stroke_width,
                decl.view_box,
                svg["rect"],
                svg["color"],
                fallback,
            )
    prs.save(pptx_path)


def postprocess_icons(html_path: str, pptx_path: str) -> None:
    """转换后调用：HTML 含 data-icon → 读 measurements → 注入矢量图标。

    无图标声明 → 原样返回。measurements.json 缺失 → RuntimeError。图标/数据非法
    → ValueError（从 parse 上抛）。
    """

    with Path(html_path).open(encoding="utf-8") as f:
        html_text = f.read()
    if "data-icon" not in html_text:
        return
    decls = parse_icon_declarations(html_text)
    if not decls:
        return
    meas_path = _measurements_path(pptx_path)
    if not Path(meas_path).exists():
        raise RuntimeError(
            f"找不到 convert 审计产物 {meas_path}——图标注入需要 measurements.json，"
            "请勿用 --no-visual-audit"
        )
    boxes = load_icon_boxes(meas_path)
    matched: dict[int, list[tuple[IconDecl, dict]]] = {}
    for decl in decls:
        svgs = boxes.get(decl.slide_index, [])
        svg = _match_svg(decl, svgs)
        if svg is None:
            raise RuntimeError(
                f"第 {decl.slide_index} 页图标 {decl.data_icon} 没测到 <svg> 占位——"
                "请确认图标容器是空 <svg data-icon=...>（不是 <i>/伪元素），且尺寸非 0"
            )
        svgs.remove(svg)  # 防同页重复匹配同一条 svg record
        matched.setdefault(decl.slide_index, []).append((decl, svg))
    fallback = _theme_accent_fallback(html_text)
    inject_icons(pptx_path, matched, fallback=fallback)
