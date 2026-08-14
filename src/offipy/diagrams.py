"""Mermaid 图 → PPTX 原生可编辑形状（结构图渲染，非原生图表）。

输入 Mermaid flowchart 源码 → vendored mermaid_extract.py 提取 IR（拓扑）
→ offipy 分层布局（layout_diagram）→ python-pptx 渲染成可编辑形状
（autoshape + connector，render_to_slide）。可独立生成整页 PPTX
（mermaid_to_pptx），也可作为 deck 后处理注入（postprocess_mermaid：
HTML <pre class="mermaid"> 块 → 替换为可编辑形状）。

vendored 提取器用 importlib 加载，保持上游文件原样；仅支持 flowchart
（TD/TB/LR/RL/BT）；sequence/state/er 及不支持语法 → ValueError。
"""

from __future__ import annotations

import importlib.util
import json
import os
import re
import sys
from collections import deque
from dataclasses import dataclass
from html.parser import HTMLParser
from pathlib import Path

_EXTRACT_REL = Path("_vendor/diagram-design/skills/diagram-design/scripts/mermaid_extract.py")
SUPPORTED_DIRECTIONS = {"TD", "TB", "LR", "RL", "BT"}
# vendored extractor 契约：graph/flowchart 必须带显式方向（否则 _kind_and_direction
# 直接 _fail("not a Mermaid file")）。裸 graph 是常见用户输入，parse_mermaid 用它做
# 预检，把误导性的 "not a Mermaid file" 换成清晰错误。
_FLOW_HEADER_WITHOUT_DIRECTION = re.compile(r"^\s*(graph|flowchart)\s*$", re.I)

_extractor = None


def _load_extractor():
    """惰性加载 vendored mermaid_extract（importlib，保持上游原样）。"""
    global _extractor
    if _extractor is None:
        script = Path(__file__).resolve().parent / _EXTRACT_REL
        if not script.exists():
            raise RuntimeError(f"vendored mermaid_extract 缺失: {script}")
        name = "offipy_vendored_mermaid_extract"
        spec = importlib.util.spec_from_file_location(name, script)
        mod = importlib.util.module_from_spec(spec)
        sys.modules[name] = mod
        assert spec.loader is not None
        spec.loader.exec_module(mod)
        _extractor = mod
    return _extractor


def parse_mermaid(text: str):
    """Mermaid 文本 → vendored Diagram IR。仅接受 flowchart。

    mermaid_extract 对坏输入/不支持语法用 SystemExit(2) 退出，这里捕获并
    归一化成 ValueError（调用方是 deck 后处理/独立 API，都按用户输入错误处理）。
    """
    # 预检裸 graph/flowchart（无方向）：extractor 会报误导性的 "not a Mermaid file"，
    # 这里先拦成清晰错误（Mermaid 标准默认 TD，但 vendor 契约要求显式方向）。
    first = text.splitlines()[0] if text.splitlines() else ""
    if _FLOW_HEADER_WITHOUT_DIRECTION.match(first):
        raise ValueError(
            "flowchart/graph 需要显式方向（TD/TB/LR/RL/BT），收到缺方向的 graph/flowchart"
        )
    ex = _load_extractor()
    block = ex.SourceBlock(0, text, 1)
    try:
        diagram = ex.parse_block(block)
    except SystemExit as e:
        code = e.code if isinstance(e.code, str) else str(e)
        raise ValueError(f"无法解析 Mermaid 源码: {code}") from None
    if diagram.kind != "flowchart":
        raise ValueError(
            f"仅支持 flowchart/graph，收到 {diagram.kind}（sequence/state/er 后续迭代）"
        )
    # vendor 对部分坏输入不抛 SystemExit，而是静默产出空图（无节点无边）。
    # 空 flowchart 无渲染意义，统一按用户输入错误处理（plan 期望：坏语法 → ValueError）。
    if not diagram.nodes and not diagram.edges:
        raise ValueError("无法解析 Mermaid 源码: 未提取到任何节点或边")
    return diagram


_NODE_W = 2.0
_NODE_H = 0.6
_GAP_X = 0.9
_GAP_Y = 1.1
_MARGIN = 0.5


@dataclass
class PlacedNode:
    id: str
    label: str
    shape: str
    x: float
    y: float
    w: float
    h: float
    is_container: bool = False
    parent: str | None = None
    fill: str = ""
    stroke: str = ""
    font_color: str = ""
    font_pt: float | None = None


@dataclass
class PlacedEdge:
    source: str
    target: str
    label: str = ""
    style: str = "solid"
    arrowhead: str = "arrow"
    undirected: bool = False
    ax1: float = 0.0
    ay1: float = 0.0
    ax2: float = 0.0
    ay2: float = 0.0
    stroke: str = ""
    waypoints: list[tuple[float, float]] | None = None


@dataclass
class DiagramLayout:
    nodes: list[PlacedNode]
    edges: list[PlacedEdge]
    canvas_w: float
    canvas_h: float


def _kahn_layers(ids: list[str], edges) -> tuple[dict[str, int], dict[int, list[str]]]:
    """Kahn 拓扑分层。返回 (node→layer, layer→node_id 列表)。环残留统一追加尾层。

    环残留（Kahn 消不掉的节点：成环节点及其下游）不参与逐层传播，全部落在
    DAG 最大层 + 1 的同一层——纯环归一层，分层结果与 ids 迭代序无关。
    """
    adj: dict[str, list[str]] = {i: [] for i in ids}
    indeg = {i: 0 for i in ids}
    for e in edges:
        if e.source in adj and e.target in adj:
            adj[e.source].append(e.target)
            indeg[e.target] += 1
    q = deque(i for i in ids if indeg[i] == 0)
    order: list[str] = []
    while q:
        i = q.popleft()
        order.append(i)
        for t in adj[i]:
            indeg[t] -= 1
            if indeg[t] == 0:
                q.append(t)
    in_order = set(order)
    residue = [i for i in ids if i not in in_order]
    layer = {i: 0 for i in ids}
    # 只在 DAG 内部传播层深（层深 = 最长路径，与遍历序无关）；环残留不参与，
    # 统一落尾层。
    for i in order:
        for t in adj[i]:
            if t in in_order:
                layer[t] = max(layer[t], layer[i] + 1)
    tail = (max((layer[i] for i in order), default=-1) + 1) if order else 0
    for i in residue:
        layer[i] = tail
    layers: dict[int, list[str]] = {}
    for i in ids:
        layers.setdefault(layer[i], []).append(i)
    return layer, layers


def layout_diagram(
    diagram, *, direction: str | None = None, max_w: float = 12.0, max_h: float = 6.75
) -> DiagramLayout:
    """把 flowchart IR 布局成坐标（inches）。方向 TD/TB/LR/RL/BT；整体 fit 到 max_w/max_h。

    - 叶子（非 container）按 Kahn 分层；layer 间距、层内列距。
    - container 节点不占独立坐标，作为背景框覆盖其子孙叶子 bbox（Task 4 画）。
    - 反向方向（BT/RL）反转层序/列序。
    """
    dr = (direction or diagram.direction or "TD").upper()
    if dr not in SUPPORTED_DIRECTIONS:
        raise ValueError(f"不支持方向 {dr}（可选 TD/TB/LR/RL/BT）")
    leaves = [n for n in diagram.nodes if not (n.container or n.shape == "container")]
    leaves_by_id = {n.id: n for n in leaves}

    layer_of, layers = _kahn_layers(list(leaves_by_id), diagram.edges)
    col_of: dict[str, int] = {}
    for ids in layers.values():
        for col, i in enumerate(ids):
            col_of[i] = col
    max_row = max(layers) if layers else 0
    max_col = max((len(v) for v in layers.values()), default=0)

    vertical = dr in ("TD", "TB", "BT")
    if vertical:
        raw_w = _MARGIN + max_col * (_NODE_W + _GAP_X) + _MARGIN
        raw_h = _MARGIN + (max_row + 1) * (_NODE_H + _GAP_Y) + _MARGIN
    else:
        raw_w = _MARGIN + (max_row + 1) * (_NODE_W + _GAP_X) + _MARGIN
        raw_h = _MARGIN + max_col * (_NODE_H + _GAP_Y) + _MARGIN
    scale = min(1.0, max_w / raw_w if raw_w > 0 else 1.0, max_h / raw_h if raw_h > 0 else 1.0)
    W, H = _NODE_W * scale, _NODE_H * scale
    GX, GY = _GAP_X * scale, _GAP_Y * scale
    M = _MARGIN * scale

    placed: list[PlacedNode] = []
    for i, n in leaves_by_id.items():
        lyr = layer_of[i]
        col = col_of[i]
        if vertical:
            row_idx = (max_row - lyr) if dr == "BT" else lyr
            px = M + col * (W + GX)
            py = M + row_idx * (H + GY)
        else:
            # RL 沿主方向（x/layer）反转：layer0 在最右；垂直分布（y/col）不镜像。
            # 对应 _layout_edges 的 RL 分支（a1=源左边缘、a2=目标右边缘）依赖此反转。
            px = M + (max_row - lyr) * (W + GX) if dr == "RL" else M + lyr * (W + GX)
            py = M + col * (H + GY)
        placed.append(PlacedNode(i, n.label, n.shape, px, py, W, H, False, n.parent))

    placed_nodes = placed + _layout_containers(diagram, placed, scale)
    # 容器框入 by_id → 容器端点边被路由（Task 2 的跳过分支转为死路径）
    placed_edges = _layout_edges(diagram, placed_nodes, vertical, dr)
    canvas_w = M + (max_col * (W + GX)) if vertical else M + ((max_row + 1) * (W + GX))
    canvas_h = M + ((max_row + 1) * (H + GY)) if vertical else M + (max_col * (H + GY))
    return DiagramLayout(placed_nodes, placed_edges, canvas_w, canvas_h)


def _layout_edges(
    diagram, placed: list[PlacedNode], vertical: bool, direction: str
) -> list[PlacedEdge]:
    by_id = {n.id: n for n in placed}
    out: list[PlacedEdge] = []
    for e in diagram.edges:
        # 可达分支：容器端点边在 Task 2 布局阶段被跳过（容器不占坐标，无框可锚）。
        # Task 4 把 placed_nodes 传给本函数后容器框入 by_id，此分支转为真正死路径。
        if e.source not in by_id or e.target not in by_id:
            continue
        s, t = by_id[e.source], by_id[e.target]
        if vertical:
            if direction == "BT":
                a1 = (s.x + s.w / 2, s.y)
                a2 = (t.x + t.w / 2, t.y + t.h)
            else:
                a1 = (s.x + s.w / 2, s.y + s.h)
                a2 = (t.x + t.w / 2, t.y)
        else:
            if direction == "RL":
                a1 = (s.x, s.y + s.h / 2)
                a2 = (t.x + t.w, t.y + t.h / 2)
            else:
                a1 = (s.x + s.w, s.y + s.h / 2)
                a2 = (t.x, t.y + t.h / 2)
        out.append(
            PlacedEdge(
                e.source,
                e.target,
                e.label,
                e.style,
                e.arrowhead,
                e.undirected,
                a1[0],
                a1[1],
                a2[0],
                a2[1],
            )
        )
    return out


_CONTAINER_PAD_X = 0.3
_CONTAINER_PAD_Y = 0.3
_CONTAINER_HDR = 0.4


def _descendants(children: dict[str, list[str]], root: str) -> set[str]:
    seen: set[str] = set()
    stack = list(children.get(root, []))
    while stack:
        i = stack.pop()
        if i in seen:
            continue
        seen.add(i)
        stack.extend(children.get(i, []))
    return seen


def _layout_containers(diagram, placed: list[PlacedNode], scale: float) -> list[PlacedNode]:
    """container 节点 → 背景框 PlacedNode（覆盖其子孙叶子 bbox + padding + 标题区）。

    嵌套容器：从最深层往上算（子先于父），外层框用已算好的子框扩展。
    """
    containers = [n for n in diagram.nodes if n.container or n.shape == "container"]
    if not containers:
        return []
    # 全节点父子树（容器→子容器/叶子），供 _descendants 跨层收集嵌套后代
    children: dict[str, list[str]] = {}
    for n in diagram.nodes:
        if n.parent:
            children.setdefault(n.parent, []).append(n.id)
    by_id = {n.id: n for n in placed}
    # 子树高度决定计算顺序：子容器先算（parent 依赖），外框再引用已算好的子框
    order = sorted(containers, key=lambda c: -c.depth)
    boxes: dict[str, PlacedNode] = {}
    for c in order:
        members = _descendants(children, c.id)
        rects = [by_id[i] for i in members if i in by_id]
        rects += [boxes[i] for i in members if i in boxes]
        if rects:
            # 顶/左缘不小于画布原点：顶层容器（首节点在 y=M=0.5）的 y0 会算成
            # 0.5-(0.3+0.4)=-0.2in，负坐标把容器标题渲到画布/幻灯片顶外被裁掉。
            x0 = max(0.0, min(r.x for r in rects) - _CONTAINER_PAD_X * scale)
            y0 = max(0.0, min(r.y for r in rects) - (_CONTAINER_PAD_Y + _CONTAINER_HDR) * scale)
            x1 = max(r.x + r.w for r in rects) + _CONTAINER_PAD_X * scale
            y1 = max(r.y + r.h for r in rects) + _CONTAINER_PAD_Y * scale
        else:
            x0 = y0 = 0.0
            x1 = y1 = 0.5 * scale
        boxes[c.id] = PlacedNode(
            c.id, c.label, "container", x0, y0, x1 - x0, y1 - y0, True, c.parent
        )
    return list(boxes.values())


_NODE_FONT = "Microsoft YaHei"


def _set_ea_font(run, name: str) -> None:
    """同时设置 latin/ea/cs typeface，保证中文在 PowerPoint 正常渲染。"""
    from pptx.oxml.ns import qn

    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", name)


def _set_tail_end(conn, kind: str) -> None:
    from pptx.oxml.ns import qn

    ln = conn.line._get_or_add_ln()
    tail = ln.find(qn("a:tailEnd"))
    if tail is None:
        tail = ln.makeelement(qn("a:tailEnd"), {})
        ln.append(tail)
    tail.set("type", kind)
    tail.set("w", "med")
    tail.set("len", "med")


def _add_text(shape, text: str, *, size_pt: float, bold: bool = False, color=None) -> None:
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Emu, Pt

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(91440)  # 0.1"
    tf.margin_top = tf.margin_bottom = Emu(45720)  # 0.05"
    tf.text = text
    for para in tf.paragraphs:
        para.alignment = PP_ALIGN.CENTER
        for run in para.runs:
            run.font.size = Pt(size_pt)
            run.font.bold = bold
            run.font.color.rgb = color if color is not None else RGBColor(0x2D, 0x31, 0x42)
            _set_ea_font(run, _NODE_FONT)


def _hex_rgb(s: str):
    """三态颜色解析：#RRGGBB/#RGB → RGBColor；空 / "none" / 非法 hex → None。

    draw.io 的 fillColor=none / strokeColor=none 表示透明；空串表示「未指定」。
    返回 None 时由调用方落到各自默认（白底 / 深线 / 深字）。
    """
    from pptx.dml.color import RGBColor

    text = (s or "").strip()
    if not text or text.lower() == "none":
        return None
    text = text.lstrip("#")
    if len(text) in (3, 4):
        text = "".join(c * 2 for c in text[:3])
    if len(text) < 6:
        return None
    try:
        return RGBColor(int(text[:2], 16), int(text[2:4], 16), int(text[4:6], 16))
    except ValueError:
        return None


def render_to_slide(
    slide, layout: DiagramLayout, *, offset_x: int = 0, offset_y: int = 0, node_font_pt: float = 14
) -> None:
    """把布局画到 slide。offset 为 EMU 偏移（deck 注入时 = bbox 左上角）。

    画序：边（连线在节点下层）→ 容器背景（Task 4）→ 节点。connector/autoshape
    全部是原生可编辑形状。所有 python-pptx import 惰性化，避免顶层 import 拖慢
    纯 deck 路径。
    """
    from pptx.dml.color import RGBColor
    from pptx.enum.dml import MSO_LINE_DASH_STYLE
    from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Emu, Inches, Pt

    # map 引用 pptx 枚举，须在惰性 import 之后构造（模块级会 import 即 NameError）。
    _SHAPE_MAP = {
        "rect": MSO_SHAPE.ROUNDED_RECTANGLE,
        "round": MSO_SHAPE.ROUNDED_RECTANGLE,
        "rectangle": MSO_SHAPE.RECTANGLE,  # 新增：drawio 直角矩形
        "ellipse": MSO_SHAPE.OVAL,  # 新增：drawio 椭圆
        "triangle": MSO_SHAPE.ISOSCELES_TRIANGLE,  # 新增：drawio 三角
        "stadium": MSO_SHAPE.ROUNDED_RECTANGLE,
        "subroutine": MSO_SHAPE.ROUNDED_RECTANGLE,
        "asymmetric": MSO_SHAPE.ROUNDED_RECTANGLE,
        "cylinder": MSO_SHAPE.ROUNDED_RECTANGLE,
        "circle": MSO_SHAPE.OVAL,
        "rhombus": MSO_SHAPE.DIAMOND,
        "hexagon": MSO_SHAPE.HEXAGON,
        "parallelogram": MSO_SHAPE.PARALLELOGRAM,
        "trapezoid": MSO_SHAPE.TRAPEZOID,
    }
    _ARROW_MAP = {"circle": "oval", "cross": "triangle", "arrow": "triangle"}
    _DASH_MAP = {
        "solid": MSO_LINE_DASH_STYLE.SOLID,
        "dashed": MSO_LINE_DASH_STYLE.DASH,
        "thick": MSO_LINE_DASH_STYLE.SOLID,
    }
    edge_color = RGBColor(0x66, 0x70, 0x85)

    for e in layout.edges:
        x1, y1, x2, y2 = e.ax1, e.ay1, e.ax2, e.ay2
        if e.waypoints:
            # waypoint 边：open polyline（build_freeform），直线段逼近正交/曲线路径
            pts = [(x1, y1)] + list(e.waypoints) + [(x2, y2)]
            fb = slide.shapes.build_freeform(
                start_x=Emu(offset_x) + Inches(x1),
                start_y=Emu(offset_y) + Inches(y1),
                scale=1.0,
            )
            fb.add_line_segments(
                vertices=[
                    (Emu(offset_x) + Inches(wx), Emu(offset_y) + Inches(wy))
                    for wx, wy in pts[1:]
                ],
                close=False,
            )
            conn = fb.convert_to_shape()
            lx, ly = (
                sum(p[0] for p in pts) / len(pts),
                sum(p[1] for p in pts) / len(pts),
            )
        else:
            conn = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT,
                Emu(offset_x) + Inches(x1),
                Emu(offset_y) + Inches(y1),
                Emu(offset_x) + Inches(x2),
                Emu(offset_y) + Inches(y2),
            )
            lx, ly = (x1 + x2) / 2, (y1 + y2) / 2
        if e.stroke == "none":
            conn.line.fill.background()  # strokeColor=none → 无线框
        else:
            conn.line.color.rgb = _hex_rgb(e.stroke) or edge_color
        conn.line.width = Pt(1.6 if e.style == "thick" else 1.2)
        conn.line.dash_style = _DASH_MAP.get(e.style, MSO_LINE_DASH_STYLE.SOLID)
        if not e.undirected:
            _set_tail_end(conn, _ARROW_MAP.get(e.arrowhead, "triangle"))
        if e.label:
            lab = slide.shapes.add_textbox(
                Emu(offset_x) + Inches(lx - 0.8),
                Emu(offset_y) + Inches(ly - 0.15),
                Inches(1.6),
                Inches(0.3),
            )
            _add_text(lab, e.label, size_pt=10)

    containers = [n for n in layout.nodes if n.is_container]
    for c in containers:
        box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Emu(offset_x) + Inches(c.x),
            Emu(offset_y) + Inches(c.y),
            Inches(c.w),
            Inches(c.h),
        )
        if c.fill == "none":
            box.fill.background()  # draw.io group 默认无 fillColor
        else:
            fill_rgb = _hex_rgb(c.fill)
            if fill_rgb is not None:
                box.fill.solid()
                box.fill.fore_color.rgb = fill_rgb
            else:
                box.fill.solid()
                box.fill.fore_color.rgb = RGBColor(0xF2, 0xF4, 0xF8)  # 空 → 浅灰兜底
        if c.stroke == "none":
            box.line.fill.background()
        else:
            stroke_rgb = _hex_rgb(c.stroke)
            if stroke_rgb is not None:
                box.line.color.rgb = stroke_rgb
            else:
                box.line.color.rgb = RGBColor(0x9A, 0xA3, 0xB2)
        box.line.width = Pt(1.0)
        tf = box.text_frame
        tf.margin_left = tf.margin_right = Emu(91440)
        tf.text = c.label
        for para in tf.paragraphs:
            para.alignment = PP_ALIGN.CENTER
            for run in para.runs:
                run.font.size = Pt(12)
                run.font.bold = True
                run.font.color.rgb = RGBColor(0x2D, 0x31, 0x42)
                _set_ea_font(run, _NODE_FONT)

    for n in layout.nodes:
        if n.is_container:
            continue
        shape = slide.shapes.add_shape(
            _SHAPE_MAP.get(n.shape, MSO_SHAPE.ROUNDED_RECTANGLE),
            Emu(offset_x) + Inches(n.x),
            Emu(offset_y) + Inches(n.y),
            Inches(n.w),
            Inches(n.h),
        )
        if n.fill == "none":
            shape.fill.background()  # fillColor=none → 透明
        else:
            fill_rgb = _hex_rgb(n.fill)
            if fill_rgb is not None:
                shape.fill.solid()
                shape.fill.fore_color.rgb = fill_rgb
            else:
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)  # 空 → 默认白
        if n.stroke == "none":
            shape.line.fill.background()  # strokeColor=none → 无线框
        else:
            stroke_rgb = _hex_rgb(n.stroke)
            if stroke_rgb is not None:
                shape.line.color.rgb = stroke_rgb
            else:
                shape.line.color.rgb = RGBColor(0x2D, 0x31, 0x42)  # 空 → 默认深线
        shape.line.width = Pt(1.0)
        _add_text(shape, n.label, size_pt=n.font_pt or node_font_pt, color=_hex_rgb(n.font_color))


_MMD_MARKER_RE = re.compile(
    r"\b(graph|flowchart|subgraph|sequenceDiagram|classDiagram|stateDiagram|"
    r"erDiagram|gantt|journey|mindmap|timeline)\b"
)


def _has_mermaid_markers(s: str) -> bool:
    """Mermaid 结构标记（graph/flowchart/--> 等）出现即视为文本而非路径。

    用于 #92 区分「含分隔符的相对路径（data/foo.mmd）」与「合法 Mermaid 源码里的 /」
    （如 graph TD; A[/决策/] 平行四边形节点）。有标记 → 文本，无标记 → 路径意图。
    """
    if _MMD_MARKER_RE.search(s):
        return True
    return any(m in s for m in ("-->", "---", "==>", "%%"))


def _looks_like_path(s: str) -> bool:
    """路径形态启发（#85/#92）：盘符（C:\\ / C:/）、前导 / 或 \\（绝对/UNC）、
    点相对段（./ ../）、无前缀相对路径（data/foo.mmd、foo\\bar.mmd）。

    str 入参「路径 or 文本」无法从类型区分，只能靠存在性试探。试探失败后按形态
    优先「路径意图」——这些形态不会是合法 Mermaid 源码的开头，误伤面为零。
    """
    if len(s) >= 3 and s[1] == ":" and s[0].isalpha() and s[2] in "/\\":
        return True
    if s.startswith(("/", "\\")):
        return True
    if s.startswith(("./", "../")):
        return True
    if "/" in s or "\\" in s:
        return not _has_mermaid_markers(s)
    return False


def _read_source(source) -> str:
    """source 是文件路径 → 读文件；否则当作 Mermaid 文本。"""
    if isinstance(source, os.PathLike):
        p = Path(source)
        # PathLike 明确表示「这是路径」：缺失/非文件直接抛 FileNotFoundError，
        # 不把路径字符串当文本解析（否则是误导性的「无法解析」ValueError）。
        if not p.is_file():
            raise FileNotFoundError(f"源文件不存在: {p}")
        return p.read_text(encoding="utf-8")
    if isinstance(source, str):
        p = Path(source)
        if p.exists() and p.is_file():
            return p.read_text(encoding="utf-8")
        # #85：路径打错时字符串不存在，旧实现静默当文本 → 抛误导性「无法解析 Mermaid
        # 源码」。存在性试探失败后按形态识别「路径意图」，给出与 PathLike 一致的
        # FileNotFoundError（可定位根因）；形态不路径化才当作文本解析。
        if _looks_like_path(source):
            raise FileNotFoundError(
                f"源文件不存在: {source}（字符串按路径处理；若本意是 Mermaid 文本，"
                "请不要让它形如路径：盘符 / / / \\ / ./ 前缀，或含 / 或 \\ 分隔符）"
            )
    return str(source)


def mermaid_to_pptx(source, out_path: str, *, direction: str | None = None) -> str:
    """Mermaid 文本或 .mmd 文件路径 → 可编辑 PPTX（16:9 整页）。返回 out_path。

    python-pptx 惰性 import（deck extra），未安装时给可操作错误。
    """
    text = _read_source(source)
    diagram = parse_mermaid(text)
    lay = layout_diagram(diagram, direction=direction)
    try:
        from pptx import Presentation
        from pptx.util import Emu
    except ImportError as e:  # pragma: no cover
        raise RuntimeError("mermaid_to_pptx 需要 python-pptx，请安装 offipy[deck]") from e
    prs = Presentation()
    prs.slide_width = Emu(12192000)
    prs.slide_height = Emu(6858000)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    render_to_slide(slide, lay)
    prs.save(out_path)
    return str(out_path)


PX_TO_EMU = 6350


class _MermaidHTMLParser(HTMLParser):
    """扫 <pre class="mermaid"> 块 → [{slide, source}]。slide 1-based 对齐 section 顺序。"""

    def __init__(self) -> None:
        super().__init__()
        self.slide_index = 0
        self.decls: list[dict] = []
        self._cur: dict | None = None

    def handle_starttag(self, tag: str, attrs) -> None:
        d = {k: (v or "") for k, v in attrs}
        if tag == "section" and "data-pptx-slide" in d:
            self.slide_index += 1
            self._cur = None
            return
        if tag == "pre" and "mermaid" in d.get("class", "").split():
            self._cur = {"slide": self.slide_index, "source": ""}
            return

    def handle_data(self, data: str) -> None:
        if self._cur is not None:
            self._cur["source"] += data

    def handle_endtag(self, tag: str) -> None:
        if tag == "pre" and self._cur is not None:
            self._cur["source"] = self._cur["source"].strip()
            self.decls.append(self._cur)
            self._cur = None


def parse_mermaid_declarations(html_text: str) -> list[dict]:
    """HTML → [{slide, source}]（slide 1-based）。pre.mermaid 在 data-pptx-slide
    <section> 之外 → ValueError（对齐 charts：声明必须落在 section 内）。"""
    p = _MermaidHTMLParser()
    p.feed(html_text)
    p.close()
    for d in p.decls:
        if d["slide"] <= 0:
            raise ValueError(
                'mermaid 块出现在 slide 之外——<pre class="mermaid"> 必须在 '
                "data-pptx-slide 的 <section> 内"
            )
    return p.decls


def load_mermaid_boxes(measurements_path: str) -> dict[int, dict]:
    """measurements.json → {slide_index(1-based): {"x","y","w","h"}}（px）。

    匹配 record className 分词含 "mermaid"（pre.mermaid 的 rect）。
    """
    with open(measurements_path, encoding="utf-8") as f:
        data = json.load(f)
    boxes: dict[int, dict] = {}
    for i, slide in enumerate(data.get("slides", []), start=1):
        for rec in slide.get("records", []):
            cls = (rec.get("className") or "").split()
            if "mermaid" in cls:
                boxes[i] = dict(rec["rect"])
                break
    return boxes


def _measurements_path(pptx_path: str) -> str:
    p = Path(pptx_path)
    return str(p.with_name(f"{p.stem}_audit") / "_cache" / "measurements.json")


def _remove_bbox_shapes(slide, box_emu: dict) -> None:
    """删除与注入矩形几何一致（同位置同尺寸）的占位形状（pre 文本块等）。

    占位矩形是确定尺寸的（box_emu = px→EMU），用户任意形状几乎不可能精确重合；
    旧「中心点落在 bbox 内」判定会误删用户内容。容差 0.01in 吸收 px→EMU 舍入。
    """
    tol = int(0.01 * 914400)  # 0.01in，EMU
    bx, by, bw, bh = box_emu["x"], box_emu["y"], box_emu["w"], box_emu["h"]
    for shape in list(slide.shapes):
        if (
            abs(shape.left - bx) <= tol
            and abs(shape.top - by) <= tol
            and abs(shape.width - bw) <= tol
            and abs(shape.height - bh) <= tol
        ):
            slide.shapes._spTree.remove(shape._element)


def inject_mermaid(pptx_path: str, decls: list[dict], boxes: dict[int, dict]) -> None:
    """把每块 mermaid 渲染成可编辑形状，替换 slide 内对应 bbox 占位。"""
    from pptx import Presentation

    prs = Presentation(pptx_path)
    for decl in decls:
        box = boxes.get(decl["slide"])
        if box is None or not decl["source"].strip():
            continue
        box_emu = {k: int(v * PX_TO_EMU) for k, v in box.items()}
        slide = prs.slides[decl["slide"] - 1]
        try:
            diagram = parse_mermaid(decl["source"])
        except ValueError as e:
            raise ValueError(f"第 {decl['slide']} 页 mermaid 块解析失败: {e}") from None
        # bbox 是 px；转换器画布 1920px = 13.333in = 12192000 EMU → 1px = 1/144 in。
        # layout 的 max_w/max_h 单位是 inches，px→inch 除以 144（不能除以 PX_TO_EMU）。
        lay = layout_diagram(diagram, max_w=box["w"] / 144, max_h=box["h"] / 144)
        _remove_bbox_shapes(slide, box_emu)
        render_to_slide(slide, lay, offset_x=box_emu["x"], offset_y=box_emu["y"])
    prs.save(pptx_path)


def postprocess_mermaid(html_path: str, pptx_path: str) -> None:
    """转换后调用（对齐 charts.postprocess_charts 签名）：HTML 含 <pre class="mermaid">
    → 读 measurements → 注入可编辑形状。无声明 → 原样返回。声明/数据非法 → ValueError。
    """
    with open(html_path, encoding="utf-8") as f:
        html_text = f.read()
    # 无子串守卫：HTMLParser 按 class 值分词，多 class（class="mermaid chart"）或
    # 属性空格（class = "mermaid"）都能命中；子串守卫会漏掉这些变体导致 mermaid
    # 块以字面文本留在 PPTX。解析本身便宜，直接走 parse 再按 decls 空判断返回。
    decls = parse_mermaid_declarations(html_text)
    if not decls:
        return
    meas_path = _measurements_path(pptx_path)
    if not os.path.exists(meas_path):
        raise RuntimeError(
            f"找不到 convert 审计产物 {meas_path}——mermaid 注入需要 measurements.json，"
            "请勿用 --no-visual-audit"
        )
    boxes = load_mermaid_boxes(meas_path)
    missing = [d["slide"] for d in decls if d["slide"] not in boxes]
    if missing:
        raise RuntimeError(
            f'第 {missing} 页没测到 <pre class="mermaid"> 容器——请确认 pre 有 '
            'class="mermaid" 且被渲染进 visual audit'
        )
    inject_mermaid(pptx_path, decls, boxes)
