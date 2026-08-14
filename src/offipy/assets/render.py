"""offipy.assets — generic asset measurement binding and renderer/z-order.

A3 Task 6/7: bind converter `kind=asset` measurements to deterministic
declaration ids, locate the transparent placeholder in an assembled deck, and
render resolved payloads into exact placeholder XML slots (replace / decorative
/ background). Internal-facing helpers for the postprocess pipeline; not part of
the public asset core surface.
"""

from __future__ import annotations

import contextlib
import json
import logging
import re
from dataclasses import dataclass
from pathlib import Path
from typing import TYPE_CHECKING, Any, cast

from offipy.assets.declarations import (
    AssetDeclaration,
    parse_injected_asset_declarations,
)
from offipy.assets.model import (
    AssetPlacement,
    AssetProviderMeta,
    AssetRect,
    AssetRenderContext,
    NativeShapePayload,
    RasterPayload,
    ResolvedAsset,
    SvgPayload,
    SvgTemplatePayload,
)
from offipy.assets.registry import get_default_registry
from offipy.assets.uri import format_asset_uri
from offipy.exceptions import InvalidArgumentError
from offipy.icons import _measurements_path

if TYPE_CHECKING:
    from collections.abc import Callable, Sequence

    from playwright.sync_api import Page
    from pptx.dml.color import RGBColor

PLACEHOLDER_PREFIX = "OFFIPY_ASSET::"


@dataclass(frozen=True)
class AssetMeasurement:
    """converter `kind=asset` 测量记录的内部视图（非公共 Asset 模型）。"""

    asset_id: str
    slide_index: int
    rect: AssetRect
    theme_vars: dict[str, str]
    html_tag: str
    color: str | None = None

    def __post_init__(self) -> None:
        # 冻结上下文的 theme_vars 与调用方可变 dict 解耦
        object.__setattr__(self, "theme_vars", dict(self.theme_vars))


def _load_slides_data(data: Any) -> list[dict[str, Any]]:
    if isinstance(data, dict) and "slides" in data:
        return cast("list[dict[str, Any]]", data["slides"])
    if isinstance(data, list):
        return data
    raise InvalidArgumentError("measurements.json 缺少 slides 数组")


def load_asset_measurements(path: str | Path) -> dict[str, AssetMeasurement]:
    """读取 measurements.json，收集 kind=asset 记录并按 assetId 建索引。

    严格校验（不静默丢弃）：assetId 非空、全 deck 内唯一、rect 必须为正。
    返回 dict[declaration_id → AssetMeasurement]。
    """
    data = json.loads(Path(path).read_text(encoding="utf-8"))
    result: dict[str, AssetMeasurement] = {}
    for i, sdata in enumerate(_load_slides_data(data)):
        for rec in sdata.get("records", []):
            if rec.get("kind") != "asset":
                continue
            asset_id = rec.get("assetId") or ""
            if not asset_id:
                raise InvalidArgumentError(f"slide {i + 1} 的 asset 测量记录缺少 assetId")
            if asset_id in result:
                raise InvalidArgumentError(f"asset 测量记录 assetId 重复: {asset_id}")
            rect = rec.get("rect") or {}
            try:
                ar = AssetRect(
                    x=float(rect["x"]),
                    y=float(rect["y"]),
                    width=float(rect["w"]),
                    height=float(rect["h"]),
                )
            except (KeyError, TypeError, ValueError) as exc:
                raise InvalidArgumentError(f"asset {asset_id} rect 无效: {rect!r}") from exc
            ar.validate_render()  # 非正 rect → 可执行错误，不静默丢弃
            color = rec.get("color")
            result[asset_id] = AssetMeasurement(
                asset_id=asset_id,
                slide_index=i + 1,
                rect=ar,
                theme_vars={k: str(v) for k, v in (rec.get("themeVars") or {}).items()},
                html_tag=str(rec.get("tag") or ""),
                color=str(color) if color is not None else None,
            )
    return result


def find_asset_placeholder(slide: Any, declaration_id: str) -> Any:
    """在已装配 slide 里按 name 精确查找透明占位符。

    必须恰好一个；0 或 >1 都是绑定错误。不做 outerHTML / 空间匹配兜底。
    """
    name = f"{PLACEHOLDER_PREFIX}{declaration_id}"
    matches = [sp for sp in slide.shapes if sp.name == name]
    if len(matches) != 1:
        raise InvalidArgumentError(
            f"asset 占位符 {declaration_id} 应为 1 个，实际 {len(matches)} 个"
        )
    return matches[0]


def bind_asset_measurements(
    declarations: list[AssetDeclaration],
    measurements: dict[str, AssetMeasurement],
) -> dict[str, AssetMeasurement]:
    """把声明按 declaration_id 绑定到测量，校验整副 deck 一致性。

    Fail on: 声明无测量、测量无声明、slide 序号不匹配。返回 dict 保持声明顺序。
    """
    bound: dict[str, AssetMeasurement] = {}
    for decl in declarations:
        measurement = measurements.get(decl.declaration_id)
        if measurement is None:
            raise InvalidArgumentError(f"asset 声明 {decl.declaration_id} 缺少对应测量记录")
        if measurement.slide_index != decl.slide_index:
            raise InvalidArgumentError(
                f"asset 声明 {decl.declaration_id} slide 序号不匹配："
                f"声明 slide {decl.slide_index} vs 测量 slide {measurement.slide_index}"
            )
        bound[decl.declaration_id] = measurement
    for asset_id in measurements:
        if asset_id not in bound:
            raise InvalidArgumentError(f"asset 测量记录 {asset_id} 无对应声明")
    return bound


# ---------------------------------------------------------------------------
# 渲染（Task 7）
# ---------------------------------------------------------------------------

_PX_TO_EMU = 6350  # 转换器画布 1920×1080 px = 12192000×6858000 EMU

_logger = logging.getLogger(__name__)

# SVG→PNG fallback 防御：screenshot 超时 + viewBox 画布上限，防畸形/恶意 SVG 挂起或
# 制造超大页面占满内存（#H8）。合法 SVG 长边不会超过 deck 画布数量级。
_SCREENSHOT_TIMEOUT_MS = 10_000
_MAX_SVG_DIM = 8192


def _px_rect(rect: AssetRect) -> dict[str, float]:
    return {"x": rect.x, "y": rect.y, "w": rect.width, "h": rect.height}


def _accent_rgb(context: AssetRenderContext) -> RGBColor | None:
    """theme_vars['accent'] → RGBColor（图标 computed color 缺失时的 v0.13.2 兜底）。

    与 legacy `_parse_color` 同源解析：接受 #RRGGBB 与 rgb(r,g,b)（computed CSS
    值可能以 rgb() 返回，旧实现只认 hex 会静默丢主题 accent）；无法解析返回 None
    （再落缺省主蓝）。
    """
    accent = context.theme_vars.get("accent")
    if not accent:
        return None
    from pptx.dml.color import RGBColor

    m = re.fullmatch(r"#([0-9a-fA-F]{6})", accent.strip())
    if m:
        # python-pptx RGBColor.from_string is unannotated → returns Any
        return cast("RGBColor", RGBColor.from_string(m.group(1)))  # type: ignore[no-untyped-call]
    m = re.search(r"rgba?\(\s*(\d+)\s*,\s*(\d+)\s*,\s*(\d+)", accent)
    if m:
        # CSS 对超界分量按最近边界钳制；RGBColor 严格 0-255，越界会抛 ValueError
        # python-pptx RGBColor.__new__ is unannotated → no-untyped-call
        return RGBColor(  # type: ignore[no-untyped-call]
            max(0, min(255, int(m.group(1)))),
            max(0, min(255, int(m.group(2)))),
            max(0, min(255, int(m.group(3)))),
        )
    return None


def render_asset(
    slide: Any,
    resolved: ResolvedAsset,
    context: AssetRenderContext,
    *,
    color: str | None = None,
    svg_to_png: Callable[[str], bytes | None] | None = None,
) -> list[Any]:
    """按 payload 类型把 ResolvedAsset 渲染进 slide，返回创建的形状/元素列表。

    - freeform_svg → 既有图标 freeform 渲染器（ph fill / lu stroke）；
    - RasterPayload → add_picture；
    - SvgPayload(svg)/SvgTemplatePayload → A1 选定 P2（OOXML SVG picture + svgBlip）；
    - NativeShapePayload → A5 前未注册渲染器，明确报错。
    color 是测量出的计算色（CSS color），None 时兜底主题 accent。
    svg_to_png 把 SVG 渲染成 PNG 字节（供不支持 SVG 的查看器回退），返回 None 时
    保持纯 SVG picture（#58）。
    """
    payload = resolved.payload
    rect = context.rect
    if isinstance(payload, SvgPayload) and payload.render_mode == "freeform_svg":
        from offipy.assets.providers.icons import icon_render_mode
        from offipy.icons import render_icon_payload

        if payload.view_box is None:
            raise InvalidArgumentError("freeform_svg payload missing view_box")
        return render_icon_payload(
            slide,
            payload.svg,
            mode=icon_render_mode(resolved.provider_meta.provider_id),
            view_box=payload.view_box,
            rect=_px_rect(rect),
            color=color,
            fallback=_accent_rgb(context),
        )
    if isinstance(payload, RasterPayload):
        return [_render_raster_picture(slide, rect, payload)]
    if isinstance(payload, SvgTemplatePayload):
        from offipy.assets.materialize import materialize_svg_template

        payload = materialize_svg_template(payload, context.theme_vars)
    if isinstance(payload, SvgPayload):
        png = svg_to_png(payload.svg) if svg_to_png is not None else None
        return [_render_svg_picture(slide, rect, payload.svg, png)]
    if isinstance(payload, NativeShapePayload):
        if context.placement == "background":
            raise InvalidArgumentError(
                "native primitives do not support background placement in v0.14"
            )
        from offipy.assets.primitives import get_native_renderer

        renderer = get_native_renderer(payload.primitive)
        shapes = renderer(slide, dict(payload.params), context)
        return list(shapes)
    raise InvalidArgumentError(f"unknown asset payload type {type(payload).__name__}")


def _render_raster_picture(slide: Any, rect: AssetRect, payload: RasterPayload) -> Any:
    from io import BytesIO

    from pptx.util import Emu

    x, y, w, h = (round(v * _PX_TO_EMU) for v in (rect.x, rect.y, rect.width, rect.height))
    return slide.shapes.add_picture(BytesIO(payload.data), Emu(x), Emu(y), Emu(w), Emu(h))


def _render_svg_picture(
    slide: Any, rect: AssetRect, svg_text: str, png_bytes: bytes | None = None
) -> Any:
    """OOXML SVG picture（asvg svgBlip + 可选 raster fallback blip）。

    结构匹配 PowerPoint 自己对 SVG 产出的形式（COM AddPicture 对照验证过）。
    提供 png_bytes 时主 <a:blip r:embed> 指向 PNG、asvg:svgBlip 指向 SVG——支持 SVG
    的宿主渲染矢量，不支持的查看器回退 PNG，两端都不空白（#58）；无 png_bytes
    保持纯 SVG（blip 不挂 embed）。
    """
    from lxml import etree
    from pptx.opc.package import Part
    from pptx.opc.packuri import PackURI

    svg_bytes = svg_text.encode("utf-8")
    pkg = slide.part.package
    slide_part = slide.part
    image_rel = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image"

    svg_partname = PackURI(str(pkg.next_partname("/ppt/media/image%d.svg")))
    svg_part = Part(svg_partname, "image/svg+xml", pkg, svg_bytes)
    svg_rid = slide_part.relate_to(svg_part, image_rel)

    png_rid = None
    if png_bytes is not None:
        png_partname = PackURI(str(pkg.next_partname("/ppt/media/image%d.png")))
        png_part = Part(png_partname, "image/png", pkg, png_bytes)
        png_rid = slide_part.relate_to(png_part, image_rel)

    a = "http://schemas.openxmlformats.org/drawingml/2006/main"
    p = "http://schemas.openxmlformats.org/presentationml/2006/main"
    r = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    asvg = "http://schemas.microsoft.com/office/drawing/2016/SVG/main"
    svg_ext_uri = "{96DAC541-7B7A-43D3-8B79-37D633B846F1}"

    etree.register_namespace("asvg", asvg)
    x, y, w, h = (round(v * _PX_TO_EMU) for v in (rect.x, rect.y, rect.width, rect.height))

    sp_tree = slide.shapes._spTree
    pic = etree.SubElement(sp_tree, f"{{{p}}}pic")
    nvPicPr = etree.SubElement(pic, f"{{{p}}}nvPicPr")
    etree.SubElement(nvPicPr, f"{{{p}}}cNvPr", id=str(_next_shape_id(slide)), name="SVG Asset")
    etree.SubElement(nvPicPr, f"{{{p}}}cNvPicPr")
    etree.SubElement(nvPicPr, f"{{{p}}}nvPr")
    blipFill = etree.SubElement(pic, f"{{{p}}}blipFill")
    blip = etree.SubElement(blipFill, f"{{{a}}}blip")
    if png_rid is not None:
        blip.set(f"{{{r}}}embed", png_rid)
    extLst = etree.SubElement(blip, f"{{{a}}}extLst")
    ext = etree.SubElement(extLst, f"{{{a}}}ext", uri=svg_ext_uri)
    svgBlip = etree.SubElement(ext, f"{{{asvg}}}svgBlip")
    svgBlip.set(f"{{{r}}}embed", svg_rid)
    stretch = etree.SubElement(blipFill, f"{{{a}}}stretch")
    etree.SubElement(stretch, f"{{{a}}}fillRect")
    spPr = etree.SubElement(pic, f"{{{p}}}spPr")
    xfrm = etree.SubElement(spPr, f"{{{a}}}xfrm")
    etree.SubElement(xfrm, f"{{{a}}}off", x=str(x), y=str(y))
    etree.SubElement(xfrm, f"{{{a}}}ext", cx=str(w), cy=str(h))
    prstGeom = etree.SubElement(spPr, f"{{{a}}}prstGeom", prst="rect")
    etree.SubElement(prstGeom, f"{{{a}}}avLst")
    return pic


def _svg_page_screenshot(page: Page, svg_text: str) -> bytes | None:
    """用浏览器页面把 SVG 渲染成 PNG；viewBox 尺寸定画布，失败返回 None。"""
    import math
    import re

    m = re.search(r'viewBox="([^"]+)"', svg_text)
    if m:
        # SVG 规范允许逗号/空白任意混用作坐标分隔（viewBox="0,0 100,100"）
        parts = [p for p in re.split(r"[,\s]+", m.group(1).strip()) if p]
        if len(parts) != 4:
            return None
        try:
            vx, vy, vw, vh = (float(v) for v in parts)
        except ValueError:
            return None
    else:
        vx, vy, vw, vh = 0.0, 0.0, 512.0, 512.0
    if not all(math.isfinite(v) for v in (vx, vy, vw, vh)) or vw <= 0 or vh <= 0:
        return None
    clip_x, clip_y = max(0.0, vx), max(0.0, vy)
    width = min(max(1, math.ceil(clip_x + vw)), _MAX_SVG_DIM)
    height = min(max(1, math.ceil(clip_y + vh)), _MAX_SVG_DIM)
    page.set_viewport_size({"width": width, "height": height})
    page.set_content(
        "<html><head><style>html,body{margin:0;padding:0}"
        "svg{width:100%;height:100%;display:block}</style></head>"
        f"<body><div style='width:{width}px;height:{height}px'>{svg_text}</div></body></html>",
        timeout=_SCREENSHOT_TIMEOUT_MS,
    )
    return page.screenshot(
        clip={
            "x": min(clip_x, float(width - 1)),
            "y": min(clip_y, float(height - 1)),
            "width": min(int(vw), width),
            "height": min(int(vh), height),
        },
        type="png",
        timeout=_SCREENSHOT_TIMEOUT_MS,
    )


def _make_svg_to_png() -> tuple[Callable[[str], bytes | None], Callable[[], None]]:
    """惰性共享 Playwright chromium 的 SVG→PNG 渲染器（#58）。

    首次 convert 调用启动一次 browser，之后复用；返回 (convert, close)。Playwright
    或 chromium 缺失、任一 SVG 渲染失败 → convert 返回 None（调用方降级为纯 SVG），
    绝不因 fallback 失败中断资产注入。
    """
    browser = None
    pw = None

    def convert(svg_text: str) -> bytes | None:
        nonlocal browser, pw
        if pw is None:
            try:
                from playwright.sync_api import sync_playwright
            except ImportError:
                return None
            pw = sync_playwright().start()
        if browser is None:
            try:
                browser = pw.chromium.launch()
            except Exception:
                return None
        page = None
        try:
            # new_page 必须在 try 内：创建失败不得让整个资产注入崩溃（#H9）
            page = browser.new_page()
            return _svg_page_screenshot(page, svg_text)
        except Exception as exc:
            _logger.warning("SVG→PNG fallback failed: %s", exc)
            return None
        finally:
            if page is not None:
                with contextlib.suppress(Exception):
                    page.close()

    def close() -> None:
        nonlocal browser, pw
        if browser is not None:
            browser.close()
            browser = None
        if pw is not None:
            pw.stop()
            pw = None

    return convert, close


def _next_shape_id(slide: Any) -> int:
    """spTree 里下一个可用 cNvPr id（当前最大 id + 1，无则 2）。"""
    ids = [int(sp.get("id") or 0) for sp in slide.shapes._spTree.iter() if sp.tag.endswith("cNvPr")]
    return max(ids, default=1) + 1


def _as_element(item: Any) -> Any:
    """python-pptx shape → 其 XML 元素；lxml 元素原样返回。"""
    return getattr(item, "_element", item)


def place_rendered_elements(
    slide: Any, placeholder: Any, rendered: Sequence[object], placement: AssetPlacement
) -> None:
    """把渲染产物精确放入占位符 XML 槽位，不使用 z_order 猜测或 send_to_back 循环。

    replace/decorative：记录占位符 XML 索引 → 产物按渲染顺序插回该槽 → 移除占位符
    并验证。background：移除占位符 → 产物移到 grpSpPr 之后、所有内容形状之前（保持
    内部顺序）。
    """
    from pptx.oxml.ns import qn

    sp_tree = slide.shapes._spTree
    ph_element = placeholder._element
    elements = [_as_element(r) for r in rendered]

    if placement == "background":
        sp_tree.remove(ph_element)
        grp_sp = sp_tree.find(qn("p:grpSpPr"))
        anchor = list(sp_tree).index(grp_sp) + 1 if grp_sp is not None else 0
        for el in elements:
            sp_tree.insert(anchor, el)
            anchor += 1
        return

    ph_index = list(sp_tree).index(ph_element)
    for el in elements:
        sp_tree.insert(ph_index, el)
        ph_index += 1
    sp_tree.remove(ph_element)
    if ph_element.getparent() is not None:
        raise RuntimeError(f"asset 占位符 {placeholder.name} 移除失败")


# ---------------------------------------------------------------------------
# 编排（Task 8）
# ---------------------------------------------------------------------------


@dataclass(frozen=True)
class AssetUsageRecord:
    """一次 asset 渲染的确定性用量记录（内部面；Task 9 序列化进 assets.json）。"""

    declaration_id: str
    slide_index: int
    request: str
    placement: AssetPlacement
    provider: AssetProviderMeta


@dataclass(frozen=True)
class AssetUsageReport:
    """一次 postprocess_assets 的用量汇总，按声明顺序保存。"""

    records: tuple[AssetUsageRecord, ...]


def _theme_name(html_text: str) -> str | None:
    m = re.search(r'<style\s+data-theme="([a-z0-9-]+)"', html_text)
    return m.group(1) if m else None


def postprocess_assets(html_path: str, pptx_path: str) -> AssetUsageReport:
    """转换后调用：注入副本 asset 声明 → 测量绑定 → 解析/渲染/落位，返回用量报告。

    无 asset 声明 → 空报告，原样返回（不打开 PPTX）。measurements.json 缺失 →
    RuntimeError（deck._postprocess 统一映射为 ConversionError）。打开/保存 PPTX
    各一次，绝不逐资产重复。
    """

    with Path(html_path).open(encoding="utf-8") as f:
        html_text = f.read()
    if "data-offipy-asset-id" not in html_text:
        return AssetUsageReport(())
    decls = parse_injected_asset_declarations(html_text)
    if not decls:
        return AssetUsageReport(())
    meas_path = _measurements_path(pptx_path)
    if not Path(meas_path).exists():
        raise RuntimeError(
            f"找不到 convert 审计产物 {meas_path}——asset 注入需要 measurements.json，"
            "请勿用 --no-visual-audit"
        )
    measurements = load_asset_measurements(meas_path)
    bound = bind_asset_measurements(decls, measurements)
    registry = get_default_registry()
    theme_name = _theme_name(html_text)

    from pptx import Presentation

    prs = Presentation(pptx_path)
    svg_to_png, close_svg_png = _make_svg_to_png()
    records: list[AssetUsageRecord] = []
    try:
        for decl in decls:
            slide = prs.slides[decl.slide_index - 1]
            measurement = bound[decl.declaration_id]
            context = AssetRenderContext(
                slide_index=decl.slide_index,
                rect=measurement.rect,
                theme_name=theme_name,
                theme_vars=measurement.theme_vars,
                placement=decl.placement,
            )
            placeholder = find_asset_placeholder(slide, decl.declaration_id)
            resolved = registry.resolve(decl.request)
            rendered = render_asset(
                slide, resolved, context, color=measurement.color, svg_to_png=svg_to_png
            )
            place_rendered_elements(slide, placeholder, rendered, decl.placement)
            records.append(
                AssetUsageRecord(
                    declaration_id=decl.declaration_id,
                    slide_index=decl.slide_index,
                    request=format_asset_uri(decl.request),
                    placement=decl.placement,
                    provider=resolved.provider_meta,
                )
            )
    finally:
        close_svg_png()
    prs.save(pptx_path)
    return AssetUsageReport(tuple(records))
