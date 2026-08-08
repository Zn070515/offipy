"""offipy.assets.primitives._common — shared native-rendering helpers (A5).

Font/theme policy: v0.14 AssetRenderContext carries color vars only. Colors
resolve from measurement theme vars via the A2 `resolve_asset_color` helper;
primitive labels use a stable common sans font (Arial), matching the deck
theme's `--font-sans`. No new public font token model is introduced.

Geometry discipline: every primitive draws inside its `AssetRect` (EMU via the
same 6350 px→EMU conversion as the converter) and returns python-pptx shapes /
XML elements in visual stacking order bottom → top. Renderers never touch
placeholder z-order themselves.
"""

from __future__ import annotations

from collections.abc import Iterable, Mapping
from typing import Any

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Pt

from offipy.assets.materialize import resolve_asset_color
from offipy.assets.model import AssetRect, AssetRenderContext
from offipy.exceptions import InvalidArgumentError

PX_TO_EMU = 6350  # same conversion as the converter (1920×1080 px canvas)

# Stable common sans for primitive labels, matching deck theme --font-sans.
PRIMITIVE_FONT = "Arial"

_EMPTY_RECT_MIN = 1.0  # px


def px_to_emu(v: float) -> int:
    return int(round(v * PX_TO_EMU))


def require_rect(ctx: AssetRenderContext) -> AssetRect:
    """Validate the render rect is positive and return it."""
    ctx.rect.validate_render()
    return ctx.rect


def _resolve_text_token(ctx: AssetRenderContext, token: str, fallback: str) -> str:
    """Resolve a theme token for internal text color, falling back on absence.

    Text tokens (ink/muted) are rendering details, not user-facing params; a
    theme that omits them must not break a primitive.
    """
    try:
        return resolve_asset_color(token, ctx.theme_vars)
    except InvalidArgumentError:
        return fallback


def resolve_native_colors(params: Mapping[str, str], ctx: AssetRenderContext) -> dict[str, str]:
    """Resolve accent/fill (user-facing) plus ink/muted text tokens.

    accent/fill come from the validated payload params; unknown theme tokens
    there raise (matches A2). ink/muted are internal and fall back to safe
    defaults when the theme does not define them.

    The ``fill`` default for some primitives is the ``accent`` token (e.g.
    label-pill); that must follow the primitive's final accent (incl. an
    explicit ``accent`` param), not resolve against theme vars independently,
    otherwise an accent override is silently lost (#59).
    """
    accent = resolve_asset_color(params.get("accent", "accent"), ctx.theme_vars)
    fill_raw = params.get("fill", "transparent")
    # fill 默认跟随本图元最终 accent（#59），不独立对 theme_vars 解析
    fill = accent if fill_raw == "accent" else resolve_asset_color(fill_raw, ctx.theme_vars)
    return {
        "accent": accent,
        "fill": fill,
        "ink": _resolve_text_token(ctx, "ink", "#222222"),
        "muted": _resolve_text_token(ctx, "muted", "#667085"),
    }


def _rgb(hex_color: str) -> RGBColor:
    return RGBColor.from_string(hex_color.lstrip("#"))


def set_fill_line(shape, fill: str | None, line: str | None = None) -> None:
    """Set solid/transparent fill and line on a python-pptx shape.

    ``None`` or ``"transparent"`` means no fill (shape.fill.background()) /
    no line (line.fill.background()).
    """
    if fill is None or fill == "transparent":
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = _rgb(fill)
    if line is None or line == "transparent":
        shape.line.fill.background()
    else:
        shape.line.color.rgb = _rgb(line)


def add_shape(
    slide,
    shape_type: MSO_SHAPE,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    fill: str | None = None,
    line: str | None = None,
) -> Any:
    """Add a native autoshape at EMU coordinates, returning the python-pptx shape."""
    sp = slide.shapes.add_shape(shape_type, px_to_emu(x), px_to_emu(y), px_to_emu(w), px_to_emu(h))
    set_fill_line(sp, fill, line)
    return sp


def add_textbox(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    text: str = "",
    font_size_pt: float = 12.0,
    color: str = "#222222",
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    anchor: MSO_ANCHOR = MSO_ANCHOR.TOP,
    word_wrap: bool = False,
) -> Any:
    """Add a textbox with a single editable run at EMU coordinates.

    Margins are zeroed so the box geometry matches the given rect exactly and
    text fills it deterministically.
    """
    tb = slide.shapes.add_textbox(px_to_emu(x), px_to_emu(y), px_to_emu(w), px_to_emu(h))
    tf = tb.text_frame
    tf.word_wrap = word_wrap
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = PRIMITIVE_FONT
    run.font.size = Pt(font_size_pt)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color)
    return tb


def set_text_style(
    run,
    *,
    text: str | None = None,
    font_size_pt: float | None = None,
    color: str | None = None,
    bold: bool | None = None,
) -> None:
    """Re-style a text run in place (mutate an already-created run)."""
    if text is not None:
        run.text = text
    run.font.name = PRIMITIVE_FONT
    if font_size_pt is not None:
        run.font.size = Pt(font_size_pt)
    if color is not None:
        run.font.color.rgb = _rgb(color)
    if bold is not None:
        run.font.bold = bold


def set_shape_text(
    sp,
    text: str,
    *,
    font_size_pt: float,
    color: str,
    align: PP_ALIGN = PP_ALIGN.CENTER,
    anchor: MSO_ANCHOR = MSO_ANCHOR.MIDDLE,
    word_wrap: bool = True,
) -> None:
    """Populate an autoshape's text frame with a single editable run.

    Margins are zeroed so text fills the shape deterministically.
    """
    tf = sp.text_frame
    tf.word_wrap = word_wrap
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = PRIMITIVE_FONT
    run.font.size = Pt(font_size_pt)
    run.font.color.rgb = _rgb(color)


def shape_elements(shapes: Iterable) -> tuple[object, ...]:
    """Normalize python-pptx shapes / XML elements to their XML elements.

    Renderers return shapes; the generic placement code inserts their elements
    into the placeholder slot, so we must hand it XML elements (python-pptx
    shapes are wrapped; their ``_element`` is the actual OOXML).
    """
    return tuple(getattr(s, "_element", s) for s in shapes)


def _char_width_em(ch: str, approx: float) -> float:
    """Per-glyph width as an em fraction.

    Full-width glyphs (CJK / 全角标点/空格) advance ≈1.0em; latin glyphs average
    ≈``approx``. Arrows (U+2190-21FF) and dashes (U+2010-2015) are half-width
    code points whose glyphs render full-width in Arial/CJK, so they also count
    as 1.0em (#62). The old uniform 0.55em estimate under-measures these by
    ~1.8×, so CJK text overflowed the measured rect (#56).
    """
    cp = ord(ch)
    if (
        cp > 0x2E80
        or cp == 0x3000
        or 0x2010 <= cp <= 0x2015
        or 0x2190 <= cp <= 0x21FF
        or 0xFF01 <= cp <= 0xFF60
        or 0xFFE0 <= cp <= 0xFFE6
    ):
        return 1.0
    return approx


def fit_font_size(
    text: str,
    w_px: float,
    h_px: float,
    *,
    start_pt: float,
    min_pt: float,
    approx_char_width: float = 0.55,
    line_height_pt: float = 1.2,
) -> float:
    """Largest font size (pt) that fits a single line in w_px×h_px.

    Uses a rough proportional-width estimate, weighting full-width (CJK) glyphs
    at 1.0em and latin glyphs at ``approx_char_width``. Never returns below
    ``min_pt``; if even ``min_pt`` overflows, raises instead of silently
    producing unreadable output.
    """
    if not text:
        return min(start_pt, max(min_pt, h_px * 72.0 / 96.0 / line_height_pt))
    for size_pt in _descending_pt(start_pt, min_pt):
        em_px = size_pt * (96.0 / 72.0)
        line_h_px = em_px * line_height_pt
        text_w_px = em_px * sum(_char_width_em(ch, approx_char_width) for ch in text)
        if text_w_px <= w_px and line_h_px <= h_px:
            return size_pt
    raise InvalidArgumentError(
        f"text {text[:24]!r} does not fit rect {w_px:.0f}x{h_px:.0f}px at min {min_pt}pt"
    )


def fit_font_size_wrapped(
    text: str,
    w_px: float,
    h_px: float,
    *,
    start_pt: float,
    min_pt: float,
    approx_char_width: float = 0.55,
    line_height_pt: float = 1.2,
) -> float:
    """Largest font size (pt) whose greedy word-wrap fits w_px×h_px.

    Whitespace is the only wrap point; a token wider than the box must shrink
    (never truncates user text). Raises when even ``min_pt`` overflows either
    axis, so no rendered line ever escapes the rect.
    """
    if not text:
        return min(start_pt, max(min_pt, h_px * 72.0 / 96.0 / line_height_pt))
    for size_pt in _descending_pt(start_pt, min_pt):
        em_px = size_pt * (96.0 / 72.0)
        line_h_px = em_px * line_height_pt
        n_lines, max_line_w = _wrap_measure(text, em_px, w_px, approx_char_width)
        if max_line_w <= w_px and n_lines * line_h_px <= h_px:
            return size_pt
    raise InvalidArgumentError(
        f"text {text[:24]!r} does not fit rect {w_px:.0f}x{h_px:.0f}px at min {min_pt}pt"
    )


def _wrap_measure(
    text: str, em_px: float, w_px: float, approx_char_width: float
) -> tuple[int, float]:
    space_w_px = em_px * approx_char_width
    lines = 1
    cur = 0.0
    max_line_w = 0.0
    for token in text.split():
        token_w = em_px * sum(_char_width_em(ch, approx_char_width) for ch in token)
        if cur == 0:
            cur = token_w
        elif cur + space_w_px + token_w <= w_px:
            cur += space_w_px + token_w
        else:
            max_line_w = max(max_line_w, cur)
            lines += 1
            cur = token_w
    return lines, max(max_line_w, cur)


def _descending_pt(start_pt: float, min_pt: float) -> Iterable[float]:
    """Font-size candidates from start_pt down to min_pt, coarse then fine.

    1pt integer steps first; once a step would cross min_pt, switch to 0.25pt
    steps so a fractional optimum near min_pt is never skipped, ending exactly
    at min_pt (the old `step = 0.25` assignment was unreachable dead code).
    """
    step = 1.0
    size = start_pt
    while size >= min_pt - 1e-9:
        yield size
        if size <= min_pt + 1e-9:
            return
        next_size = size - step
        if next_size < min_pt:
            if step >= 1.0:
                size = size - 0.25
                step = 0.25
            else:
                size = min_pt
        else:
            size = next_size
