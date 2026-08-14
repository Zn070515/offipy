"""deck 集成：HTML <div class="drawio" data-drawio="..."> → measurements bbox → 注入。"""

from __future__ import annotations

import json

import pytest
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu, Inches

from offipy.drawio import parse_drawio_declarations, postprocess_drawio

HTML = """\
<!doctype html>
<html><head><meta charset="utf-8"></head><body>
<section data-pptx-slide>
  <h1>架构图</h1>
  <div class="drawio-box">
    <div class="drawio" data-drawio="arch.drawio"></div>
  </div>
</section>
</body></html>
"""

DRAWIO = """\
<mxfile><diagram name="Page 1"><mxGraphModel><root>
  <mxCell id="0"/><mxCell id="1" parent="0"/>
  <mxCell id="a" value="模块A"
          style="rounded=0;fillColor=#dae8fc;strokeColor=#6c8ebf;"
          vertex="1" parent="1">
    <mxGeometry x="40" y="40" width="120" height="60" as="geometry"/>
  </mxCell>
  <mxCell id="b" value="模块B"
          style="rounded=0;fillColor=#d5e8d4;strokeColor=#82b366;"
          vertex="1" parent="1">
    <mxGeometry x="220" y="40" width="120" height="60" as="geometry"/>
  </mxCell>
</root></mxGraphModel></diagram></mxfile>
"""

MULTI_DRAWIO = """\
<mxfile>
  <diagram name="第一页"><mxGraphModel><root>
    <mxCell id="0"/><mxCell id="1" parent="0"/>
    <mxCell id="a" value="第一页节点"
            style="rounded=0;fillColor=#dae8fc;strokeColor=#6c8ebf;"
            vertex="1" parent="1">
      <mxGeometry x="40" y="40" width="120" height="60" as="geometry"/>
    </mxCell>
  </root></mxGraphModel></diagram>
  <diagram name="第二页"><mxGraphModel><root>
    <mxCell id="0"/><mxCell id="1" parent="0"/>
    <mxCell id="b" value="第二页节点"
            style="rounded=0;fillColor=#dae8fc;strokeColor=#6c8ebf;"
            vertex="1" parent="1">
      <mxGeometry x="40" y="40" width="120" height="60" as="geometry"/>
    </mxCell>
  </root></mxGraphModel></diagram>
</mxfile>
"""


def _measurements(slide_index=1):
    return {
        "slides": [
            {
                "index": slide_index,
                "records": [
                    {"className": "drawio-box", "rect": {"x": 0, "y": 0, "w": 0, "h": 0}},
                    {"className": "drawio", "rect": {"x": 60, "y": 120, "w": 400, "h": 240}},
                ],
            }
        ]
    }


def _write(tmp_path):
    html = tmp_path / "deck.html"
    html.write_text(HTML, encoding="utf-8")
    (tmp_path / "arch.drawio").write_text(DRAWIO, encoding="utf-8")
    meas = tmp_path / "deck_audit" / "_cache" / "measurements.json"
    meas.parent.mkdir(parents=True)
    meas.write_text(json.dumps(_measurements()), encoding="utf-8")
    return str(html)


def _write_multi(tmp_path, page_attr=None):
    html = tmp_path / "deck.html"
    attr = 'data-drawio="arch.drawio"'
    if page_attr is not None:
        attr += f' data-drawio-page="{page_attr}"'
    html.write_text(HTML.replace('data-drawio="arch.drawio"', attr), encoding="utf-8")
    (tmp_path / "arch.drawio").write_text(MULTI_DRAWIO, encoding="utf-8")
    meas = tmp_path / "deck_audit" / "_cache" / "measurements.json"
    meas.parent.mkdir(parents=True)
    meas.write_text(json.dumps(_measurements()), encoding="utf-8")
    return str(html)


def test_parse_drawio_declarations():
    decls = parse_drawio_declarations(HTML)
    assert decls == [{"slide": 1, "path": "arch.drawio", "page": None}]


def test_parse_drawio_declarations_page_attr():
    html = HTML.replace(
        'data-drawio="arch.drawio"',
        'data-drawio="arch.drawio" data-drawio-page="2"',
    )
    decls = parse_drawio_declarations(html)
    assert decls == [{"slide": 1, "path": "arch.drawio", "page": "2"}]


def test_parse_drawio_declarations_skips_non_drawio():
    html = HTML.replace('class="drawio"', 'class="code"')
    assert parse_drawio_declarations(html) == []


def test_parse_drawio_declarations_multi_class_tokens():
    html = HTML.replace('class="drawio"', 'class="drawio chart"')
    decls = parse_drawio_declarations(html)
    assert len(decls) == 1
    assert decls[0]["path"] == "arch.drawio"


def test_parse_drawio_declarations_rejects_outside_section():
    html = '<html><body><div class="drawio" data-drawio="a.drawio"></div></body></html>'
    with pytest.raises(ValueError, match="slide 之外"):
        parse_drawio_declarations(html)


def test_parse_drawio_declarations_rejects_multiple_per_slide():
    html = HTML.replace(
        '<div class="drawio" data-drawio="arch.drawio"></div>',
        '<div class="drawio" data-drawio="arch.drawio"></div>\n'
        '    <div class="drawio" data-drawio="other.drawio"></div>',
    )
    with pytest.raises(ValueError, match="每页仅支持一个 drawio"):
        parse_drawio_declarations(html)


def test_postprocess_drawio_injects_shapes(tmp_path):
    html = _write(tmp_path)
    # 模拟 convert 产物：drawio 占位容器渲染成占位文本框（bbox px 60,120,400,240 → inch
    # 0.42,0.83,2.78,1.67，中心落在 bbox 内）。拿 drawio_to_pptx 产物当占位会让
    # "postprocess 没注入也通过" 的假阳性逃过断言，故手摆占位文本框。
    pptx_path = str(tmp_path / "deck.pptx")
    prs = Presentation()
    prs.slide_width = Emu(12192000)
    prs.slide_height = Emu(6858000)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(0.42), Inches(0.83), Inches(2.78), Inches(1.67))
    tb.text = "占位"
    prs.save(pptx_path)

    postprocess_drawio(html, pptx_path)
    prs = Presentation(pptx_path)
    auto = [
        sh for sl in prs.slides for sh in sl.shapes if sh.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
    ]
    texts = {sh.text_frame.text for sh in auto}
    assert "模块A" in texts  # 注入的可编辑形状出现
    assert "模块B" in texts
    # 占位文本被移除——证明 postprocess 真的替换了
    assert not any(
        sh.has_text_frame and "占位" in sh.text_frame.text for sh in prs.slides[0].shapes
    )


def test_postprocess_drawio_multi_page_requires_page(tmp_path):
    # #99：多页 .drawio 未指定 data-drawio-page → 明确报错，不静默取第一页
    html = _write_multi(tmp_path)
    pptx_path = str(tmp_path / "deck.pptx")
    prs = Presentation()
    prs.slide_width = Emu(12192000)
    prs.slide_height = Emu(6858000)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(0.42), Inches(0.83), Inches(2.78), Inches(1.67))
    tb.text = "占位"
    prs.save(pptx_path)

    with pytest.raises(ValueError, match="data-drawio-page"):
        postprocess_drawio(html, pptx_path)


def test_postprocess_drawio_selects_page(tmp_path):
    # #99：data-drawio-page="2"（1 基）→ 注入第二页内容
    html = _write_multi(tmp_path, page_attr="2")
    pptx_path = str(tmp_path / "deck.pptx")
    prs = Presentation()
    prs.slide_width = Emu(12192000)
    prs.slide_height = Emu(6858000)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(0.42), Inches(0.83), Inches(2.78), Inches(1.67))
    tb.text = "占位"
    prs.save(pptx_path)

    postprocess_drawio(html, pptx_path)
    prs = Presentation(pptx_path)
    texts = {sh.text_frame.text for sh in prs.slides[0].shapes if sh.has_text_frame}
    assert "第二页节点" in texts
    assert "第一页节点" not in texts
    assert "占位" not in texts


def test_postprocess_drawio_selects_first_page(tmp_path):
    # #99：data-drawio-page="1" → 注入第一页内容
    html = _write_multi(tmp_path, page_attr="1")
    pptx_path = str(tmp_path / "deck.pptx")
    prs = Presentation()
    prs.slide_width = Emu(12192000)
    prs.slide_height = Emu(6858000)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(0.42), Inches(0.83), Inches(2.78), Inches(1.67))
    tb.text = "占位"
    prs.save(pptx_path)

    postprocess_drawio(html, pptx_path)
    prs = Presentation(pptx_path)
    texts = {sh.text_frame.text for sh in prs.slides[0].shapes if sh.has_text_frame}
    assert "第一页节点" in texts
    assert "第二页节点" not in texts
    assert "占位" not in texts


def test_postprocess_drawio_missing_measurements(tmp_path):
    html = tmp_path / "deck.html"
    html.write_text(HTML, encoding="utf-8")
    (tmp_path / "arch.drawio").write_text(DRAWIO, encoding="utf-8")
    # 不写 measurements.json
    with pytest.raises(RuntimeError, match="measurements.json"):
        postprocess_drawio(str(html), str(tmp_path / "deck.pptx"))


def test_postprocess_drawio_missing_source(tmp_path):
    html = _write(tmp_path)
    (tmp_path / "arch.drawio").unlink()
    pptx_path = str(tmp_path / "deck.pptx")
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[6])
    prs.save(pptx_path)
    with pytest.raises(RuntimeError, match="arch.drawio"):
        postprocess_drawio(html, pptx_path)


def test_postprocess_drawio_zero_height_box(tmp_path):
    # #94 守卫：measurements 里容器 h=0 → 可操作报错，而非零缩放静默坏渲染
    html = _write(tmp_path)
    meas = tmp_path / "deck_audit" / "_cache" / "measurements.json"
    m = json.loads(meas.read_text(encoding="utf-8"))
    m["slides"][0]["records"][1]["rect"]["h"] = 0
    meas.write_text(json.dumps(m), encoding="utf-8")
    pptx_path = str(tmp_path / "deck.pptx")
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[6])
    prs.save(pptx_path)
    with pytest.raises(RuntimeError, match="尺寸为 0"):
        postprocess_drawio(html, pptx_path)


def test_postprocess_drawio_accepts_file_uri(tmp_path):
    # #93：data-drawio 直接用 file:// URI（deck 管线把相对路径改写成的产物）
    # → 不再被当相对路径拼到 base，注入成功
    src = tmp_path / "arch.drawio"
    src.write_text(DRAWIO, encoding="utf-8")
    html = tmp_path / "deck.html"
    html.write_text(
        HTML.replace('data-drawio="arch.drawio"', f'data-drawio="{src.as_uri()}"'),
        encoding="utf-8",
    )
    meas = tmp_path / "deck_audit" / "_cache" / "measurements.json"
    meas.parent.mkdir(parents=True)
    meas.write_text(json.dumps(_measurements()), encoding="utf-8")
    pptx_path = str(tmp_path / "deck.pptx")
    prs = Presentation()
    prs.slide_width = Emu(12192000)
    prs.slide_height = Emu(6858000)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(0.42), Inches(0.83), Inches(2.78), Inches(1.67))
    tb.text = "占位"
    prs.save(pptx_path)

    postprocess_drawio(str(html), pptx_path)  # 不得抛「源文件缺失」
    prs = Presentation(pptx_path)
    texts = {sh.text_frame.text for sl in prs.slides for sh in sl.shapes if sh.has_text_frame}
    assert "模块A" in texts and "模块B" in texts


def test_postprocess_drawio_preserves_non_placeholder_shape_inside_box(tmp_path):
    # #96：bbox 内不同尺寸的用户形状不被误删（旧「中心点在内」判定会删掉它）
    html = _write(tmp_path)
    pptx_path = str(tmp_path / "deck.pptx")
    prs = Presentation()
    prs.slide_width = Emu(12192000)
    prs.slide_height = Emu(6858000)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(0.42), Inches(0.83), Inches(2.78), Inches(1.67))
    tb.text = "占位"
    user = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(1.0), Inches(0.5))
    user.text = "用户形状"  # 中心 (1.0, 1.25) 落在 bbox 内，但尺寸与注入矩形不同
    prs.save(pptx_path)

    postprocess_drawio(html, pptx_path)
    prs = Presentation(pptx_path)
    texts = {sh.text_frame.text for sh in prs.slides[0].shapes if sh.has_text_frame}
    assert "模块A" in texts and "模块B" in texts  # 注入的可编辑形状出现
    assert "用户形状" in texts  # bbox 内用户形状保留
    assert "占位" not in texts  # 占位被替换


def test_no_visual_audit_rejects_drawio_before_browser(tmp_path, monkeypatch):
    html = tmp_path / "deck.html"
    html.write_text(HTML, encoding="utf-8")
    from offipy import deck

    browser_calls: list = []
    monkeypatch.setattr(deck, "_preflight_browser", lambda *a, **k: browser_calls.append(a))
    monkeypatch.setattr(deck, "_run_convert", lambda *a, **k: None)

    with pytest.raises(deck.InvalidArgumentError, match="drawio"):
        deck.render(
            str(html), out=str(tmp_path / "deck.pptx"), overwrite=True, no_visual_audit=True
        )
    assert browser_calls == []  # chromium 从未启动
