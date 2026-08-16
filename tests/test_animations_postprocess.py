"""postprocess_animations：读 measurements.json → OFFIPY_ELEM::<elem_id> 定位 → 注入。

集成测试用构造的 measurements.json + 预打名 shape 的 pptx，验证 <p:timing>
落位、多形状单元全命中、p:sld 顺序、报告计数、未命中告警。
"""

import json

import pytest

from offipy.animations.apply import postprocess_animations


def _qn(tag):
    return tag.split("}")[-1]


def _measurements_path(pptx):
    from pathlib import Path
    p = Path(pptx)
    return p.with_name(f"{p.stem}_audit") / "_cache" / "measurements.json"


def _build_fixture(tmp_path):
    """构造 1 页 pptx（title shape 名为 OFFIPY_ELEM::1.0，卡片拆成 fill+2 边框线同 elem_id）
    + measurements.json（record 带 anim_decl + elem_id）。返回 (html, pptx)。"""
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    sp = slide.shapes.add_shape(1, Inches(1), Inches(1), Inches(3), Inches(0.8))
    sp.name = "OFFIPY_ELEM::1.0"
    line1 = slide.shapes.add_shape(1, Inches(1), Inches(1.9), Inches(3), Inches(0.05))
    line1.name = "OFFIPY_ELEM::1.0"
    line2 = slide.shapes.add_shape(1, Inches(1), Inches(1.95), Inches(3), Inches(0.05))
    line2.name = "OFFIPY_ELEM::1.0"

    html = tmp_path / "t.audited.html"
    html.write_text("<section data-pptx-slide>card</section>", encoding="utf-8")
    pptx = tmp_path / "t.pptx"
    prs.save(str(pptx))

    meas_path = _measurements_path(str(pptx))
    meas_path.parent.mkdir(parents=True, exist_ok=True)
    meas = {
        "_total": 1,
        "slides": [
            {
                "slide": {"background": "rgb(255,255,255)"},
                "records": [
                    {"id": 1, "kind": "shape", "className": "",
                     "anim_decl": {"anim": "fade"}, "elem_id": "1.0",
                     "rect": {"x": 0, "y": 0, "w": 10, "h": 10}},
                ],
            }
        ],
    }
    meas_path.write_text(json.dumps(meas), encoding="utf-8")
    return str(html), str(pptx), meas_path


def test_postprocess_injects_timing_and_counts(tmp_path):
    html, pptx, _ = _build_fixture(tmp_path)
    report = postprocess_animations(html, pptx)
    assert report["animations_applied"] == 1
    assert report["unmatched"] == []
    from pptx import Presentation
    sld = Presentation(pptx).slides[0]
    timings = sld._element.findall(".//{*}timing")
    assert len(timings) == 1


def test_postprocess_multi_shape_unit_all_shapes_animated(tmp_path):
    html, pptx, _ = _build_fixture(tmp_path)
    postprocess_animations(html, pptx)
    from pptx import Presentation
    sld = Presentation(pptx).slides[0]
    spids = [t.get("spid") for t in sld._element.findall(".//{*}spTgt")]
    # fill + 2 边框线 = 3 个形状，各一组 set+effect → 6 个 spTgt
    assert len(spids) == 6


def test_postprocess_sld_child_order(tmp_path):
    html, pptx, _ = _build_fixture(tmp_path)
    postprocess_animations(html, pptx)
    from pptx import Presentation
    sld = Presentation(pptx).slides[0]
    names = [_qn(c.tag) for c in sld._element]
    assert names.index("cSld") < names.index("timing")


def test_postprocess_transition_injected_and_before_timing(tmp_path):
    html, pptx, meas_path = _build_fixture(tmp_path)
    # 给 slide 元数据加 transition_decl
    data = json.loads(meas_path.read_text(encoding="utf-8"))
    data["slides"][0]["slide"]["transition_decl"] = {"kind": "push", "speed": "medium"}
    meas_path.write_text(json.dumps(data), encoding="utf-8")
    report = postprocess_animations(html, pptx)
    assert report["transitions_applied"] == 1
    from pptx import Presentation
    sld = Presentation(pptx).slides[0]
    names = [_qn(c.tag) for c in sld._element]
    assert names.index("transition") < names.index("timing")


def test_postprocess_unmatched_elem_records_warning_not_fail(tmp_path):
    html, pptx, meas_path = _build_fixture(tmp_path)
    data = json.loads(meas_path.read_text(encoding="utf-8"))
    # 加一个 anim_decl 但 assemble 没产出对应形状（无 OFFIPY_ELEM::9.9 shape）
    data["slides"][0]["records"].append(
        {"id": 2, "kind": "shape", "className": "",
         "anim_decl": {"anim": "fade"}, "elem_id": "9.9",
         "rect": {"x": 0, "y": 0, "w": 5, "h": 5}},
    )
    meas_path.write_text(json.dumps(data), encoding="utf-8")
    report = postprocess_animations(html, pptx)
    assert report["unmatched"] == [{"elem_id": "9.9", "slide": 1}]
    # 告警追加进 measurements.json 的 _warnings
    updated = json.loads(meas_path.read_text(encoding="utf-8"))
    assert any("动画" in str(w) for w in updated.get("_warnings", []))


def test_postprocess_no_animations_noop(tmp_path):
    html, pptx, _ = _build_fixture(tmp_path)
    # 清掉 anim_decl
    from pathlib import Path
    meas_path = _measurements_path(pptx)
    data = json.loads(meas_path.read_text(encoding="utf-8"))
    for rec in data["slides"][0]["records"]:
        rec.pop("anim_decl", None)
        rec.pop("elem_id", None)
    meas_path.write_text(json.dumps(data), encoding="utf-8")
    report = postprocess_animations(html, pptx)
    assert report["animations_applied"] == 0
    assert report["transitions_applied"] == 0
    from pptx import Presentation
    sld = Presentation(pptx).slides[0]
    assert not sld._element.findall(".//{*}timing")


def test_postprocess_missing_measurements_noop(tmp_path):
    html = tmp_path / "x.html"
    html.write_text("x", encoding="utf-8")
    pptx = tmp_path / "x.pptx"
    from pptx import Presentation
    Presentation().save(str(pptx))
    report = postprocess_animations(str(html), str(pptx))
    assert report["animations_applied"] == 0


def test_postprocess_invalid_decl_warning_persisted_even_without_injection(tmp_path):
    # 唯一动画声明非法（未知效果 flip）→ parse None 产 spec 告警；无任何注入，
    # 但告警必须落进 measurements.json 的 _warnings（P1 dict 不崩 + P2 不丢）。
    html, pptx, meas_path = _build_fixture(tmp_path)
    data = json.loads(meas_path.read_text(encoding="utf-8"))
    data["slides"][0]["records"][0]["anim_decl"] = {"anim": "flip"}
    meas_path.write_text(json.dumps(data), encoding="utf-8")
    report = postprocess_animations(html, pptx)
    assert report["animations_applied"] == 0
    assert report["transitions_applied"] == 0
    updated = json.loads(meas_path.read_text(encoding="utf-8"))
    ws = updated.get("_warnings", [])
    assert ws and any(isinstance(w, dict) and w.get("kind") == "anim" for w in ws)
    assert any("flip" in str(w) for w in ws)


def test_postprocess_invalid_transition_speed_warns_not_raises(tmp_path):
    # 过渡声明 speed 非法（fastest）→ 不硬失败（kind 合法但 speed 越界），
    # 记告警跳过、不构造 TransitionSpec；measurements.json _warnings 提到该 speed。
    html, pptx, meas_path = _build_fixture(tmp_path)
    data = json.loads(meas_path.read_text(encoding="utf-8"))
    data["slides"][0]["slide"]["transition_decl"] = {"kind": "push", "speed": "fastest"}
    meas_path.write_text(json.dumps(data), encoding="utf-8")
    report = postprocess_animations(html, pptx)
    assert report["transitions_applied"] == 0
    updated = json.loads(meas_path.read_text(encoding="utf-8"))
    ws = updated.get("_warnings", [])
    assert ws and any(isinstance(w, dict) and w.get("kind") == "anim" for w in ws)
    assert any("fastest" in str(w) for w in ws)
