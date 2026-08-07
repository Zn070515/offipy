"""A4 Tasks 12–14 — procedural patterns through the deck asset pipeline.

HTML declarations (data-asset / data-asset-param-* / data-asset-placement)
flow into the real deck.render → postprocess_assets chain: measurements drive
theme materialization (not THEMES.base_vars), each declaration becomes exactly
one P2 svgBlip picture, background placement lands below content while
decorative stays in place, and provenance lands in the audit assets.json.
"""

import json
import subprocess as _sp
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE

from offipy import deck
from offipy.assets.declarations import preprocess_asset_declarations
from offipy.assets.materialize import materialize_svg_template
from offipy.assets.model import AssetRef, AssetRequest
from offipy.assets.providers.procedural import ProceduralProvider
from offipy.envcheck import _check_browser
from offipy.exceptions import InvalidArgumentError

_DEFAULT_VARS = {"bg": "#ffffff", "surface": "#f3f4f6", "accent": "#2251ff"}

_HTML_RINGS = """<!doctype html>
<html><head></head><body>
<section data-pptx-slide>
  <div data-asset="asset://procedural/pattern/rings?count=7"></div>
</section>
</body></html>"""

_HTML_TOPO_BG = """<!doctype html>
<html><head></head><body>
<section data-pptx-slide>
  <div data-asset="asset://procedural/pattern/topography?seed=42"
       data-asset-placement="background"></div>
</section>
</body></html>"""

_HTML_PARAM_FORM = """<!doctype html>
<html><head></head><body>
<section data-pptx-slide>
  <div data-asset="asset://procedural/pattern/wave"
       data-asset-param-seed="4" data-asset-param-foreground="accent-soft"></div>
</section>
</body></html>"""

_HTML_CONFLICT = """<!doctype html>
<html><head></head><body>
<section data-pptx-slide>
  <div data-asset="asset://procedural/pattern/wave?seed=4" data-asset-param-seed="5"></div>
</section>
</body></html>"""

_HTML_HEX_FG = """<!doctype html>
<html><head></head><body>
<section data-pptx-slide>
  <div data-asset="asset://procedural/pattern/rings" data-asset-param-foreground="#00FF00"></div>
</section>
</body></html>"""

_HTML_MISSING_TOKEN = """<!doctype html>
<html><head></head><body>
<section data-pptx-slide>
  <div data-asset="asset://procedural/pattern/wave" data-asset-param-foreground="muted"></div>
</section>
</body></html>"""

_HTML_ZORDER = """<!doctype html>
<html><head></head><body>
<section data-pptx-slide>
  <div data-asset="asset://procedural/pattern/rings?count=4"></div>
  <div data-asset="asset://procedural/pattern/wave?seed=9" data-asset-placement="background"></div>
</section>
</body></html>"""

_HTML_PROVENANCE = """<!doctype html>
<html><head></head><body>
<section data-pptx-slide><div data-asset="asset://procedural/pattern/wave"></div></section>
<section data-pptx-slide><div data-asset="asset://procedural/pattern/wave"></div></section>
<section data-pptx-slide><div data-asset="asset://procedural/pattern/rings?count=3"></div></section>
</body></html>"""


@pytest.fixture(autouse=True)
def _no_browser(monkeypatch):
    monkeypatch.setattr(deck, "_preflight_browser", lambda: None)


def _asset_meas(asset_id, x=20, y=30, w=120, h=60, color="rgb(255, 0, 0)", theme_vars=None):
    return {
        "id": 1,
        "kind": "asset",
        "tag": "div",
        "assetId": asset_id,
        "rect": {"x": x, "y": y, "w": w, "h": h},
        "themeVars": theme_vars if theme_vars is not None else _DEFAULT_VARS,
        "color": color,
    }


def _make_fake_convert(decls, *, meas_override=None):
    """fake subprocess.run：造带 OFFIPY_ASSET:: 占位符的 pptx + measurements.json。"""

    def fake_run(cmd, *a, **k):
        out = cmd[cmd.index("--out") + 1]
        prs = Presentation()
        for d in decls:
            while len(prs.slides) < d.slide_index:
                prs.slides.add_slide(prs.slide_layouts[6])
            slide = prs.slides[d.slide_index - 1]
            sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, 100, 100)
            sp.name = f"OFFIPY_ASSET::{d.declaration_id}"
        prs.save(out)
        max_slide = max((d.slide_index for d in decls), default=0)
        slides = [{"records": []} for _ in range(max_slide)]
        for d in decls:
            rec = _asset_meas(d.declaration_id)
            if meas_override:
                rec = meas_override(rec, d.declaration_id)
            slides[d.slide_index - 1]["records"].append(rec)
        audit = Path(out).with_name(f"{Path(out).stem}_audit")
        meas_dir = audit / "_cache"
        meas_dir.mkdir(parents=True)
        (meas_dir / "measurements.json").write_text(
            json.dumps({"slides": slides}, ensure_ascii=False), encoding="utf-8"
        )
        return _sp.CompletedProcess(args=[], returncode=0, stdout="", stderr="")

    return fake_run


def _render_deck(tmp_path, html_text, *, theme=None):
    html = tmp_path / "d.html"
    html.write_text(html_text, encoding="utf-8")
    pptx = tmp_path / "d.pptx"
    pptx.write_bytes(b"placeholder")
    return deck.render(str(html), out=str(pptx), overwrite=True, theme=theme), pptx


def _pic_blobs(slide):
    """spTree 顺序下每个 p:pic 关联的 SVG 文本（P2 svgBlip 媒体部件）。"""
    from pptx.oxml.ns import qn

    r = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    asvg = "http://schemas.microsoft.com/office/drawing/2016/SVG/main"
    blobs: list[str] = []
    for el in slide.shapes._spTree:
        if el.tag != qn("p:pic"):
            continue
        svg_blip = el.find(f".//{{{asvg}}}svgBlip")
        rid = svg_blip.get(f"{{{r}}}embed")
        blobs.append(slide.part.related_part(rid).blob.decode("utf-8"))
    return blobs


def _shape_names_in_order(slide):
    """spTree 里 p:pic / p:sp 的 (kind, cNvPr name) 顺序。"""
    from pptx.oxml.ns import qn

    out: list[tuple[str, str]] = []
    for el in slide.shapes._spTree:
        if el.tag not in (qn("p:pic"), qn("p:sp")):
            continue
        cNvPr = next(c for c in el.iter() if c.tag.endswith("cNvPr"))
        kind = "pic" if el.tag == qn("p:pic") else "sp"
        out.append((kind, cNvPr.get("name", "")))
    return out


def _load_manifest(tmp_path):
    p = tmp_path / "d_audit" / "assets.json"
    if not p.exists():
        return None
    return json.loads(p.read_text(encoding="utf-8"))


# ---------------------------------------------------------------------------
# Task 12 — HTML integration / placement / theme materialization
# ---------------------------------------------------------------------------


def test_procedural_background_renders_single_p2_picture(tmp_path, monkeypatch):
    _, decls = preprocess_asset_declarations(_HTML_TOPO_BG)
    monkeypatch.setattr(deck.subprocess, "run", _make_fake_convert(decls))
    _render_deck(tmp_path, _HTML_TOPO_BG)

    slide = Presentation(str(tmp_path / "d.pptx")).slides[0]
    blobs = _pic_blobs(slide)
    assert len(blobs) == 1  # 一个声明 → 恰好一个 P2 picture
    svg = blobs[0]
    assert "__OFFIPY_ASSET" not in svg  # 模板已物化
    assert "#2251FF" in svg  # 测量 themeVars accent 进入前景
    order = _shape_names_in_order(slide)
    assert order == [("pic", "SVG Asset")]  # background 置于 grpSpPr 后、内容之前


def test_rings_count_flows_into_materialized_svg(tmp_path, monkeypatch):
    _, decls = preprocess_asset_declarations(_HTML_RINGS)
    monkeypatch.setattr(deck.subprocess, "run", _make_fake_convert(decls))
    _render_deck(tmp_path, _HTML_RINGS)

    svg = _pic_blobs(Presentation(str(tmp_path / "d.pptx")).slides[0])[0]
    assert svg.count("<circle") == 7
    assert 'stroke="#2251FF"' in svg


def test_param_attr_form_matches_direct_resolve(tmp_path, monkeypatch):
    theme_vars = {**_DEFAULT_VARS, "accent-soft": "#E9EDFF"}
    _, decls = preprocess_asset_declarations(_HTML_PARAM_FORM)
    override = lambda rec, aid: {**rec, "themeVars": theme_vars}  # noqa: E731
    monkeypatch.setattr(deck.subprocess, "run", _make_fake_convert(decls, meas_override=override))
    _render_deck(tmp_path, _HTML_PARAM_FORM)

    svg = _pic_blobs(Presentation(str(tmp_path / "d.pptx")).slides[0])[0]
    assert 'stroke="#E9EDFF"' in svg  # foreground=accent-soft 已解析
    resolved = ProceduralProvider().resolve(
        AssetRequest(
            AssetRef("procedural", "pattern", "wave"),
            (("seed", "4"), ("foreground", "accent-soft")),
        )
    )
    expected = materialize_svg_template(resolved.payload, theme_vars).svg
    assert svg == expected  # attr 参数 → URI 等价 resolve → 几何逐字节一致


def test_uri_and_attr_param_conflict_rejected_before_convert(tmp_path, monkeypatch):
    html = tmp_path / "d.html"
    html.write_text(_HTML_CONFLICT, encoding="utf-8")
    pptx = tmp_path / "d.pptx"
    pptx.write_bytes(b"placeholder")
    original = pptx.read_bytes()

    def _never(cmd, *a, **k):
        raise AssertionError("convert 不应被调用")

    monkeypatch.setattr(deck.subprocess, "run", _never)
    with pytest.raises(InvalidArgumentError, match="duplicate"):
        deck.render(str(html), out=str(pptx), overwrite=True)

    assert pptx.read_bytes() == original  # 旧最终输出原样保留
    assert not list(tmp_path.glob("offipy-deck-*"))  # 无残留注入副本


def test_explicit_hex_foreground_overrides_theme_token(tmp_path, monkeypatch):
    _, decls = preprocess_asset_declarations(_HTML_HEX_FG)
    monkeypatch.setattr(deck.subprocess, "run", _make_fake_convert(decls))
    _render_deck(tmp_path, _HTML_HEX_FG)

    svg = _pic_blobs(Presentation(str(tmp_path / "d.pptx")).slides[0])[0]
    assert 'stroke="#00FF00"' in svg
    assert "#2251FF" not in svg  # 显式 hex 胜过主题 token


def test_missing_theme_token_fails_postprocess(tmp_path, monkeypatch):
    theme_vars = _DEFAULT_VARS  # 不含 muted
    _, decls = preprocess_asset_declarations(_HTML_MISSING_TOKEN)
    override = lambda rec, aid: {**rec, "themeVars": theme_vars}  # noqa: E731
    monkeypatch.setattr(deck.subprocess, "run", _make_fake_convert(decls, meas_override=override))

    with pytest.raises(InvalidArgumentError, match="muted"):
        _render_deck(tmp_path, _HTML_MISSING_TOKEN)


@pytest.mark.parametrize(
    ("theme_vars", "expected"),
    [
        (_DEFAULT_VARS, "#2251FF"),  # mckinsey base accent
        ({"bg": "#051C2C", "surface": "#0C1C23", "accent": "#5B8CFF"}, "#5B8CFF"),  # dark variant
        ({"bg": "#ffffff", "surface": "#ffffff", "accent": "#FF00AA"}, "#FF00AA"),  # 自定义
        ({"bg": "#ffffff", "surface": "#ffffff", "--accent": "#FF00AA"}, "#FF00AA"),  # --key 形式
    ],
)
def test_measurement_theme_vars_drive_materialization(tmp_path, monkeypatch, theme_vars, expected):
    """deck 主题参数只注入 CSS；materialization 以测量 themeVars 为准，不落回 THEMES.base_vars。"""
    _, decls = preprocess_asset_declarations(_HTML_RINGS)
    override = lambda rec, aid: {**rec, "themeVars": theme_vars}  # noqa: E731
    monkeypatch.setattr(deck.subprocess, "run", _make_fake_convert(decls, meas_override=override))
    _render_deck(tmp_path, _HTML_RINGS, theme="mckinsey")

    svg = _pic_blobs(Presentation(str(tmp_path / "d.pptx")).slides[0])[0]
    assert f'stroke="{expected}"' in svg


def test_zero_size_rect_fails_postprocess(tmp_path, monkeypatch):
    """CSS 量到 0×60 的 rect → 绑定阶段 InvalidArgumentError，deck 报错不产出。"""
    _, decls = preprocess_asset_declarations(_HTML_RINGS)
    override = lambda rec, aid: {**rec, "rect": {"x": 20, "y": 30, "w": 0, "h": 60}}  # noqa: E731
    monkeypatch.setattr(deck.subprocess, "run", _make_fake_convert(decls, meas_override=override))

    with pytest.raises(InvalidArgumentError, match="positive"):
        _render_deck(tmp_path, _HTML_RINGS)


# ---------------------------------------------------------------------------
# Task 13 — z-order 结构：背景沉底、装饰居中、占位符全消
# ---------------------------------------------------------------------------


def _zorder_fake_convert():
    def fake_run(cmd, *a, **k):
        out = cmd[cmd.index("--out") + 1]
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[6])
        slide = prs.slides[0]
        for name, decl in [
            ("A", None),
            (None, "asset-s01-001"),
            ("B", None),
            (None, "asset-s01-002"),
        ]:
            sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, 100, 100)
            sp.name = name if decl is None else f"OFFIPY_ASSET::{decl}"
        prs.save(out)
        audit = Path(out).with_name(f"{Path(out).stem}_audit")
        meas_dir = audit / "_cache"
        meas_dir.mkdir(parents=True)
        (meas_dir / "measurements.json").write_text(
            json.dumps(
                {
                    "slides": [
                        {
                            "records": [
                                _asset_meas("asset-s01-001"),
                                _asset_meas("asset-s01-002"),
                            ]
                        }
                    ]
                },
                ensure_ascii=False,
            ),
            encoding="utf-8",
        )
        return _sp.CompletedProcess(args=[], returncode=0, stdout="", stderr="")

    return fake_run


def test_zorder_background_below_decorative_between(tmp_path, monkeypatch):
    monkeypatch.setattr(deck.subprocess, "run", _zorder_fake_convert())
    _render_deck(tmp_path, _HTML_ZORDER)

    slide = Presentation(str(tmp_path / "d.pptx")).slides[0]
    assert not any(s.name.startswith("OFFIPY_ASSET::") for s in slide.shapes)
    order = _shape_names_in_order(slide)
    assert order == [
        ("pic", "SVG Asset"),  # 背景 wave：A 之前（grpSpPr 之后）
        ("sp", "A"),
        ("pic", "SVG Asset"),  # 装饰 rings：A 与 B 之间
        ("sp", "B"),
    ]
    blobs = _pic_blobs(slide)
    assert len(blobs) == 2  # 两个声明各恰好一个 P2 picture
    assert blobs[0].count("<path") > 0 and "<circle" not in blobs[0]  # wave 背景
    assert blobs[1].count("<circle") == 4  # rings?count=4 装饰


# ---------------------------------------------------------------------------
# Task 14 — assets.json 来源清单（procedural 元数据）
# ---------------------------------------------------------------------------


def test_assets_json_procedural_provenance(tmp_path, monkeypatch):
    _, decls = preprocess_asset_declarations(_HTML_PROVENANCE)
    monkeypatch.setattr(deck.subprocess, "run", _make_fake_convert(decls))
    _render_deck(tmp_path, _HTML_PROVENANCE)

    manifest = _load_manifest(tmp_path)
    assert manifest["schema"] == 1
    assets = manifest["assets"]
    assert [a["declaration_id"] for a in assets] == [
        "asset-s01-001",
        "asset-s02-001",
        "asset-s03-001",
    ]
    assert [a["slide_index"] for a in assets] == [1, 2, 3]
    assert assets[0]["request"] == "asset://procedural/pattern/wave"
    assert assets[1]["request"] == "asset://procedural/pattern/wave"  # 同一 ref 二次声明
    assert assets[2]["request"] == "asset://procedural/pattern/rings?count=3"
    assert all(a["placement"] == "replace" for a in assets)
    for a in assets:
        p = a["provider"]
        assert p["id"] == "procedural"
        assert p["license"] == "MIT"
        assert p["first_party"] is True
        assert p["redistributable"] is True
        assert p["source_url"] == "https://github.com/Zn070515/offipy"
        assert p["source_commit"] is None  # 不捏造 commit
        assert p["attribution"] is None


# ---------------------------------------------------------------------------
# Task 16 — 50 个 procedural 资产：确定性 + 占位符全消费
# ---------------------------------------------------------------------------


def _many_decl_html(n):
    divs = "\n".join(
        f'<div data-asset="asset://procedural/pattern/rings?count={1 + (i % 12)}"></div>'
        for i in range(n)
    )
    return (
        "<!doctype html>\n<html><head></head><body>\n"
        f"<section data-pptx-slide>\n{divs}\n</section>\n</body></html>"
    )


def test_50_procedural_assets_deterministic(tmp_path, monkeypatch):
    html_text = _many_decl_html(50)
    _, decls = preprocess_asset_declarations(html_text)
    assert len(decls) == 50
    monkeypatch.setattr(deck.subprocess, "run", _make_fake_convert(decls))
    _render_deck(tmp_path, html_text)

    slide = Presentation(str(tmp_path / "d.pptx")).slides[0]
    assert not any(s.name.startswith("OFFIPY_ASSET::") for s in slide.shapes)
    blobs = _pic_blobs(slide)
    assert len(blobs) == 50
    assert all("__OFFIPY_ASSET" not in b for b in blobs)
    assert sorted({b.count("<circle") for b in blobs}) == list(range(1, 13))

    # 同一输入再渲一次 → 几何与顺序逐字节一致
    out2 = tmp_path / "d2.pptx"
    out2.write_bytes(b"placeholder")
    deck.render(str(tmp_path / "d.html"), out=str(out2), overwrite=True)
    blobs2 = _pic_blobs(Presentation(str(out2)).slides[0])
    assert blobs2 == blobs


# ---------------------------------------------------------------------------
# Task 11 — 真实 Chromium 渲染（deck_render 门禁，CI 无浏览器则 skip）
# ---------------------------------------------------------------------------

_HTML_REAL = """<!doctype html>
<html><head><style>
  :root { --bg:#ffffff; --surface:#f3f4f6; --ink:#222222; --muted:#667085;
          --accent:#2251ff; --accent-soft:#e9edff; --divider:#e5e7eb; }
</style></head><body>
<section data-pptx-slide>
  <h1>Probe</h1>
  <div data-asset="asset://procedural/pattern/topography?seed=42"
       data-asset-placement="background"
       style="position:absolute;left:0;top:0;width:100%;height:100%"></div>
  <div data-asset="asset://procedural/pattern/rings?count=7"
       style="position:absolute;left:120px;top:160px;width:500px;height:500px"></div>
</section>
</body></html>"""


@pytest.mark.deck_render
@pytest.mark.skipif(
    not deck.CONVERT_PY.exists() or not _check_browser().ok,
    reason="vendored 转换器缺失或 chromium 不可用",
)
def test_real_render_procedural_background_and_decorative(tmp_path):
    """真实浏览器渲染：procedural 资产按声明进 P2 图片，占位符全消。"""
    html = tmp_path / "real.html"
    html.write_text(_HTML_REAL, encoding="utf-8")
    out = tmp_path / "real.pptx"
    deck.render(str(html), out=str(out), overwrite=True, timeout=240)

    slide = Presentation(str(out)).slides[0]
    assert not any(s.name.startswith("OFFIPY_ASSET::") for s in slide.shapes)
    blobs = _pic_blobs(slide)
    assert len(blobs) == 2
    assert all("__OFFIPY_ASSET" not in b for b in blobs)
    assert all("#2251FF" in b for b in blobs)  # 主题 accent 物化
    # 背景 topography 在装饰 rings 之前（spTree 顺序）
    assert "<circle" not in blobs[0] and blobs[0].count("<path") > 0
    assert blobs[1].count("<circle") == 7
