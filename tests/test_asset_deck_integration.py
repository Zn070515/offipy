"""deck 资产管线集成：目标注入 / no_visual_audit 前置 / only_slides 一致性。

A3 Task 3：asset 声明（data-asset/data-primitive/legacy data-icon）与图表一样
需要注入副本 target（converter 量到注入后的确定性 ID），且 no_visual_audit 的
前置检查要在 chromium / convert 子进程之前拦下非法组合。单测 monkeypatch 浏览器
与转换子进程，任何环境可跑。
"""

import json
import subprocess as _sp
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE

from offipy import charts, deck
from offipy.assets.declarations import preprocess_asset_declarations


@pytest.fixture(autouse=True)
def _no_browser(monkeypatch):
    # 单测不真启动 chromium：render 的浏览器前置检查换成 no-op（重测 preflight 顺序时覆盖）
    monkeypatch.setattr(deck, "_preflight_browser", lambda: None)


def _fake_run(recorded: dict):
    """mock subprocess.run：按 cmd 的 --out 产出临时 .pptx，记录 cmd。"""

    def fake_run(cmd, **kw):
        recorded["cmd"] = cmd
        out = cmd[cmd.index("--out") + 1]
        Path(out).write_bytes(b"fake pptx")
        return _sp.CompletedProcess(args=[], returncode=0, stdout="", stderr="")

    return fake_run


def _record_target(recorded: dict):
    """mock postprocess_charts：记录收到的 html target 路径 + 调用瞬间内容。

    注入副本在 render 返回时已被 finally 清理，内容必须在 postprocess 调用瞬间读。
    """

    def record(html_path, pptx_path):
        recorded["html"] = html_path
        recorded["exists_at_call"] = Path(html_path).exists()
        recorded["content_at_call"] = Path(html_path).read_text(encoding="utf-8")

    return record


def _patch_assets_noop(monkeypatch):
    """本批用例只关心 charts 收 target；asset 主链另由本文件 manifest 生命周期
    用例覆盖（真实 postprocess_assets），这里打桩避免 measurements 缺失抛错。"""
    monkeypatch.setattr(
        "offipy.assets.render.postprocess_assets", lambda html_path, pptx_path: None
    )


_HTML = """<!doctype html>
<html><head></head><body>
<section data-pptx-slide><h1>一</h1></section>
<section data-pptx-slide><div data-asset="asset://ph/icon/check"></div></section>
</body></html>
"""


def _write(tmp_path, text=_HTML):
    html = tmp_path / "d.html"
    html.write_text(text, encoding="utf-8")
    pptx = tmp_path / "d.pptx"
    pptx.write_bytes(b"placeholder")
    return html, pptx


# ---------------------------------------------------------------------------
# no_visual_audit 前置：声明 → chromium / convert 之前 fail-fast
# ---------------------------------------------------------------------------


def _assert_nova_rejected_before_browser(tmp_path, monkeypatch, html_text, marker):
    html, pptx = _write(tmp_path, html_text)
    browser_calls: list = []
    run_calls: list = []

    monkeypatch.setattr(deck, "_preflight_browser", lambda *a, **k: browser_calls.append(a))
    monkeypatch.setattr(deck.subprocess, "run", lambda *a, **k: run_calls.append(a))

    with pytest.raises(deck.InvalidArgumentError) as exc:
        deck.render(str(html), out=str(pptx), overwrite=True, no_visual_audit=True)
    assert marker in str(exc.value)  # 信息列出具体声明类型
    assert browser_calls == []  # chromium 从未启动
    assert run_calls == []  # convert 子进程从未调用


def test_no_visual_audit_rejects_data_asset_before_browser(tmp_path, monkeypatch):
    _assert_nova_rejected_before_browser(tmp_path, monkeypatch, _HTML, "data-asset")


def test_no_visual_audit_rejects_data_primitive_before_browser(tmp_path, monkeypatch):
    _assert_nova_rejected_before_browser(
        tmp_path,
        monkeypatch,
        '<section data-pptx-slide><div data-primitive="quote-mark"></div></section>',
        "data-primitive",
    )


def test_no_visual_audit_still_rejects_legacy_data_icon_before_browser(tmp_path, monkeypatch):
    _assert_nova_rejected_before_browser(
        tmp_path,
        monkeypatch,
        '<section data-pptx-slide><svg data-icon="ph:check"></svg></section>',
        "data-icon",
    )


def test_no_visual_audit_still_rejects_data_chart_before_browser(tmp_path, monkeypatch):
    _assert_nova_rejected_before_browser(
        tmp_path,
        monkeypatch,
        '<section data-pptx-slide><div data-chart="bar"></div></section>',
        "data-chart",
    )


def test_no_visual_audit_message_lists_declaration_type(tmp_path, monkeypatch):
    """data-asset 单独出现时，信息点名的就是它（不混入无关类型）。"""
    html, pptx = _write(tmp_path)
    monkeypatch.setattr(deck.subprocess, "run", lambda *a, **k: None)

    with pytest.raises(deck.InvalidArgumentError) as exc:
        deck.render(str(html), out=str(pptx), overwrite=True, no_visual_audit=True)
    msg = str(exc.value)
    assert "data-asset" in msg
    assert "data-chart" not in msg  # 未声明图表就不误报图表
    assert "data-icon" not in msg
    assert "data-primitive" not in msg


def test_no_visual_audit_plain_deck_still_ok(tmp_path, monkeypatch):
    """回归：无任何声明的纯 deck 配 no_visual_audit 不受前置校验影响。"""
    html, pptx = _write(
        tmp_path, "<html><body><section data-pptx-slide><p>ok</p></section></body></html>"
    )
    recorded = {}
    monkeypatch.setattr(deck.subprocess, "run", _fake_run(recorded))

    out = deck.render(str(html), out=str(pptx), overwrite=True, no_visual_audit=True)
    assert out == str(pptx)
    assert "--no-visual-audit" in recorded["cmd"]  # 转换器确实收到该开关


# ---------------------------------------------------------------------------
# 注入副本 target：asset 声明强制临时 target，无 theme/layout 也如此
# ---------------------------------------------------------------------------


def test_asset_only_passes_temp_injected_target(tmp_path, monkeypatch):
    """无 theme/layout 时，asset 声明也强制临时注入副本；源文件不动、用后即清。"""
    html, pptx = _write(tmp_path)
    original = html.read_text(encoding="utf-8")

    recorded = {}
    monkeypatch.setattr(charts, "postprocess_charts", _record_target(recorded))
    _patch_assets_noop(monkeypatch)
    monkeypatch.setattr(deck.subprocess, "run", _fake_run(recorded))

    out = deck.render(str(html), out=str(pptx), overwrite=True)
    assert out == str(pptx)

    target = Path(recorded["html"])
    assert str(target) != str(html)  # 注入副本而非源文件
    assert target.name == "d.audited.html"  # 命名让 convert 跳过 work-copy 分支
    assert target.parent.name.startswith("offipy-deck-")  # 临时目录前缀
    assert recorded["exists_at_call"] is True  # 后处理瞬间副本存在
    assert 'data-offipy-asset-id="asset-s02-001"' in recorded["content_at_call"]
    # 用户源 HTML 未被就地编辑
    assert html.read_text(encoding="utf-8") == original
    # 注入副本随 TemporaryDirectory 用后即清
    assert not target.exists()
    assert not target.parent.exists()


def test_asset_only_no_visual_audit_rejected_before_temp_target(tmp_path, monkeypatch):
    """no_visual_audit + asset 声明在临时 target 生成前就拦下（无残留副本）。"""
    html, pptx = _write(tmp_path)
    monkeypatch.setattr(deck.subprocess, "run", lambda *a, **k: None)

    with pytest.raises(deck.InvalidArgumentError):
        deck.render(str(html), out=str(pptx), overwrite=True, no_visual_audit=True)
    assert not list(tmp_path.glob("offipy-deck-*"))


def test_theme_plus_asset_target_has_theme_and_ids(tmp_path, monkeypatch):
    """theme + asset：注入副本同时带主题 token 与确定性 ID（顺序：先注入后解析）。"""
    html, pptx = _write(tmp_path)
    recorded = {}
    monkeypatch.setattr(charts, "postprocess_charts", _record_target(recorded))
    _patch_assets_noop(monkeypatch)
    monkeypatch.setattr(deck.subprocess, "run", _fake_run(recorded))

    out = deck.render(str(html), out=str(pptx), overwrite=True, theme="mckinsey")
    assert out == str(pptx)

    content = recorded["content_at_call"]
    assert 'data-theme="mckinsey"' in content  # 主题注入在 asset 解析前完成
    assert 'data-offipy-asset-id="asset-s02-001"' in content


def test_asset_target_cleaned_on_failure(tmp_path, monkeypatch):
    """后处理抛错时，注入副本随 finally 清理，不留残留。"""
    html, pptx = _write(tmp_path)
    target_seen: dict[str, str] = {}

    def boom(html_path, pptx_path):
        target_seen["html"] = html_path
        raise RuntimeError("boom")

    monkeypatch.setattr(charts, "postprocess_charts", boom)
    monkeypatch.setattr(deck.subprocess, "run", _fake_run({}))

    with pytest.raises(deck.ConversionError):
        deck.render(str(html), out=str(pptx), overwrite=True)
    assert not Path(target_seen["html"]).exists()
    assert not Path(target_seen["html"]).parent.exists()


# ---------------------------------------------------------------------------
# only_slides：声明页码与 converter 一致
# ---------------------------------------------------------------------------


def test_only_slides_keeps_full_html_declaration_ordinals(tmp_path, monkeypatch):
    """only_slides 只重渲子集，但声明 ID 仍按全 HTML 页码编号（converter 对齐）。"""
    html, pptx = _write(tmp_path)
    recorded = {}
    monkeypatch.setattr(charts, "postprocess_charts", _record_target(recorded))
    _patch_assets_noop(monkeypatch)
    monkeypatch.setattr(deck.subprocess, "run", _fake_run(recorded))

    out = deck.render(str(html), out=str(pptx), overwrite=True, only_slides=[2])
    assert out == str(pptx)

    cmd = recorded["cmd"]
    assert cmd[cmd.index("--only-slides") + 1] == "2"
    content = recorded["content_at_call"]
    # 声明按全 HTML 第 2 页编号（asset-s02-001），不被 only_slides 重排
    assert 'data-offipy-asset-id="asset-s02-001"' in content


# ---------------------------------------------------------------------------
# Task 10 — manifest 生命周期：postprocess_assets 报告 → assets.json 落盘
# ---------------------------------------------------------------------------


_ASSET_ORDER_HTML = """<!doctype html>
<html><head></head><body>
<section data-pptx-slide><div data-asset="asset://ph/icon/check"></div></section>
<section data-pptx-slide>
  <div data-asset="asset://lu/icon/settings" data-asset-placement="background"></div>
</section>
</body></html>"""

_PLAIN_HTML = "<html><head></head><body><section data-pptx-slide><p>x</p></section></body></html>"


def _asset_meas(asset_id, x=20, y=30, w=120, h=60, color="rgb(255, 0, 0)"):
    return {
        "id": 1,
        "kind": "asset",
        "tag": "div",
        "assetId": asset_id,
        "rect": {"x": x, "y": y, "w": w, "h": h},
        "themeVars": {"bg": "#ffffff", "surface": "#f3f4f6", "accent": "#2251ff"},
        "color": color,
    }


def _make_fake_convert(decls, *, create_audit=True):
    """fake subprocess.run：读 --out，造带 OFFIPY_ASSET:: 占位符的真实 pptx，
    并按需产出 <stem>_audit/_cache/measurements.json（与真实 convert 审计同形）。"""

    def fake_run(cmd, *a, **k):
        out = cmd[cmd.index("--out") + 1]
        prs = Presentation()
        for d in decls:
            while len(prs.slides) < d.slide_index:
                prs.slides.add_slide(prs.slide_layouts[6])
            slide = prs.slides[d.slide_index - 1]
            sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, 100, 100)
            sp.name = f"OFFIPY_ASSET::{d.declaration_id}"
        prs.save(out)
        if create_audit:
            max_slide = max((d.slide_index for d in decls), default=0)
            slides = [{"records": []} for _ in range(max_slide)]
            for d in decls:
                slides[d.slide_index - 1]["records"].append(_asset_meas(d.declaration_id))
            audit = Path(out).with_name(f"{Path(out).stem}_audit")
            meas_dir = audit / "_cache"
            meas_dir.mkdir(parents=True)
            (meas_dir / "measurements.json").write_text(
                json.dumps({"slides": slides}, ensure_ascii=False), encoding="utf-8"
            )
        return _sp.CompletedProcess(args=[], returncode=0, stdout="", stderr="")

    return fake_run


def _load_manifest(tmp_path):
    p = tmp_path / "d_audit" / "assets.json"
    if not p.exists():
        return None
    return json.loads(p.read_text(encoding="utf-8"))


def test_visual_audit_zero_assets_writes_empty_manifest(tmp_path, monkeypatch):
    """纯文本 deck + visual audit：审计目录存在 → 空用量报告落成空 assets 列表。"""
    html, pptx = _write(tmp_path, _PLAIN_HTML)
    monkeypatch.setattr(deck.subprocess, "run", _make_fake_convert([]))

    deck.render(str(html), out=str(pptx), overwrite=True)

    assert _load_manifest(tmp_path) == {"schema": 1, "assets": []}


def test_assets_written_in_declaration_order_with_provider_meta(tmp_path, monkeypatch):
    """真实 asset 主链：用量清单按声明顺序落盘，provider 元数据（ph→MIT / lu→ISC）正确。"""
    html, pptx = _write(tmp_path, _ASSET_ORDER_HTML)
    _, decls = preprocess_asset_declarations(_ASSET_ORDER_HTML)
    monkeypatch.setattr(deck.subprocess, "run", _make_fake_convert(decls))

    deck.render(str(html), out=str(pptx), overwrite=True)

    manifest = _load_manifest(tmp_path)
    assert manifest["schema"] == 1
    assets = manifest["assets"]
    assert [a["declaration_id"] for a in assets] == ["asset-s01-001", "asset-s02-001"]
    assert [a["provider"]["license"] for a in assets] == ["MIT", "ISC"]
    assert assets[0]["placement"] == "replace"
    assert assets[1]["placement"] == "background"


def test_no_visual_audit_writes_no_audit_dir_or_manifest(tmp_path, monkeypatch):
    """no_visual_audit：无审计目录 → 不写 assets.json，纯 deck 正常产出。"""
    html, pptx = _write(tmp_path, _PLAIN_HTML)
    monkeypatch.setattr(deck.subprocess, "run", _make_fake_convert([], create_audit=False))

    deck.render(str(html), out=str(pptx), overwrite=True, no_visual_audit=True)

    assert pptx.exists()
    assert not (tmp_path / "d_audit").exists()
    assert not list(tmp_path.rglob("assets.json"))


def test_failed_postprocess_preserves_old_pptx_and_manifest(tmp_path, monkeypatch):
    """后处理失败：不产生新的最终清单，旧 PPTX 与旧 d_audit/assets.json 原样保留。"""
    html, pptx = _write(tmp_path, _ASSET_ORDER_HTML)
    _, decls = preprocess_asset_declarations(_ASSET_ORDER_HTML)
    monkeypatch.setattr(deck.subprocess, "run", _make_fake_convert(decls))

    # 第一次成功渲染 → 产出 d.pptx + d_audit/assets.json
    deck.render(str(html), out=str(pptx), overwrite=True)
    old_pptx = pptx.read_bytes()
    old_manifest = (tmp_path / "d_audit" / "assets.json").read_bytes()

    # 第二次渲染：postprocess_assets 抛错 → 失败不碰旧产物
    def boom(html_path, pptx_path):
        raise RuntimeError("boom")

    monkeypatch.setattr("offipy.assets.render.postprocess_assets", boom)
    with pytest.raises(deck.ConversionError) as exc:
        deck.render(str(html), out=str(pptx), overwrite=True)
    assert "资源" in str(exc.value)

    assert pptx.read_bytes() == old_pptx
    assert (tmp_path / "d_audit" / "assets.json").read_bytes() == old_manifest
    assert not list(tmp_path.glob(".*_audit"))  # 本次失败渲染的临时审计目录已清理


# ---------------------------------------------------------------------------
# Task 12 — 失败原子性补全：非法 URI / 缺失占位符
# ---------------------------------------------------------------------------


def test_invalid_asset_uri_preserves_final_output(tmp_path, monkeypatch):
    """data-asset 非法 URI：preprocess 阶段 InvalidArgumentError，最终输出不被替换。"""
    html, pptx = _write(
        tmp_path, '<section data-pptx-slide><div data-asset="not-a-uri"></div></section>'
    )
    original = pptx.read_bytes()
    monkeypatch.setattr(deck.subprocess, "run", lambda *a, **k: None)

    with pytest.raises(deck.InvalidArgumentError, match="asset://"):
        deck.render(str(html), out=str(pptx), overwrite=True)

    assert pptx.read_bytes() == original  # 旧最终输出原样保留
    assert not list(tmp_path.glob("offipy-deck-*"))  # 无残留注入副本


def test_missing_placeholder_preserves_final_output(tmp_path, monkeypatch):
    """measurements 有记录但转换器漏放 OFFIPY_ASSET 占位符：绑定失败，旧产物保留。"""
    html, pptx = _write(tmp_path, _ASSET_ORDER_HTML)
    _, decls = preprocess_asset_declarations(_ASSET_ORDER_HTML)
    monkeypatch.setattr(deck.subprocess, "run", _make_fake_convert(decls))

    # 第一次成功渲染 → 产出 d.pptx + d_audit/assets.json
    deck.render(str(html), out=str(pptx), overwrite=True)
    old_pptx = pptx.read_bytes()
    old_manifest = (tmp_path / "d_audit" / "assets.json").read_bytes()

    # 第二次：fake convert 写 measurements.json 但不放占位符（converter 产物缺损）
    def fake_no_placeholder(cmd, *a, **k):
        out = cmd[cmd.index("--out") + 1]
        prs = Presentation()
        while len(prs.slides) < 2:
            prs.slides.add_slide(prs.slide_layouts[6])
        prs.save(out)
        audit = Path(out).with_name(f"{Path(out).stem}_audit")
        meas_dir = audit / "_cache"
        meas_dir.mkdir(parents=True)
        (meas_dir / "measurements.json").write_text(
            json.dumps(
                {
                    "slides": [
                        {"records": [_asset_meas("asset-s01-001")]},
                        {"records": [_asset_meas("asset-s02-001")]},
                    ]
                },
                ensure_ascii=False,
            ),
            encoding="utf-8",
        )
        return _sp.CompletedProcess(args=[], returncode=0, stdout="", stderr="")

    monkeypatch.setattr(deck.subprocess, "run", fake_no_placeholder)
    with pytest.raises(deck.InvalidArgumentError, match="占位符"):
        deck.render(str(html), out=str(pptx), overwrite=True)

    assert pptx.read_bytes() == old_pptx  # 失败不替换最终输出
    assert (tmp_path / "d_audit" / "assets.json").read_bytes() == old_manifest  # 旧清单原样
    assert not list(tmp_path.glob(".*_audit"))  # 本次失败渲染的临时审计目录已清理
