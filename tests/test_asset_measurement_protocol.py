"""A3 Task 6 — converter kind=asset 测量绑定 / 占位符查找 / deck 一致性校验。

load_asset_measurements 只做测量内部校验（assetId 非空/唯一/rect 为正）；
bind_asset_measurements 做声明↔测量交叉校验（缺测量/缺声明/slide 不匹配）。
占位符查找按 name 精确匹配，0 或 >1 都是绑定错误，不做空间/outerHTML 兜底。
"""

import json

import pytest

from offipy.assets.declarations import AssetDeclaration
from offipy.assets.model import AssetRect, AssetRef, AssetRequest
from offipy.assets.render import (
    bind_asset_measurements,
    find_asset_placeholder,
    load_asset_measurements,
)
from offipy.exceptions import InvalidArgumentError


def _decl(declaration_id, slide_index):
    return AssetDeclaration(
        declaration_id=declaration_id,
        slide_index=slide_index,
        request=AssetRequest(AssetRef("ph", "icon", "check")),
        placement="replace",
        html_tag="div",
    )


def _asset_rec(asset_id, x=20, y=30, w=120, h=60, tag="div"):
    return {
        "id": 1,
        "kind": "asset",
        "tag": tag,
        "assetId": asset_id,
        "rect": {"x": x, "y": y, "w": w, "h": h},
        "themeVars": {"bg": "#ffffff", "surface": "#f3f4f6", "accent": "#2251ff"},
        "color": "rgb(255, 0, 0)",
    }


def _write_meas(tmp_path, slides):
    p = tmp_path / "measurements.json"
    p.write_text(json.dumps({"slides": slides}, ensure_ascii=False), encoding="utf-8")
    return p


# ---------------------------------------------------------------------------
# load_asset_measurements
# ---------------------------------------------------------------------------


class TestLoadMeasurements:
    def test_loads_asset_records_with_fields(self, tmp_path):
        meas = load_asset_measurements(
            _write_meas(
                tmp_path,
                [
                    {"records": [_asset_rec("asset-s01-001", 20, 30, 120, 60)]},
                    {"records": [_asset_rec("asset-s02-001", 320, 20, 80, 80, tag="svg")]},
                ],
            )
        )
        assert list(meas) == ["asset-s01-001", "asset-s02-001"]
        m1 = meas["asset-s01-001"]
        assert m1.slide_index == 1
        assert m1.rect == AssetRect(20, 30, 120, 60)
        assert m1.theme_vars["accent"] == "#2251ff"
        assert m1.html_tag == "div"
        assert m1.color == "rgb(255, 0, 0)"
        m2 = meas["asset-s02-001"]
        assert m2.slide_index == 2
        assert m2.html_tag == "svg"

    def test_ignores_non_asset_records(self, tmp_path):
        meas = load_asset_measurements(
            _write_meas(
                tmp_path,
                [
                    {
                        "records": [
                            {"id": 1, "kind": "shape", "rect": {"x": 0, "y": 0, "w": 10, "h": 10}},
                            {"id": 2, "kind": "text", "rect": {"x": 0, "y": 0, "w": 10, "h": 10}},
                            _asset_rec("asset-s01-001"),
                        ]
                    },
                ],
            )
        )
        assert list(meas) == ["asset-s01-001"]

    def test_missing_asset_id_rejected(self, tmp_path):
        rec = _asset_rec("asset-s01-001")
        rec["assetId"] = ""
        with pytest.raises(InvalidArgumentError, match="assetId"):
            load_asset_measurements(_write_meas(tmp_path, [{"records": [rec]}]))

    def test_duplicate_asset_id_rejected(self, tmp_path):
        with pytest.raises(InvalidArgumentError, match="asset-s01-001"):
            load_asset_measurements(
                _write_meas(
                    tmp_path,
                    [
                        {"records": [_asset_rec("asset-s01-001")]},
                        {"records": [_asset_rec("asset-s01-001")]},
                    ],
                )
            )

    def test_nonpositive_rect_rejected(self, tmp_path):
        rec = _asset_rec("asset-s01-001", w=0)
        with pytest.raises(InvalidArgumentError, match="positive"):
            load_asset_measurements(_write_meas(tmp_path, [{"records": [rec]}]))

    def test_bad_rect_shape_rejected(self, tmp_path):
        rec = _asset_rec("asset-s01-001")
        rec["rect"] = {"x": 0, "y": 0}
        with pytest.raises(InvalidArgumentError, match="rect"):
            load_asset_measurements(_write_meas(tmp_path, [{"records": [rec]}]))


# ---------------------------------------------------------------------------
# find_asset_placeholder
# ---------------------------------------------------------------------------


def _slide_with_placeholder(*names):
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    for n in names:
        sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, 100, 100)
        sp.name = n
    return slide


class TestFindPlaceholder:
    def test_exact_match_returns_shape(self):
        slide = _slide_with_placeholder("OFFIPY_ASSET::asset-s01-001")
        sp = find_asset_placeholder(slide, "asset-s01-001")
        assert sp.name == "OFFIPY_ASSET::asset-s01-001"

    def test_missing_placeholder_rejected(self):
        slide = _slide_with_placeholder("OFFIPY_ASSET::asset-s01-001")
        with pytest.raises(InvalidArgumentError, match="asset-s02-001"):
            find_asset_placeholder(slide, "asset-s02-001")

    def test_duplicate_placeholder_rejected(self):
        slide = _slide_with_placeholder(
            "OFFIPY_ASSET::asset-s01-001", "OFFIPY_ASSET::asset-s01-001"
        )
        with pytest.raises(InvalidArgumentError, match="2 个"):
            find_asset_placeholder(slide, "asset-s01-001")


# ---------------------------------------------------------------------------
# bind_asset_measurements
# ---------------------------------------------------------------------------


class TestBindMeasurements:
    def test_binds_by_id_preserving_declaration_order(self):
        decls = [_decl("asset-s01-001", 1), _decl("asset-s01-002", 1), _decl("asset-s02-001", 2)]
        meas = {
            "asset-s01-001": _measurement("asset-s01-001", 1),
            "asset-s02-001": _measurement("asset-s02-001", 2),
            "asset-s01-002": _measurement("asset-s01-002", 1),
        }
        bound = bind_asset_measurements(decls, meas)
        assert list(bound) == ["asset-s01-001", "asset-s01-002", "asset-s02-001"]

    def test_declaration_without_measurement_rejected(self):
        decls = [_decl("asset-s01-001", 1), _decl("asset-s01-002", 1)]
        meas = {"asset-s01-001": _measurement("asset-s01-001", 1)}
        with pytest.raises(InvalidArgumentError, match="asset-s01-002"):
            bind_asset_measurements(decls, meas)

    def test_measurement_without_declaration_rejected(self):
        decls = [_decl("asset-s01-001", 1)]
        meas = {
            "asset-s01-001": _measurement("asset-s01-001", 1),
            "asset-s02-001": _measurement("asset-s02-001", 2),
        }
        with pytest.raises(InvalidArgumentError, match="asset-s02-001"):
            bind_asset_measurements(decls, meas)

    def test_slide_mismatch_rejected(self):
        decls = [_decl("asset-s01-001", 1)]
        meas = {"asset-s01-001": _measurement("asset-s01-001", 2)}
        with pytest.raises(InvalidArgumentError, match="slide"):
            bind_asset_measurements(decls, meas)

    def test_same_uri_ten_times_all_bind_by_id(self):
        # 同 URI 重复 10 次：每份声明都有确定性 ID，必须按 ID 绑定而非按 URI/空间
        decls = [_decl(f"asset-s01-{i:03d}", 1) for i in range(1, 11)]
        meas = {f"asset-s01-{i:03d}": _measurement(f"asset-s01-{i:03d}", 1) for i in range(1, 11)}
        bound = bind_asset_measurements(decls, meas)
        assert len(bound) == 10
        assert [m.asset_id for m in bound.values()] == [f"asset-s01-{i:03d}" for i in range(1, 11)]


def _measurement(asset_id, slide_index):
    from offipy.assets.render import AssetMeasurement

    return AssetMeasurement(
        asset_id=asset_id,
        slide_index=slide_index,
        rect=AssetRect(20, 30, 120, 60),
        theme_vars={"bg": "#ffffff", "accent": "#2251ff"},
        html_tag="div",
        color="rgb(0, 0, 0)",
    )
