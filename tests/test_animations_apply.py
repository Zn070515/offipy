"""animations/apply.py：独立 API——按形状名定位 spid → 注入 timing/transition → 报告。"""

import pytest

from offipy.animations.apply import apply_animations, apply_transitions
from offipy.animations.spec import AnimationSpec, TransitionSpec
from offipy.exceptions import InvalidArgumentError


def _build_pptx(tmp_path, shapes, *, existing_timing=False, existing_transition=False):
    """造一个 1 页 .pptx，shapes=[(name,)]。返回路径。"""
    from pptx import Presentation
    from pptx.oxml.ns import qn
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    for i, (name,) in enumerate(shapes):
        sp = slide.shapes.add_shape(1, Inches(1), Inches(1 + i), Inches(2), Inches(0.6))
        sp.name = name
    if existing_timing:
        # 给 slide 塞一个假 <p:timing>（构造字符串插入）
        t = slide._element.makeelement(qn("p:timing"), {})
        slide._element.append(t)
    if existing_transition:
        # 给 slide 塞一个真 <p:transition>（复用真实构建器）
        from offipy.animations.transition import build_transition

        slide._element.append(build_transition("fade", "medium"))
    p = tmp_path / "a.pptx"
    prs.save(str(p))
    return str(p)


def _qn(tag):
    return tag.split("}")[-1]


def test_apply_animations_injects_timing(tmp_path):
    p = _build_pptx(tmp_path, [("title",)])
    report = apply_animations(
        p,
        animations=[AnimationSpec(slide=1, target="title", effect="fade")],
    )
    assert report["animations_applied"] == 1
    assert report["transitions_applied"] == 0
    assert report["unmatched"] == []
    from pptx import Presentation

    sld = Presentation(p).slides[0]
    timing = sld._element.findall(".//{*}timing")
    assert len(timing) == 1


def test_apply_animations_target_not_found_reported_not_fatal(tmp_path):
    p = _build_pptx(tmp_path, [("title",)])
    report = apply_animations(
        p,
        animations=[AnimationSpec(slide=1, target="missing", effect="fade")],
        raise_on_all_unmatched=False,
    )
    assert report["unmatched"] == [{"slide": 1, "target": "missing"}]
    assert report["animations_applied"] == 0


def test_apply_animations_all_unmatched_raises(tmp_path):
    p = _build_pptx(tmp_path, [("title",)])
    with pytest.raises(InvalidArgumentError):
        apply_animations(
            p,
            animations=[
                AnimationSpec(slide=1, target="a", effect="fade"),
                AnimationSpec(slide=1, target="b", effect="fade"),
            ],
        )


def test_apply_animations_transition_only(tmp_path):
    p = _build_pptx(tmp_path, [("title",)])
    report = apply_animations(
        p,
        transitions=[TransitionSpec(slide=1, kind="push", speed="medium")],
    )
    assert report["transitions_applied"] == 1
    assert report["animations_applied"] == 0
    from pptx import Presentation

    sld = Presentation(p).slides[0]
    transitions = sld._element.findall(".//{*}transition")
    assert len(transitions) == 1


def test_apply_animations_order_transition_before_timing(tmp_path):
    p = _build_pptx(tmp_path, [("title",)])
    apply_animations(
        p,
        animations=[AnimationSpec(slide=1, target="title", effect="fade")],
        transitions=[TransitionSpec(slide=1, kind="fade", speed="fast")],
    )
    from pptx import Presentation

    sld = Presentation(p).slides[0]
    names = [_qn(c.tag) for c in sld._element]
    assert names.index("transition") < names.index("timing")
    # 且都在 clrMapOvr 之后
    assert names.index("cSld") < names.index("transition")
    assert names.index("clrMapOvr") < names.index("transition")


def test_apply_animations_idempotency_raises(tmp_path):
    p = _build_pptx(tmp_path, [("title",)], existing_timing=True)
    with pytest.raises(InvalidArgumentError, match="已有动画"):
        apply_animations(
            p,
            animations=[AnimationSpec(slide=1, target="title", effect="fade")],
        )


def test_apply_transitions_wrapper(tmp_path):
    p = _build_pptx(tmp_path, [("title",)])
    report = apply_transitions(p, [TransitionSpec(slide=1, kind="cover", speed="slow")])
    assert report["transitions_applied"] == 1


def test_apply_animations_existing_transition_skips(tmp_path):
    p = _build_pptx(tmp_path, [("title",)], existing_transition=True)
    report = apply_animations(
        p,
        transitions=[TransitionSpec(slide=1, kind="push", speed="medium")],
    )
    # slide 已有过渡 → 不注入第二个，也不计 applied
    assert report["transitions_applied"] == 0
    assert report["skipped"] == [{"slide": 1, "kind": "push", "reason": "slide 已有过渡"}]
    from pptx import Presentation

    sld = Presentation(p).slides[0]
    assert len(sld._element.findall(".//{*}transition")) == 1


def test_apply_animations_multi_transition_same_slide_skip(tmp_path):
    p = _build_pptx(tmp_path, [("title",)])
    report = apply_animations(
        p,
        transitions=[
            TransitionSpec(slide=1, kind="push", speed="medium"),
            TransitionSpec(slide=1, kind="cover", speed="medium"),
        ],
    )
    assert report["transitions_applied"] == 1
    assert report["skipped"] == [{"slide": 1, "kind": "cover", "reason": "每页仅一个过渡"}]
    from pptx import Presentation

    sld = Presentation(p).slides[0]
    assert len(sld._element.findall(".//{*}transition")) == 1


def test_apply_animations_multi_shape_same_name(tmp_path):
    p = _build_pptx(tmp_path, [("title",), ("title",)])
    report = apply_animations(
        p,
        animations=[AnimationSpec(slide=1, target="title", effect="fade")],
    )
    assert report["animations_applied"] == 1
    from pptx import Presentation

    sld = Presentation(p).slides[0]
    bld_spids = {b.get("spid") for b in sld._element.findall(".//{*}bldP")}
    shape_ids = {str(sh.shape_id) for sh in sld.shapes if sh.name == "title"}
    assert bld_spids == shape_ids == {str(s.shape_id) for s in sld.shapes}


def test_apply_animations_all_unmatched_preserves_file(tmp_path):
    p = _build_pptx(tmp_path, [("title",)])
    from pathlib import Path

    before = Path(p).read_bytes()
    with pytest.raises(InvalidArgumentError):
        apply_animations(
            p,
            animations=[
                AnimationSpec(slide=1, target="missing-a", effect="fade"),
                AnimationSpec(slide=1, target="missing-b", effect="fade"),
            ],
        )
    assert Path(p).read_bytes() == before


def test_apply_animations_dict_coercion(tmp_path):
    p = _build_pptx(tmp_path, [("title",)])
    report = apply_animations(
        p,
        animations=[{"slide": 1, "target": "title", "effect": "fade"}],
    )
    assert report["animations_applied"] == 1

    p2 = _build_pptx(tmp_path, [("title",)])
    with pytest.raises(InvalidArgumentError):
        apply_animations(
            p2,
            animations=[{"slide": 1, "target": "title", "effect": "fade", "bogus": 1}],
        )
