"""office-real：PowerPoint COM 读 MainSequence.Count / transition 存在性。

真机验证 apply_animations 注入的 <p:timing>/<p:transition> 被 PowerPoint 识别为
可播动画（MainSequence 效果数 = 注入单元数）与页面过渡。无存活 PowerPoint 时
`live_ppt` fixture 跳过（session server 8890 持有）。标记 `com` 让 office-real
workflow（-m "com or deck_render"）作为真机门禁拾取。
"""

import pytest

from offipy.animations.apply import apply_animations
from offipy.animations.spec import AnimationSpec, TransitionSpec

pytestmark = [pytest.mark.com]


def test_com_reads_main_sequence_count(tmp_path, live_ppt):
    from pptx import Presentation
    from pptx.util import Inches

    from offipy.api import Ppt

    p = tmp_path / "a.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    for i, name in enumerate(("title", "body")):
        sp = slide.shapes.add_shape(1, Inches(1), Inches(1 + i), Inches(3), Inches(0.8))
        sp.name = name
    prs.save(str(p))

    apply_animations(
        str(p),
        animations=[
            AnimationSpec(slide=1, target="title", effect="fade"),
            AnimationSpec(
                slide=1,
                target="body",
                effect="fly_in",
                direction="left",
                trigger="after",
                delay=0.2,
            ),
        ],
        transitions=[TransitionSpec(slide=1, kind="push", speed="medium")],
    )

    with Ppt() as ppt:
        doc_id = ppt.open_pres(str(p))
        try:
            # 本地直连 Ppt()：同一进程持有 COM Presentation，读原始对象模型
            pres = ppt._app._docs[doc_id]
            slide1 = pres.Slides(1)
            count = slide1.TimeLine.MainSequence.Count
            # title(fade,click) + body(fly_in,after) = 2 个 MainSequence 效果
            assert count == 2, f"PowerPoint 识别到 {count} 个动画效果，期望 2"
            # 过渡：SlideShowTransition.EntryEffect 非 0
            assert slide1.SlideShowTransition.EntryEffect != 0, "应存在页面过渡"
        finally:
            ppt.close_pres(doc_id=doc_id, save=False)
