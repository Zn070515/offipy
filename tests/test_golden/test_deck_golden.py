"""Deck golden（P2-6）：固定 starter deck 渲染结果指纹，防行为漂移。

断言与 2026-08-05 基线一致：5 页、文件大小带、每页首行文本。
任何 HTML→PPTX 管线的输出结构变化都会在这里被抓住。
"""

from pathlib import Path

import pytest

from offipy.deck import CONVERT_PY
from offipy.envcheck import _check_browser

STARTER = (
    Path(__file__).resolve().parent.parent.parent / "examples" / "decks" / "starter" / "deck.html"
)

# 2026-08-05 基线：render(starter) → 173,700 B，5 页
MIN_SIZE = 100_000
MAX_SIZE = 400_000
EXPECTED_SLIDES = 5
# 每页第一个文本形状（kicker 行）——稳定指纹
EXPECTED_KICKERS = [
    "QUARTERLY REVIEW · Q2 2026",
    "AGENDA",
    "HIGHLIGHT",
    "PLAN",
    "THANK YOU",
]

pytestmark = [
    pytest.mark.deck_render,
    pytest.mark.skipif(
        not CONVERT_PY.exists() or not _check_browser().ok,
        reason="vendored 转换器缺失或 chromium 不可用",
    ),
]


def _first_texts(path: str) -> list[str]:
    pytest.importorskip("pptx")
    from pptx import Presentation

    out = []
    for slide in Presentation(path).slides:
        first = ""
        for sh in slide.shapes:
            if sh.has_text_frame and sh.text_frame.text.strip():
                first = sh.text_frame.text.strip().split("\n")[0].strip()
                break
        out.append(first)
    return out


def test_starter_deck_golden_fingerprint(tmp_path):
    from offipy.deck import render

    out = tmp_path / "golden.pptx"
    render(str(STARTER), out=str(out), no_visual_audit=True)
    assert out.stat().st_size >= MIN_SIZE
    assert out.stat().st_size <= MAX_SIZE


def test_starter_deck_golden_slide_count(tmp_path):
    from offipy.deck import render

    pytest.importorskip("pptx")
    from pptx import Presentation

    out = tmp_path / "golden.pptx"
    render(str(STARTER), out=str(out), no_visual_audit=True)
    assert len(Presentation(str(out)).slides) == EXPECTED_SLIDES


def test_starter_deck_golden_kickers(tmp_path):
    from offipy.deck import render

    out = tmp_path / "golden.pptx"
    render(str(STARTER), out=str(out), no_visual_audit=True)
    assert _first_texts(str(out)) == EXPECTED_KICKERS


def test_starter_deck_golden_no_elem_anchor_names(tmp_path):
    pytest.importorskip("pptx")
    from offipy.deck import render

    out = tmp_path / "golden.pptx"
    render(str(STARTER), out=str(out), no_visual_audit=True)
    from pptx import Presentation

    for slide in Presentation(str(out)).slides:
        for sh in slide.shapes:
            assert not sh.name.startswith("OFFIPY_ELEM::"), (
                f"animations=False 不应戳 OFFIPY_ELEM 锚点名（shape={sh.name}）"
            )
