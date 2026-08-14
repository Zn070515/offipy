"""offipy.diagrams — Mermaid 解析（vendored mermaid_extract 包装）。"""

from __future__ import annotations

import pytest
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.oxml.ns import qn
from pptx.util import Emu

from offipy.diagrams import (
    _kahn_layers,
    _read_source,
    layout_diagram,
    mermaid_to_pptx,
    parse_mermaid,
    render_to_slide,
)


def test_parse_flowchart_basic():
    diagram = parse_mermaid("graph TD\n    A[开始] --> B[处理]")
    assert diagram.kind == "flowchart"
    assert diagram.direction == "TD"
    assert [n.id for n in diagram.nodes] == ["A", "B"]
    assert [e.source for e in diagram.edges] == ["A"]


def test_parse_chinese_label_preserved():
    diagram = parse_mermaid("graph LR\n    X[数据管道] --> Y[结果]")
    labels = {n.id: n.label for n in diagram.nodes}
    assert labels["X"] == "数据管道"


def test_parse_direction_lr():
    diagram = parse_mermaid("graph LR\n    A --> B")
    assert diagram.direction == "LR"


def test_unsupported_kind_rejected():
    with pytest.raises(ValueError, match="仅支持 flowchart"):
        parse_mermaid("sequenceDiagram\n    A->>B: hi")


def test_bad_syntax_raises_value_error():
    with pytest.raises(ValueError):
        parse_mermaid("graph TD\n    this is not mermaid !!!")


def test_extractor_systemexit_normalized_to_value_error():
    # gantt 不是 extractor 支持的 kind，_kind_and_direction 直接 _fail → SystemExit(2)。
    # 消息判别器「无法解析 Mermaid 源码」证明走的是 SystemExit 归一化分支
    # （区别于空图守卫的「未提取到任何节点或边」与 kind 检查的「仅支持 flowchart」）。
    with pytest.raises(ValueError, match="无法解析 Mermaid 源码"):
        parse_mermaid("gantt\ntitle x")


def test_mermaid_missing_direction_clear_error():
    # extractor 契约：graph/flowchart 必须带显式方向，裸 graph 会报 "not a Mermaid file"；
    # offipy 包装预检后给清晰 ValueError（含"方向"），不透传误导性消息。
    with pytest.raises(ValueError, match="方向"):
        parse_mermaid("graph\n    A[开始] --> B[处理]")


def test_layout_td_layers_and_coords():
    diagram = parse_mermaid(
        "graph TD\n    A[开始] --> B[处理]\n    B --> C[输出]\n    B --> D[结束]"
    )
    lay = layout_diagram(diagram)
    # A 在 (0,0)，B 在第二层，C/D 在第三层（各占一层内的列）
    a = next(n for n in lay.nodes if n.id == "A")
    b = next(n for n in lay.nodes if n.id == "B")
    c = next(n for n in lay.nodes if n.id == "C")
    d = next(n for n in lay.nodes if n.id == "D")
    assert b.y > a.y  # 分层：A/B 不同行
    assert c.y > b.y
    assert c.y == pytest.approx(d.y)  # C/D 同层同 y
    assert c.x != d.x  # 同层不同列
    assert all(n.w > 0 and n.h > 0 for n in lay.nodes)


def test_layout_lr_swaps_axis():
    diagram = parse_mermaid("graph LR\n    A --> B")
    lay = layout_diagram(diagram)
    a = next(n for n in lay.nodes if n.id == "A")
    b = next(n for n in lay.nodes if n.id == "B")
    assert b.x > a.x  # LR: 水平前进
    assert b.y == a.y  # 同层同列


def test_layout_fits_max():
    diagram = parse_mermaid("graph TD\n    A --> B\n    B --> C\n    C --> D")
    lay = layout_diagram(diagram, max_w=3.0, max_h=2.0)
    assert lay.canvas_w <= 3.0 + 1e-9
    assert lay.canvas_h <= 2.0 + 1e-9


def test_layout_cycle_tolerated():
    diagram = parse_mermaid("graph TD\n    A --> B\n    B --> A")
    lay = layout_diagram(diagram)
    assert {n.id for n in lay.nodes} == {"A", "B"}


def test_kahn_cycle_residue_unified_tail_layer():
    # #84：纯环归一层，分层结果与 ids 迭代序无关（不再散到多层依赖输入顺序）
    d = parse_mermaid("graph TD\n    A --> B\n    B --> C\n    C --> A")
    leaves = [n for n in d.nodes if not n.container]
    layer, layers = _kahn_layers([n.id for n in leaves], d.edges)
    assert set(layer.values()) == {0}  # 纯环统一一层
    assert sorted(layers[0]) == ["A", "B", "C"]
    # 倒序 ids：layer 分配不变
    layer2, _ = _kahn_layers(list(reversed([n.id for n in leaves])), d.edges)
    assert dict(layer2) == dict(layer)


def test_kahn_dag_feeding_cycle_tail_layer():
    # #84：DAG 节点正常分层；环残留（成环节点及其下游）统一落 DAG 最大层 + 1
    d = parse_mermaid("graph TD\n    D --> A\n    A --> B\n    B --> C\n    C --> A")
    leaves = [n for n in d.nodes if not n.container]
    layer, _ = _kahn_layers([n.id for n in leaves], d.edges)
    assert layer["D"] == 0
    assert layer["A"] == layer["B"] == layer["C"]
    assert layer["A"] == 1  # = max(DAG layer) + 1


def test_layout_routes_container_endpoint_edges():
    # 容器框落地后，容器端点边 A→Z 被路由到容器框（源锚底部中心，TD 方向）。
    diagram = parse_mermaid("graph TD\n    subgraph A\n        X --> Y\n    end\n    A --> Z")
    lay = layout_diagram(diagram)
    box = next(n for n in lay.nodes if n.is_container)
    assert box.id == "A"
    assert box.x >= 0 and box.y >= 0
    x = next(n for n in lay.nodes if n.id == "X")
    y = next(n for n in lay.nodes if n.id == "Y")
    assert box.x <= x.x and box.y <= y.y
    assert box.x + box.w >= x.x + x.w
    assert box.y + box.h >= y.y + y.h
    e = next(e for e in lay.edges if e.source == "A" and e.target == "Z")
    assert e.ax1 == pytest.approx(box.x + box.w / 2)
    assert e.ay1 == pytest.approx(box.y + box.h)
    assert len(lay.edges) == 2


def test_layout_subgraph_container_bbox():
    diagram = parse_mermaid(
        "graph TD\n    subgraph 数据处理\n        A[清洗] --> B[聚合]\n    end\n    B --> C[输出]"
    )
    lay = layout_diagram(diagram)
    cont = [n for n in lay.nodes if n.is_container]
    assert len(cont) == 1
    box = cont[0]
    # 容器框必须包住其子节点 A/B
    a = next(n for n in lay.nodes if n.id == "A")
    b = next(n for n in lay.nodes if n.id == "B")
    assert box.x <= a.x and box.y <= a.y
    assert box.x + box.w >= b.x + b.w
    assert box.y + box.h >= b.y + b.h
    # 顶/左缘不越过画布原点（防容器标题被裁到画布顶外）
    assert box.x >= 0 and box.y >= 0


def test_render_subgraph_container_shape():
    diagram = parse_mermaid("graph TD\n    subgraph 数据处理\n        A[清洗] --> B[聚合]\n    end")
    prs, slide = _render(diagram)
    texts = {sh.text_frame.text for sh in slide.shapes if sh.has_text_frame}
    assert "数据处理" in texts
    assert texts >= {"清洗", "聚合"}
    # 容器背景框（RECTANGLE）与叶节点都是 autoshape
    auto = [sh for sh in slide.shapes if sh.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE]
    assert len(auto) == 3


def test_render_edge_label_text():
    diagram = parse_mermaid("graph TD\n    A[发] -->|是| B[收]")
    prs, slide = _render(diagram)
    texts = {sh.text_frame.text for sh in slide.shapes if sh.has_text_frame}
    assert "是" in texts


def test_layout_nested_subgraph_bbox():
    # 嵌套 subgraph：外层框必须包住内层框（内层是外层的跨层后代），内层框包住叶子。
    diagram = parse_mermaid(
        "graph TD\n"
        "    subgraph 外层\n"
        "        subgraph 内层\n"
        "            A[开始] --> B[结束]\n"
        "        end\n"
        "    end"
    )
    lay = layout_diagram(diagram)
    boxes = {n.label: n for n in lay.nodes if n.is_container}
    assert set(boxes) == {"外层", "内层"}
    outer, inner = boxes["外层"], boxes["内层"]
    assert outer.x <= inner.x and outer.y <= inner.y
    assert outer.x + outer.w >= inner.x + inner.w
    assert outer.y + outer.h >= inner.y + inner.h
    a = next(n for n in lay.nodes if n.id == "A")
    b = next(n for n in lay.nodes if n.id == "B")
    assert inner.x <= a.x and inner.y <= a.y
    assert inner.x + inner.w >= b.x + b.w
    assert inner.y + inner.h >= b.y + b.h


@pytest.mark.parametrize(
    "direction",
    ["BT", "RL"],
)
def test_layout_reversed_directions_anchors(direction):
    # 反向方向：BT 层 1 在层 0 上方（y 更小），RL 层 1 在层 0 左侧（x 更小）。
    # 同时锁边锚点（BT: 源顶部/目标底部；RL: 源左缘/目标右缘）。
    diagram = parse_mermaid(f"graph {direction}\n    A --> B")
    lay = layout_diagram(diagram)
    a = next(n for n in lay.nodes if n.id == "A")
    b = next(n for n in lay.nodes if n.id == "B")
    if direction == "BT":
        assert b.y < a.y
        assert b.x == a.x
    else:
        assert b.x < a.x
        assert b.y == a.y
    edge = lay.edges[0]
    if direction == "BT":
        assert edge.ax1 == pytest.approx(a.x + a.w / 2)
        assert edge.ay1 == pytest.approx(a.y)
        assert edge.ax2 == pytest.approx(b.x + b.w / 2)
        assert edge.ay2 == pytest.approx(b.y + b.h)
    else:  # RL
        assert edge.ax1 == pytest.approx(a.x)
        assert edge.ay1 == pytest.approx(a.y + a.h / 2)
        assert edge.ax2 == pytest.approx(b.x + b.w)
        assert edge.ay2 == pytest.approx(b.y + b.h / 2)


def _render(diagram, max_w=12.0, max_h=6.75):
    lay = layout_diagram(diagram, max_w=max_w, max_h=max_h)
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    render_to_slide(slide, lay)
    return prs, slide


def _shape_map(slide):
    return {sh.shape_type: sh for sh in slide.shapes}


def test_render_creates_editable_autoshapes():
    diagram = parse_mermaid("graph TD\n    A[开始] --> B{判断} --> C[结束]")
    prs, slide = _render(diagram)
    auto = [sh for sh in slide.shapes if sh.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE]
    conn = [sh for sh in slide.shapes if sh.shape_type == MSO_SHAPE_TYPE.LINE]
    assert len(auto) == 3
    assert len(conn) == 2
    texts = {sh.text_frame.text for sh in slide.shapes if sh.has_text_frame}
    assert texts >= {"开始", "判断", "结束"}


def test_render_shape_mapping_rhombus():
    diagram = parse_mermaid("graph TD\n    A{决策} --> B[结果]")
    prs, slide = _render(diagram)
    diamond = [
        sh
        for sh in slide.shapes
        if sh.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and sh.auto_shape_type == MSO_SHAPE.DIAMOND
    ]
    assert len(diamond) == 1
    assert diamond[0].text_frame.text == "决策"


def test_render_arrowheads_present():
    diagram = parse_mermaid("graph TD\n    A --> B")
    prs, slide = _render(diagram)
    conns = [sh for sh in slide.shapes if sh.shape_type == MSO_SHAPE_TYPE.LINE]
    assert len(conns) == 1
    ln = conns[0].line._get_or_add_ln()
    tail = ln.find(qn("a:tailEnd"))
    assert tail is not None
    assert tail.get("type") == "triangle"


def test_render_dashed_edge_when_dotted():
    diagram = parse_mermaid("graph TD\n    A -.-> B")
    prs, slide = _render(diagram)
    conns = [sh for sh in slide.shapes if sh.shape_type == MSO_SHAPE_TYPE.LINE]
    assert conns and conns[0].line.dash_style == MSO_LINE_DASH_STYLE.DASH


def test_render_chinese_font_set():
    diagram = parse_mermaid("graph TD\n    A[数据] --> B[管道]")
    prs, slide = _render(diagram)
    assert slide.shapes
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        for para in sh.text_frame.paragraphs:
            assert para.runs
            for run in para.runs:
                rPr = run._r.find(qn("a:rPr"))
                assert rPr is not None
                for tag in ("a:latin", "a:ea", "a:cs"):
                    el = rPr.find(qn(tag))
                    assert el is not None, f"{tag} 未设置"
                    assert el.get("typeface") == "Microsoft YaHei"


def test_render_no_arrowhead_when_undirected():
    diagram = parse_mermaid("graph TD\n    A --- B")
    prs, slide = _render(diagram)
    conns = [sh for sh in slide.shapes if sh.shape_type == MSO_SHAPE_TYPE.LINE]
    assert conns
    ln = conns[0].line._get_or_add_ln()
    assert ln.find(qn("a:tailEnd")) is None


def _write_source_file(tmp_path, text):
    p = tmp_path / "input.mmd"
    p.write_text(text, encoding="utf-8")
    return str(p)


def test_mermaid_to_pptx_from_text(tmp_path):
    out = str(tmp_path / "out.pptx")
    result = mermaid_to_pptx("graph TD\n    A[开始] --> B[结束]", out)
    assert result == out
    prs = Presentation(out)
    assert prs.slide_width == Emu(12192000)  # 16:9
    texts = {sh.text_frame.text for sl in prs.slides for sh in sl.shapes if sh.has_text_frame}
    assert texts >= {"开始", "结束"}


def test_mermaid_to_pptx_from_path(tmp_path):
    src = _write_source_file(tmp_path, "graph LR\n    A --> B")
    out = str(tmp_path / "out.pptx")
    mermaid_to_pptx(src, out)
    prs = Presentation(out)
    assert len(prs.slides[0].shapes) >= 2


def test_mermaid_to_pptx_bad_source(tmp_path):
    with pytest.raises(ValueError, match="无法解析"):
        mermaid_to_pptx("not mermaid at all", str(tmp_path / "x.pptx"))


def test_mermaid_to_pptx_missing_file_raises(tmp_path):
    # PathLike 明确表示路径：缺失 → FileNotFoundError，而非当 Mermaid 文本解析
    src = tmp_path / "nope.mmd"
    with pytest.raises(FileNotFoundError):
        mermaid_to_pptx(src, str(tmp_path / "out.pptx"))


def test_read_source_missing_path_like_str_raises():
    # #85：str 入参含路径形态（盘符 / 绝对 / 点相对段）但不存在 → FileNotFoundError，
    # 与 PathLike 分支一致；不再静默当文本解析（否则是误导性「无法解析 Mermaid 源码」）
    for raw in (r"C:\nope\dir\flow.mmd", "/nope/dir/flow.mmd", "./nope/flow.mmd"):
        with pytest.raises(FileNotFoundError, match="源文件不存在"):
            _read_source(raw)


def test_mermaid_to_pptx_missing_path_str_raises(tmp_path):
    # #85：mermaid_to_pptx 传不存在的路径字符串 → FileNotFoundError（可定位「路径打错」）
    with pytest.raises(FileNotFoundError, match="源文件不存在"):
        mermaid_to_pptx(r"C:\nope\dir\flow.mmd", str(tmp_path / "out.pptx"))


def test_mermaid_to_pptx_bad_text_without_path_shape_is_value_error(tmp_path):
    # #85：无路径形态的坏文本仍当作文本 → 「无法解析」ValueError（路径启发不误伤）
    with pytest.raises(ValueError, match="无法解析"):
        mermaid_to_pptx("not mermaid at all", str(tmp_path / "x.pptx"))


def test_looks_like_path_relative_with_separator():
    # #92：无前缀相对路径（data/foo.mmd、foo\bar.mmd）识别为路径意图；
    # 含 Mermaid 结构标记的文本（graph/-->）不当路径，避免误伤 [/.../] 平行四边形节点
    from offipy.diagrams import _looks_like_path

    for raw in ("data/foo.mmd", "foo\\bar.mmd", "sub/dir/flow.mmd"):
        assert _looks_like_path(raw), raw
    assert not _looks_like_path("graph TD\n    A[/输入/] --> B[处理]")
    assert not _looks_like_path("graph TD; A --> B")
    assert not _looks_like_path("plain text no separator")


def test_read_source_relative_missing_path_raises():
    # #92：#85 只覆盖前缀形态，data/foo.mmd、foo\bar.mmd 这类含分隔符的相对路径
    # 文件不存在时不再静默当 Mermaid 文本 → FileNotFoundError（可定位「路径打错」）
    for raw in ("data/foo.mmd", "foo\\bar.mmd"):
        with pytest.raises(FileNotFoundError, match="源文件不存在"):
            _read_source(raw)


def test_read_source_mermaid_text_with_slashes_is_text():
    # #92：合法 Mermaid 源码含 /（[/.../] 平行四边形节点、label 斜杠）仍按文本处理，
    # 不因含分隔符被误判为路径 → 原样返回
    text = "graph TD\n    A[/输入/] --> B[处理]"
    assert _read_source(text) == text


def test_hex_rgb_three_states():
    from offipy.diagrams import _hex_rgb

    assert _hex_rgb("#dae8fc") == RGBColor(0xDA, 0xE8, 0xFC)
    assert _hex_rgb("#dae8fcff") == RGBColor(0xDA, 0xE8, 0xFC)  # 带 alpha 取前 6 位
    assert _hex_rgb("#abc") == RGBColor(0xAA, 0xBB, 0xCC)  # #RGB 容错
    assert _hex_rgb("none") is None  # 透明
    assert _hex_rgb("") is None  # 未指定
    assert _hex_rgb("not-a-color") is None  # 非法 hex 宽容


def test_render_fill_stroke_font_color_branches():
    from offipy.diagrams import DiagramLayout, PlacedEdge, PlacedNode

    lay = DiagramLayout(
        nodes=[
            PlacedNode(
                "n1",
                "彩色",
                "rectangle",
                0.5,
                0.5,
                2.0,
                0.8,
                fill="#dae8fc",
                stroke="#6c8ebf",
                font_color="#ff0000",
            ),
            PlacedNode("n2", "透明", "rectangle", 3.5, 0.5, 2.0, 0.8, fill="none", stroke="none"),
            PlacedNode("n3", "默认", "rectangle", 0.5, 2.5, 2.0, 0.8),
        ],
        edges=[PlacedEdge("n1", "n3", "连", stroke="#333333")],
        canvas_w=6.0,
        canvas_h=4.0,
    )
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    render_to_slide(slide, lay)
    by_text = {sh.text_frame.text: sh for sh in slide.shapes if sh.has_text_frame}
    c = by_text["彩色"]
    assert c.fill.fore_color.rgb == RGBColor(0xDA, 0xE8, 0xFC)
    assert c.line.color.rgb == RGBColor(0x6C, 0x8E, 0xBF)
    assert c.text_frame.paragraphs[0].runs[0].font.color.rgb == RGBColor(0xFF, 0x00, 0x00)
    assert c.auto_shape_type == MSO_SHAPE.RECTANGLE
    t = by_text["透明"]
    spPr = t._element.find(qn("p:spPr"))
    assert spPr is not None
    assert spPr.find(qn("a:noFill")) is not None
    assert spPr.find(qn("a:ln")).find(qn("a:noFill")) is not None
    d = by_text["默认"]
    assert d.fill.fore_color.rgb == RGBColor(0xFF, 0xFF, 0xFF)
    conns = [sh for sh in slide.shapes if sh.shape_type == MSO_SHAPE_TYPE.LINE]
    assert conns and conns[0].line.color.rgb == RGBColor(0x33, 0x33, 0x33)
