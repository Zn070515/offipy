"""audit autofit 规则：normAutofit→shrink / spAutoFit→grow 分开；fontScale/估算字号记录。"""

from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Inches, Pt

from offipy.audit.extract import extract_presentation
from offipy.audit.models import (
    RULE_AUTOFIT_GROW,
    RULE_AUTOFIT_SHRINK,
    AuditConfig,
    Severity,
)
from offipy.audit.rules import run_rules


def _run(prs, tmp_path, config=None):
    path = tmp_path / "x.pptx"
    prs.save(path)
    ext = extract_presentation(path)
    records = [r for slide in ext.slides for r in slide.shapes]
    findings, suppressed = run_rules(records, ext.slide_size, config or AuditConfig())
    return findings, suppressed, records


def _add_tb(slide, x, y, text, w=1, h=0.5, font_size=None, autofit=None, font_scale=None):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tb.text = text
    tf = tb.text_frame
    if font_size is not None:
        tf.paragraphs[0].runs[0].font.size = Pt(font_size)
    if autofit is not None:
        tf.auto_size = autofit
    if font_scale is not None:
        norm = tf._txBody.xpath("./a:bodyPr/a:normAutofit")
        assert norm, "设置 TEXT_TO_FIT_SHAPE 后应存在 normAutofit"
        norm[0].set("fontScale", str(font_scale))
    return tb


def _by_rule(findings, rule_id):
    return [f for f in findings if f.rule_id == rule_id]


def test_shrink_mid_records_scale(tmp_path):
    # normAutofit + fontScale 0.5，18pt → 估算 9pt ≥ 8 → MID
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_tb(
        slide,
        1,
        1,
        "hello world",
        w=2,
        h=1,
        font_size=18,
        autofit=MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE,
        font_scale=50000,
    )
    findings, _, _ = _run(prs, tmp_path)
    shrinks = _by_rule(findings, RULE_AUTOFIT_SHRINK)
    assert len(shrinks) == 1
    f = shrinks[0]
    assert f.severity == Severity.MID
    assert f.confidence == 1.0
    assert f.details["original_font_size_pt"] == 18
    assert f.details["font_scale"] == 0.5
    assert f.details["estimated_font_size_pt"] == 9.0


def test_shrink_below_readable_high(tmp_path):
    # 36pt × 0.2 = 7.2pt < 8pt 最小可读 → HIGH
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_tb(
        slide,
        1,
        1,
        "hello",
        w=2,
        h=1,
        font_size=36,
        autofit=MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE,
        font_scale=20000,
    )
    findings, _, _ = _run(prs, tmp_path)
    shrinks = _by_rule(findings, RULE_AUTOFIT_SHRINK)
    assert len(shrinks) == 1
    assert shrinks[0].severity == Severity.HIGH
    assert shrinks[0].details["estimated_font_size_pt"] == 7.2


def test_shrink_no_scale_no_overflow_no_finding(tmp_path):
    # normAutofit 但无 fontScale 且文本不溢出 → 无风险不报
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_tb(slide, 1, 1, "hi", w=2, h=1, autofit=MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE)
    findings, _, _ = _run(prs, tmp_path)
    assert _by_rule(findings, RULE_AUTOFIT_SHRINK) == []


def test_grow_mid_off_page_false(tmp_path):
    # 默认 spAutoFit，文本溢出撑大后仍在页内 → MID
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_tb(slide, 1, 1, "W" * 20, w=1, h=0.5)
    findings, _, _ = _run(prs, tmp_path)
    grows = _by_rule(findings, RULE_AUTOFIT_GROW)
    assert len(grows) == 1
    f = grows[0]
    assert f.severity == Severity.MID
    assert f.confidence == 1.0
    assert f.details["text_width_in"] > f.details["avail_width_in"]


def test_grow_off_page_high(tmp_path):
    # 靠右文本框撑大后越过页面右边界 → HIGH
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_tb(slide, 8, 1, "W" * 20, w=1, h=0.5)
    findings, _, _ = _run(prs, tmp_path)
    grows = _by_rule(findings, RULE_AUTOFIT_GROW)
    assert len(grows) == 1
    assert grows[0].severity == Severity.HIGH


def test_grow_no_overflow_no_finding(tmp_path):
    # 现有框能装下 → 不会撑大 → 不报
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_tb(slide, 1, 1, "hi", w=2, h=1)
    findings, _, _ = _run(prs, tmp_path)
    assert _by_rule(findings, RULE_AUTOFIT_GROW) == []
