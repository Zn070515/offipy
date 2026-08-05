"""audit bounds 规则：部分越界/完全出画布/容差内不报/背景不误报/不与 margin 双报。"""

from pptx import Presentation
from pptx.util import Inches

from offipy.audit.extract import extract_presentation
from offipy.audit.models import (
    RULE_BOUNDS_OFF_CANVAS,
    RULE_BOUNDS_PARTIAL,
    AuditConfig,
    Severity,
)
from offipy.audit.rules import run_rules


def _run(prs, tmp_path, config=None):
    path = tmp_path / "x.pptx"
    prs.save(path)
    ext = extract_presentation(path)
    records = [r for slide in ext.slides for r in slide.shapes]
    findings, suppressed = run_rules(records, ext.slide_size, config or AuditConfig())
    return findings, suppressed, records


def _add_shape(slide, x, y, w, h):
    return slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))


def test_partial_right_overshoot_mid(tmp_path):
    # 右边界越出 0.1 英寸，未到 HIGH 阈值 → MID partial
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_shape(slide, 9.6, 1, 0.5, 1)
    findings, suppressed, _ = _run(prs, tmp_path)
    assert len(findings) == 1
    f = findings[0]
    assert f.rule_id == RULE_BOUNDS_PARTIAL
    assert f.severity == Severity.MID
    assert f.details["edges"] == ["right"]
    assert f.details["max_overshoot_in"] == 0.1
    assert suppressed == []


def test_partial_large_overshoot_high(tmp_path):
    # 右边界越出 3.0 英寸 → HIGH partial
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_shape(slide, 9, 1, 4, 1)
    findings, _, _ = _run(prs, tmp_path)
    assert len(findings) == 1
    assert findings[0].rule_id == RULE_BOUNDS_PARTIAL
    assert findings[0].severity == Severity.HIGH


def test_off_canvas_mid_for_large_shape(tmp_path):
    # 完全在画布外且面积 >=1 → MID off_canvas
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_shape(slide, 11, 2, 1, 1)
    findings, _, _ = _run(prs, tmp_path)
    assert len(findings) == 1
    assert findings[0].rule_id == RULE_BOUNDS_OFF_CANVAS
    assert findings[0].severity == Severity.MID


def test_off_canvas_low_for_tiny_shape(tmp_path):
    # 完全在画布外但面积 <1 → LOW off_canvas（暂存/装饰残留候选）
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_shape(slide, 11, 2, 0.5, 0.5)
    findings, _, _ = _run(prs, tmp_path)
    assert len(findings) == 1
    assert findings[0].rule_id == RULE_BOUNDS_OFF_CANVAS
    assert findings[0].severity == Severity.LOW


def test_within_tolerance_no_finding(tmp_path):
    # 越界 0.005 < bounds_tolerance_in(0.01) → 不报
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_shape(slide, 9.995, 1, 0.01, 1)
    findings, suppressed, _ = _run(prs, tmp_path)
    assert findings == []
    assert suppressed == []


def test_full_bleed_background_no_false_positive(tmp_path):
    # 全页背景无文本 → 不产生 finding，贴边 margin 走 full_bleed 豁免进 suppressed
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_shape(slide, 0, 0, 10, 7.5)
    findings, suppressed, _ = _run(prs, tmp_path)
    assert findings == []
    assert suppressed and all(s.reason == "full_bleed" for s in suppressed)


def test_ignored_shape_goes_to_suppressed(tmp_path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    sh = _add_shape(slide, 9.6, 1, 0.5, 1)
    path = tmp_path / "x.pptx"
    prs.save(path)
    ext = extract_presentation(path)
    records = [r for slide in ext.slides for r in slide.shapes]
    findings, suppressed = run_rules(
        records, ext.slide_size, AuditConfig(ignored_shapes={(1, sh.shape_id)})
    )
    assert findings == []
    assert len(suppressed) == 1
    assert suppressed[0].reason == "user_shape"


def test_bounds_no_margin_double_report(tmp_path):
    # 右边界越出只报 bounds，不额外报同方向 margin
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_shape(slide, 9.6, 1, 0.5, 1)
    findings, _, _ = _run(prs, tmp_path)
    assert len(findings) == 1
    assert findings[0].rule_id == RULE_BOUNDS_PARTIAL
