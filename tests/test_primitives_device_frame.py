"""A5 Task 8 — device-frame renderer structural contract."""

from __future__ import annotations

import pytest
from pptx import Presentation
from pptx.enum.dml import MSO_FILL_TYPE
from pptx.enum.shapes import MSO_SHAPE_TYPE

from offipy.assets.model import (
    AssetRect,
    AssetRef,
    AssetRenderContext,
    AssetRequest,
    NativeShapePayload,
)
from offipy.assets.primitives import get_native_renderer
from offipy.assets.providers.primitives import PrimitivesProvider
from offipy.exceptions import InvalidArgumentError

_THEME = {
    "accent": "#2251ff",
    "surface": "#F2F4F7",
    "ink": "#222222",
    "muted": "#667085",
    "bg": "#ffffff",
}

_DEVICES = ("phone", "tablet", "desktop")


def _ctx(**overrides) -> AssetRenderContext:
    kwargs = {
        "slide_index": 1,
        "rect": AssetRect(x=0, y=0, width=400, height=200),
        "theme_name": "mckinsey",
        "theme_vars": dict(_THEME),
        "placement": "replace",
    }
    kwargs.update(overrides)
    return AssetRenderContext(**kwargs)


def _blank_slide():
    prs = Presentation()
    return prs.slides.add_slide(prs.slide_layouts[6])


def _render(params: dict[str, str], **ctx_overrides):
    resolved = PrimitivesProvider().resolve(
        AssetRequest(
            ref=AssetRef("primitives", "primitive", "device-frame"),
            params=tuple(params.items()),
        )
    )
    assert isinstance(resolved.payload, NativeShapePayload)
    slide = _blank_slide()
    get_native_renderer("device-frame")(slide, dict(resolved.payload.params), _ctx(**ctx_overrides))
    return slide


def _solid_fill_shapes(slide, hex_color: str) -> list:
    target = hex_color.lstrip("#").upper()
    out = []
    for sp in slide.shapes:
        if sp.shape_type != MSO_SHAPE_TYPE.AUTO_SHAPE:
            continue
        if sp.fill.type != MSO_FILL_TYPE.SOLID:
            continue
        if str(sp.fill.fore_color.rgb) == target:
            out.append(sp)
    return out


def _assert_within_rect(slide, rect: AssetRect, tol_px: float = 1.0) -> None:
    tol = round(tol_px * 6350)
    x0, y0 = int(rect.x * 6350), int(rect.y * 6350)
    x1, y1 = int((rect.x + rect.width) * 6350), int((rect.y + rect.height) * 6350)
    for sp in slide.shapes:
        assert sp.left >= x0 - tol, f"left overflow: {sp.left}"
        assert sp.top >= y0 - tol, f"top overflow: {sp.top}"
        assert sp.left + sp.width <= x1 + tol, f"right overflow: {sp.left + sp.width}"
        assert sp.top + sp.height <= y1 + tol, f"bottom overflow: {sp.top + sp.height}"


@pytest.mark.parametrize("device", _DEVICES)
def test_device_frame_shape_count(device: str) -> None:
    slide = _render({"device": device})
    expected = 4 if device in ("phone", "desktop") else 3
    assert len(list(slide.shapes)) == expected


@pytest.mark.parametrize("device", _DEVICES)
def test_device_frame_within_rect(device: str) -> None:
    slide = _render({"device": device})
    _assert_within_rect(slide, _ctx().rect)


@pytest.mark.parametrize("device", _DEVICES)
def test_device_frame_no_pictures(device: str) -> None:
    slide = _render({"device": device})
    assert not any(sp.shape_type == MSO_SHAPE_TYPE.PICTURE for sp in slide.shapes)


@pytest.mark.parametrize("device", _DEVICES)
def test_device_frame_custom_fill_screen_color(device: str) -> None:
    slide = _render({"device": device, "fill": "#123456"})
    screens = _solid_fill_shapes(slide, "#123456")
    assert len(screens) == 1, f"expected exactly one screen with fill 123456, got {len(screens)}"


@pytest.mark.parametrize("device", ["phone", "desktop"])
def test_device_frame_accent_highlight_color(device: str) -> None:
    slide = _render({"device": device, "accent": "#FF0000"})
    highlights = _solid_fill_shapes(slide, "#FF0000")
    assert len(highlights) == 1, f"expected exactly one accent shape, got {len(highlights)}"


@pytest.mark.parametrize("device", ["phone", "desktop"])
def test_device_frame_tiny_rect_raises(device: str) -> None:
    with pytest.raises(InvalidArgumentError, match="too small"):
        _render({"device": device}, rect=AssetRect(x=0, y=0, width=30, height=30))
