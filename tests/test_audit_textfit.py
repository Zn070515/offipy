"""audit text-fit 规则：Pillow 字体度量 vs 字符权重回退、无可用空间、横/纵溢出。"""

from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches, Pt

from offipy.audit.extract import extract_presentation
from offipy.audit.models import (
    RULE_TEXT_FIT_HORIZONTAL,
    RULE_TEXT_FIT_VERTICAL,
    AuditConfig,
    Severity,
)
from offipy.audit.rules import _DEFAULT_FONT_SIZE_PT, _text_height_in, run_rules


def _run(prs, tmp_path, config=None):
    path = tmp_path / "x.pptx"
    prs.save(path)
    ext = extract_presentation(path)
    records = [r for slide in ext.slides for r in slide.shapes]
    findings, suppressed = run_rules(records, ext.slide_size, config or AuditConfig())
    return findings, suppressed, records


def _add_tb(slide, x, y, text, w=1, h=0.5, font_name=None, font_size=None):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tb.text = text
    run = tb.text_frame.paragraphs[0].runs[0]
    if font_name is not None:
        run.font.name = font_name
    if font_size is not None:
        run.font.size = Pt(font_size)
    return tb


def _by_rule(findings, rule_id):
    return [f for f in findings if f.rule_id == rule_id]


def _any_font_path() -> str | None:
    for path in (
        r"C:\Windows\Fonts\arial.ttf",
        "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
    ):
        if Path(path).exists():
            return path
    return None


def test_fallback_char_weight_low_confidence(tmp_path):
    # 找不到字体文件 → 字符权重回退，confidence 0.4，消息标注低置信
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_tb(slide, 1, 1, "WWWWWWWWWW", font_name="NonexistentFontXYZ", font_size=18)
    findings, _, _ = _run(prs, tmp_path)
    fits = _by_rule(findings, RULE_TEXT_FIT_HORIZONTAL)
    assert len(fits) == 1
    f = fits[0]
    assert f.severity == Severity.LOW
    assert f.confidence == 0.4
    assert "字符估算低置信" in f.message
    assert f.details["text_width_in"] == 1.25


def test_pillow_metrics_high_confidence(tmp_path):
    font_path = _any_font_path()
    if font_path is None:
        pytest.skip("无可用系统字体，跳过 Pillow 度量测试")
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_tb(slide, 1, 1, "WWWWWWWWWWWWWWWWWWWW", font_name=font_path, font_size=18)
    findings, _, _ = _run(prs, tmp_path)
    fits = _by_rule(findings, RULE_TEXT_FIT_HORIZONTAL)
    assert len(fits) == 1
    f = fits[0]
    assert f.severity == Severity.LOW
    assert f.confidence == 0.8
    assert "字符估算低置信" not in f.message


def test_vertical_overflow_low(tmp_path):
    # 5 行文字在 0.5 英寸高框内 → 纵向溢出 LOW
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_tb(slide, 1, 1, "line1\nline2\nline3\nline4\nline5", w=2, h=0.5)
    findings, _, _ = _run(prs, tmp_path)
    verts = _by_rule(findings, RULE_TEXT_FIT_VERTICAL)
    assert len(verts) == 1
    f = verts[0]
    assert f.severity == Severity.LOW
    assert f.details["text_height_in"] > f.details["avail_height_in"]


def test_no_usable_space_mid(tmp_path):
    # 内边距吃尽可用区域 → 无可用空间 MID
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_tb(slide, 1, 1, "x", w=0.1, h=0.1)
    findings, _, _ = _run(prs, tmp_path)
    fits = _by_rule(findings, RULE_TEXT_FIT_HORIZONTAL)
    assert len(fits) == 1
    assert fits[0].severity == Severity.MID


def _add_seg_tb(slide, x, y, w, h, segs, wrap, font_name="NonexistentFontXYZ", font_size=18):
    """多段（a:br 软换行）文本框；Nonexistent 字体强制走字符权重，宽度可预测。"""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    for i, seg in enumerate(segs):
        if i > 0:
            p.add_line_break()
        r = p.add_run()
        r.text = seg
        r.font.name = font_name
        r.font.size = Pt(font_size)
    return tb


def test_multi_segment_horizontal_uses_max_not_sum(tmp_path):
    # 段落含 a:br：多段求和（旧 S2#15 16.25in 假象）会误报，应取最长段。
    # 两段各 10 个 W（18pt × 0.5 权重 = 1.25in），求和 2.5in > avail 2.3in，
    # 但最长段 1.25in < avail 2.3in → 不报横溢。
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_seg_tb(slide, 1, 1, 2.5, 1.0, ["WWWWWWWWWW", "WWWWWWWWWW"], wrap=False)
    findings, _, _ = _run(prs, tmp_path)
    assert _by_rule(findings, RULE_TEXT_FIT_HORIZONTAL) == []
    assert _by_rule(findings, RULE_TEXT_FIT_VERTICAL) == []  # 2 段 0.6in < avail 0.8in


def test_multi_segment_genuine_overflow_still_reported(tmp_path):
    # 最长段确实超过框宽 → 仍报横溢，且宽度是段宽（2.5in）不是求和（2.75in）
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_seg_tb(slide, 1, 1, 1.0, 1.0, ["WWWWWWWWWWWWWWWWWWWW", "WW"], wrap=False)
    findings, _, _ = _run(prs, tmp_path)
    fits = _by_rule(findings, RULE_TEXT_FIT_HORIZONTAL)
    assert len(fits) == 1
    assert fits[0].details["text_width_in"] == pytest.approx(2.5, abs=0.02)


def test_unset_wrap_default_square_no_horizontal(tmp_path):
    # bodyPr@wrap 未设 → PowerPoint 默认 square（自动折行），不报横溢（旧 not None 误报）
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4.0), Inches(1.0))
    tf = tb.text_frame
    tf.paragraphs[0].add_run().text = "WWWWWWWWWWWWWWWWWWWW"
    bodyPr = tf._txBody.xpath("./a:bodyPr")[0]
    bodyPr.attrib.pop("wrap", None)  # 移除 wrap → 未设
    findings, _, records = _run(prs, tmp_path)
    assert records[0].word_wrap is None
    assert _by_rule(findings, RULE_TEXT_FIT_HORIZONTAL) == []


def test_explicit_wrap_none_horizontal_reported(tmp_path):
    # 显式 wrap=none（不折行）单行超宽 → 报横溢
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_seg_tb(slide, 1, 1, 1.0, 1.0, ["WWWWWWWWWWWWWWWWWWWW"], wrap=False)
    findings, _, _ = _run(prs, tmp_path)
    fits = _by_rule(findings, RULE_TEXT_FIT_HORIZONTAL)
    assert len(fits) == 1
    assert fits[0].severity == Severity.LOW


def test_vertical_counts_no_wrap_segments(tmp_path):
    # wrap=none 多段：纵向按段数计行（3 段 × 18pt × 1.2 / 72 = 0.9in），非 1 行
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_seg_tb(slide, 1, 1, 2.0, 0.5, ["short", "line", "three"], wrap=False)
    findings, _, _ = _run(prs, tmp_path)
    verts = _by_rule(findings, RULE_TEXT_FIT_VERTICAL)
    assert len(verts) == 1
    assert verts[0].details["text_height_in"] == pytest.approx(0.9, abs=0.01)


def test_vertical_uses_spcpts_line_spacing(tmp_path):
    # a:lnSpc spcPts 绝对值决定行高（回归 S13 定位结论：行距 14.85pt 而非字号×1.2）
    # 2 行 × 14.85pt = 0.4125in 装进 0.55in 框 → 不报纵向溢出；
    # 旧 字号×1.2（2×18×1.2=0.6in）则会误报。
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = _add_seg_tb(slide, 1, 1, 2.0, 0.55, ["lineA", "lineB"], wrap=False)
    tb.text_frame.margin_top = 0
    tb.text_frame.margin_bottom = 0
    tb.text_frame.paragraphs[0].line_spacing = Pt(14.85)
    findings, _, records = _run(prs, tmp_path)
    rec = records[0]
    assert rec.paragraphs[0].line_spacing_pts == pytest.approx(14.85)
    text_h, _ = _text_height_in(rec, rec.width, _DEFAULT_FONT_SIZE_PT)
    assert text_h == pytest.approx(2 * 14.85 / 72.0)  # 0.4125in，非 0.6in
    assert _by_rule(findings, RULE_TEXT_FIT_VERTICAL) == []


def test_vertical_uses_spcpct_line_spacing(tmp_path):
    # a:lnSpc spcPct 百分比（150%）→ 行高 = 字号×1.5；3 段 × 18×1.5 = 1.125in
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = _add_seg_tb(slide, 1, 1, 2.0, 1.0, ["a", "b", "c"], wrap=False)
    tb.text_frame.margin_top = 0
    tb.text_frame.margin_bottom = 0
    tb.text_frame.paragraphs[0].line_spacing = 1.5
    _, _, records = _run(prs, tmp_path)
    rec = records[0]
    assert rec.paragraphs[0].line_spacing_pct == pytest.approx(150.0)
    text_h, _ = _text_height_in(rec, rec.width, _DEFAULT_FONT_SIZE_PT)
    assert text_h == pytest.approx(3 * 18 * 1.5 / 72.0)


def test_vertical_trailing_soft_break_not_extra_line(tmp_path):
    # 尾部 a:br 不渲染额外一行：2 实线 + 尾部 br（3 段含空段）→ 按 2 行计高
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2.0), Inches(0.7))
    tf = tb.text_frame
    tf.word_wrap = False
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.add_run().text = "lineA"
    p.add_line_break()
    p.add_run().text = "lineB"
    p.add_line_break()  # 尾部 br → 末尾空段
    findings, _, records = _run(prs, tmp_path)
    rec = records[0]
    assert len(rec.paragraphs[0].segments) == 3  # 提取层保留空段
    text_h, _ = _text_height_in(rec, rec.width, _DEFAULT_FONT_SIZE_PT)
    assert text_h == pytest.approx(2 * 18 * 1.2 / 72.0)  # 0.6in，非 0.9in
    assert _by_rule(findings, RULE_TEXT_FIT_VERTICAL) == []  # 0.6 < 0.7+1pt


# ---------------------------------------------------------------- 字体定位


def test_font_candidates_include_msyh_ttc():
    # 微软雅黑是 TTC 集合，base 拼 .ttf 找不到 → 必须返回 msyh.ttc 候选
    from offipy.audit.rules import _font_candidates

    cands = _font_candidates("Microsoft YaHei", False)
    assert any(p.name.lower() == "msyh.ttc" for p in cands)
    assert any(p.name.lower() == "msyhbd.ttc" for p in _font_candidates("Microsoft YaHei", True))


def test_msyh_metric_high_confidence_when_available(tmp_path):
    # 本机有 msyh.ttc → 真实 Pillow 度量 confidence 0.8（回归：旧版找不到字体回退 0.4）
    from offipy.audit.rules import _font_candidates

    if not any(p.exists() for p in _font_candidates("Microsoft YaHei", False)):
        pytest.skip("本机无 msyh.ttc，跳过 Pillow 实测")
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_tb(slide, 1, 1, "WWWWWWWWWW", font_name="Microsoft YaHei", font_size=18)
    findings, _, _ = _run(prs, tmp_path)
    fits = _by_rule(findings, RULE_TEXT_FIT_HORIZONTAL)
    assert len(fits) == 1
    assert fits[0].confidence == 0.8
    assert "字符估算低置信" not in fits[0].message


def test_horizontal_overflow_below_floor_not_reported(tmp_path):
    # 溢出 < 1pt（引擎度量噪声下限）：Pillow(FreeType) 与 PowerPoint(DirectWrite)
    # 度量 msyh 有亚 pt 级差，低于下限的「超出」在 PowerPoint 渲染刚好放下。
    # 10 W @18pt 字符权重 = 1.25in；框 1.443in → avail 1.243in，溢出 0.007in(<1pt)
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_seg_tb(slide, 1, 1, 1.443, 1.0, ["WWWWWWWWWW"], wrap=False)
    findings, _, _ = _run(prs, tmp_path)
    assert _by_rule(findings, RULE_TEXT_FIT_HORIZONTAL) == []


def test_run_width_pt_applies_kerning():
    # msyhbd "FPS · TensorRT FP16" 22pt bold：fontTools hmtx+kern ≈ DirectWrite 3.1261in
    # （Pillow getlength 无 kern 会高估 ~1.3% 到 3.16in）——治理拉丁+粗体高估的根因
    from offipy.audit.rules import _font_candidates, _run_width_pt

    if not any(p.exists() for p in _font_candidates("Microsoft YaHei", True)):
        pytest.skip("本机无 msyhbd.ttc，跳过 kerning 度量测试")
    w_pt, used = _run_width_pt("FPS · TensorRT FP16", 22, True, "Microsoft YaHei")
    assert used is True
    assert w_pt / 72.0 == pytest.approx(3.1264, abs=0.01)


def test_horizontal_overflow_above_floor_reported(tmp_path):
    # 溢出 > 1pt 仍报：框 1.3in → avail 1.1in，溢出 0.15in(10.8pt)
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_seg_tb(slide, 1, 1, 1.3, 1.0, ["WWWWWWWWWW"], wrap=False)
    findings, _, _ = _run(prs, tmp_path)
    fits = _by_rule(findings, RULE_TEXT_FIT_HORIZONTAL)
    assert len(fits) == 1
    assert fits[0].severity == Severity.LOW
