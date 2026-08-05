"""audit 角色：背景/页码/页眉页脚/重复装饰识别，不能全局忽略纯数字短文本。"""

from pptx import Presentation
from pptx.util import Inches

from offipy.audit.extract import extract_presentation
from offipy.audit.roles import classify_presentation


def _extract(prs: Presentation, tmp_path) -> list:
    path = tmp_path / "x.pptx"
    prs.save(path)
    ext = extract_presentation(path)
    records = [r for slide in ext.slides for r in slide.shapes]
    classify_presentation(records, ext.slide_size)
    return records


def _by_name(records: list, name: str):
    return next(r for r in records if r.name == name)


# ---------------------------------------------------------------- 单页角色


def test_full_bleed_rectangle_is_background(tmp_path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(7.5))
    records = _extract(prs, tmp_path)
    assert records[0].role == "background"


def test_full_page_text_is_content_not_background(tmp_path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(0), Inches(0), Inches(10), Inches(7.5))
    tb.text = "hello content"
    records = _extract(prs, tmp_path)
    assert records[0].role == "content"


def test_partial_cover_shape_not_background(tmp_path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_shape(1, Inches(1), Inches(1), Inches(3), Inches(3))
    records = _extract(prs, tmp_path)
    assert records[0].role == "unknown"


def test_page_number_bottom_pure_digit(tmp_path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(9), Inches(7.0), Inches(0.8), Inches(0.3))
    tb.text = "12"
    records = _extract(prs, tmp_path)
    assert records[0].role == "page_number"


def test_pure_digit_top_is_content(tmp_path):
    # 不能全局忽略所有纯数字短文本：顶部纯数字不是页码
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(0.8), Inches(0.3))
    tb.text = "12"
    records = _extract(prs, tmp_path)
    assert records[0].role == "content"


def test_pure_digit_bottom_large_is_content(tmp_path):
    # 底部但尺寸大的纯数字不是页码
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(9), Inches(6.6), Inches(1), Inches(1))
    tb.text = "12"
    records = _extract(prs, tmp_path)
    assert records[0].role == "content"


def test_title_placeholder_role(tmp_path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title Slide
    slide.shapes.title.text = "T"
    records = _extract(prs, tmp_path)
    title = _by_name(records, "Title 1")
    assert title.role == "title"


def test_subtitle_placeholder_role(tmp_path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    sub = slide.placeholders[1]
    sub.text = "sub"
    records = _extract(prs, tmp_path)
    assert _by_name(records, sub.name).role == "content"


# ---------------------------------------------------------------- 跨页重复


def _repeated_pages(n: int, build):
    prs = Presentation()
    for _ in range(n):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        build(slide)
    return prs


def _add_rect(slide, x, y, w=0.5, h=0.5):
    return slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))


def _add_tb(slide, x, y, text, w=2, h=0.3):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tb.text = text
    return tb


def test_repeated_logo_is_decoration(tmp_path):
    prs = _repeated_pages(3, lambda s: _add_rect(s, 0.1, 0.1))
    records = _extract(prs, tmp_path)
    assert records[0].role == "decoration"


def test_repeated_header_text(tmp_path):
    prs = _repeated_pages(3, lambda s: _add_tb(s, 4, 0.1, "ACME CORP"))
    records = _extract(prs, tmp_path)
    assert records[0].role == "header"


def test_repeated_footer_text(tmp_path):
    def build(s):
        tb = s.shapes.add_textbox(Inches(4), Inches(7.1), Inches(2), Inches(0.3))
        tb.text = "FOOTER NOTE"

    prs = _repeated_pages(3, build)
    records = _extract(prs, tmp_path)
    assert records[0].role == "footer"


def test_repeated_middle_text_is_decoration(tmp_path):
    def build(s):
        tb = s.shapes.add_textbox(Inches(4), Inches(3.5), Inches(2), Inches(0.3))
        tb.text = "LOGOMARK"

    prs = _repeated_pages(3, build)
    records = _extract(prs, tmp_path)
    assert records[0].role == "decoration"


def test_single_shape_not_repeated_stays_unknown(tmp_path):
    prs = _repeated_pages(3, lambda s: None)
    prs.slides[0].shapes.add_shape(1, Inches(0.1), Inches(0.1), Inches(0.5), Inches(0.5))
    records = _extract(prs, tmp_path)
    lonely = [r for r in records if r.shape_type == "AUTO_SHAPE"][0]
    assert lonely.role == "unknown"


def test_same_shape_different_positions_not_decoration(tmp_path):
    # 同形状但每页位置不同 → 指纹不同 → 不抑制
    prs = Presentation()
    for x, y in [(0.1, 0.1), (3.0, 2.0), (6.0, 5.0)]:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(0.5), Inches(0.5))
    records = _extract(prs, tmp_path)
    assert all(r.role == "unknown" for r in records)
