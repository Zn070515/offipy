"""audit overlap 规则：部分重叠/完全覆盖分类、卡片容器抑制、线条/隐藏/极小点跳过。"""

from io import BytesIO

from pptx import Presentation
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.util import Inches

from offipy.audit.extract import extract_presentation
from offipy.audit.models import (
    RULE_OVERLAP_COVERED_TEXT,
    RULE_OVERLAP_PARTIAL,
    AuditConfig,
    Severity,
)
from offipy.audit.rules import run_rules

_PNG_1PX = bytes.fromhex(
    "89504e470d0a1a0a0000000d49484452000000010000000108060000001f15c489"
    "0000000d4944415478da63f8cfc0f01f000500018002f1b8d5f10000000049454e44ae426082"
)


def _run(prs, tmp_path, config=None):
    path = tmp_path / "x.pptx"
    prs.save(path)
    ext = extract_presentation(path)
    records = [r for slide in ext.slides for r in slide.shapes]
    findings, suppressed = run_rules(records, ext.slide_size, config or AuditConfig())
    return findings, suppressed, records


def _add_rect(slide, x, y, w, h):
    return slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))


def _add_picture(slide, x, y, w, h):
    return slide.shapes.add_picture(BytesIO(_PNG_1PX), Inches(x), Inches(y), Inches(w), Inches(h))


def test_partial_overlap_mid(tmp_path):
    # A(2,2,4,2) 与 B(3,2,4,2)：交集 3×2，ratio=0.75 → MID partial
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    a = _add_rect(slide, 2, 2, 4, 2)
    b = _add_rect(slide, 3, 2, 4, 2)
    findings, suppressed, _ = _run(prs, tmp_path)
    assert len(findings) == 1
    f = findings[0]
    assert f.rule_id == RULE_OVERLAP_PARTIAL
    assert f.severity == Severity.MID
    assert f.primary.shape_id == a.shape_id
    assert f.secondary.shape_id == b.shape_id
    assert f.details["overlap_ratio"] == 0.75
    assert suppressed == []


def test_no_overlap_no_finding(tmp_path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(slide, 1, 1, 2, 2)
    _add_rect(slide, 5, 5, 2, 2)
    findings, suppressed, _ = _run(prs, tmp_path)
    assert findings == []
    assert suppressed == []


def test_text_in_autoshape_suppressed_container(tmp_path):
    # 文本位于填充 AutoShape 内且在上层 → 卡片容器，不报 overlap
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(slide, 1, 1, 4, 3)  # card 先加（在下层）
    tb = slide.shapes.add_textbox(Inches(1.5), Inches(1.5), Inches(2), Inches(1))
    tb.text = "card text"
    findings, suppressed, _ = _run(prs, tmp_path)
    assert findings == []
    assert len(suppressed) == 1
    assert suppressed[0].reason == "intentional_containment"


def test_picture_covers_text_high(tmp_path):
    # 图片在上层完全盖住文本 → HIGH covered_text
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(3), Inches(3))
    tb.text = "hidden text"
    pic = _add_picture(slide, 0.5, 0.5, 3, 3)
    findings, suppressed, _ = _run(prs, tmp_path)
    assert len(findings) == 1
    f = findings[0]
    assert f.rule_id == RULE_OVERLAP_COVERED_TEXT
    assert f.severity == Severity.HIGH
    assert f.primary.shape_id == tb.shape_id
    assert f.secondary.shape_id == pic.shape_id
    assert suppressed == []


def test_caption_over_picture_no_finding(tmp_path):
    # 文本题注盖在图片上 → 正常配图，不报 overlap
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_picture(slide, 0.5, 0.5, 3, 3)
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(3), Inches(3))
    tb.text = "caption"
    findings, suppressed, _ = _run(prs, tmp_path)
    assert findings == []
    assert suppressed == []


def test_same_group_partial_low(tmp_path):
    # 同一 group 内两子形状部分重叠 → LOW partial
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    g = slide.shapes.add_group_shape()
    g.shapes.add_shape(1, Inches(0), Inches(0), Inches(2), Inches(2))
    g.shapes.add_shape(1, Inches(1), Inches(0), Inches(2), Inches(2))
    g.left, g.top = Inches(2), Inches(2)
    findings, suppressed, _ = _run(prs, tmp_path)
    assert len(findings) == 1
    assert findings[0].rule_id == RULE_OVERLAP_PARTIAL
    assert findings[0].severity == Severity.LOW
    assert suppressed == []


def test_connector_skipped(tmp_path):
    # 连接线与形状重叠 → 连接线不参与 overlap
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(slide, 1, 1, 4, 4)
    slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1), Inches(2), Inches(5), Inches(2))
    findings, suppressed, _ = _run(prs, tmp_path)
    assert findings == []
    assert suppressed == []


def test_hidden_shape_skipped(tmp_path):
    # 隐藏形状与可见形状重叠 → 隐藏的不参与
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    h = _add_rect(slide, 0, 0, 3, 3)
    h._element.xpath("./p:nvSpPr/p:cNvPr")[0].set("hidden", "1")
    _add_rect(slide, 1, 1, 3, 3)
    findings, suppressed, _ = _run(prs, tmp_path)
    assert findings == []
    assert suppressed == []


def test_tiny_shape_skipped(tmp_path):
    # 极小装饰点（面积 <0.0025）与正常形状重叠 → 极小点不参与
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(slide, 1, 1, 3, 3)
    _add_rect(slide, 2, 2, 0.04, 0.04)  # 面积 0.0016
    findings, suppressed, _ = _run(prs, tmp_path)
    assert findings == []
    assert suppressed == []
