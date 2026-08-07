"""A5 Task 5 — metric-badge renderer structural contract."""

from __future__ import annotations

import pytest
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from offipy.assets.model import (
    AssetRect,
    AssetRef,
    AssetRenderContext,
    AssetRequest,
    NativeShapePayload,
)
from offipy.assets.primitives import get_native_renderer
from offipy.assets.providers.primitives import PrimitivesProvider
from offipy.exceptions import InvalidArgumentError

_THEME = {
    "accent": "#2251ff",
    "surface": "#F2F4F7",
    "ink": "#222222",
    "muted": "#667085",
    "bg": "#ffffff",
}


def _ctx(**overrides) -> AssetRenderContext:
    kwargs = dict(
        slide_index=1,
        rect=AssetRect(x=0, y=0, width=400, height=200),
        theme_name="mckinsey",
        theme_vars=dict(_THEME),
        placement="replace",
    )
    kwargs.update(overrides)
    return AssetRenderContext(**kwargs)


def _blank_slide():
    prs = Presentation()
    return prs.slides.add_slide(prs.slide_layouts[6])


def _render(params: dict[str, str], **ctx_overrides):
    resolved = PrimitivesProvider().resolve(
        AssetRequest(
            ref=AssetRef("primitives", "primitive", "metric-badge"),
            params=tuple(params.items()),
        )
    )
    assert isinstance(resolved.payload, NativeShapePayload)
    slide = _blank_slide()
    get_native_renderer("metric-badge")(slide, dict(resolved.payload.params), _ctx(**ctx_overrides))
    return slide


def _runs(sp):
    return [r.text for p in sp.text_frame.paragraphs for r in p.runs]


def _texts(slide) -> list[str]:
    return [t for sp in slide.shapes if hasattr(sp, "text_frame") for t in _runs(sp)]


def _autoshapes(slide):
    return [sp for sp in slide.shapes if sp.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE]


def _assert_within_rect(slide, rect: AssetRect, tol_px: float = 1.0) -> None:
    tol = int(round(tol_px * 6350))
    x0, y0 = int(rect.x * 6350), int(rect.y * 6350)
    x1, y1 = int((rect.x + rect.width) * 6350), int((rect.y + rect.height) * 6350)
    for sp in slide.shapes:
        assert sp.left >= x0 - tol, f"left overflow: {sp.left}"
        assert sp.top >= y0 - tol, f"top overflow: {sp.top}"
        assert sp.left + sp.width <= x1 + tol, f"right overflow: {sp.left + sp.width}"
        assert sp.top + sp.height <= y1 + tol, f"bottom overflow: {sp.top + sp.height}"


def test_metric_badge_only_value() -> None:
    slide = _render({"value": "42"})
    assert len(list(slide.shapes)) == 2  # surface card + value
    assert "42" in _texts(slide)


def test_metric_badge_value_and_label() -> None:
    slide = _render({"value": "42", "label": "YoY"})
    assert len(list(slide.shapes)) == 3
    assert "42" in _texts(slide)
    assert "YoY" in _texts(slide)


def test_metric_badge_all_three() -> None:
    slide = _render({"value": "24%", "label": "YoY", "delta": "+3.2%"})
    assert len(list(slide.shapes)) == 4
    for text in ("24%", "YoY", "+3.2%"):
        assert text in _texts(slide)


def test_metric_badge_custom_fill() -> None:
    slide = _render({"value": "42", "fill": "#123456"})
    card = _autoshapes(slide)[0]
    assert str(card.fill.fore_color.rgb) == "123456"


def test_metric_badge_shapes_within_rect() -> None:
    slide = _render({"value": "24%", "label": "YoY", "delta": "+3.2%"})
    _assert_within_rect(slide, _ctx().rect)


def test_metric_badge_long_value_stays_in_rect() -> None:
    rect = AssetRect(x=0, y=0, width=800, height=200)
    slide = _render({"value": "x" * 80}, rect=rect)
    _assert_within_rect(slide, rect)
    assert "x" * 80 in _texts(slide)


def test_metric_badge_value_too_long_raises() -> None:
    rect = AssetRect(x=0, y=0, width=200, height=200)
    with pytest.raises(InvalidArgumentError, match="does not fit"):
        _render({"value": "x" * 80}, rect=rect)


def test_metric_badge_no_pictures() -> None:
    slide = _render({"value": "42", "label": "L", "delta": "+1%"})
    assert not any(sp.shape_type == MSO_SHAPE_TYPE.PICTURE for sp in slide.shapes)
