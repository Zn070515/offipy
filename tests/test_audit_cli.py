"""offipy audit CLI：参数/配置映射、门槛退出码 0/1/2/3、main() 不劫持 audit 的码。"""

import argparse
import json
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches

from offipy.cli import _build_audit_config, _main, build_parser, main


def _rect(slide, x, y, w, h):
    return slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))


def _deck_high(tmp_path):
    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(s, 9.5, 3, 3, 1)  # bounds HIGH
    path = tmp_path / "high.pptx"
    prs.save(path)
    return str(path)


def _deck_low(tmp_path):
    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(s, 9.85, 5, 0.1, 1)  # margin LOW
    path = tmp_path / "low.pptx"
    prs.save(path)
    return str(path)


def _pair(tmp_path):
    base = Presentation()
    s = base.slides.add_slide(base.slide_layouts[6])
    _rect(s, 9.85, 5, 0.1, 1)  # margin LOW
    base_path = tmp_path / "base.pptx"
    base.save(base_path)

    cand = Presentation()
    s = cand.slides.add_slide(cand.slide_layouts[6])
    _rect(s, 5, 5, 0.1, 1)  # margin 已解决
    _rect(s, 9.5, 3, 3, 1)  # 新增 bounds HIGH
    cand_path = tmp_path / "cand.pptx"
    cand.save(cand_path)
    return str(base_path), str(cand_path)


# ---------------------------------------------------------------- 参数解析


def test_audit_parser_accepts_flags():
    args = build_parser().parse_args(
        [
            "audit",
            "deck.pptx",
            "--format",
            "markdown",
            "--fail-on",
            "MID",
            "--slides-dir",
            "imgs",
            "--no-page-number-ignore",
            "--safe-margin",
            "0.5",
        ]
    )
    assert args.app == "audit"
    assert args.file == "deck.pptx"
    assert args.format == "markdown"
    assert args.fail_on == "MID"
    assert args.slides_dir == "imgs"
    assert args.no_page_number_ignore is True
    assert args.safe_margin == 0.5


def test_audit_bad_format_parser_exit2():
    with pytest.raises(SystemExit) as exc:
        build_parser().parse_args(["audit", "x.pptx", "--format", "bogus"])
    assert exc.value.code == 2


def test_build_audit_config_flags():
    ns = argparse.Namespace(
        safe_margin=0.5,
        bounds_tolerance=0.05,
        no_full_bleed_ignore=True,
        no_repeated_decoration_ignore=True,
        no_page_number_ignore=True,
        no_header_footer_ignore=True,
    )
    cfg = _build_audit_config(ns)
    assert cfg.safe_margin_in == 0.5
    assert cfg.bounds_tolerance_in == 0.05
    assert cfg.ignore_full_bleed_shapes is False
    assert cfg.ignore_repeated_decorations is False
    assert cfg.ignore_page_numbers is False
    assert cfg.ignore_headers_footers is False


# ---------------------------------------------------------------- 门槛退出码


def test_audit_fail_on_gate(tmp_path):
    assert _main(["audit", _deck_high(tmp_path), "--format", "text"]) == 1  # HIGH≥HIGH
    assert _main(["audit", _deck_low(tmp_path), "--format", "text"]) == 0  # LOW<HIGH
    assert _main(["audit", _deck_low(tmp_path), "--fail-on", "LOW"]) == 1
    assert _main(["audit", _deck_high(tmp_path), "--fail-on", "LOW"]) == 1


def test_audit_fail_on_new_gate(tmp_path):
    base, cand = _pair(tmp_path)
    assert _main(["audit", cand, "--baseline", base, "--fail-on-new", "MID"]) == 1
    assert _main(["audit", cand, "--baseline", base, "--fail-on-new", "HIGH"]) == 1


def test_audit_fail_on_new_not_triggered(tmp_path):
    base = _deck_low(tmp_path)
    cand = Presentation()
    s = cand.slides.add_slide(cand.slide_layouts[6])
    _rect(s, 5, 5, 0.1, 1)  # 修复 margin
    cand_path = tmp_path / "fixed.pptx"
    cand.save(cand_path)
    assert _main(["audit", str(cand_path), "--baseline", base, "--fail-on-new", "LOW"]) == 0


# ---------------------------------------------------------------- 输入/解析错误


def test_audit_missing_file_returns_2(tmp_path):
    assert _main(["audit", str(tmp_path / "nope.pptx")]) == 2


def test_audit_wrong_extension_returns_2(tmp_path):
    bad = tmp_path / "bad.ppt"
    bad.write_bytes(b"x")
    assert _main(["audit", str(bad)]) == 2


def test_audit_corrupted_pptx_returns_3(tmp_path):
    bad = tmp_path / "bad.pptx"
    bad.write_bytes(b"not a zip")
    assert _main(["audit", str(bad)]) == 3


def test_audit_fail_on_new_without_baseline_returns_2(tmp_path):
    assert _main(["audit", _deck_high(tmp_path), "--fail-on-new", "MID"]) == 2


# ---------------------------------------------------------------- main() 不劫持


def test_main_does_not_hijack_audit_codes(tmp_path):
    assert main(["audit", str(tmp_path / "nope.pptx")]) == 2  # 输入错，不是 1
    assert main(["audit", _deck_high(tmp_path)]) == 1  # 门槛命中仍是 1


# ---------------------------------------------------------------- 输出


def test_audit_json_format(capsys, tmp_path):
    assert _main(["audit", _deck_high(tmp_path), "--format", "json"]) == 1
    data = json.loads(capsys.readouterr().out)
    assert data["max_severity"] == "HIGH"


def test_audit_html_writes_default_out(tmp_path):
    path = _deck_high(tmp_path)
    code = _main(["audit", path, "--format", "html"])
    out = Path(path).with_name("high.audit.html")
    assert out.exists()
    assert "<svg" in out.read_text(encoding="utf-8")
    assert code == 1


def test_audit_out_writes_file(tmp_path, capsys):
    out = tmp_path / "report.md"
    _main(["audit", _deck_high(tmp_path), "--format", "markdown", "--out", str(out)])
    assert out.exists()
    assert "# 审计报告" in out.read_text(encoding="utf-8")
