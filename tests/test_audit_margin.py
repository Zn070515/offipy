"""audit margin 规则：贴边 LOW；角色豁免（背景/页码/页眉页脚）进 suppressed；用户豁免。"""

from pptx import Presentation
from pptx.util import Inches

from offipy.audit.extract import extract_presentation
from offipy.audit.models import (
    RULE_MARGIN_BOTTOM,
    RULE_MARGIN_LEFT,
    RULE_MARGIN_RIGHT,
    RULE_MARGIN_TOP,
    AuditConfig,
    Severity,
)
from offipy.audit.rules import run_rules


def _run(prs, tmp_path, config=None):
    path = tmp_path / "x.pptx"
    prs.save(path)
    ext = extract_presentation(path)
    records = [r for slide in ext.slides for r in slide.shapes]
    findings, suppressed = run_rules(records, ext.slide_size, config or AuditConfig())
    return findings, suppressed, records


def _add_tb(slide, x, y, text, w=2, h=1):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tb.text = text
    return tb


def test_left_edge_low(tmp_path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_tb(slide, 0.05, 1, "left")
    findings, suppressed, _ = _run(prs, tmp_path)
    assert len(findings) == 1
    f = findings[0]
    assert f.rule_id == RULE_MARGIN_LEFT
    assert f.severity == Severity.LOW
    assert f.details["gap_in"] == 0.05
    assert f.details["required_in"] == 0.2
    assert suppressed == []


def test_right_edge_low(tmp_path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_tb(slide, 8.0, 1, "right")
    findings, _, _ = _run(prs, tmp_path)
    assert len(findings) == 1
    assert findings[0].rule_id == RULE_MARGIN_RIGHT
    assert findings[0].severity == Severity.LOW


def test_top_edge_low(tmp_path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_tb(slide, 1, 0.05, "top")
    findings, _, _ = _run(prs, tmp_path)
    assert len(findings) == 1
    assert findings[0].rule_id == RULE_MARGIN_TOP
    assert findings[0].severity == Severity.LOW


def test_bottom_edge_low(tmp_path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_tb(slide, 1, 6.55, "bottom", w=2, h=0.8)
    findings, _, _ = _run(prs, tmp_path)
    assert len(findings) == 1
    assert findings[0].rule_id == RULE_MARGIN_BOTTOM
    assert findings[0].severity == Severity.LOW


def test_safe_margin_ok_no_finding(tmp_path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_tb(slide, 0.3, 1, "safe")
    findings, suppressed, _ = _run(prs, tmp_path)
    assert findings == []
    assert suppressed == []


def test_background_suppressed_full_bleed(tmp_path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(7.5))
    findings, suppressed, _ = _run(prs, tmp_path)
    assert findings == []
    assert len(suppressed) >= 1
    assert all(s.reason == "full_bleed" for s in suppressed)


def test_page_number_suppressed(tmp_path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(9.8), Inches(7.1), Inches(0.2), Inches(0.2))
    tb.text = "12"
    findings, suppressed, _ = _run(prs, tmp_path)
    assert findings == []
    assert len(suppressed) == 1
    assert suppressed[0].reason == "page_number"


def test_repeated_header_suppressed(tmp_path):
    prs = Presentation()
    for _ in range(3):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _add_tb(slide, 4, 0.1, "ACME CORP", w=2, h=0.3)
    findings, suppressed, _ = _run(prs, tmp_path)
    assert findings == []
    assert len(suppressed) == 3
    assert all(s.reason == "header_footer" for s in suppressed)


def test_ignored_region_suppressed_user_region(tmp_path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_tb(slide, 0.05, 1, "x")
    path = tmp_path / "x.pptx"
    prs.save(path)
    ext = extract_presentation(path)
    records = [r for slide in ext.slides for r in slide.shapes]
    findings, suppressed = run_rules(
        records, ext.slide_size, AuditConfig(ignored_regions=[(0.0, 0.0, 2.0, 3.0)])
    )
    assert findings == []
    assert len(suppressed) == 1
    assert suppressed[0].reason == "user_region"


def test_ignored_shape_suppressed_user_shape(tmp_path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = _add_tb(slide, 0.05, 1, "x")
    path = tmp_path / "x.pptx"
    prs.save(path)
    ext = extract_presentation(path)
    records = [r for slide in ext.slides for r in slide.shapes]
    findings, suppressed = run_rules(
        records, ext.slide_size, AuditConfig(ignored_shapes={(1, tb.shape_id)})
    )
    assert findings == []
    assert len(suppressed) == 1
    assert suppressed[0].reason == "user_shape"
