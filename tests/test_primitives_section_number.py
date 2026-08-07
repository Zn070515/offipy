"""A5 Task 4 — section-number renderer structural contract."""

from __future__ import annotations

import pytest
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from offipy.assets.model import (
    AssetRect,
    AssetRef,
    AssetRenderContext,
    AssetRequest,
    NativeShapePayload,
)
from offipy.assets.primitives import get_native_renderer
from offipy.assets.providers.primitives import PrimitivesProvider

_THEME = {
    "accent": "#2251ff",
    "surface": "#F2F4F7",
    "ink": "#222222",
    "muted": "#667085",
    "bg": "#ffffff",
}


def _ctx(**overrides) -> AssetRenderContext:
    kwargs = dict(
        slide_index=1,
        rect=AssetRect(x=0, y=0, width=400, height=200),
        theme_name="mckinsey",
        theme_vars=dict(_THEME),
        placement="replace",
    )
    kwargs.update(overrides)
    return AssetRenderContext(**kwargs)


def _blank_slide():
    prs = Presentation()
    return prs.slides.add_slide(prs.slide_layouts[6])


def _render(params: dict[str, str], **ctx_overrides):
    resolved = PrimitivesProvider().resolve(
        AssetRequest(
            ref=AssetRef("primitives", "primitive", "section-number"),
            params=tuple(params.items()),
        )
    )
    assert isinstance(resolved.payload, NativeShapePayload)
    slide = _blank_slide()
    get_native_renderer("section-number")(
        slide, dict(resolved.payload.params), _ctx(**ctx_overrides)
    )
    return slide


def _runs(sp):
    return [r.text for p in sp.text_frame.paragraphs for r in p.runs]


def _texts(slide) -> list[str]:
    return [t for sp in slide.shapes if hasattr(sp, "text_frame") for t in _runs(sp)]


def _assert_within_rect(slide, rect: AssetRect, tol_px: float = 1.0) -> None:
    tol = int(round(tol_px * 6350))
    x0, y0 = int(rect.x * 6350), int(rect.y * 6350)
    x1, y1 = int((rect.x + rect.width) * 6350), int((rect.y + rect.height) * 6350)
    for sp in slide.shapes:
        assert sp.left >= x0 - tol, f"left overflow: {sp.left}"
        assert sp.top >= y0 - tol, f"top overflow: {sp.top}"
        assert sp.left + sp.width <= x1 + tol, f"right overflow: {sp.left + sp.width}"
        assert sp.top + sp.height <= y1 + tol, f"bottom overflow: {sp.top + sp.height}"


@pytest.mark.parametrize("num", ["0", "1", "9999"])
def test_section_number_digits_round_trip(num: str) -> None:
    slide = _render({"number": num})
    assert num in _texts(slide)


def test_section_number_label_absent() -> None:
    slide = _render({"number": "3"})
    assert len(list(slide.shapes)) == 2  # accent bar + number
    assert "Intro" not in _texts(slide)


def test_section_number_label_present() -> None:
    slide = _render({"number": "3", "label": "Intro"})
    assert len(list(slide.shapes)) == 3
    assert "Intro" in _texts(slide)


def test_section_number_text_editable() -> None:
    slide = _render({"number": "42"})
    boxes = [sp for sp in slide.shapes if hasattr(sp, "text_frame")]
    assert any("42" in _runs(sp) for sp in boxes)


def test_section_number_shapes_within_rect() -> None:
    slide = _render({"number": "9999", "label": "A label"})
    _assert_within_rect(slide, _ctx().rect)


def test_section_number_no_pictures() -> None:
    slide = _render({"number": "7", "label": "L"})
    assert not any(sp.shape_type == MSO_SHAPE_TYPE.PICTURE for sp in slide.shapes)
