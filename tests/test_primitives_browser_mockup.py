"""A5 Task 9 — browser-mockup renderer structural contract."""

from __future__ import annotations

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from offipy.assets.model import (
    AssetRect,
    AssetRef,
    AssetRenderContext,
    AssetRequest,
    NativeShapePayload,
)
from offipy.assets.primitives import get_native_renderer
from offipy.assets.providers.primitives import PrimitivesProvider

_THEME = {
    "accent": "#2251ff",
    "surface": "#F2F4F7",
    "ink": "#222222",
    "muted": "#667085",
    "bg": "#ffffff",
}


def _ctx(**overrides) -> AssetRenderContext:
    kwargs = {
        "slide_index": 1,
        "rect": AssetRect(x=0, y=0, width=400, height=200),
        "theme_name": "mckinsey",
        "theme_vars": dict(_THEME),
        "placement": "replace",
    }
    kwargs.update(overrides)
    return AssetRenderContext(**kwargs)


def _blank_slide():
    prs = Presentation()
    return prs.slides.add_slide(prs.slide_layouts[6])


def _render(params: dict[str, str], **ctx_overrides):
    resolved = PrimitivesProvider().resolve(
        AssetRequest(
            ref=AssetRef("primitives", "primitive", "browser-mockup"),
            params=tuple(params.items()),
        )
    )
    assert isinstance(resolved.payload, NativeShapePayload)
    slide = _blank_slide()
    get_native_renderer("browser-mockup")(
        slide, dict(resolved.payload.params), _ctx(**ctx_overrides)
    )
    return slide


def _runs(sp):
    return [r.text for p in sp.text_frame.paragraphs for r in p.runs]


def _texts(slide) -> list[str]:
    return [t for sp in slide.shapes if hasattr(sp, "text_frame") for t in _runs(sp)]


def _assert_within_rect(slide, rect: AssetRect, tol_px: float = 1.0) -> None:
    tol = round(tol_px * 6350)
    x0, y0 = int(rect.x * 6350), int(rect.y * 6350)
    x1, y1 = int((rect.x + rect.width) * 6350), int((rect.y + rect.height) * 6350)
    for sp in slide.shapes:
        assert sp.left >= x0 - tol, f"left overflow: {sp.left}"
        assert sp.top >= y0 - tol, f"top overflow: {sp.top}"
        assert sp.left + sp.width <= x1 + tol, f"right overflow: {sp.left + sp.width}"
        assert sp.top + sp.height <= y1 + tol, f"bottom overflow: {sp.top + sp.height}"


def test_browser_mockup_empty_window() -> None:
    slide = _render({})
    assert len(list(slide.shapes)) == 7
    assert _texts(slide) == []


def test_browser_mockup_title_only() -> None:
    slide = _render({"title": "Q3 Report"})
    assert len(list(slide.shapes)) == 8
    assert "Q3 Report" in _texts(slide)


def test_browser_mockup_url_only() -> None:
    slide = _render({"url": "https://offipy.dev"})
    assert len(list(slide.shapes)) == 9
    assert "https://offipy.dev" in _texts(slide)


def test_browser_mockup_title_and_url() -> None:
    slide = _render({"title": "Q3 Report", "url": "https://offipy.dev"})
    assert len(list(slide.shapes)) == 10
    texts = _texts(slide)
    assert "Q3 Report" in texts
    assert "https://offipy.dev" in texts


def test_browser_mockup_long_content_stays_in_rect() -> None:
    slide = _render(
        {"title": "x" * 120, "url": "y" * 240},
        rect=AssetRect(x=0, y=0, width=1700, height=240),
    )
    texts = _texts(slide)
    assert "x" * 120 in texts
    assert "y" * 240 in texts
    _assert_within_rect(slide, AssetRect(x=0, y=0, width=1700, height=240))


def test_browser_mockup_no_pictures() -> None:
    slide = _render({"title": "T", "url": "https://offipy.dev"})
    assert not any(sp.shape_type == MSO_SHAPE_TYPE.PICTURE for sp in slide.shapes)
