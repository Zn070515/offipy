"""A5 Task 2 — native rendering foundation helpers.

px→EMU conversion, color resolution from theme vars, font fitting, and the
renderer dispatch guard (unknown primitive / background placement rejection).
"""

from __future__ import annotations

import pytest
from pptx import Presentation
from pptx.enum.dml import MSO_FILL_TYPE
from pptx.enum.shapes import MSO_SHAPE

from offipy.assets.model import (
    AssetRect,
    AssetRenderContext,
    NativeShapePayload,
)
from offipy.assets.primitives import get_native_renderer
from offipy.assets.primitives._common import (
    _char_width_em,
    _descending_pt,
    add_shape,
    add_textbox,
    fit_font_size,
    fit_font_size_wrapped,
    px_to_emu,
    resolve_native_colors,
    shape_elements,
)
from offipy.assets.registry import get_default_registry
from offipy.assets.render import render_asset
from offipy.exceptions import InvalidArgumentError

_THEME = {
    "accent": "#2251ff",
    "surface": "#F2F4F7",
    "ink": "#222222",
    "muted": "#667085",
    "bg": "#ffffff",
}


def _ctx(**overrides) -> AssetRenderContext:
    kwargs = dict(
        slide_index=1,
        rect=AssetRect(x=0, y=0, width=400, height=200),
        theme_name="mckinsey",
        theme_vars=dict(_THEME),
        placement="replace",
    )
    kwargs.update(overrides)
    return AssetRenderContext(**kwargs)


def _blank_slide():
    prs = Presentation()
    return prs.slides.add_slide(prs.slide_layouts[6])


# -- conversion ------------------------------------------------------------


def test_px_to_emu_exact() -> None:
    assert px_to_emu(0) == 0
    assert px_to_emu(1) == 6350
    assert px_to_emu(1920) == 12192000
    assert px_to_emu(1080) == 6858000
    assert px_to_emu(0.5) == 3175
    assert px_to_emu(1.5) == 9525


def test_px_to_emu_rounds_to_nearest() -> None:
    assert px_to_emu(0.0001) == 1  # round(0.635) -> 1
    assert px_to_emu(0.7) == 4445  # round(4445.0)


# -- color resolution ------------------------------------------------------


def test_resolve_native_colors_exact() -> None:
    colors = resolve_native_colors({"accent": "accent", "fill": "surface"}, _ctx())
    assert colors == {
        "accent": "#2251FF",
        "fill": "#F2F4F7",
        "ink": "#222222",
        "muted": "#667085",
    }


def test_resolve_native_colors_hex_uppercases() -> None:
    colors = resolve_native_colors({"accent": "#123456", "fill": "#abcdef"}, _ctx())
    assert colors["accent"] == "#123456"
    assert colors["fill"] == "#ABCDEF"


def test_resolve_native_colors_transparent() -> None:
    colors = resolve_native_colors({"accent": "accent", "fill": "transparent"}, _ctx())
    assert colors["fill"] == "transparent"


def test_resolve_native_colors_missing_token_raises() -> None:
    with pytest.raises(InvalidArgumentError, match="accent"):
        resolve_native_colors({"accent": "accent"}, _ctx(theme_vars={"bg": "#fff"}))


def test_resolve_native_colors_text_fallback_on_missing_ink() -> None:
    colors = resolve_native_colors(
        {"accent": "accent", "fill": "surface"},
        _ctx(theme_vars={"accent": "#2251ff", "surface": "#eee"}),
    )
    assert colors["ink"] == "#222222"
    assert colors["muted"] == "#667085"


# -- font fitting ----------------------------------------------------------


def test_fit_font_size_returns_start_for_short_text() -> None:
    assert fit_font_size("Hi", 400, 200, start_pt=20, min_pt=8) == pytest.approx(20)


def test_fit_font_size_shrinks_for_long_text() -> None:
    size = fit_font_size("x" * 40, 400, 200, start_pt=48, min_pt=6)
    assert 6 <= size < 48
    # 40 chars * size * (96/72) * 0.55 <= 400px  ->  size <= ~13.6
    assert size <= 14


def test_fit_font_size_raises_when_impossible() -> None:
    with pytest.raises(InvalidArgumentError, match="does not fit"):
        fit_font_size("x" * 500, 100, 40, start_pt=24, min_pt=8)


def test_fit_font_size_empty_text() -> None:
    assert fit_font_size("", 400, 200, start_pt=20, min_pt=8) > 0


def test_fit_font_size_wrapped_short_text_stays_single_line() -> None:
    size = fit_font_size_wrapped("Hi there", 400, 200, start_pt=40, min_pt=8)
    assert size == pytest.approx(40)


def test_fit_font_size_wrapped_shrinks_for_long_multiline_text() -> None:
    size = fit_font_size_wrapped(" ".join(["word"] * 20), 400, 80, start_pt=48, min_pt=6)
    assert 6 <= size < 48


def test_fit_font_size_wrapped_unbreakable_token_raises() -> None:
    with pytest.raises(InvalidArgumentError, match="does not fit"):
        fit_font_size_wrapped("x" * 500, 100, 40, start_pt=24, min_pt=8)


def test_fit_font_size_wrapped_empty_text() -> None:
    assert fit_font_size_wrapped("", 400, 200, start_pt=20, min_pt=8) > 0


def test_descending_pt_reaches_exact_fractional_min() -> None:
    # 0.25 分数兜底此前是死代码：非整数 min_pt 永远不会被精确尝试。
    # 修复后序列含 min_pt 与附近 0.25 步进，严格递减、不含 min_pt 以下值。
    sizes = list(_descending_pt(16.0, 12.5))
    assert sizes == [16.0, 15.0, 14.0, 13.0, 12.75, 12.5]
    assert sizes == sorted(sizes, reverse=True)
    assert all(s >= 12.5 - 1e-9 for s in sizes)


def test_descending_pt_integer_min_stays_integer() -> None:
    assert list(_descending_pt(16.0, 12.0)) == [16.0, 15.0, 14.0, 13.0, 12.0]


def test_descending_pt_terminates_at_exact_start() -> None:
    assert list(_descending_pt(20.0, 20.0)) == [20.0]


def test_fit_font_size_uses_fractional_size_between_integers() -> None:
    # 13pt 放不下、12.75pt 放得下的文本：旧实现只试整数并早退，会误报
    # 「does not fit at min 12.5pt」；修复后应精确命中 12.75pt。
    text = "x" * 21  # 13pt → 17.33em*21*0.55 ≈ 200.2px > 199；12.75pt → 196.3px ≤ 199
    size = fit_font_size(text, 199, 100, start_pt=16, min_pt=12.5)
    assert size == pytest.approx(12.75)


def _full_width_cp(cp: int) -> bool:
    """Arial/CJK 下全角（≈1.0em）的码点：#56/#62 的「真实字形」模型。

    CJK 表意字/假名/谚文、全角标点，以及半角码位里的全角字形——箭头
    （U+2190-21FF）、连字符/破折号（U+2010-2015）在 Arial 下也是 1.0em。
    """
    return (
        cp > 0x2E80
        or cp == 0x3000
        or 0x2010 <= cp <= 0x2015
        or 0x2190 <= cp <= 0x21FF
        or 0xFF01 <= cp <= 0xFF60
        or 0xFFE0 <= cp <= 0xFFE6
    )


def _real_wrapped_max_width_px(text: str, size_pt: float, w_px: float) -> float:
    """真实字形宽度（CJK/全角字形 ≈1.0em，拉丁/空格 ≈0.55em）下 greedy wrap 的最宽行。"""
    em_px = size_pt * (96.0 / 72.0)
    space_w_px = em_px * 0.55
    cur = 0.0
    max_w = 0.0
    for token in text.split():
        token_w = em_px * sum(1.0 if _full_width_cp(ord(ch)) else 0.55 for ch in token)
        if cur == 0 or cur + space_w_px + token_w <= w_px:
            cur += (space_w_px if cur else 0.0) + token_w
        else:
            max_w = max(max_w, cur)
            cur = token_w
    return max(max_w, cur)


def test_char_width_em_fullwidth_halfwidth_code_points() -> None:
    # #62：箭头/破折号是半角码位里的全角字形，Arial 下 1.0em，此前被按 0.55em 低估
    assert _char_width_em("→", 0.55) == 1.0
    assert _char_width_em("←", 0.55) == 1.0
    assert _char_width_em("—", 0.55) == 1.0
    assert _char_width_em("–", 0.55) == 1.0
    assert _char_width_em("A", 0.55) == pytest.approx(0.55)
    assert _char_width_em("让", 0.55) == 1.0


def test_fit_font_size_cjk_does_not_overflow() -> None:
    # #56：0.55em 统一估算低估 CJK 全角字形（实际 ~1.0em）→ 中文文本溢出测量矩形。
    # 修复后返回字号下按真实字形宽度模拟 wrapped 折行，最宽行必须 ≤ 矩形宽。
    cases = [
        ("操作点 漏检/误检", 420.0, 240.0),
        ("城市环境智能治理", 400.0, 100.0),
        ("让每处垃圾都被看见", 420.0, 200.0),
    ]
    for text, w, h in cases:
        pt = fit_font_size_wrapped(text, w, h, start_pt=60.0, min_pt=8.0)
        real_w = _real_wrapped_max_width_px(text, pt, w)
        assert real_w <= w + 1e-6, f"{text!r} pt={pt:.1f} real_w={real_w:.1f} > box {w}"


def test_fit_font_size_wrapped_arrow_and_dash_do_not_overflow() -> None:
    # #62 回归：含 →/—— 的图元文本此前按 0.55em 低估 → 溢出测量矩形
    cases = [
        ("FP32 → FP16", 420.0, 200.0),
        ("——让每处垃圾都被看见", 420.0, 200.0),
    ]
    for text, w, h in cases:
        pt = fit_font_size_wrapped(text, w, h, start_pt=60.0, min_pt=8.0)
        real_w = _real_wrapped_max_width_px(text, pt, w)
        assert real_w <= w + 1e-6, f"{text!r} pt={pt:.1f} real_w={real_w:.1f} > box {w}"


# -- shape helpers ---------------------------------------------------------


def test_add_shape_and_textbox_are_editable_shapes() -> None:
    slide = _blank_slide()
    sp = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 10, 20, 100, 50, fill="#2251FF")
    tb = add_textbox(slide, 10, 70, 100, 50, text="Hello", font_size_pt=14, color="#222222")
    assert sp._element.prst == MSO_SHAPE.ROUNDED_RECTANGLE
    assert str(sp.fill.fore_color.rgb) == "2251FF"
    assert tb.text_frame.paragraphs[0].runs[0].text == "Hello"
    assert tb.text_frame.paragraphs[0].runs[0].font.name == "Arial"


def test_set_fill_transparent_gives_no_fill() -> None:
    slide = _blank_slide()
    sp = add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, 10, 10, fill="transparent", line="transparent")
    assert sp.fill.type == MSO_FILL_TYPE.BACKGROUND
    assert sp.line.fill.type == MSO_FILL_TYPE.BACKGROUND


def test_shape_elements_returns_xml_elements() -> None:
    slide = _blank_slide()
    shapes = [
        add_shape(slide, 1, 0, 0, 10, 10),
        add_textbox(slide, 0, 20, 10, 10, text="t"),
    ]
    elements = shape_elements(shapes)
    assert len(elements) == 2
    for el in elements:
        assert el.tag.endswith("sp")


# -- renderer dispatch guards ----------------------------------------------


def test_unknown_primitive_renderer_rejected() -> None:
    with pytest.raises(InvalidArgumentError, match="unknown native primitive"):
        get_native_renderer("does-not-exist")


def test_background_placement_rejected_for_primitives() -> None:
    slide = _blank_slide()
    resolved = get_default_registry().resolve(
        "asset://primitives/primitive/label-pill?text=Onboard"
    )
    with pytest.raises(InvalidArgumentError, match="background"):
        render_asset(slide, resolved, _ctx(placement="background"))


def test_replace_placement_resolves_native_payload() -> None:
    # renderers are added in A5 Tasks 3-9; the registry wiring (resolve ->
    # NativeShapePayload) is asserted here, renderer presence is covered per
    # primitive in its own structural tests.
    resolved = get_default_registry().resolve("asset://primitives/primitive/quote-mark?text=hi")
    assert isinstance(resolved.payload, NativeShapePayload)
