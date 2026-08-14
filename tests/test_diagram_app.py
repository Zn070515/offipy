# tests/test_diagram_app.py
"""diagram app 测试：build 格式识别 / 参数透传 / 冒烟 / 惰性 import 红线。
install_skill 与三入口接线测试在后续 Task 追加。"""

import subprocess
import sys
from pathlib import Path

import pytest

from offipy.diagram import DiagramApp

_MMD = "graph LR\n  A --> B\n"
_DRAWIO = """\
<mxfile>
  <diagram name="Page 1">
    <mxGraphModel>
      <root>
        <mxCell id="0"/>
        <mxCell id="1" parent="0"/>
        <mxCell id="a" value="A" vertex="1" parent="1">
          <mxGeometry x="20" y="20" width="120" height="40" as="geometry"/>
        </mxCell>
        <mxCell id="b" value="B" vertex="1" parent="1">
          <mxGeometry x="200" y="20" width="120" height="40" as="geometry"/>
        </mxCell>
        <mxCell id="e" value="" style="edgeStyle=orthogonalEdgeStyle;" edge="1"
                source="a" target="b" parent="1">
          <mxGeometry relative="1" as="geometry"/>
        </mxCell>
      </root>
    </mxGraphModel>
  </diagram>
</mxfile>
"""


def _write(tmp_path, name, text):
    p = tmp_path / name
    p.write_text(text, encoding="utf-8")
    return p


def test_build_mermaid_smoke(tmp_path):
    src = _write(tmp_path, "flow.mmd", _MMD)
    out = tmp_path / "out.pptx"
    assert DiagramApp().build(str(src), str(out)) == {"pptx": str(out)}
    assert out.is_file()
    from pptx import Presentation

    prs = Presentation(str(out))
    assert len(prs.slides) == 1
    assert list(prs.slides[0].shapes)


def test_build_drawio_smoke(tmp_path):
    src = _write(tmp_path, "graph.drawio", _DRAWIO)
    out = tmp_path / "out.pptx"
    assert DiagramApp().build(str(src), str(out)) == {"pptx": str(out)}
    assert out.is_file()
    from pptx import Presentation

    prs = Presentation(str(out))
    assert len(prs.slides) == 1


def test_build_content_sniff_mermaid(tmp_path):
    src = _write(tmp_path, "noext", _MMD)  # 无扩展名 → 内容嗅探
    out = tmp_path / "out.pptx"
    assert DiagramApp().build(str(src), str(out)) == {"pptx": str(out)}


def test_build_content_sniff_drawio(tmp_path):
    src = _write(tmp_path, "noext", _DRAWIO)  # 无扩展名 → 内容嗅探
    out = tmp_path / "out.pptx"
    assert DiagramApp().build(str(src), str(out)) == {"pptx": str(out)}


def test_build_missing_source_raises_file_not_found(tmp_path):
    with pytest.raises(FileNotFoundError):
        DiagramApp().build(str(tmp_path / "nope.mmd"), str(tmp_path / "o.pptx"))


def test_build_unknown_format_raises_value_error(tmp_path):
    src = _write(tmp_path, "data.txt", "no diagram markers at all")
    with pytest.raises(ValueError, match="无法识别"):
        DiagramApp().build(str(src), str(tmp_path / "o.pptx"))


def test_build_direction_forwarded_to_mermaid(tmp_path, monkeypatch):
    from offipy import diagrams

    src = _write(tmp_path, "flow.mmd", _MMD)
    captured = {}

    def fake(source, out_path, *, direction=None):
        captured["direction"] = direction
        return str(out_path)

    monkeypatch.setattr(diagrams, "mermaid_to_pptx", fake)
    DiagramApp().build(str(src), str(tmp_path / "o.pptx"), direction="LR")
    assert captured == {"direction": "LR"}


def test_build_page_forwarded_to_drawio(tmp_path, monkeypatch):
    from offipy import drawio

    src = _write(tmp_path, "graph.drawio", _DRAWIO)
    captured = {}

    def fake(source, out_path, *, page=None):
        captured["page"] = page
        return str(out_path)

    monkeypatch.setattr(drawio, "drawio_to_pptx", fake)
    DiagramApp().build(str(src), str(tmp_path / "o.pptx"), page=2)
    assert captured == {"page": 2}


def test_build_irrelevant_param_not_forwarded(tmp_path, monkeypatch):
    # direction 对 drawio 不适用：fake 无 direction 参数，传了会 TypeError → 测试即失败
    from offipy import drawio

    src = _write(tmp_path, "graph.drawio", _DRAWIO)

    def fake(source, out_path, *, page=None):
        return str(out_path)

    monkeypatch.setattr(drawio, "drawio_to_pptx", fake)
    DiagramApp().build(str(src), str(tmp_path / "o.pptx"), direction="LR")


def test_build_extension_priority(tmp_path, monkeypatch):
    # .mmd 扩展名优先：内容即使像 drawio 也走 mermaid 分支
    from offipy import diagrams

    src = _write(tmp_path, "weird.mmd", _DRAWIO)
    captured = {}

    def fake(source, out_path, *, direction=None):
        captured["kind"] = "mermaid"
        return str(out_path)

    monkeypatch.setattr(diagrams, "mermaid_to_pptx", fake)
    DiagramApp().build(str(src), str(tmp_path / "o.pptx"))
    assert captured["kind"] == "mermaid"


def test_diagram_lazy_import_no_pptx():
    """diagram.py 顶层不得 import pptx（惰性 import 红线，与 diagrams/drawio 一致）。"""
    code = (
        "import sys\n"
        "import offipy.diagram\n"
        "assert 'pptx' not in sys.modules, 'diagram.py 顶层不得 import pptx'\n"
    )
    subprocess.run([sys.executable, "-c", code], check=True)


def test_install_skill_target_dir(tmp_path):
    target = tmp_path / "skills"
    res = DiagramApp().install_skill(target_dir=str(target))
    assert len(res["installed"]) == 2
    assert res["skipped"] == []
    assert (target / "diagram-design" / "SKILL.md").is_file()
    assert (target / "offipy-diagram" / "SKILL.md").is_file()


def test_install_skill_idempotent_no_overwrite(tmp_path):
    target = tmp_path / "skills"
    first = DiagramApp().install_skill(target_dir=str(target))
    second = DiagramApp().install_skill(target_dir=str(target))
    assert len(first["installed"]) == 2
    assert second["installed"] == []
    assert len(second["skipped"]) == 2
    # 用户对已存在 skill 的手动编辑不被覆盖
    marker = target / "offipy-diagram" / "user_edit.txt"
    marker.write_text("keep", encoding="utf-8")
    DiagramApp().install_skill(target_dir=str(target))
    assert marker.read_text(encoding="utf-8") == "keep"


def test_install_skill_force_overwrites(tmp_path):
    target = tmp_path / "skills"
    DiagramApp().install_skill(target_dir=str(target))
    marker = target / "offipy-diagram" / "user_edit.txt"
    marker.write_text("keep", encoding="utf-8")
    res = DiagramApp().install_skill(target_dir=str(target), force=True)
    assert res["skipped"] == []
    assert len(res["installed"]) == 2
    assert not marker.exists()  # force 删除并重建目标目录


def test_install_skill_default_dir(tmp_path, monkeypatch):
    fake_home = tmp_path / "home"
    monkeypatch.setattr(Path, "home", lambda: fake_home)
    res = DiagramApp().install_skill()
    assert res["skipped"] == []
    assert (fake_home / ".claude" / "skills" / "offipy-diagram" / "SKILL.md").is_file()


def test_wrapper_skill_contract_keywords():
    skill = (
        Path(__file__).resolve().parent.parent
        / "src"
        / "offipy"
        / "_skills"
        / "offipy-diagram"
        / "SKILL.md"
    )
    text = skill.read_text(encoding="utf-8")
    for kw in ("Mermaid", "draw.io", "offipy diagram build", "mxGraphModel"):
        assert kw in text
