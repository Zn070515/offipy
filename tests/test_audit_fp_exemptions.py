"""使用反馈三类 overlap/margin 误报豁免的行为测试（真实 deck 实测反哺）。

治本设计：extract 提取显式填充（a:noFill 透明），overlap 统一按「上层是否真的
遮挡下方内容」判定，而非面积比的补丁式启发：

1. transparent_overlay —— 透明（a:noFill）无文本上层 → 不遮挡 → 豁免；
2. decorative_overlay —— 实心小装饰/色条浮卡片（尺寸判别）→ 豁免 covered_text；
3. text_on_background —— 文字浮无文本背景 shape → 豁免 partial；
4. 全宽贴边条 → background/header/footer → 豁免 margin（full_bleed/header_footer）。

并断言不误伤真实问题：无文本大形状盖文字仍报、无文本矩形盖文字仍报、
文本叠文本仍报（透明有文本也不豁免）、关闭对应 ignore 开关后同一 finding 恢复。
"""

from pptx import Presentation
from pptx.util import Inches

from offipy.audit.extract import extract_presentation
from offipy.audit.models import (
    RULE_OVERLAP_COVERED_TEXT,
    RULE_OVERLAP_PARTIAL,
    AuditConfig,
)
from offipy.audit.rules import run_rules


def _run(prs, tmp_path, config=None):
    path = tmp_path / "x.pptx"
    prs.save(path)
    ext = extract_presentation(path)
    records = [r for slide in ext.slides for r in slide.shapes]
    findings, suppressed = run_rules(records, ext.slide_size, config or AuditConfig())
    return findings, suppressed, records


def _new_deck():
    prs = Presentation()  # 默认 10×7.5 in
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    return prs, slide


# ---------------------------------------------------------------- 豁免 A：装饰浮卡片


def test_overlap_decorative_marker_on_card_suppressed(tmp_path):
    prs, slide = _new_deck()
    card = slide.shapes.add_shape(1, Inches(1), Inches(1), Inches(6), Inches(3))
    card.text = "Card title"
    slide.shapes.add_shape(1, Inches(2), Inches(1.5), Inches(0.3), Inches(0.3))  # 12px 级 marker
    findings, suppressed, records = _run(prs, tmp_path)
    assert len(records) == 2
    assert not any(f.kind == "overlap" for f in findings)
    assert any(s.reason == "decorative_overlay" for s in suppressed)


def test_overlap_decorative_bar_on_card_suppressed(tmp_path):
    prs, slide = _new_deck()
    card = slide.shapes.add_shape(1, Inches(1), Inches(1), Inches(6), Inches(3))
    card.text = "Card"
    slide.shapes.add_shape(1, Inches(1), Inches(1), Inches(6), Inches(0.4))  # 顶部色条，无文本
    findings, suppressed, _ = _run(prs, tmp_path)
    assert not any(f.kind == "overlap" for f in findings)
    assert any(s.reason == "decorative_overlay" for s in suppressed)


def test_overlap_large_shape_covering_text_not_suppressed(tmp_path):
    """无文本大形状盖住有文本内容（面积占优）→ 仍报 covered_text，不误伤。"""
    prs, slide = _new_deck()
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(0.5))
    textbox.text = "Important"
    slide.shapes.add_shape(1, Inches(1), Inches(1), Inches(4), Inches(0.8))  # 完全盖住文字
    findings, suppressed, _ = _run(prs, tmp_path)
    assert any(f.rule_id == RULE_OVERLAP_COVERED_TEXT for f in findings)
    assert not any(s.reason == "decorative_overlay" for s in suppressed)


# ---------------------------------------------------------------- 豁免 B：文字浮背景


def test_overlap_text_on_background_suppressed(tmp_path):
    prs, slide = _new_deck()
    slide.shapes.add_shape(1, Inches(1), Inches(2), Inches(5), Inches(0.5))  # 无文本行背景条
    tb = slide.shapes.add_textbox(Inches(1), Inches(2.05), Inches(5), Inches(0.5))
    tb.text = "Row text"  # 浮在背景条上，覆盖比 ≈0.9
    findings, suppressed, records = _run(prs, tmp_path)
    assert len(records) == 2
    assert not any(f.kind == "overlap" for f in findings)
    assert any(s.reason == "text_on_background" for s in suppressed)


def test_overlap_shape_over_text_not_suppressed(tmp_path):
    """无文本矩形（上层）盖住文字一部分 → 仍报 partial，不误伤。"""
    prs, slide = _new_deck()
    tb = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(3), Inches(0.5))
    tb.text = "Important"
    slide.shapes.add_shape(1, Inches(1), Inches(2.1), Inches(4), Inches(0.5))
    findings, suppressed, _ = _run(prs, tmp_path)
    assert any(f.rule_id == RULE_OVERLAP_PARTIAL for f in findings)
    assert not any(s.reason == "text_on_background" for s in suppressed)


# ---------------------------------------------------------------- 豁免 D：卡片+行条纹装饰分层


def test_overlap_decorative_stripe_poking_card_suppressed(tmp_path):
    """S13 卡片+行条纹：条纹仅越出卡片底边 0.1in（partial ratio 0.8）→ 装饰分层豁免。"""
    prs, slide = _new_deck()
    slide.shapes.add_shape(1, Inches(1), Inches(1), Inches(6), Inches(2))  # 卡片 aspect 3
    slide.shapes.add_shape(1, Inches(1), Inches(2.6), Inches(5.7), Inches(0.5))  # 条纹 poking 0.1
    findings, suppressed, records = _run(prs, tmp_path)
    assert len(records) == 2
    assert not any(f.kind == "overlap" for f in findings)
    assert any(s.reason == "decorative_layering" for s in suppressed)


def test_overlap_non_strip_textless_partial_still_reported(tmp_path):
    """非长条空框（aspect < 3）部分重叠 → 仍报 partial，不误伤。"""
    prs, slide = _new_deck()
    slide.shapes.add_shape(1, Inches(1), Inches(1), Inches(4), Inches(3))  # card aspect 1.33
    slide.shapes.add_shape(1, Inches(1), Inches(2.7), Inches(1.5), Inches(1.5))  # 方形 poking
    findings, suppressed, _ = _run(prs, tmp_path)
    assert any(f.rule_id == RULE_OVERLAP_PARTIAL for f in findings)
    assert not any(s.reason == "decorative_layering" for s in suppressed)


# ---------------------------------------------------------------- 豁免 A2：透明装饰（治本根信号）


def test_overlap_transparent_marker_on_card_suppressed(tmp_path):
    """显式 a:noFill（透明）无文本 marker 浮卡片 → transparent_overlay。"""
    prs, slide = _new_deck()
    card = slide.shapes.add_shape(1, Inches(1), Inches(1), Inches(6), Inches(3))
    card.text = "Card title"
    marker = slide.shapes.add_shape(1, Inches(2), Inches(1.5), Inches(0.3), Inches(0.3))
    marker.fill.background()
    findings, suppressed, _ = _run(prs, tmp_path)
    assert not any(f.kind == "overlap" for f in findings)
    assert any(s.reason == "transparent_overlay" for s in suppressed)


def test_overlap_transparent_text_shape_not_suppressed(tmp_path):
    """透明但有文本 → 仍是内容叠内容，必须报 covered_text（透明不改变内容冲突）。"""
    prs, slide = _new_deck()
    big = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(3))
    big.text = "Big text"
    small = slide.shapes.add_textbox(Inches(2), Inches(2), Inches(1), Inches(1))
    small.text = "Small"
    small.fill.background()
    findings, suppressed, _ = _run(prs, tmp_path)
    assert any(f.rule_id == RULE_OVERLAP_COVERED_TEXT for f in findings)
    assert not any(s.reason == "transparent_overlay" for s in suppressed)


# ---------------------------------------------------------------- 豁免 C：全宽贴边条


def test_margin_full_width_bar_suppressed(tmp_path):
    prs, slide = _new_deck()
    slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.6))  # 全宽贴顶条
    findings, suppressed, records = _run(prs, tmp_path)
    assert records[0].role == "background"
    assert not any(f.kind == "margin" for f in findings)
    assert any(s.reason == "full_bleed" for s in suppressed)


def test_margin_full_width_header_bar_suppressed(tmp_path):
    prs, slide = _new_deck()
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.6))
    bar.text = "ACME CORP"  # 全宽有文本 → header
    findings, suppressed, records = _run(prs, tmp_path)
    assert records[0].role == "header"
    assert not any(f.kind == "margin" for f in findings)
    assert any(s.reason == "header_footer" for s in suppressed)


def test_margin_full_width_bar_recovered_when_ignored_off(tmp_path):
    prs, slide = _new_deck()
    slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.6))
    findings, suppressed, _ = _run(prs, tmp_path, AuditConfig(ignore_full_bleed_shapes=False))
    assert any(f.kind == "margin" for f in findings)
    assert not any(s.reason == "full_bleed" for s in suppressed)

    prs2, slide2 = _new_deck()
    bar = slide2.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.6))
    bar.text = "ACME CORP"
    findings, suppressed, _ = _run(prs2, tmp_path, AuditConfig(ignore_headers_footers=False))
    assert any(f.kind == "margin" for f in findings)
    assert not any(s.reason == "header_footer" for s in suppressed)


# ---------------------------------------------------------------- 回归：文本叠文本仍报


def test_overlap_text_over_text_not_suppressed(tmp_path):
    prs, slide = _new_deck()
    big = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(3))
    big.text = "Big text"
    small = slide.shapes.add_textbox(Inches(2), Inches(2), Inches(1), Inches(1))
    small.text = "Small"  # 完全在 big 内 → 文本叠文本
    findings, suppressed, _ = _run(prs, tmp_path)
    assert any(f.rule_id == RULE_OVERLAP_COVERED_TEXT for f in findings)
    assert not any(s.reason in ("decorative_overlay", "text_on_background") for s in suppressed)
