"""审计/回归报告渲染：text/markdown/json/html 输出；HTML 无外网依赖、筛选齐全。"""

import json

from pptx import Presentation
from pptx.util import Inches

from offipy.audit import (
    AuditConfig,
    audit_pptx,
    compare_pptx,
    render_markdown,
    render_text,
)


def _rect(slide, x, y, w, h):
    return slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))


def _deck(tmp_path, name):
    """bounds HIGH + margin LOW + overlap MID 三页不混，单页合成。"""
    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(s, 9.5, 3, 3, 1)  # 部分越界 → bounds.partial HIGH
    _rect(s, 9.85, 5, 0.1, 1)  # 贴右边缘 → margin.right LOW
    _rect(s, 1, 1, 2, 2)  # overlap 对 1/2 → overlap.partial MID
    _rect(s, 1.5, 1.5, 2, 2)
    path = tmp_path / name
    prs.save(path)
    return path


def _diff_pair(tmp_path):
    base = Presentation()
    s = base.slides.add_slide(base.slide_layouts[6])
    _rect(s, 9.85, 5, 0.1, 1)  # margin.right LOW
    base_path = tmp_path / "base.pptx"
    base.save(base_path)

    cand = Presentation()
    s = cand.slides.add_slide(cand.slide_layouts[6])
    _rect(s, 5, 5, 0.1, 1)  # margin 已解决
    _rect(s, 9.5, 3, 3, 1)  # 新增越界 → bounds HIGH
    cand_path = tmp_path / "cand.pptx"
    cand.save(cand_path)
    return base_path, cand_path


# ---------------------------------------------------------------- text


def test_render_text_audit(tmp_path):
    report = audit_pptx(_deck(tmp_path, "deck.pptx"))
    text = render_text(report)
    assert "审计报告" in text
    assert "[HIGH]" in text
    assert "geometry.bounds.partial" in text
    assert "geometry.margin.right" in text
    assert "概要" in text
    assert f"sha256 {report.source_sha256}" in text


def test_render_text_diff(tmp_path):
    base_path, cand_path = _diff_pair(tmp_path)
    diff = compare_pptx(base_path, cand_path)
    text = render_text(diff)
    assert "回归对比" in text
    assert "基线:" in text
    assert "已解决" in text
    assert "新增/恶化最高严重度: HIGH" in text


# ---------------------------------------------------------------- markdown


def test_render_markdown_audit(tmp_path):
    report = audit_pptx(_deck(tmp_path, "deck.pptx"))
    md = report.to_markdown()
    assert md.startswith("# 审计报告")
    assert "| 严重度 | rule_id | 形状 | 描述 |" in md
    assert "## 概要" in md
    assert "geometry.overlap.partial" in md


def test_render_markdown_diff(tmp_path):
    base_path, cand_path = _diff_pair(tmp_path)
    diff = compare_pptx(base_path, cand_path)
    md = render_markdown(diff)
    assert md.startswith("# 回归对比")
    assert "## 新增问题" in md
    assert "## 已解决问题" in md
    assert "## 形状变化" in md


# ---------------------------------------------------------------- json


def test_render_json_audit(tmp_path):
    report = audit_pptx(_deck(tmp_path, "deck.pptx"))
    data = json.loads(report.to_json())
    assert data["max_severity"] == "HIGH"
    assert len(data["findings"]) >= 3
    assert data["shapes"]  # HTML/SVG 渲染几何快照齐备


def test_render_json_diff(tmp_path):
    base_path, cand_path = _diff_pair(tmp_path)
    diff = compare_pptx(base_path, cand_path)
    data = json.loads(diff.to_json())
    assert data["max_new_severity"] == "HIGH"
    assert [f["rule_id"] for f in data["resolved_findings"]] == ["geometry.margin.right"]
    assert [f["rule_id"] for f in data["added_findings"]] == ["geometry.bounds.partial"]
    assert data["added_shapes"]


# ---------------------------------------------------------------- html：审计


def test_render_html_audit_no_external(tmp_path):
    report = audit_pptx(_deck(tmp_path, "deck.pptx"))
    html = report.to_html()
    assert "http://" not in html
    assert "https://" not in html
    assert "<svg" in html
    assert "filterAudit" in html
    assert 'class="shape sev-high"' in html
    assert 'class="shape sev-mid"' in html
    assert "overlap-area" in html
    for label in ("全部", "HIGH", "MID", "LOW", "豁免", "警告"):
        assert label in html


def test_render_html_audit_suppressed(tmp_path):
    deck = _deck(tmp_path, "deck.pptx")
    report = audit_pptx(deck)
    margin_id = next(
        f.primary.shape_id for f in report.findings if f.rule_id == "geometry.margin.right"
    )
    report2 = audit_pptx(deck, AuditConfig(ignored_shapes={(1, margin_id)}))
    assert [s.reason for s in report2.suppressed] == ["user_shape"]
    html = report2.to_html()
    assert 'class="shape sup"' in html
    assert "豁免" in html


def test_render_html_audit_slides_dir(tmp_path):
    report = audit_pptx(_deck(tmp_path, "deck.pptx"))
    html = report.to_html(slides_dir=str(tmp_path))
    assert "slide-1.png" in html
    assert "href=" in html


# ---------------------------------------------------------------- html：对比


def test_render_html_diff(tmp_path):
    base_path, cand_path = _diff_pair(tmp_path)
    diff = compare_pptx(base_path, cand_path)
    html = diff.to_html()
    assert "http://" not in html
    assert "https://" not in html
    assert "<table>" in html
    assert "<h2>形状变化</h2>" in html
    assert "已解决" in html
    assert "geometry.bounds.partial" in html
    assert "新增/恶化最高严重度" in html
