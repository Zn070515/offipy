"""audit compare 基线回归：shape 增删/移动/缩放/文本 + finding 新增/已解决/变化 + 匹配链。"""

import pytest
from pptx import Presentation
from pptx.util import Inches

from offipy.audit import compare_pptx
from offipy.audit.models import RULE_BOUNDS_PARTIAL, RULE_MARGIN_RIGHT, Severity
from offipy.exceptions import InvalidArgumentError


def _force_id(shape, value: int) -> None:
    shape._element.xpath(".//p:cNvPr")[0].set("id", str(value))


def _rect(slide, x, y, w, h):
    return slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))


def _tb(slide, x, y, text, w=2, h=1):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tb.text = text
    return tb


def test_identical_decks_no_diff(tmp_path):
    base = Presentation()
    cand = Presentation()
    for prs in (base, cand):
        s = prs.slides.add_slide(prs.slide_layouts[6])
        _tb(s, 2, 2, "hello")
    base_path = tmp_path / "base.pptx"
    cand_path = tmp_path / "cand.pptx"
    base.save(base_path)
    cand.save(cand_path)

    diff = compare_pptx(base_path, cand_path)
    assert diff.added_shapes == []
    assert diff.removed_shapes == []
    assert diff.moved_shapes == []
    assert diff.resized_shapes == []
    assert diff.text_changes == []
    assert diff.added_findings == []
    assert diff.resolved_findings == []
    assert diff.changed_findings == []
    assert diff.gate_severity() is None
    assert len(diff.baseline_sha256) == 64
    assert len(diff.candidate_sha256) == 64
    assert diff.baseline_slide_count == diff.candidate_slide_count == 1


def test_shape_moved_resized_text_removed(tmp_path):
    base = Presentation()
    s = base.slides.add_slide(base.slide_layouts[6])
    _tb(s, 2, 2, "AAA")
    _tb(s, 2, 4, "BBB")
    _tb(s, 5, 4, "CCC")
    _tb(s, 2, 6, "EEE")
    base_path = tmp_path / "base.pptx"
    base.save(base_path)

    cand = Presentation()
    s = cand.slides.add_slide(cand.slide_layouts[6])
    _tb(s, 4, 2, "AAA")  # moved
    _tb(s, 2, 4, "BBB", w=3)  # resized
    _tb(s, 5, 4, "CCC changed")  # text
    cand_path = tmp_path / "cand.pptx"
    cand.save(cand_path)

    diff = compare_pptx(base_path, cand_path)
    assert [x.name for x in diff.added_shapes] == []
    assert [x.name for x in diff.removed_shapes] == ["TextBox 4"]
    moved = [x.name for x in diff.moved_shapes]
    resized = [x.name for x in diff.resized_shapes]
    texts = [x.name for x in diff.text_changes]
    assert "TextBox 1" in moved
    assert "TextBox 2" in resized
    assert "TextBox 3" in texts
    assert diff.added_findings == []
    assert diff.resolved_findings == []
    assert diff.changed_findings == []


def test_diff_added_resolved_changed(tmp_path):
    # baseline: T(9.85,1)→margin LOW；S(9.5,3,1,1)→bounds MID
    base = Presentation()
    s = base.slides.add_slide(base.slide_layouts[6])
    _rect(s, 9.85, 1, 0.1, 1)
    _rect(s, 9.5, 3, 1, 1)
    base_path = tmp_path / "base.pptx"
    base.save(base_path)

    # candidate: T→安全位置（margin 已解决）；S 撑大→bounds HIGH；新增 C→margin LOW
    cand = Presentation()
    s = cand.slides.add_slide(cand.slide_layouts[6])
    _rect(s, 5, 1, 0.1, 1)
    _rect(s, 9.5, 3, 3, 1)
    _rect(s, 9.85, 5, 0.1, 1)
    cand_path = tmp_path / "cand.pptx"
    cand.save(cand_path)

    diff = compare_pptx(base_path, cand_path)
    assert [f.rule_id for f in diff.resolved_findings] == [RULE_MARGIN_RIGHT]
    assert [f.rule_id for f in diff.added_findings] == [RULE_MARGIN_RIGHT]
    assert len(diff.changed_findings) == 1
    ch = diff.changed_findings[0]
    assert ch.rule_id == RULE_BOUNDS_PARTIAL
    assert ch.old_severity == Severity.MID
    assert ch.new_severity == Severity.HIGH
    assert ch.worsened
    assert diff.gate_severity() == Severity.HIGH

    assert [x.name for x in diff.added_shapes] == ["Rectangle 3"]
    moved_names = [x.name for x in diff.moved_shapes]
    resized_names = [x.name for x in diff.resized_shapes]
    assert "Rectangle 1" in moved_names
    assert "Rectangle 2" in resized_names


def test_match_fallback_text_geometry_when_id_differs(tmp_path):
    # 候选里逻辑相同的形状换了 shape_id + name → 走 text+geometry 兜底匹配
    base = Presentation()
    s = base.slides.add_slide(base.slide_layouts[6])
    tb = _tb(s, 3, 3, "UNIQUE TEXT")
    tb.name = "Marker"
    _force_id(tb, 100)
    base_path = tmp_path / "base.pptx"
    base.save(base_path)

    cand = Presentation()
    s = cand.slides.add_slide(cand.slide_layouts[6])
    filler = _rect(s, 1, 1, 1, 1)
    _force_id(filler, 200)
    tb2 = _tb(s, 3, 3, "UNIQUE TEXT")
    tb2.name = "Renamed"
    _force_id(tb2, 300)
    cand_path = tmp_path / "cand.pptx"
    cand.save(cand_path)

    diff = compare_pptx(base_path, cand_path)
    added_names = [x.name for x in diff.added_shapes]
    assert "Renamed" not in added_names  # 已通过 text+geometry 匹配，不误报新增
    assert any(n == "Rectangle 1" for n in added_names)
    assert not any(x.name == "Renamed" for x in diff.unmatched_candidate)


def test_compare_missing_baseline_raises(tmp_path):
    cand = Presentation()
    s = cand.slides.add_slide(cand.slide_layouts[6])
    _tb(s, 2, 2, "x")
    cand_path = tmp_path / "cand.pptx"
    cand.save(cand_path)
    with pytest.raises(InvalidArgumentError):
        compare_pptx(tmp_path / "nope.pptx", cand_path)


def test_compare_wrong_extension_raises(tmp_path):
    base = Presentation()
    s = base.slides.add_slide(base.slide_layouts[6])
    _tb(s, 2, 2, "x")
    base_path = tmp_path / "base.pptx"
    base.save(base_path)
    bad = tmp_path / "cand.ppt"
    bad.write_bytes(b"not a pptx")
    with pytest.raises(InvalidArgumentError):
        compare_pptx(base_path, bad)
