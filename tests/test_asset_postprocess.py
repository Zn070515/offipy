"""A3 Task 8 — postprocess_assets 编排 + 注入副本声明读回。

postprocess_assets 从注入副本读回 data-offipy-asset-id 声明 → 绑定测量 →
registry.resolve → render_asset → place_rendered_elements，PPTX 只打开/保存各一次，
返回按声明顺序的用量报告。无声明 → 空报告不碰 PPTX；测量缺失 → RuntimeError；
非法 rect / 缺绑定 → InvalidArgumentError。
"""

import json

import pytest

from offipy.assets.declarations import (
    parse_injected_asset_declarations,
    preprocess_asset_declarations,
)
from offipy.assets.model import AssetRef
from offipy.assets.render import AssetUsageReport, postprocess_assets
from offipy.exceptions import InvalidArgumentError

_ASVG_NS = "http://schemas.microsoft.com/office/drawing/2016/SVG/main"
_R = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed"
_A = "http://schemas.openxmlformats.org/drawingml/2006/main"

HTML_TWO = """<!doctype html>
<html><head><style data-theme="dark"></style></head>
<body>
<section data-pptx-slide>
  <div data-asset="asset://ph/icon/check"></div>
</section>
<section data-pptx-slide>
  <div data-asset="asset://lu/icon/settings" data-asset-placement="background"></div>
</section>
</body></html>"""

HTML_NONE = """<!doctype html>
<html><body><section data-pptx-slide><p>plain text</p></section></body></html>"""

HTML_LEGACY = """<!doctype html>
<html><body>
<section data-pptx-slide><svg data-icon="ph:check" width="10" height="10"></svg></section>
</body></html>"""


def _asset_meas(asset_id, x=20, y=30, w=120, h=60, color="rgb(255, 0, 0)"):
    return {
        "id": 1,
        "kind": "asset",
        "tag": "div",
        "assetId": asset_id,
        "rect": {"x": x, "y": y, "w": w, "h": h},
        "themeVars": {"bg": "#ffffff", "surface": "#f3f4f6", "accent": "#2251ff"},
        "color": color,
    }


def _fixture(tmp_path, html_source, *, drop_ids=(), override=None):
    """构造注入副本 + 占位符 PPTX + measurements.json，返回 (html, decls, pptx)。

    drop_ids: 不写入测量的声明 id（模拟缺失测量）；override: (rec, decl_id) → 改后的 rec。
    """
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE

    injected, decls = preprocess_asset_declarations(html_source)
    html = tmp_path / "t.audited.html"
    html.write_text(injected, encoding="utf-8")

    prs = Presentation()
    for d in decls:
        while len(prs.slides) < d.slide_index:
            prs.slides.add_slide(prs.slide_layouts[6])
        slide = prs.slides[d.slide_index - 1]
        sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, 100, 100)
        sp.name = f"OFFIPY_ASSET::{d.declaration_id}"
    pptx = tmp_path / "t.pptx"
    prs.save(pptx)

    slides = [{"records": []} for _ in range(max(d.slide_index for d in decls) or 0)]
    for d in decls:
        if d.declaration_id in drop_ids:
            continue
        rec = _asset_meas(d.declaration_id)
        if override:
            rec = override(rec, d.declaration_id)
        slides[d.slide_index - 1]["records"].append(rec)
    meas_dir = tmp_path / "t_audit" / "_cache"
    meas_dir.mkdir(parents=True)
    (meas_dir / "measurements.json").write_text(
        json.dumps({"slides": slides}, ensure_ascii=False), encoding="utf-8"
    )
    return str(html), decls, str(pptx)


def _reload(pptx_path):
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(pptx_path)
    slides = list(prs.slides)
    for slide in slides:
        assert not any(s.name.startswith("OFFIPY_ASSET::") for s in slide.shapes)
    freeforms = [
        len([s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.FREEFORM])
        for slide in slides
    ]
    return slides, freeforms


# ---------------------------------------------------------------------------
# postprocess_assets
# ---------------------------------------------------------------------------


class TestPostprocessAssets:
    def test_renders_all_declarations_and_reports_in_order(self, tmp_path):
        html, decls, pptx = _fixture(tmp_path, HTML_TWO)
        report = postprocess_assets(html, pptx)
        assert isinstance(report, AssetUsageReport)
        assert len(report.records) == 2
        r0, r1 = report.records
        assert r0.declaration_id == "asset-s01-001"
        assert r0.slide_index == 1
        assert r0.request == "asset://ph/icon/check"
        assert r0.placement == "replace"
        assert r0.provider.provider_id == "ph"
        assert r1.declaration_id == "asset-s02-001"
        assert r1.slide_index == 2
        assert r1.request == "asset://lu/icon/settings"
        assert r1.placement == "background"
        assert r1.provider.provider_id == "lu"

        slides, freeforms = _reload(pptx)
        assert len(slides) == 2
        assert freeforms == [2, 2]  # ph:check 与 lu:settings 各 2 子路径 → 各 2 个 freeform

    def test_no_assets_returns_empty_report_without_touching_pptx(self, tmp_path):
        injected, _ = preprocess_asset_declarations(HTML_NONE)
        html = tmp_path / "plain.audited.html"
        html.write_text(injected, encoding="utf-8")
        pptx = tmp_path / "plain.pptx"
        pptx.write_bytes(b"not-a-pptx")  # 原样返回，不应被打开
        report = postprocess_assets(str(html), str(pptx))
        assert report.records == ()
        assert pptx.read_bytes() == b"not-a-pptx"

    def test_missing_measurements_raises_runtime_error(self, tmp_path):
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE

        injected, decls = preprocess_asset_declarations(HTML_TWO)
        html = tmp_path / "t.audited.html"
        html.write_text(injected, encoding="utf-8")
        prs = Presentation()
        for d in decls:
            while len(prs.slides) < d.slide_index:
                prs.slides.add_slide(prs.slide_layouts[6])
            slide = prs.slides[d.slide_index - 1]
            sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, 100, 100)
            sp.name = f"OFFIPY_ASSET::{d.declaration_id}"
        pptx = tmp_path / "t.pptx"
        prs.save(pptx)
        with pytest.raises(RuntimeError, match="measurements.json"):
            postprocess_assets(html, str(pptx))

    def test_nonpositive_rect_raises_invalid_argument(self, tmp_path):
        def _bad(rec, aid):
            rec["rect"] = {"x": 20, "y": 30, "w": 0, "h": 60}
            return rec

        html, _, pptx = _fixture(tmp_path, HTML_TWO, override=_bad)
        with pytest.raises(InvalidArgumentError, match="positive"):
            postprocess_assets(html, pptx)

    def test_declaration_without_measurement_raises(self, tmp_path):
        html, _, pptx = _fixture(tmp_path, HTML_TWO, drop_ids=("asset-s02-001",))
        with pytest.raises(InvalidArgumentError, match="asset-s02-001"):
            postprocess_assets(html, pptx)

    def test_svg_asset_gets_png_fallback(self, tmp_path, monkeypatch):
        # #58：postprocess 用共享 PNG 渲染器给 SVG picture 补 raster fallback。
        # Playwright 渲染器本身 mock 掉，只验证接线与最终 blip 结构：主 <a:blip>
        # 指向 PNG、svgBlip 仍指向 SVG。
        import offipy.assets.render as render_mod

        calls = []
        monkeypatch.setattr(
            render_mod,
            "_make_svg_to_png",
            lambda: (lambda svg: (calls.append(svg) or b"fake-png"), lambda: None),
        )
        html_src = (
            "<html><body><section data-pptx-slide>"
            '<div data-asset="asset://procedural/pattern/wave"></div>'
            "</section></body></html>"
        )
        html, _, pptx = _fixture(tmp_path, html_src)
        postprocess_assets(html, pptx)
        assert calls, "PNG 渲染器应被 SVG asset 调用"

        from pptx import Presentation

        prs = Presentation(pptx)
        slide = prs.slides[0]
        blip = slide.shapes._spTree.find(f".//{{{_A}}}blip")
        assert blip is not None
        png_rid = blip.get(_R)
        assert png_rid
        assert slide.part.related_part(png_rid).partname.endswith(".png")
        svg_blip = slide.shapes._spTree.find(f".//{{{_ASVG_NS}}}svgBlip")
        assert svg_blip is not None
        assert slide.part.related_part(svg_blip.get(_R)).partname.endswith(".svg")


# ---------------------------------------------------------------------------
# parse_injected_asset_declarations
# ---------------------------------------------------------------------------


class TestParseInjected:
    def test_roundtrip_matches_preprocess(self, tmp_path):
        injected, decls = preprocess_asset_declarations(HTML_TWO)
        parsed = parse_injected_asset_declarations(injected)
        assert len(parsed) == len(decls)
        for a, b in zip(parsed, decls, strict=True):
            assert a.declaration_id == b.declaration_id
            assert a.slide_index == b.slide_index
            assert a.request == b.request
            assert a.placement == b.placement
            assert a.html_tag == b.html_tag

    def test_legacy_icon_injected_reads_back_canonical(self, tmp_path):
        injected, _ = preprocess_asset_declarations(HTML_LEGACY)
        parsed = parse_injected_asset_declarations(injected)
        assert len(parsed) == 1
        decl = parsed[0]
        assert decl.declaration_id == "asset-s01-001"
        assert decl.slide_index == 1
        assert decl.request.ref == AssetRef("ph", "icon", "check")
        assert decl.placement == "replace"
        assert decl.html_tag == "svg"

    def test_missing_data_asset_rejected(self, tmp_path):
        bad = '<section data-pptx-slide><div data-offipy-asset-id="asset-s01-001"></div></section>'
        with pytest.raises(InvalidArgumentError, match="missing data-asset"):
            parse_injected_asset_declarations(bad)

    def test_invalid_placement_rejected(self, tmp_path):
        bad = (
            '<section data-pptx-slide><div data-asset="asset://ph/icon/check" '
            'data-offipy-asset-id="asset-s01-001" data-asset-placement="bogus"></div></section>'
        )
        with pytest.raises(InvalidArgumentError, match="placement"):
            parse_injected_asset_declarations(bad)
