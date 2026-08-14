"""offipy.drawio — draw.io 解析（vendored drawio_extract 包装）。"""

from __future__ import annotations

import subprocess
import sys

import pytest
from pptx import Presentation
from pptx.util import Emu

from offipy.drawio import _remove_bbox_shapes, drawio_to_pptx, layout_drawio, parse_drawio

SINGLE_PAGE = """\
<mxfile>
  <diagram name="Page 1">
    <mxGraphModel>
      <root>
        <mxCell id="0"/>
        <mxCell id="1" parent="0"/>
        <mxCell id="g" value="容器"
                style="rounded=1;fillColor=#fff2cc;strokeColor=#d6b656;container=1;"
                vertex="1" parent="1">
          <mxGeometry x="20" y="20" width="340" height="200" as="geometry"/>
        </mxCell>
        <mxCell id="a" value="节点A"
                style="rounded=0;fillColor=#dae8fc;strokeColor=#6c8ebf;fontColor=#000000;"
                vertex="1" parent="g">
          <mxGeometry x="40" y="40" width="120" height="60" as="geometry"/>
        </mxCell>
        <mxCell id="b" value="节点B"
                style="rounded=0;fillColor=#d5e8d4;strokeColor=#82b366;"
                vertex="1" parent="g">
          <mxGeometry x="200" y="40" width="120" height="60" as="geometry"/>
        </mxCell>
        <mxCell id="c" value="节点C"
                style="rounded=1;fillColor=#ffe6cc;strokeColor=#d79b00;"
                vertex="1" parent="1">
          <mxGeometry x="60" y="260" width="140" height="60" as="geometry"/>
        </mxCell>
        <mxCell id="e1" value="连接"
                style="edgeStyle=orthogonalEdgeStyle;rounded=1;strokeColor=#333333;dashed=1;"
                edge="1" source="a" target="b" parent="1">
          <mxGeometry relative="1" as="geometry"/>
        </mxCell>
      </root>
    </mxGraphModel>
  </diagram>
</mxfile>
"""

MULTI_PAGE = """\
<mxfile>
  <diagram name="第一页">
    <mxGraphModel><root>
      <mxCell id="0"/><mxCell id="1" parent="0"/>
      <mxCell id="x" value="X" vertex="1" parent="1">
        <mxGeometry x="10" y="10" width="100" height="50" as="geometry"/>
      </mxCell>
    </root></mxGraphModel>
  </diagram>
  <diagram name="第二页">
    <mxGraphModel><root>
      <mxCell id="0"/><mxCell id="1" parent="0"/>
      <mxCell id="y" value="Y" vertex="1" parent="1">
        <mxGeometry x="30" y="30" width="100" height="50" as="geometry"/>
      </mxCell>
    </root></mxGraphModel>
  </diagram>
</mxfile>
"""

WAYPOINTS = """\
<mxfile><diagram name="P"><mxGraphModel><root>
  <mxCell id="0"/><mxCell id="1" parent="0"/>
  <mxCell id="a" value="A" vertex="1" parent="1">
    <mxGeometry x="0" y="0" width="50" height="30" as="geometry"/></mxCell>
  <mxCell id="b" value="B" vertex="1" parent="1">
    <mxGeometry x="200" y="100" width="50" height="30" as="geometry"/></mxCell>
  <mxCell id="e1" style="edgeStyle=orthogonalEdgeStyle;strokeColor=#333333;"
          edge="1" source="a" target="b" parent="1">
    <mxGeometry relative="1" as="geometry">
      <Array as="points">
        <mxPoint x="80" y="60"/>
        <mxPoint x="170" y="60"/>
      </Array>
    </mxGeometry>
  </mxCell>
</root></mxGraphModel></diagram></mxfile>
"""

STYLED = """\
<mxfile><diagram name="P"><mxGraphModel><root>
  <mxCell id="0"/><mxCell id="1" parent="0"/>
  <mxCell id="a" value="转"
          style="strokeWidth=3;rotation=45;dashPattern=3 2;fillColor=#dae8fc;strokeColor=#6c8ebf;"
          vertex="1" parent="1">
    <mxGeometry x="0" y="0" width="100" height="50" as="geometry"/>
  </mxCell>
  <mxCell id="b" value="自定义"
          style="dashPattern=8 3 2 2;fillColor=#d5e8d4;strokeColor=#82b366;"
          vertex="1" parent="1">
    <mxGeometry x="120" y="0" width="100" height="50" as="geometry"/>
  </mxCell>
  <mxCell id="e1" style="edgeStyle=orthogonalEdgeStyle;strokeColor=#333333;strokeWidth=4;"
          edge="1" source="a" target="b" parent="1">
    <mxGeometry relative="1" as="geometry"/>
  </mxCell>
</root></mxGraphModel></diagram></mxfile>
"""

SHAPES = """\
<mxfile><diagram name="P"><mxGraphModel><root>
  <mxCell id="0"/><mxCell id="1" parent="0"/>
  <mxCell id="ra" value="RA" style="rounded=0;" vertex="1" parent="1">
    <mxGeometry x="0" y="0" width="50" height="30" as="geometry"/>
  </mxCell>
  <mxCell id="rb" value="RB" style="rounded=1;" vertex="1" parent="1">
    <mxGeometry x="80" y="0" width="50" height="30" as="geometry"/>
  </mxCell>
  <mxCell id="el" value="EL" style="ellipse;" vertex="1" parent="1">
    <mxGeometry x="160" y="0" width="50" height="30" as="geometry"/>
  </mxCell>
  <mxCell id="tr" value="TR" style="triangle;" vertex="1" parent="1">
    <mxGeometry x="240" y="0" width="50" height="30" as="geometry"/>
  </mxCell>
</root></mxGraphModel></diagram></mxfile>
"""


def _write(tmp_path, text=SINGLE_PAGE, name="input.drawio"):
    p = tmp_path / name
    p.write_text(text, encoding="utf-8")
    return str(p)


def test_parse_drawio_basic(tmp_path):
    d = parse_drawio(_write(tmp_path))
    by_id = {n.id: n for n in d.nodes}
    assert set(by_id) == {"g", "a", "b", "c"}
    a = by_id["a"]
    assert a.x == 60.0 and a.y == 60.0  # 父链累加：g(20,20) + a(40,40)
    assert a.w == 120.0 and a.h == 60.0
    assert a.fill == "#dae8fc"
    assert a.stroke == "#6c8ebf"
    assert a.font_color == "#000000"
    assert a.rounded is False  # rounded=0 → 直角
    g = by_id["g"]
    assert g.container is True
    assert g.label == "容器"
    assert g.x == 20.0 and g.y == 20.0
    c = by_id["c"]
    assert c.rounded is True  # rounded=1 → 圆角
    assert len(d.edges) == 1
    e = d.edges[0]
    assert e.source == "a" and e.target == "b"
    assert e.label == "连接"
    assert e.dashed is True
    assert e.stroke == "#333333"


def test_parse_drawio_default_first_page(tmp_path):
    d = parse_drawio(_write(tmp_path, MULTI_PAGE, "multi.drawio"))
    assert [n.id for n in d.nodes] == ["x"]


def test_parse_drawio_page_by_index(tmp_path):
    src = _write(tmp_path, MULTI_PAGE, "multi.drawio")
    assert [n.id for n in parse_drawio(src, page=0).nodes] == ["x"]
    assert [n.id for n in parse_drawio(src, page=1).nodes] == ["y"]
    assert [n.id for n in parse_drawio(src, page="0").nodes] == ["x"]


def test_parse_drawio_page_by_name(tmp_path):
    src = _write(tmp_path, MULTI_PAGE, "multi.drawio")
    assert [n.id for n in parse_drawio(src, page="第二页").nodes] == ["y"]


def test_parse_drawio_missing_file(tmp_path):
    with pytest.raises(FileNotFoundError):
        parse_drawio(str(tmp_path / "nope.drawio"))


def test_parse_drawio_bad_source(tmp_path):
    p = tmp_path / "bad.drawio"
    p.write_text("not draw.io at all", encoding="utf-8")
    with pytest.raises(ValueError, match=r"无法解析 draw.io 源码"):
        parse_drawio(str(p))


def test_parse_drawio_empty_diagram(tmp_path):
    xml = (
        '<mxfile><diagram name="P"><mxGraphModel><root>'
        '<mxCell id="0"/><mxCell id="1" parent="0"/>'
        "</root></mxGraphModel></diagram></mxfile>"
    )
    with pytest.raises(ValueError, match="未提取到任何节点或边"):
        parse_drawio(_write(tmp_path, xml, "empty.drawio"))


def test_parse_drawio_dangling_edge_filtered(tmp_path):
    xml = (
        '<mxfile><diagram name="P"><mxGraphModel><root>'
        '<mxCell id="0"/><mxCell id="1" parent="0"/>'
        '<mxCell id="a" value="A" vertex="1" parent="1">'
        '<mxGeometry x="0" y="0" width="50" height="30" as="geometry"/></mxCell>'
        '<mxCell id="e1" value="dangle" style="edgeStyle=orthogonalEdgeStyle;" '
        'edge="1" source="a" target="ghost" parent="1">'
        '<mxGeometry relative="1" as="geometry"/></mxCell>'
        "</root></mxGraphModel></diagram></mxfile>"
    )
    d = parse_drawio(_write(tmp_path, xml, "dangle.drawio"))
    assert [n.id for n in d.nodes] == ["a"]
    assert d.edges == []  # 悬空边（target 指向不存在的节点）被过滤


def test_drawio_lazy_import_no_pptx():
    """drawio.py 顶层不得 import pptx（惰性 import 红线）。"""
    code = (
        "import sys\n"
        "import offipy.drawio\n"
        "assert 'pptx' not in sys.modules, 'drawio.py 顶层不得 import pptx'\n"
    )
    subprocess.run([sys.executable, "-c", code], check=True)


def test_layout_drawio_centers_in_box(tmp_path):
    # #95：非绑定轴居中留白对称（SINGLE_PAGE raw 340×300、max 12×6.75 → scale=0.0225、
    # content 7.65×6.75、高度为绑定轴 → off_y=0、off_x=2.175、左右边距对称）
    d = parse_drawio(_write(tmp_path))
    lay = layout_drawio(d)
    assert lay.canvas_w == pytest.approx(7.65)
    assert lay.canvas_h == pytest.approx(6.75)
    min_x = min(n.x for n in lay.nodes)
    max_x = max(n.x + n.w for n in lay.nodes)
    assert min_x == pytest.approx(2.175)
    assert max_x == pytest.approx(9.825)
    assert 12.0 - max_x == pytest.approx(min_x)  # 左右留白对称
    assert min(n.y for n in lay.nodes) == pytest.approx(0.0)
    assert all(n.w > 0 and n.h > 0 for n in lay.nodes if not n.is_container)


def test_layout_drawio_fits_max(tmp_path):
    d = parse_drawio(_write(tmp_path))
    lay = layout_drawio(d, max_w=3.0, max_h=2.0)
    assert lay.canvas_w <= 3.0 + 1e-9
    assert lay.canvas_h <= 2.0 + 1e-9


def test_layout_drawio_color_passthrough(tmp_path):
    d = parse_drawio(_write(tmp_path))
    lay = layout_drawio(d)
    by_id = {n.id: n for n in lay.nodes}
    assert by_id["a"].fill == "#dae8fc"
    assert by_id["a"].stroke == "#6c8ebf"
    assert by_id["a"].font_color == "#000000"
    assert by_id["a"].is_container is False
    assert by_id["g"].is_container is True


def test_layout_drawio_shape_synthesis(tmp_path):
    d = parse_drawio(_write(tmp_path, SHAPES, "shapes.drawio"))
    lay = layout_drawio(d)
    by_id = {n.id: n for n in lay.nodes}
    assert by_id["ra"].shape == "rectangle"  # rect + rounded=0 → 直角
    assert by_id["rb"].shape == "round"  # rect + rounded=1 → 圆角
    assert by_id["el"].shape == "ellipse"
    assert by_id["tr"].shape == "triangle"


def test_layout_drawio_edge_anchors_on_node_edges(tmp_path):
    d = parse_drawio(_write(tmp_path))
    lay = layout_drawio(d)
    assert len(lay.edges) == 1
    e = lay.edges[0]
    by_id = {n.id: n for n in lay.nodes}
    a, b = by_id["a"], by_id["b"]
    # b 在 a 右侧且同 y → 水平主导：源取右缘中点、目标取左缘中点
    # （无方向概念，固定规则，不做 BT/RL 反转）。精确钉死锚点，防回归成 (0,0)。
    assert e.ax1 == pytest.approx(a.x + a.w)
    assert e.ay1 == pytest.approx(a.y + a.h / 2)
    assert e.ax2 == pytest.approx(b.x)
    assert e.ay2 == pytest.approx(b.y + b.h / 2)
    assert e.stroke == "#333333"
    assert e.style == "dashed"


def test_parse_drawio_edge_waypoints(tmp_path):
    d = parse_drawio(_write(tmp_path, WAYPOINTS, "waypoints.drawio"))
    assert len(d.edges) == 1
    e = d.edges[0]
    assert e.waypoints == [(80.0, 60.0), (170.0, 60.0)]  # mxPoint 坐标（排除 source/target）
    assert e.edge_style == "orthogonalEdgeStyle"  # edgeStyle 不再坍缩成 "orthogonal"


def test_layout_drawio_transforms_waypoints(tmp_path):
    d = parse_drawio(_write(tmp_path, WAYPOINTS, "waypoints.drawio"))
    lay = layout_drawio(d)
    assert len(lay.edges) == 1
    e = lay.edges[0]
    assert e.waypoints is not None
    # raw 包络 250×130（a(0,0,50,30)+b(200,100,50,30)）；max 12×6.75
    # scale=min(1,12/250,6.75/130)=0.048 → content 12×6.24，宽度绑定轴 → off_x=0、off_y=0.255
    assert e.waypoints[0] == pytest.approx((3.84, 3.135))
    assert e.waypoints[1] == pytest.approx((8.16, 3.135))


def test_render_drawio_waypoint_polyline(tmp_path):
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.oxml.ns import qn

    from offipy.diagrams import render_to_slide

    d = parse_drawio(_write(tmp_path, WAYPOINTS, "waypoints.drawio"))
    lay = layout_drawio(d)
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    render_to_slide(slide, lay)
    freeforms = [sh for sh in slide.shapes if sh.shape_type == MSO_SHAPE_TYPE.FREEFORM]
    assert len(freeforms) == 1  # waypoint 边 → build_freeform polyline（非单直线 connector）
    ff = freeforms[0]
    assert ff.line.color.rgb == RGBColor(0x33, 0x33, 0x33)
    assert ff.line.width == Emu(15240)  # Pt(1.2)
    assert ff.line._get_or_add_ln().find(qn("a:tailEnd")) is not None  # 箭头保留


def test_parse_drawio_stroke_rotation_dash(tmp_path):
    d = parse_drawio(_write(tmp_path, STYLED, "styled.drawio"))
    by_id = {n.id: n for n in d.nodes}
    a = by_id["a"]
    assert a.stroke_width == 3.0
    assert a.rotation == 45.0
    assert a.dash_pattern == "3 2"
    b = by_id["b"]
    assert b.dash_pattern == "8 3 2 2"
    assert b.stroke_width is None  # 未指定 → None（默认线宽）
    assert len(d.edges) == 1
    assert d.edges[0].stroke_width == 4.0


def test_parse_drawio_rotation_from_geometry(tmp_path):
    # rotation 的真实出口是 mxGeometry 属性（非 style key）——验证 vendored 兜底
    xml = (
        '<mxfile><diagram name="P"><mxGraphModel><root>'
        '<mxCell id="0"/><mxCell id="1" parent="0"/>'
        '<mxCell id="r" value="R" style="strokeWidth=2;" vertex="1" parent="1">'
        '<mxGeometry x="0" y="0" width="50" height="30" rotation="90" as="geometry"/>'
        "</mxCell>"
        "</root></mxGraphModel></diagram></mxfile>"
    )
    n = parse_drawio(_write(tmp_path, xml, "rot.drawio")).nodes[0]
    assert n.rotation == 90.0
    assert n.stroke_width == 2.0


def test_render_drawio_styles_propagate(tmp_path):
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.oxml.ns import qn
    from pptx.util import Pt

    from offipy.diagrams import render_to_slide

    d = parse_drawio(_write(tmp_path, STYLED, "styled.drawio"))
    lay = layout_drawio(d)
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    render_to_slide(slide, lay)
    by_text = {sh.text_frame.text: sh for sh in slide.shapes if sh.has_text_frame}
    sh = by_text["转"]
    assert sh.line.width == Pt(3)  # strokeWidth=3 透传（drawio 1px≈1pt）
    assert sh.rotation == 45  # rotation 透传
    ln = sh.line._get_or_add_ln()
    # dashPattern="3 2" → 预设 DASH（prstDash），非 custDash
    assert ln.find(qn("a:prstDash")) is not None
    assert ln.find(qn("a:custDash")) is None
    custom = by_text["自定义"]
    cln = custom.line._get_or_add_ln()
    cust = cln.find(qn("a:custDash"))  # "8 3 2 2" 不在预设 → 自定义
    assert cust is not None
    assert len(cust.findall(qn("a:ds"))) == 2
    # 边 strokeWidth=4 透传（无 waypoints → 单直线 connector）
    lines = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.LINE]
    assert lines and lines[0].line.width == Pt(4)


def test_render_drawio_preserves_colors(tmp_path):
    # spec「render_to_slide(use_colors=True)：读回 PPTX 断言形状类型、fill/stroke 颜色、
    # label 文本」的 drawio fixture 覆盖。用函数内 import（Task 4 才把 pptx 提为顶层）。
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.dml import MSO_LINE_DASH_STYLE
    from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE

    from offipy.diagrams import render_to_slide

    d = parse_drawio(_write(tmp_path))
    lay = layout_drawio(d)
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    render_to_slide(slide, lay)
    by_text = {sh.text_frame.text: sh for sh in slide.shapes if sh.has_text_frame}
    a = by_text["节点A"]
    assert a.auto_shape_type == MSO_SHAPE.RECTANGLE  # rounded=0 → 直角
    assert a.fill.fore_color.rgb == RGBColor(0xDA, 0xE8, 0xFC)
    assert a.line.color.rgb == RGBColor(0x6C, 0x8E, 0xBF)
    c = by_text["节点C"]
    assert c.auto_shape_type == MSO_SHAPE.ROUNDED_RECTANGLE  # rounded=1 → 圆角
    assert c.fill.fore_color.rgb == RGBColor(0xFF, 0xE6, 0xCC)
    g = by_text["容器"]
    assert g.fill.fore_color.rgb == RGBColor(0xFF, 0xF2, 0xCC)  # 容器 fillColor 保留
    conns = [sh for sh in slide.shapes if sh.shape_type == MSO_SHAPE_TYPE.LINE]
    assert conns and conns[0].line.color.rgb == RGBColor(0x33, 0x33, 0x33)
    assert conns[0].line.dash_style == MSO_LINE_DASH_STYLE.DASH  # dashed=1 → 虚线


def test_layout_drawio_no_upscale(tmp_path):
    # 图比画布小 → scale 必须封顶 1.0（绝不放大）
    xml = (
        '<mxfile><diagram name="P"><mxGraphModel><root>'
        '<mxCell id="0"/><mxCell id="1" parent="0"/>'
        '<mxCell id="s" value="S" vertex="1" parent="1">'
        '<mxGeometry x="0" y="0" width="5" height="3" as="geometry"/></mxCell>'
        "</root></mxGraphModel></diagram></mxfile>"
    )
    lay = layout_drawio(parse_drawio(_write(tmp_path, xml, "small.drawio")))
    assert lay.canvas_w == pytest.approx(5.0)  # 未被放大到 12.0
    assert lay.canvas_h == pytest.approx(3.0)
    s = next(n for n in lay.nodes if n.id == "s")
    assert s.w == pytest.approx(5.0) and s.h == pytest.approx(3.0)


def test_layout_drawio_uniform_scale(tmp_path):
    # 等比 fit：canvas 宽高比 = raw 宽高比（SINGLE_PAGE raw 340x300），
    # x/y 各自独立缩放会被此断言抓出。
    lay = layout_drawio(parse_drawio(_write(tmp_path)))
    assert lay.canvas_w / lay.canvas_h == pytest.approx(340 / 300)


def test_drawio_to_pptx_smoke(tmp_path):
    out = str(tmp_path / "out.pptx")
    result = drawio_to_pptx(_write(tmp_path), out)
    assert result == out
    prs = Presentation(out)
    assert prs.slide_width == Emu(12192000)  # 16:9
    texts = {sh.text_frame.text for sl in prs.slides for sh in sl.shapes if sh.has_text_frame}
    assert texts >= {"节点A", "节点B", "节点C", "容器", "连接"}


def test_drawio_to_pptx_bad_source(tmp_path):
    p = tmp_path / "bad.drawio"
    p.write_text("not draw.io at all", encoding="utf-8")
    with pytest.raises(ValueError, match="无法解析"):
        drawio_to_pptx(str(p), str(tmp_path / "x.pptx"))


def test_drawio_to_pptx_missing_file_raises(tmp_path):
    with pytest.raises(FileNotFoundError):
        drawio_to_pptx(tmp_path / "nope.drawio", str(tmp_path / "out.pptx"))


FONT_SIZE = """\
<mxfile><diagram name="P"><mxGraphModel><root>
  <mxCell id="0"/><mxCell id="1" parent="0"/>
  <mxCell id="a" value="大" style="fontSize=14;fillColor=#dae8fc;" vertex="1" parent="1">
    <mxGeometry x="0" y="0" width="100" height="50" as="geometry"/></mxCell>
  <mxCell id="b" value="小" style="rounded=0;" vertex="1" parent="1">
    <mxGeometry x="120" y="0" width="100" height="50" as="geometry"/></mxCell>
  <mxCell id="c" value="坏" style="fontSize=abc;" vertex="1" parent="1">
    <mxGeometry x="240" y="0" width="100" height="50" as="geometry"/></mxCell>
</root></mxGraphModel></diagram></mxfile>
"""


FONT_DIAGRAM = """\
<mxfile><diagram name="P"><mxGraphModel><root>
  <mxCell id="0"/><mxCell id="1" parent="0"/>
  <mxCell id="big" value="大" style="fontSize=28;fillColor=#dae8fc;" vertex="1" parent="1">
    <mxGeometry x="0" y="0" width="100" height="50" as="geometry"/></mxCell>
  <mxCell id="small" value="小" style="fontSize=12;fillColor=#dae8fc;" vertex="1" parent="1">
    <mxGeometry x="120" y="0" width="100" height="50" as="geometry"/></mxCell>
</root></mxGraphModel></diagram></mxfile>
"""


def test_layout_drawio_font_scales(tmp_path):
    # #97：#fontSize 按 scale 换算成 pt（缺省 fontSize → 12 默认），单位 in/px × pt/in
    d = parse_drawio(_write(tmp_path, FONT_DIAGRAM, "font2.drawio"))
    lay = layout_drawio(d)
    scale = min(1.0, 12.0 / 220, 6.75 / 50)  # raw 220×50
    by_id = {n.id: n for n in lay.nodes}
    assert by_id["big"].font_pt == pytest.approx(28.0 * scale * 72)
    assert by_id["small"].font_pt == pytest.approx(12.0 * scale * 72)


def test_render_drawio_respects_font_size(tmp_path):
    # #97：render 层字号 = font_pt（层级不被拍平），相对比较避免依赖绝对 scale 数学
    from pptx import Presentation

    from offipy.diagrams import render_to_slide

    d = parse_drawio(_write(tmp_path, FONT_DIAGRAM, "font3.drawio"))
    lay = layout_drawio(d)
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    render_to_slide(slide, lay)
    by_text = {sh.text_frame.text: sh for sh in slide.shapes if sh.has_text_frame}
    big = by_text["大"].text_frame.paragraphs[0].runs[0].font.size
    small = by_text["小"].text_frame.paragraphs[0].runs[0].font.size
    assert big is not None and small is not None
    assert big > small


def test_remove_bbox_shapes_only_matches_injected():
    # #96：占位删除改几何一致匹配——中心点落入 bbox 但尺寸不同的用户形状保留，
    # 与注入矩形同位置同尺寸的占位才删
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    prs.slide_width = Emu(12192000)
    prs.slide_height = Emu(6858000)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    user = slide.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(2.4), Inches(1.4))
    user.text = "用户形状"  # 中心 (2.0, 1.5) 恰与 bbox 中心重合，但尺寸不同
    placeholder = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
    placeholder.text = "占位"
    box_emu = {
        "x": int(Inches(1)),
        "y": int(Inches(1)),
        "w": int(Inches(2)),
        "h": int(Inches(1)),
    }
    _remove_bbox_shapes(slide, box_emu)
    remaining = list(slide.shapes)
    assert len(remaining) == 1
    assert remaining[0].text_frame.text == "用户形状"


def test_parse_drawio_font_size(tmp_path):
    # #97：fontSize 从 style 提取；缺省 / 非数值 → None（走 12pt 默认）
    d = parse_drawio(_write(tmp_path, FONT_SIZE, "font.drawio"))
    by_id = {n.id: n for n in d.nodes}
    assert by_id["a"].font_size == 14.0
    assert by_id["b"].font_size is None
    assert by_id["c"].font_size is None
