"""生成 tests/fixtures/audit/ 下的固定审计资产（纯 python-pptx，无 COM/无 chromium）。

产出（默认页面 10×7.5，4:3，与 python-pptx Presentation() 默认一致）：
- synthetic.pptx      一条规则一个场景：越界/贴边/部分重叠/文本被覆盖/nowrap/
                       显式多行/页码/全页背景/off-canvas/normAutofit/spAutoFit
- edge_cases.pptx     缩放 group/嵌套 group/rotation/flipH/flipV/hidden/
                       connector/table/normAutofit/spAutoFit
- baseline.pptx       对比基线：贴边 margin.right LOW + 常规 shape
- candidate.pptx      对比候选：修复贴边 + 新增越界 + 移动/缩放 + 文本变化 + 增删 shape

deck_generated.pptx 由 offipy deck 渲染示例 starter 产生（需 chromium），见 README.md。
这些 .pptx 是一次性资产入库；CI 不要求每次字节级重建（改动审计逻辑时手动重跑本脚本）。
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Inches, Pt

HERE = Path(__file__).resolve().parent


def _page():
    prs = Presentation()
    return prs, prs.slides.add_slide(prs.slide_layouts[6])


def _box(slide, x, y, w, h, text="", wrap=True):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    if text:
        tb.text_frame.text = text
    tb.text_frame.word_wrap = wrap
    return tb


def _rect(slide, x, y, w, h, fill=1):
    return slide.shapes.add_shape(fill, Inches(x), Inches(y), Inches(w), Inches(h))


def _synthetic() -> Presentation:
    prs, slide = _page()
    # 全页背景（z=0，覆盖 100% 页面 → background 角色，豁免 margin/不参与 overlap）
    _rect(slide, 0, 0, 10, 7.5, fill=1)
    # 部分越界 HIGH：out_ratio=(11.2-10)/2=0.6>0.5
    _rect(slide, 9.2, 3, 2, 1)
    # 完全在画布外 LOW：area 0.4in²<1
    _rect(slide, 10.2, 0.5, 0.8, 0.5)
    # 贴边 margin.right LOW：右缘 9.9，距 10 边 0.1<0.2
    _rect(slide, 9.7, 5, 0.2, 1)
    # 部分重叠 MID：ratio=(1.5*1.5)/4=0.5625>0.5
    _rect(slide, 1, 1, 2, 2)
    _rect(slide, 1.5, 1.5, 2, 2)
    # 文本被填充形状完全覆盖 → covered_text HIGH
    _rect(slide, 6, 1, 3, 2)
    _box(slide, 6.5, 1.2, 2, 1, text="被覆盖的文本")
    # nowrap 单行超宽 → text.fit.horizontal
    _box(slide, 4, 3, 1, 0.4, text="W" * 20, wrap=False)
    # 显式多行超高 → text.fit.vertical
    _box(slide, 5, 4.5, 1.5, 0.4, text="line1\nline2\nline3\nline4")
    # 页码候选：纯数字 + 底部 15% 区域 + 小尺寸 → page_number，豁免 margin
    _box(slide, 9.3, 7.05, 0.6, 0.3, text="3")
    # normAutofit → text.autofit.shrink（注入 fontScale=40%：18pt→7.2pt < 8pt 最小可读 → HIGH）
    nf = slide.shapes.add_textbox(Inches(7), Inches(5), Inches(1.5), Inches(0.4))
    nf.text_frame.text = "缩小的文本"
    if nf.text_frame.paragraphs[0].runs:
        nf.text_frame.paragraphs[0].runs[0].font.size = Pt(18)
    nf.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    norm_el = nf.text_frame._txBody.xpath("./a:bodyPr/a:normAutofit")
    if norm_el:
        norm_el[0].set("fontScale", "40000")
    # spAutoFit → text.autofit.grow
    sf = slide.shapes.add_textbox(Inches(0.3), Inches(5.5), Inches(1), Inches(0.3))
    sf.text_frame.text = "扩张的文本"
    sf.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    return prs


def _edge_cases() -> Presentation:
    prs, slide = _page()
    # 缩放 group：ext=(3,2) ≠ chExt=(1.5,1) → 缩放 2×2；
    # 子 "gs" 局部 (0,0,0.75,0.5) → 绝对 (1,1,1.5,1)。
    # 陷阱：python-pptx 往嵌套 group 里加子 shape 会重置外层 group 的
    # p:grpSpPr/a:xfrm 为 off=(0,0) ext=(1,1) → 外层 position 必须在所有
    # 子（含嵌套）都加完之后再设，chExt 最后强制。
    g = slide.shapes.add_group_shape()
    g.shapes.add_textbox(Inches(0), Inches(0), Inches(0.75), Inches(0.5)).text_frame.text = "gs"
    g2 = g.shapes.add_group_shape()
    g2.shapes.add_textbox(Inches(0), Inches(0), Inches(0.5), Inches(0.5)).text_frame.text = "deep"
    g2.left, g2.top = Inches(0.75), Inches(0.5)
    g2.width, g2.height = Inches(0.75), Inches(0.5)
    g.left, g.top = Inches(1), Inches(1)
    g.width, g.height = Inches(3), Inches(2)
    g_xfrm = g._element.xpath(".//a:xfrm")[0]
    g_xfrm.find("{http://schemas.openxmlformats.org/drawingml/2006/main}chExt").set(
        "cx", str(int(1.5 * 914400))
    )
    g_xfrm.find("{http://schemas.openxmlformats.org/drawingml/2006/main}chExt").set(
        "cy", str(int(1.0 * 914400))
    )
    # 旋转 45° 的 shape（AABB 近似，处于安全区无越界/贴边/重叠）
    rot = _rect(slide, 5, 5, 1, 0.6)
    rot.rotation = 45
    # flipH group（绕组中心镜像，chExt 强约束观察位移；与其它形状不相邻）
    fg = slide.shapes.add_group_shape()
    fg.shapes.add_textbox(Inches(0), Inches(0), Inches(0.5), Inches(1)).text_frame.text = "fl"
    fg.left, fg.top = Inches(7), Inches(2)
    fg.width, fg.height = Inches(1), Inches(1)
    fg_xfrm = fg._element.xpath(".//a:xfrm")[0]
    fg_xfrm.set("flipH", "1")
    fg_xfrm.find("{http://schemas.openxmlformats.org/drawingml/2006/main}chExt").set(
        "cx", str(int(1.0 * 914400))
    )
    fg_xfrm.find("{http://schemas.openxmlformats.org/drawingml/2006/main}chExt").set(
        "cy", str(int(1.0 * 914400))
    )
    # flipV shape（普通 shape 的 xfrm 翻转；bbox 不变，只影响内容朝向）
    fv = _rect(slide, 0.5, 3.2, 1, 1)
    fv._element.xpath(".//a:xfrm")[0].set("flipV", "1")
    # hidden shape（不参与任何规则）
    h = _rect(slide, 8, 6, 1, 1)
    h._element.xpath("./p:nvSpPr/p:cNvPr")[0].set("hidden", "1")
    # connector（不参与面积 overlap）
    slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(2), Inches(6), Inches(4), Inches(7))
    # table（含文本；TextFit/Autofit 跳过表格单元格）
    tbl = slide.shapes.add_table(2, 2, Inches(3), Inches(5), Inches(1.5), Inches(0.8))
    tbl.table.cell(0, 0).text = "a"
    tbl.table.cell(1, 1).text = "b"
    return prs


def _compare_pair() -> tuple[Presentation, Presentation]:
    base, bs = _page()
    cand, cs = _page()
    # 构造顺序保持一致 → 前四个 shape_id 对应（2/3/4/5）
    # s1：base 贴边 → cand 修复（移回安全区）
    _rect(bs, 9.7, 5, 0.2, 1)
    _rect(cs, 5, 5, 0.2, 1)
    # s2：常规 shape，两版不变
    _rect(bs, 1, 1, 2, 2)
    _rect(cs, 1, 1, 2, 2)
    # s3：文本框，cand 改文本
    _box(bs, 4, 1, 2, 1, text="hello")
    _box(cs, 4, 1, 2, 1, text="world")
    # s4：cand 移动 + 缩放
    _rect(bs, 6, 1, 2, 2)
    _rect(cs, 6, 2, 3, 2)
    # s5/s6：仅 cand 新增（越界 + 常规）
    _rect(cs, 8, 3, 3, 1)  # 新增越界 → bounds.partial MID
    _rect(cs, 0.5, 6, 1, 1)
    return base, cand


def main() -> None:
    Path(HERE).mkdir(exist_ok=True, parents=True)
    _synthetic().save(HERE / "synthetic.pptx")
    _edge_cases().save(HERE / "edge_cases.pptx")
    base, cand = _compare_pair()
    base.save(HERE / "baseline.pptx")
    cand.save(HERE / "candidate.pptx")


if __name__ == "__main__":
    main()
