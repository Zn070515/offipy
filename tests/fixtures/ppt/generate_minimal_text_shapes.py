"""重建最小 PPT 固定 fixture：tests/fixtures/ppt/minimal_text_shapes.pptx。

维护工具，非测试——CI（Linux）不跑本脚本，只在 Windows 本机有 python-pptx
时手工执行以重建/更新 fixture。fixture 结构验证见 tests/test_ppt_fixture_structure.py。

5 页：
1. 标题+正文占位符页（title=1 / body=2；body 由 layout 1 的 object=7 重打标签）
2. 纯文本框页（无占位符，考验摘要阅读顺序回退）
3. group 嵌套文本页（外层 grpSp + 内层 grpSp，chOff=0/chExt=ext 恒等变换）
4. 页码/页脚/日期占位符页（sldNum=13 / ftr=15 / dt=16，均为摘要豁免集）
5. 无文本图片页（PICTURE=13，read_slide_texts 不应返回）

占位符/group 无法用 python-pptx 公开 API 建（无 add_placeholder；新建页不物化
date/footer/sldNum；group 需手工 grpSp XML），故以下用手工 OOXML 插入，并在
注释标注与微软官方枚举/OOXML type 的对应（不引用项目自身常量，防自证）。
"""

from __future__ import annotations

import io
import struct
import zlib
from pathlib import Path

from pptx import Presentation
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls
from pptx.util import Emu, Inches

HERE = Path(__file__).resolve().parent
OUT = HERE / "minimal_text_shapes.pptx"

# 页尺寸 10in x 7.5in（720 x 540pt）。1in = 914400 EMU = 72pt。
PAGE_W_IN = 10.0
PAGE_H_IN = 7.5


def _emu(inches: float) -> int:
    """英寸 → EMU（OOXML a:off/a:ext 单位）。"""
    return int(inches * 914400)


def _png_bytes() -> bytes:
    """1x1 红色 PNG（无外部依赖，python-pptx add_picture 用）。"""

    def chunk(tag: bytes, data: bytes) -> bytes:
        c = struct.pack(">I", len(data)) + tag + data
        return c + struct.pack(">I", zlib.crc32(tag + data) & 0xFFFFFFFF)

    ihdr = chunk(b"IHDR", struct.pack(">IIBBBBB", 1, 1, 8, 6, 0, 0, 0))
    idat = chunk(b"IDAT", zlib.compress(b"\x00\xff\x00\x00"))
    iend = chunk(b"IEND", b"")
    return b"\x89PNG\r\n\x1a\n" + ihdr + idat + iend


def _sp_xml(
    sid: int,
    name: str,
    text: str,
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    ph_type: str | None = None,
    ph_idx: int | None = None,
) -> str:
    """单个 <p:sp> 的 OOXML。ph_type 非 None 时带上占位符属性（type 值见注释）。"""
    nv_pr = ""
    if ph_type is not None:
        # ph type 官方值（OOXML 小写）：title=1 body=2 sldNum=13 ftr=15 dt=16 ...
        idx = f' idx="{ph_idx}"' if ph_idx is not None else ""
        nv_pr = f'<p:nvPr><p:ph type="{ph_type}"{idx}/></p:nvPr>'
    xfrm = (
        f'<a:xfrm><a:off x="{_emu(left)}" y="{_emu(top)}"/>'
        f'<a:ext cx="{_emu(width)}" cy="{_emu(height)}"/></a:xfrm>'
    )
    return f"""<p:sp {nsdecls("p", "a")}>
  <p:nvSpPr>
    <p:cNvPr id="{sid}" name="{name}"/>
    <p:cNvSpPr txBox="1"/>
    {nv_pr}
  </p:nvSpPr>
  <p:spPr>
    {xfrm}
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
  </p:spPr>
  <p:txBody><a:bodyPr/><a:lstStyle/><a:p><a:r><a:t>{text}</a:t></a:r></a:p></p:txBody>
</p:sp>"""


def _grp_sp_xml(
    sid: int,
    name: str,
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    children: list[str],
) -> str:
    """<p:grpSp>：chOff=(0,0)、chExt=ext 恒等变换——子坐标即幻灯片绝对坐标。"""
    children_xml = "\n".join(children)
    return f"""<p:grpSp {nsdecls("p", "a")}>
  <p:nvGrpSpPr>
    <p:cNvPr id="{sid}" name="{name}"/>
    <p:cNvGrpSpPr/>
    <p:nvPr/>
  </p:nvGrpSpPr>
  <p:grpSpPr>
    <a:xfrm>
      <a:off x="{_emu(left)}" y="{_emu(top)}"/>
      <a:ext cx="{_emu(width)}" cy="{_emu(height)}"/>
      <a:chOff x="0" y="0"/>
      <a:chExt cx="{_emu(width)}" cy="{_emu(height)}"/>
    </a:xfrm>
  </p:grpSpPr>
  {children_xml}
</p:grpSp>"""


def _blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _make():
    prs = Presentation()
    prs.slide_width = Emu(_emu(PAGE_W_IN))
    prs.slide_height = Emu(_emu(PAGE_H_IN))

    # ---- 页 1：标题+正文占位符 ----
    s1 = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
    s1.shapes.title.text = "产品发布计划"
    body = s1.placeholders[1]
    # layout 1 的内容占位符是 object(7)，重打标签为 body(2)（p:ph type 属性）
    body._element.ph.set("type", "body")
    body.text_frame.text = "第一季度\n第二季度\n第三季度"
    s1.notes_slide.notes_text_frame.text = "演讲者备注：本页为公司级发布计划"

    # ---- 页 2：纯文本框（无占位符） ----
    s2 = _blank_slide(prs)
    tb_title = s2.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(6), Inches(0.5))
    tb_title.name = "TextBoxTitle"
    tb_title.text_frame.text = "纯文本框标题"
    tb_body = s2.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(6), Inches(1.2))
    tb_body.name = "TextBoxBody"
    tb_body.text_frame.text = "纯文本框正文第一行\n纯文本框正文第二行"

    # ---- 页 3：group 嵌套文本（恒等变换，子坐标 = 幻灯片绝对坐标） ----
    s3 = _blank_slide(prs)
    nested_title = _sp_xml(
        302, "NestedTitle", "嵌套组内标题", left=1.0, top=1.0, width=4.0, height=0.5
    )
    nested_body = _sp_xml(
        303, "NestedBody", "嵌套组内正文", left=1.0, top=1.8, width=4.0, height=0.6
    )
    inner_grp = _grp_sp_xml(
        301,
        "GroupInner",
        left=0.5,
        top=0.5,
        width=5.0,
        height=2.0,
        children=[nested_title, nested_body],
    )
    outer_text = _sp_xml(304, "OuterText", "外层组内文本", left=0.3, top=2.6, width=4.0, height=0.5)
    outer_grp = _grp_sp_xml(
        300,
        "GroupOuter",
        left=1.0,
        top=1.0,
        width=6.0,
        height=3.0,
        children=[inner_grp, outer_text],
    )
    s3.shapes._spTree.append(parse_xml(outer_grp))

    # ---- 页 4：页码/页脚/日期占位符（豁免集 13/15/16）+ 角上纯数字文本框 ----
    s4 = _blank_slide(prs)
    s4_title = s4.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(6), Inches(0.5))
    s4_title.name = "TitleText"
    s4_title.text_frame.text = "季度财报"
    # 角上纯数字文本框（非占位符）→ 走页码候选豁免路径（P1-2）
    s4_num = s4.shapes.add_textbox(Inches(8.5), Inches(6.9), Inches(0.6), Inches(0.3))
    s4_num.name = "TextBoxPageNum"
    s4_num.text_frame.text = "7"
    # 页脚(15)/日期(16)/页码(13)占位符，手工 XML
    for sid, name, ph_type, text, x in (
        (400, "Footer 1", "ftr", "机密", 1.0),
        (401, "Date 1", "dt", "2026-08-05", 7.0),
        (402, "SlideNumberPlaceholder 1", "sldNum", "4", 9.0),
    ):
        s4.shapes._spTree.append(
            parse_xml(
                _sp_xml(
                    sid,
                    name,
                    text,
                    left=x,
                    top=7.0,
                    width=1.5,
                    height=0.3,
                    ph_type=ph_type,
                    ph_idx=10,
                )
            )
        )

    # ---- 页 5：无文本图片 + 说明文本框 ----
    s5 = _blank_slide(prs)
    s5.shapes.add_picture(io.BytesIO(_png_bytes()), Inches(1), Inches(1), Inches(2), Inches(2))
    cap = s5.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(3), Inches(0.5))
    cap.name = "PictureCaption"
    cap.text_frame.text = "图片说明"

    prs.save(OUT)
    print(f"written {OUT} ({OUT.stat().st_size} bytes)")


if __name__ == "__main__":
    _make()
