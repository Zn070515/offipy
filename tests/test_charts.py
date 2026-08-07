# tests/test_charts.py
"""原生图表声明解析测试（纯 Python，不依赖 Office）。"""

import json
import re

import pytest
from pptx import Presentation

from offipy import charts
from offipy.charts import (
    ChartData,
    ChartDecl,
    ChartSeries,
    _cycle_colors,
    _decl_hex_palette,
    _effective_accent,
    _theme_name_from_html,
    derive_chart_palette,
    inject_native_charts,
    load_chart_boxes,
    parse_chart_colors_override,
    parse_chart_declarations,
)
from offipy.design import THEMES, Theme

ATTR_HTML = (
    "<!DOCTYPE html>\n"
    "<html><body>\n"
    '<section class="slide" data-pptx-slide>\n'
    '  <h2 class="title">营收增长</h2>\n'
    '  <div class="chart" data-chart="bar" '
    'data-chart-data=\'{"categories":["Q1","Q2","Q3"],'
    '"series":[{"name":"营收","values":[40,55,70]}]}\'></div>\n'
    "</section>\n"
    '<section class="slide" data-pptx-slide>\n'
    '  <h2 class="title">趋势</h2>\n'
    '  <div class="chart" data-chart="line"></div>\n'
    "</section>\n"
    "</body></html>"
)

SCRIPT_HTML = """<!DOCTYPE html>
<html><body>
<section class="slide" data-pptx-slide>
  <h2 class="title">占比</h2>
  <div class="chart" id="share-chart" data-chart="pie"></div>
  <script type="application/json" data-chart-target="#share-chart">
    {"categories":["A","B","C"],"series":[{"name":"份额","values":[45,35,20]}]}
  </script>
</section>
</body></html>"""


def test_parse_attribute_format():
    # 属性格式正常 → 正常解析返回声明（验证 data-chart-data 解析）
    decls = parse_chart_declarations(
        '<section class="slide" data-pptx-slide>'
        '<div class="chart" data-chart="bar" '
        'data-chart-data=\'{"categories":["Q1","Q2","Q3"],'
        '"series":[{"name":"营收","values":[40,55,70]}]}\'></div>'
        "</section>"
    )
    assert len(decls) == 1
    assert decls[0].slide_index == 1
    assert decls[0].chart_type == "bar"
    assert decls[0].data.categories == ["Q1", "Q2", "Q3"]
    assert decls[0].data.series[0].name == "营收"
    assert decls[0].data.series[0].values == [40.0, 55.0, 70.0]
    # ATTR_HTML 第二页 line 无数据 → 抛 ValueError（契约：声明了图表就必须给数据）
    with pytest.raises(ValueError):
        parse_chart_declarations(ATTR_HTML)


def test_parse_script_format():
    decls = parse_chart_declarations(SCRIPT_HTML)
    assert len(decls) == 1
    assert decls[0].slide_index == 1
    assert decls[0].chart_type == "pie"
    assert decls[0].data.categories == ["A", "B", "C"]


def test_parse_no_chart_returns_empty():
    assert parse_chart_declarations("<section data-pptx-slide><h2>x</h2></section>") == []


def test_parse_invalid_type_raises():
    with pytest.raises(ValueError):
        parse_chart_declarations(
            '<section data-pptx-slide><div class="chart" data-chart="radar" '
            'data-chart-data=\'{"categories":["a"],"series":[{"name":"s","values":[1]]}\'></div></section>'
        )


def test_parse_bad_data_json_raises():
    with pytest.raises(ValueError):
        parse_chart_declarations(
            '<section data-pptx-slide><div class="chart" data-chart="bar" '
            'data-chart-data="{oops}"></div></section>'
        )


def test_parse_series_values_must_be_numeric():
    with pytest.raises(ValueError):
        parse_chart_declarations(
            '<section data-pptx-slide><div class="chart" data-chart="bar" '
            'data-chart-data=\'{"categories":["a"],"series":[{"name":"s","values":["x"]}]}\'></div></section>'
        )


def test_parse_chart_without_data_raises():
    # 契约：打了 data-chart 的容器必须有数据（属性或 script），否则报错
    with pytest.raises(ValueError):
        parse_chart_declarations(
            '<section data-pptx-slide><div class="chart" data-chart="line"></div></section>'
        )


def _blank_pptx(tmp_path):
    from pptx.util import Emu

    prs = Presentation()
    prs.slide_width = Emu(12192000)
    prs.slide_height = Emu(6858000)
    prs.slides.add_slide(prs.slide_layouts[6])
    return prs


def test_load_chart_boxes(tmp_path):
    meas = tmp_path / "measurements.json"
    meas.write_text(
        json.dumps(
            {
                "slides": [
                    {
                        "slide": {"theme": "mckinsey"},
                        "records": [
                            {
                                "id": 1,
                                "kind": "text",
                                "tag": "h2",
                                "className": "title",
                                "rect": {"x": 96, "y": 96, "w": 800, "h": 80},
                                "text": "营收",
                            },
                            {
                                "id": 2,
                                "kind": "shape",
                                "tag": "div",
                                "className": "chart",
                                "rect": {"x": 96, "y": 260, "w": 900, "h": 500},
                                "text": "",
                            },
                            {
                                "id": 3,
                                "kind": "text",
                                "tag": "div",
                                "className": "chart-note",
                                "rect": {"x": 1020, "y": 260, "w": 360, "h": 500},
                                "text": "来源",
                            },
                        ],
                    }
                ]
            }
        ),
        encoding="utf-8",
    )
    boxes = load_chart_boxes(str(meas))
    assert boxes[1] == {"x": 96, "y": 260, "w": 900, "h": 500}
    assert 2 not in boxes  # 只认 className 恰为 "chart" 的容器，chart-note 不匹配


def test_inject_native_chart_bar(tmp_path):
    prs = _blank_pptx(tmp_path)
    slide = prs.slides[0]
    from pptx.util import Emu

    # 放一个占位矩形（模拟 convert 渲染的图表区），再放一个右上角说明（不该被删）
    slide.shapes.add_shape(1, Emu(96 * 6350), Emu(260 * 6350), Emu(900 * 6350), Emu(500 * 6350))
    note = slide.shapes.add_textbox(
        Emu(1020 * 6350), Emu(260 * 6350), Emu(360 * 6350), Emu(500 * 6350)
    )
    note.text = "来源"
    out = tmp_path / "out.pptx"
    prs.save(str(out))

    decls = [
        ChartDecl(
            slide_index=1,
            chart_type="bar",
            data=ChartData(
                categories=["Q1", "Q2"],
                series=[ChartSeries(name="营收", values=[40, 70])],
            ),
        )
    ]
    boxes = {1: {"x": 96, "y": 260, "w": 900, "h": 500}}
    inject_native_charts(str(out), decls, boxes)

    prs2 = Presentation(str(out))
    slide2 = prs2.slides[0]
    charts = [s for s in slide2.shapes if s.has_chart]
    assert len(charts) == 1
    ch = charts[0].chart
    assert ch.chart_type is not None
    assert ch.has_title is False
    # 占位矩形被移除，说明文本保留
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    rects = [s for s in slide2.shapes if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE]
    assert len(rects) == 0
    texts = [s.text for s in slide2.shapes if s.has_text_frame]
    assert "来源" in texts


def test_inject_native_chart_pie(tmp_path):
    prs = _blank_pptx(tmp_path)
    out = tmp_path / "pie.pptx"
    prs.save(str(out))
    decls = [
        ChartDecl(
            slide_index=1,
            chart_type="pie",
            data=ChartData(
                categories=["A", "B"],
                series=[ChartSeries(name="份额", values=[45, 55])],
            ),
        )
    ]
    inject_native_charts(str(out), decls, {1: {"x": 100, "y": 200, "w": 800, "h": 500}})
    prs2 = Presentation(str(out))
    charts = [s for s in prs2.slides[0].shapes if s.has_chart]
    assert len(charts) == 1


def test_parse_top_level_non_dict_raises():
    with pytest.raises(ValueError):
        parse_chart_declarations(
            '<section data-pptx-slide><div class="chart" data-chart="bar" '
            'data-chart-data="[1,2,3]"></div></section>'
        )


def test_parse_chart_outside_slide_raises():
    # 图表容器出现在第一个 <section data-pptx-slide> 之前 → slide_index=0 → 报错而非静默
    with pytest.raises(ValueError):
        parse_chart_declarations(
            '<div class="chart" data-chart="bar" '
            'data-chart-data=\'{"categories":["a"],"series":[{"name":"s","values":[1]}]}\'></div>'
            "<section data-pptx-slide><h2>x</h2></section>"
        )


def test_postprocess_skips_without_charts(tmp_path, monkeypatch):
    # HTML 无 data-chart → 不读 measurements、不改 pptx
    called = {}
    monkeypatch.setattr(charts, "load_chart_boxes", lambda *a, **k: called.setdefault("load", True))
    html = tmp_path / "d.html"
    html.write_text("<section data-pptx-slide><h2>x</h2></section>", encoding="utf-8")
    pptx = tmp_path / "d.pptx"
    pptx.write_bytes(b"not a real pptx")
    charts.postprocess_charts(str(html), str(pptx))  # 不抛异常、load 不被调用
    assert "load" not in called


def test_postprocess_missing_measurements_raises(tmp_path):
    html = tmp_path / "d.html"
    html.write_text(
        '<section data-pptx-slide><div class="chart" data-chart="bar" '
        'data-chart-data=\'{"categories":["a"],"series":[{"name":"s","values":[1]}]}\'></div></section>',
        encoding="utf-8",
    )
    pptx = tmp_path / "d.pptx"
    pptx.write_bytes(b"x")
    with pytest.raises(RuntimeError, match="measurements"):
        charts.postprocess_charts(str(html), str(pptx))


def test_postprocess_calls_inject(monkeypatch, tmp_path):
    from pptx.util import Emu

    prs = Presentation()
    prs.slide_width = Emu(12192000)
    prs.slide_height = Emu(6858000)
    prs.slides.add_slide(prs.slide_layouts[6])
    pptx = tmp_path / "d.pptx"
    prs.save(str(pptx))
    html = tmp_path / "d.html"
    html.write_text(
        '<section data-pptx-slide><div class="chart" data-chart="bar" '
        'data-chart-data=\'{"categories":["a"],"series":[{"name":"s","values":[1,2]}]}\'></div></section>',
        encoding="utf-8",
    )
    meas_dir = tmp_path / "d_audit" / "_cache"
    meas_dir.mkdir(parents=True)
    (meas_dir / "measurements.json").write_text(
        json.dumps(
            {
                "slides": [
                    {
                        "slide": {},
                        "records": [
                            {
                                "id": 1,
                                "kind": "shape",
                                "tag": "div",
                                "className": "chart",
                                "rect": {"x": 100, "y": 200, "w": 800, "h": 500},
                                "text": "",
                            }
                        ],
                    }
                ]
            }
        ),
        encoding="utf-8",
    )

    charts.postprocess_charts(str(html), str(pptx))

    prs2 = Presentation(str(pptx))
    assert any(s.has_chart for s in prs2.slides[0].shapes)


# ---------------------------------------------------------------------------
# Task 0: theme-aware palette + chart color overrides 元数据
# ---------------------------------------------------------------------------


def test_derive_palette_known_accent_golden():
    # 冻结的确定性调色板快照（首色 = accent 本身）
    assert derive_chart_palette("#2251FF") == (
        "#2251FF",
        "#614BFF",
        "#41B8FF",
        "#9418FF",
        "#45F5D0",
        "#EA8F6A",
    )
    palette = derive_chart_palette("#2251FF")
    assert palette[0] == "#2251FF"
    assert len(palette) == 6
    assert all(re.match(r"^#[0-9A-F]{6}$", c) for c in palette)


def test_derive_palette_deterministic():
    # 同 accent（含小写输入）→ 字节级一致；不同 accent → 首色不同
    a = derive_chart_palette("#2251FF")
    b = derive_chart_palette("#2251ff")
    assert a == b
    c = derive_chart_palette("#0E9387")
    assert c[0] == "#0E9387"
    assert a[0] != c[0]


@pytest.mark.parametrize("accent", ["#FFFFFF", "#000000", "#FF0000"])
def test_derive_palette_extremes(accent):
    # 黑白/高饱和极端值：不崩、六色全部合法大写 #RRGGBB、确定性
    palette = derive_chart_palette(accent)
    assert len(palette) == 6
    assert all(re.match(r"^#[0-9A-F]{6}$", c) for c in palette)
    assert palette == derive_chart_palette(accent)


@pytest.mark.parametrize("accent", ["notacolor", "#FFF", "#GGGGGG"])
def test_derive_palette_invalid_accent_raises(accent):
    # 只认 6 位 #RRGGBB；3 位 #FFF / 非 hex 均拒绝（6 位无 # 形式按设计接受，不在本测试断言）
    with pytest.raises(ValueError):
        derive_chart_palette(accent)


def test_parse_colors_override_valid():
    assert parse_chart_colors_override('["#2251FF","#0E9387"]') == ("#2251FF", "#0E9387")
    # 小写输入 → 归一化为大写
    assert parse_chart_colors_override('["#2251ff","#0e9387"]') == ("#2251FF", "#0E9387")


@pytest.mark.parametrize(
    "raw",
    [
        "[]",  # 空列表
        '["#FF00"]',  # 非法 hex
        '["red"]',  # 非 hex
        '"#2251FF"',  # 不是列表
        "[1,2]",  # 非字符串
        '["#2251FF","oops"]',  # 混入非法项
        "",  # 空字符串（JSON 都解析不了）
    ],
)
def test_parse_colors_override_invalid(raw):
    with pytest.raises(ValueError):
        parse_chart_colors_override(raw)


def test_parse_decl_preserves_slide_classes():
    html = (
        '<section class="slide dark" data-pptx-slide>'
        '<div class="chart" data-chart="bar" '
        'data-chart-data=\'{"categories":["Q1"],"series":[{"name":"s","values":[1]}]}\'></div>'
        "</section>"
    )
    decls = parse_chart_declarations(html)
    assert decls[0].slide_classes == frozenset({"slide", "dark"})
    # 无 class 的 section → 空集合
    no_cls = parse_chart_declarations(
        "<section data-pptx-slide>"
        '<div class="chart" data-chart="bar" '
        'data-chart-data=\'{"categories":["Q1"],"series":[{"name":"s","values":[1]}]}\'></div>'
        "</section>"
    )
    assert no_cls[0].slide_classes == frozenset()


def test_parse_decl_slide_classes_reset_per_section():
    # 两页 section：第二页无 class，不得继承第一页的 {"slide","dark"}
    html = (
        '<section class="slide dark" data-pptx-slide>'
        '<div class="chart" data-chart="bar" '
        'data-chart-data=\'{"categories":["Q1"],"series":[{"name":"s","values":[1]}]}\'></div>'
        "</section>"
        "<section data-pptx-slide>"
        '<div class="chart" data-chart="line" '
        'data-chart-data=\'{"categories":["Q1"],"series":[{"name":"s","values":[1]}]}\'></div>'
        "</section>"
    )
    decls = parse_chart_declarations(html)
    assert decls[0].slide_classes == frozenset({"slide", "dark"})
    assert decls[1].slide_classes == frozenset()


def test_parse_decl_preserves_colors_override():
    html = (
        '<section class="slide" data-pptx-slide>'
        '<div class="chart" data-chart="bar" '
        'data-chart-colors=\'["#2251FF","#0E9387"]\' '
        'data-chart-data=\'{"categories":["Q1"],"series":[{"name":"s","values":[1]}]}\'></div>'
        "</section>"
    )
    decls = parse_chart_declarations(html)
    assert decls[0].colors_override == ("#2251FF", "#0E9387")
    # 无该属性 → None
    plain = parse_chart_declarations(
        "<section data-pptx-slide>"
        '<div class="chart" data-chart="bar" '
        'data-chart-data=\'{"categories":["Q1"],"series":[{"name":"s","values":[1]}]}\'></div>'
        "</section>"
    )
    assert plain[0].colors_override is None


def test_parse_decl_empty_colors_override_raises():
    # data-chart-colors="" 也算「存在」→ 按非法 JSON 列表报错
    with pytest.raises(ValueError):
        parse_chart_declarations(
            "<section data-pptx-slide>"
            '<div class="chart" data-chart="bar" data-chart-colors="" '
            'data-chart-data=\'{"categories":["Q1"],"series":[{"name":"s","values":[1]}]}\'></div>'
            "</section>"
        )


def test_cycle_colors():
    assert _cycle_colors(("#a", "#b"), 5) == ("#a", "#b", "#a", "#b", "#a")
    assert _cycle_colors(("#a",), 2) == ("#a", "#a")
    # 截断：n < len(colors)
    assert _cycle_colors(("#a", "#b", "#c"), 2) == ("#a", "#b")


# ---------------------------------------------------------------------------
# Task 1: theme-aware per-slide chart colors
# ---------------------------------------------------------------------------


def test_theme_name_from_html_injected_full_block():
    html = (
        "<!DOCTYPE html><html><head>"
        '<style data-theme="mckinsey">\n:root { --accent: #2251FF; }\n'
        ".slide.dark { --accent: #5B8CFF; }\n</style>"
        "</head><body></body></html>"
    )
    assert _theme_name_from_html(html) == "mckinsey"


def test_theme_name_from_html_empty_placeholder():
    # Claude 约定：<head> 里写空占位 <style data-theme="<name>"></style>
    html = '<html><head><style data-theme="mckinsey"></style></head><body></body></html>'
    assert _theme_name_from_html(html) == "mckinsey"


def test_theme_name_from_html_none_when_absent():
    assert _theme_name_from_html("<html><body><p>no theme here</p></body></html>") is None


def test_effective_accent_base_page():
    assert _effective_accent(THEMES["mckinsey"], {"slide"}) == "#2251FF"


def test_effective_accent_dark_variant_page():
    assert _effective_accent(THEMES["mckinsey"], {"slide", "dark"}) == "#5B8CFF"


def test_effective_accent_dark_tech_light_variant_page():
    assert _effective_accent(THEMES["dark-tech"], {"slide", "light"}) == "#0284C7"


def test_effective_accent_non_variant_class_falls_back_to_base():
    # mckinsey 的 variant 是 .slide.dark；{"slide","light"} 不命中 → base accent
    assert _effective_accent(THEMES["mckinsey"], {"slide", "light"}) == "#2251FF"


def test_effective_accent_unparseable_selector_falls_back_to_base():
    theme = Theme(
        name="custom",
        title="custom",
        description="custom",
        base_vars={"--accent": "#123456"},
        variant_selector=".card.active",  # 非冻结内置 `.slide.<class>` 形式
        variant_vars={"--accent": "#ABCDEF"},
    )
    assert _effective_accent(theme, {"slide", "active"}) == "#123456"


def _chart_decl(
    *,
    slide_classes=frozenset(),
    colors_override=None,
) -> ChartDecl:
    return ChartDecl(
        slide_index=1,
        chart_type="bar",
        data=ChartData(
            categories=["Q1"],
            series=[ChartSeries(name="s", values=[1.0])],
        ),
        slide_classes=slide_classes,
        colors_override=colors_override,
    )


def test_decl_hex_palette_override_wins():
    decl = _chart_decl(colors_override=("#FFFFFF", "#000000"))
    assert _decl_hex_palette(decl, THEMES["mckinsey"]) == ("#FFFFFF", "#000000")


def test_decl_hex_palette_theme_base_accent():
    decl = _chart_decl(slide_classes=frozenset({"slide"}))
    assert _decl_hex_palette(decl, THEMES["mckinsey"]) == derive_chart_palette("#2251FF")


def test_decl_hex_palette_theme_variant_accent():
    decl = _chart_decl(slide_classes=frozenset({"slide", "dark"}))
    assert _decl_hex_palette(decl, THEMES["mckinsey"]) == derive_chart_palette("#5B8CFF")


def test_decl_hex_palette_no_theme_returns_none():
    assert _decl_hex_palette(_chart_decl(), None) is None


def test_decl_hex_palette_unknown_theme_returns_none():
    # THEMES.get("nope") 是 None → 与 theme=None 同路径，回退到固定调色板
    assert _decl_hex_palette(_chart_decl(), THEMES.get("nope")) is None


def test_parse_declarations_regression_task1():
    # Task 1 只加调色板解析，不改声明解析：slide_index/chart_type/data/
    # slide_classes/colors_override 字段逐项对齐
    html = (
        '<section class="slide dark" data-pptx-slide>'
        '<div class="chart" data-chart="bar" '
        'data-chart-colors=\'["#2251FF","#0E9387"]\' '
        'data-chart-data=\'{"categories":["Q1"],"series":[{"name":"s","values":[1]}]}\'></div>'
        "</section>"
        '<section class="slide" data-pptx-slide>'
        '<div class="chart" id="c2" data-chart="pie"></div>'
        '<script type="application/json" data-chart-target="#c2">'
        '{"categories":["A"],"series":[{"name":"p","values":[42]}]}'
        "</script>"
        "</section>"
    )
    decls = parse_chart_declarations(html)
    assert [d.slide_index for d in decls] == [1, 2]
    assert [d.chart_type for d in decls] == ["bar", "pie"]
    assert decls[0].slide_classes == frozenset({"slide", "dark"})
    assert decls[1].slide_classes == frozenset({"slide"})
    assert decls[0].colors_override == ("#2251FF", "#0E9387")
    assert decls[1].colors_override is None
    assert decls[0].data.categories == ["Q1"]
    assert decls[1].data.categories == ["A"]
    assert decls[0].data.series[0].values == [1.0]
    assert decls[1].data.series[0].values == [42.0]
