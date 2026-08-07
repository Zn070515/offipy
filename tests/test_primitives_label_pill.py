"""A5 Task 4 — label-pill renderer structural contract."""

from __future__ import annotations

import pytest
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN

from offipy.assets.model import (
    AssetRect,
    AssetRef,
    AssetRenderContext,
    AssetRequest,
    NativeShapePayload,
)
from offipy.assets.primitives import get_native_renderer
from offipy.assets.providers.primitives import PrimitivesProvider
from offipy.exceptions import InvalidArgumentError

_THEME = {
    "accent": "#2251ff",
    "surface": "#F2F4F7",
    "ink": "#222222",
    "muted": "#667085",
    "bg": "#ffffff",
}


def _ctx(**overrides) -> AssetRenderContext:
    kwargs = dict(
        slide_index=1,
        rect=AssetRect(x=0, y=0, width=400, height=200),
        theme_name="mckinsey",
        theme_vars=dict(_THEME),
        placement="replace",
    )
    kwargs.update(overrides)
    return AssetRenderContext(**kwargs)


def _blank_slide():
    prs = Presentation()
    return prs.slides.add_slide(prs.slide_layouts[6])


def _render(params: dict[str, str], **ctx_overrides):
    resolved = PrimitivesProvider().resolve(
        AssetRequest(
            ref=AssetRef("primitives", "primitive", "label-pill"),
            params=tuple(params.items()),
        )
    )
    assert isinstance(resolved.payload, NativeShapePayload)
    slide = _blank_slide()
    get_native_renderer("label-pill")(slide, dict(resolved.payload.params), _ctx(**ctx_overrides))
    return slide


def _runs(sp):
    return [r.text for p in sp.text_frame.paragraphs for r in p.runs]


def _find_textbox(slide, text):
    for sp in slide.shapes:
        if hasattr(sp, "text_frame") and text in _runs(sp):
            return sp
    return None


def _run_color(slide, text) -> str:
    box = _find_textbox(slide, text)
    assert box is not None
    return str(box.text_frame.paragraphs[0].runs[0].font.color.rgb)


def _autoshapes(slide):
    return [sp for sp in slide.shapes if sp.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE]


def _assert_within_rect(slide, rect: AssetRect, tol_px: float = 1.0) -> None:
    tol = int(round(tol_px * 6350))
    x0, y0 = int(rect.x * 6350), int(rect.y * 6350)
    x1, y1 = int((rect.x + rect.width) * 6350), int((rect.y + rect.height) * 6350)
    for sp in slide.shapes:
        assert sp.left >= x0 - tol, f"left overflow: {sp.left}"
        assert sp.top >= y0 - tol, f"top overflow: {sp.top}"
        assert sp.left + sp.width <= x1 + tol, f"right overflow: {sp.left + sp.width}"
        assert sp.top + sp.height <= y1 + tol, f"bottom overflow: {sp.top + sp.height}"


def test_label_pill_structure() -> None:
    slide = _render({"text": "Onboard"})
    assert len(_autoshapes(slide)) == 1
    textboxes = [sp for sp in slide.shapes if sp.shape_type == MSO_SHAPE_TYPE.TEXT_BOX]
    assert len(textboxes) == 1


def test_label_pill_text_round_trip() -> None:
    slide = _render({"text": "Onboard"})
    assert _find_textbox(slide, "Onboard") is not None


def test_label_pill_text_centered() -> None:
    slide = _render({"text": "Onboard"})
    box = _find_textbox(slide, "Onboard")
    assert box is not None
    tf = box.text_frame
    assert tf.paragraphs[0].alignment == PP_ALIGN.CENTER
    assert tf.vertical_anchor == MSO_ANCHOR.MIDDLE


def test_label_pill_default_accent_fill() -> None:
    slide = _render({"text": "Onboard"})  # default fill is the accent token
    pill = _autoshapes(slide)[0]
    assert str(pill.fill.fore_color.rgb) == "2251FF"
    assert _run_color(slide, "Onboard") == "FFFFFF"


def test_label_pill_accent_override() -> None:
    slide = _render({"text": "Onboard", "accent": "#FF0000"})
    pill = _autoshapes(slide)[0]
    assert str(pill.fill.fore_color.rgb) == "FF0000"


def test_label_pill_light_accent_override_dark_text() -> None:
    slide = _render({"text": "Onboard", "accent": "#FFFFFF"})
    pill = _autoshapes(slide)[0]
    assert str(pill.fill.fore_color.rgb) == "FFFFFF"
    assert _run_color(slide, "Onboard") == "222222"


def test_label_pill_font_shrinks_in_rect() -> None:
    slide = _render({"text": "x" * 50})
    box = _find_textbox(slide, "x" * 50)
    size_pt = box.text_frame.paragraphs[0].runs[0].font.size.pt
    assert 8 <= size_pt < 120  # start was h*0.6 = 120pt


def test_label_pill_shapes_within_rect() -> None:
    slide = _render({"text": "x" * 50})
    _assert_within_rect(slide, _ctx().rect)


def test_label_pill_too_long_raises() -> None:
    with pytest.raises(InvalidArgumentError, match="does not fit"):
        _render({"text": "x" * 120}, rect=AssetRect(x=0, y=0, width=100, height=20))


def test_label_pill_no_pictures() -> None:
    slide = _render({"text": "Onboard"})
    assert not any(sp.shape_type == MSO_SHAPE_TYPE.PICTURE for sp in slide.shapes)
