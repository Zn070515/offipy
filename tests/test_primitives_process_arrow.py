"""A5 Task 7 — process-arrow renderer structural contract."""

from __future__ import annotations

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE

from offipy.assets.model import (
    AssetRect,
    AssetRef,
    AssetRenderContext,
    AssetRequest,
    NativeShapePayload,
)
from offipy.assets.primitives import get_native_renderer
from offipy.assets.providers.primitives import PrimitivesProvider

_THEME = {
    "accent": "#2251ff",
    "surface": "#F2F4F7",
    "ink": "#222222",
    "muted": "#667085",
    "bg": "#ffffff",
}


def _ctx(**overrides) -> AssetRenderContext:
    kwargs = {
        "slide_index": 1,
        "rect": AssetRect(x=0, y=0, width=400, height=200),
        "theme_name": "mckinsey",
        "theme_vars": dict(_THEME),
        "placement": "replace",
    }
    kwargs.update(overrides)
    return AssetRenderContext(**kwargs)


def _blank_slide():
    prs = Presentation()
    return prs.slides.add_slide(prs.slide_layouts[6])


def _render(params: dict[str, str], **ctx_overrides):
    resolved = PrimitivesProvider().resolve(
        AssetRequest(
            ref=AssetRef("primitives", "primitive", "process-arrow"),
            params=tuple(params.items()),
        )
    )
    assert isinstance(resolved.payload, NativeShapePayload)
    slide = _blank_slide()
    get_native_renderer("process-arrow")(
        slide, dict(resolved.payload.params), _ctx(**ctx_overrides)
    )
    return slide


def _runs(sp):
    return [r.text for p in sp.text_frame.paragraphs for r in p.runs]


def _texts(slide) -> list[str]:
    return [t for sp in slide.shapes if hasattr(sp, "text_frame") for t in _runs(sp)]


def _by_type(slide, shape_type):
    return [
        sp
        for sp in slide.shapes
        if sp.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and sp.auto_shape_type == shape_type
    ]


def _assert_within_rect(slide, rect: AssetRect, tol_px: float = 1.0) -> None:
    tol = round(tol_px * 6350)
    x0, y0 = int(rect.x * 6350), int(rect.y * 6350)
    x1, y1 = int((rect.x + rect.width) * 6350), int((rect.y + rect.height) * 6350)
    for sp in slide.shapes:
        assert sp.left >= x0 - tol, f"left overflow: {sp.left}"
        assert sp.top >= y0 - tol, f"top overflow: {sp.top}"
        assert sp.left + sp.width <= x1 + tol, f"right overflow: {sp.left + sp.width}"
        assert sp.top + sp.height <= y1 + tol, f"bottom overflow: {sp.top + sp.height}"


def test_process_arrow_horizontal_two_steps() -> None:
    slide = _render({"steps": "a,b"})
    assert len(_by_type(slide, MSO_SHAPE.CHEVRON)) == 2


def test_process_arrow_horizontal_eight_steps() -> None:
    slide = _render({"steps": "a,b,c,d,e,f,g,h"})
    assert len(_by_type(slide, MSO_SHAPE.CHEVRON)) == 8


def test_process_arrow_vertical_two_steps() -> None:
    slide = _render({"steps": "a,b", "direction": "vertical"})
    assert len(_by_type(slide, MSO_SHAPE.ROUNDED_RECTANGLE)) == 2
    assert len(_by_type(slide, MSO_SHAPE.DOWN_ARROW)) == 1


def test_process_arrow_vertical_eight_steps() -> None:
    slide = _render({"steps": "a,b,c,d,e,f,g,h", "direction": "vertical"})
    assert len(_by_type(slide, MSO_SHAPE.ROUNDED_RECTANGLE)) == 8
    assert len(_by_type(slide, MSO_SHAPE.DOWN_ARROW)) == 7


def test_process_arrow_duplicate_labels_ok() -> None:
    slide = _render({"steps": "a,a"})
    assert _texts(slide).count("a") == 2


def test_process_arrow_step_text_editable() -> None:
    slide = _render({"steps": "Plan,Build"})
    for step in ("Plan", "Build"):
        assert step in _texts(slide)


def test_process_arrow_horizontal_within_rect() -> None:
    slide = _render({"steps": "a,b,c,d"})
    _assert_within_rect(slide, _ctx().rect)


def test_process_arrow_vertical_within_rect() -> None:
    slide = _render({"steps": "a,b,c,d", "direction": "vertical"})
    _assert_within_rect(slide, _ctx().rect)


def test_process_arrow_no_pictures() -> None:
    slide = _render({"steps": "a,b,c,d"})
    assert not any(sp.shape_type == MSO_SHAPE_TYPE.PICTURE for sp in slide.shapes)
