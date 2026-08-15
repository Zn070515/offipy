"""audit 提取：递归 shape 提取（group 嵌套）、隐藏、连接线、文本/autofit 读取。"""

import builtins
import zipfile

import pytest
from lxml import etree
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

from offipy.audit.extract import (
    _GroupTransform,
    _PresentationExtract,
    _TextRun,
    extract_presentation,
)


def _build_pptx() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    # 文本框：多 run + wrap/autofit
    tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
    tf = tb.text_frame
    tf.text = "hello"
    r = tf.paragraphs[0].add_run()
    r.text = " world"
    r.font.size = Pt(18)
    r.font.bold = True
    r.font.name = "Arial"
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE  # normAutofit
    # 自动形状：有文字，无 run 级字号（继承）
    sh = slide.shapes.add_shape(1, Inches(5), Inches(5), Inches(2), Inches(1))
    sh.text = "shape"
    # 连接线
    slide.shapes.add_connector(1, Inches(0), Inches(5), Inches(2), Inches(6))
    # 表格
    slide.shapes.add_table(2, 2, Inches(6), Inches(5), Inches(3), Inches(1))
    return prs


def _make_extract(tmp_path, prs) -> _PresentationExtract:
    path = tmp_path / "x.pptx"
    prs.save(path)
    return extract_presentation(path)


_DGM_NS = "http://schemas.openxmlformats.org/drawingml/2006/diagram"
_DSP_NS = "http://schemas.openxmlformats.org/drawingml/2006/diagramData"
_R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
_P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
_A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"


def _build_smartart_pptx(tmp_path):
    """zip-inject 一个含 SmartArt 的 pptx，返回路径。

    空白 pptx → 重写 slide1.xml 注入 diagram graphicFrame → 补 diagramData /
    quickStyle part + slide rels + [Content_Types] override（python-pptx 无
    diagram part 写入 API，故走 zip 重写）。
    """
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[6])
    out = tmp_path / "smartart.pptx"
    prs.save(out)

    gf_xml = (
        f'<p:graphicFrame xmlns:p="{_P_NS}" xmlns:a="{_A_NS}" xmlns:r="{_R_NS}">'
        '<p:nvGraphicFramePr><p:cNvPr id="100" name="SmartArt 图形"/>'
        "<p:cNvGraphicFramePr/><p:nvPr/></p:nvGraphicFramePr>"
        '<p:xfrm><a:off x="914400" y="914400"/><a:ext cx="3657600" cy="1828800"/></p:xfrm>'
        f'<a:graphic><a:graphicData uri="{_DGM_NS}">'
        f'<dgm:relIds xmlns:dgm="{_DGM_NS}" r:dm="rId5" r:lo="rId5" r:qs="rId6" r:cs="rId5"/>'
        "</a:graphicData></a:graphic></p:graphicFrame>"
    )
    data_xml = (
        f'<dgm:dataModel xmlns:dgm="{_DSP_NS}" xmlns:a="{_A_NS}"><dgm:ptLst>'
        '<dgm:pt modelId="0"><dgm:prSet/><dgm:spPr/></dgm:pt>'
        '<dgm:pt modelId="1"><dgm:prSet/><dgm:spPr/><dgm:t><dgm:txBody>'
        "<a:p><a:r><a:t>节点一</a:t></a:r></a:p></dgm:txBody></dgm:t></dgm:pt>"
        '<dgm:pt modelId="2"><dgm:prSet/><dgm:spPr/><dgm:t><dgm:txBody>'
        "<a:p><a:r><a:t>节点二</a:t></a:r></a:p></dgm:txBody></dgm:t></dgm:pt>"
        "</dgm:ptLst><dgm:cxnLst/></dgm:dataModel>"
    )
    qs_xml = f'<dgm:quickStyle xmlns:dgm="{_DGM_NS}"/>'

    with zipfile.ZipFile(out, "r") as zin:
        src = {n: zin.read(n) for n in zin.namelist()}

    slide_el = etree.fromstring(src["ppt/slides/slide1.xml"])
    sp_tree = slide_el.find(f"{{{_P_NS}}}cSld/{{{_P_NS}}}spTree")
    sp_tree.append(etree.fromstring(gf_xml.encode("utf-8")))
    src["ppt/slides/slide1.xml"] = etree.tostring(
        slide_el, xml_declaration=True, encoding="UTF-8", standalone=True
    )

    rels = src["ppt/slides/_rels/slide1.xml.rels"].decode("utf-8")
    rels = rels.replace(
        "</Relationships>",
        f'<Relationship Id="rId5" Type="{_R_NS}/diagramData" Target="../diagrams/data1.xml"/>'
        f'<Relationship Id="rId6" Type="{_R_NS}/diagramQuickStyle" '
        f'Target="../diagrams/quickStyle1.xml"/>'
        "</Relationships>",
    )
    src["ppt/slides/_rels/slide1.xml.rels"] = rels.encode("utf-8")

    src["ppt/diagrams/data1.xml"] = data_xml.encode("utf-8")
    src["ppt/diagrams/quickStyle1.xml"] = qs_xml.encode("utf-8")

    ct = src["[Content_Types].xml"].decode("utf-8")
    ct = ct.replace(
        "</Types>",
        '<Override PartName="/ppt/diagrams/data1.xml" '
        'ContentType="application/vnd.openxmlformats-officedocument.drawingml.diagramData+xml"/>'
        '<Override PartName="/ppt/diagrams/quickStyle1.xml" '
        'ContentType="application/vnd.openxmlformats-officedocument.drawingml.diagramStyle+xml"/>'
        "</Types>",
    )
    src["[Content_Types].xml"] = ct.encode("utf-8")

    with zipfile.ZipFile(out, "w", zipfile.ZIP_DEFLATED) as zout:
        for name, blob in src.items():
            zout.writestr(name, blob)
    return str(out)


def test_smartart_text_and_node_count(tmp_path):
    ext = extract_presentation(_build_smartart_pptx(tmp_path))
    rec = next(s for s in ext.slides[0].shapes if s.smartart_node_count is not None)
    assert rec.smartart_node_count == 2  # 含文本节点（modelId 1/2），根节点 modelId 0 不计
    assert "节点一" in rec.text and "节点二" in rec.text
    assert rec.shape_type == "UNKNOWN"  # SmartArt graphicFrame → python-pptx 不识别类型


def test_table_cell_text_extracted(tmp_path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tbl = slide.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(3), Inches(1))
    tbl.table.cell(0, 0).text = "单元格A"
    tbl.table.cell(1, 1).text = "单元格B"
    ext = _make_extract(tmp_path, prs)
    rec = next(s for s in ext.slides[0].shapes if s.has_table)
    assert "单元格A" in rec.text and "单元格B" in rec.text


def test_chart_text_extracted(tmp_path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    cd = CategoryChartData()
    cd.categories = ["甲", "乙"]
    cd.add_series("系列", (1, 2))
    slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1), Inches(1), Inches(3), Inches(2), cd
    )
    ext = _make_extract(tmp_path, prs)
    rec = next(s for s in ext.slides[0].shapes if s.shape_type == "CHART")
    assert rec.has_chart is True
    assert "系列" in rec.text and "甲" in rec.text and "乙" in rec.text


# ---------------------------------------------------------------- 顶层 shape


def test_extract_textbox_text_and_runs(tmp_path):
    ext = _make_extract(tmp_path, _build_pptx())
    assert ext.slide_size[0] == pytest.approx(13.333, abs=1e-3)
    assert ext.slide_size[1] == pytest.approx(7.5)
    assert len(ext.slides) == 1
    assert ext.slides[0].slide_index == 1

    rec = next(s for s in ext.slides[0].shapes if s.shape_type == "TEXT_BOX")
    assert rec.text == "hello world"
    assert rec.has_text_frame is True
    assert rec.word_wrap is True
    assert rec.autofit_mode == "TEXT_TO_FIT_SHAPE"
    assert rec.rotation == 0.0
    assert rec.left == pytest.approx(1.0)
    assert rec.width == pytest.approx(3.0)
    # run 级字号/加粗/字体
    run: _TextRun = rec.paragraphs[0].runs[1]
    assert run.text == " world"
    assert run.font_size == pytest.approx(18.0)
    assert run.bold is True
    assert run.font_name == "Arial"
    # TextFrame 默认 margin 0.1 英寸
    assert rec.tf_margin_left == pytest.approx(0.1)


def test_extract_autoshape_inherited_font_size(tmp_path):
    ext = _make_extract(tmp_path, _build_pptx())
    rec = next(s for s in ext.slides[0].shapes if s.text == "shape")
    assert rec.text == "shape"
    assert rec.has_text_frame is True
    assert rec.paragraphs[0].runs[0].font_size is None  # 继承


def test_paragraph_segments_split_at_soft_break(tmp_path):
    # a:br 软换行把段落切成视觉行分组；runs 不含 a:br，段切分由元素序遍历还原
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
    tf = tb.text_frame
    tf.text = "seg one"
    p = tf.paragraphs[0]
    p.add_line_break()
    r = p.add_run()
    r.text = "seg two"
    r.font.size = Pt(14)
    p.add_line_break()
    p.add_run().text = "seg three"
    ext = _make_extract(tmp_path, prs)
    rec = next(s for s in ext.slides[0].shapes if s.shape_type == "TEXT_BOX")
    para = rec.paragraphs[0]
    assert para.text == "seg one\vseg two\vseg three"
    assert len(para.segments) == 3
    assert [run.text for run in para.segments[0]] == ["seg one"]
    assert [run.text for run in para.segments[1]] == ["seg two"]
    assert para.segments[1][0].font_size == pytest.approx(14.0)
    assert [run.text for run in para.segments[2]] == ["seg three"]


def test_paragraph_no_break_single_segment(tmp_path):
    # 无 a:br → 整段一组
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
    tf = tb.text_frame
    tf.text = "plain"
    tf.paragraphs[0].add_run().text = " extra"
    ext = _make_extract(tmp_path, prs)
    rec = next(s for s in ext.slides[0].shapes if s.shape_type == "TEXT_BOX")
    para = rec.paragraphs[0]
    assert len(para.segments) == 1
    assert [run.text for run in para.segments[0]] == ["plain", " extra"]


def test_paragraph_line_spacing_extracted(tmp_path):
    # a:lnSpc 解析：spcPts 绝对点值 / spcPct 百分比，未设时为 None
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
    tf = tb.text_frame
    p1 = tf.paragraphs[0]
    p1.add_run().text = "line one"
    p1.line_spacing = Pt(14.85)  # spcPts val=1485 → 14.85pt
    p2 = tf.add_paragraph()
    p2.add_run().text = "line two"
    p2.line_spacing = 1.5  # spcPct val=150000 → 150%
    p3 = tf.add_paragraph()
    p3.add_run().text = "line three"  # 未设 → None
    ext = _make_extract(tmp_path, prs)
    rec = next(s for s in ext.slides[0].shapes if s.shape_type == "TEXT_BOX")
    p1r, p2r, p3r = rec.paragraphs[0], rec.paragraphs[1], rec.paragraphs[2]
    assert p1r.line_spacing_pts == pytest.approx(14.85)
    assert p1r.line_spacing_pct is None
    assert p2r.line_spacing_pts is None
    assert p2r.line_spacing_pct == pytest.approx(150.0)
    assert p3r.line_spacing_pts is None
    assert p3r.line_spacing_pct is None


# ---------------------------------------------------------------- 损坏 XML 容错


def test_corrupt_font_scale_graceful(tmp_path):
    # 原始 XML normAutofit fontScale 非数字 → autofit_font_scale None，不崩 ValueError
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
    tf = tb.text_frame
    tf.text = "hi"
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf._txBody.xpath("./a:bodyPr/a:normAutofit")[0].set("fontScale", "abc")
    ext = _make_extract(tmp_path, prs)
    rec = next(s for s in ext.slides[0].shapes if s.shape_type == "TEXT_BOX")
    assert rec.autofit_font_scale is None


def test_corrupt_line_spacing_graceful(tmp_path):
    # 原始 XML spcPts val 非数字 → line_spacing 回退 None，不崩
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.add_run().text = "line"
    p.line_spacing = Pt(14.85)
    lnSpc = p._p.xpath("./a:pPr/a:lnSpc")[0]
    spcPts = lnSpc.find("{http://schemas.openxmlformats.org/drawingml/2006/main}spcPts")
    spcPts.set("val", "abc")
    ext = _make_extract(tmp_path, prs)
    rec = next(s for s in ext.slides[0].shapes if s.shape_type == "TEXT_BOX")
    assert rec.paragraphs[0].line_spacing_pts is None
    assert rec.paragraphs[0].line_spacing_pct is None


def test_corrupt_group_transform_rot_graceful(tmp_path):
    # 原始 XML group xfrm rot 非数字 → python-pptx 属性解析 ValueError；
    # 该 group 跳过 + 告警，整个提取不崩（兄弟姐妹不受影响）
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    g = slide.shapes.add_group_shape()
    c = g.shapes.add_shape(1, Inches(0), Inches(0), Inches(1), Inches(1))
    c.text = "child"
    g._element.xpath(".//a:xfrm")[0].set("rot", "abc")
    ext = _make_extract(tmp_path, prs)
    grecs = [s for s in ext.slides[0].shapes if s.is_group]
    assert grecs == []  # 损坏 group 跳过
    assert any(w.code == "audit.extract.shape_skip_corrupt" for w in ext.warnings)


def test_corrupt_slide_size_graceful(tmp_path):
    # #70：b47421e 只兜 per-shape 损坏；演示文稿级 sldSz@cx/@cy 非数字在顶层
    # _to_inches(prs.slide_width) 抛 ValueError 整文件 audit 崩。降级 0.0 + 告警。
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[6])
    sld_sz = prs.part.presentation._element.find(qn("p:sldSz"))
    sld_sz.set("cx", "abc")
    sld_sz.set("cy", "def")
    path = tmp_path / "x.pptx"
    prs.save(path)
    ext = extract_presentation(path)
    assert ext.slide_size == (0.0, 0.0)
    assert any(w.code == "audit.extract.slidesize_corrupt" for w in ext.warnings)


# ---------------------------------------------------------------- group 嵌套


def test_extract_group_nesting_absolutized(tmp_path):
    # group off=0 且 ext==chExt（恒等缩放）时，子元素绝对坐标 = 局部坐标
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    g = slide.shapes.add_group_shape()
    c1 = g.shapes.add_shape(1, Inches(0.5), Inches(0.25), Inches(1), Inches(1))
    c1.text = "child"
    g2 = g.shapes.add_group_shape()  # 嵌套 group
    c2 = g2.shapes.add_textbox(Inches(0), Inches(0), Inches(1), Inches(0.5))
    c2.text = "deep"
    ext = _make_extract(tmp_path, prs)

    group_rec = next(s for s in ext.slides[0].shapes if s.is_group and s.shape_id == g.shape_id)
    assert group_rec.shape_type == "GROUP"
    assert group_rec.parent_shape_id is None
    assert group_rec.group_path == ()
    assert isinstance(group_rec.transform, _GroupTransform)
    assert group_rec.transform.ch_ext_cx == pytest.approx(1.5)  # 坐标空间随子元素扩展

    c1_rec = next(s for s in ext.slides[0].shapes if s.shape_id == c1.shape_id)
    assert c1_rec.parent_shape_id == g.shape_id
    assert c1_rec.group_path == (g.shape_id,)
    # 已绝对化到幻灯片坐标（本用例 ext==chExt，故 = 局部坐标）
    assert c1_rec.left == pytest.approx(0.5)
    assert c1_rec.top == pytest.approx(0.25)
    assert c1_rec.is_rotated is False
    assert c1_rec.geometry_unknown is False

    g2_rec = next(s for s in ext.slides[0].shapes if s.shape_id == g2.shape_id)
    assert g2_rec.group_path == (g.shape_id,)
    c2_rec = next(s for s in ext.slides[0].shapes if s.shape_id == c2.shape_id)
    assert c2_rec.group_path == (g.shape_id, g2.shape_id)
    assert c2_rec.text == "deep"


def test_group_scale_absolutizes_child(tmp_path):
    # group off=(1,0) ext=(2,2)，子 1×1 局部 (0,0) → 绝对 (1,0) 2×2
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    g = slide.shapes.add_group_shape()
    c = g.shapes.add_shape(1, Inches(0), Inches(0), Inches(1), Inches(1))
    c.text = "child"
    g.left, g.top = Inches(1), Inches(0)
    g.width, g.height = Inches(2), Inches(2)
    ext = _make_extract(tmp_path, prs)
    rec = next(s for s in ext.slides[0].shapes if s.shape_id == c.shape_id)
    assert rec.left == pytest.approx(1.0)
    assert rec.top == pytest.approx(0.0)
    assert rec.width == pytest.approx(2.0)
    assert rec.height == pytest.approx(2.0)


def test_group_fliph_mirrors_child_about_center(tmp_path):
    # group off=(2,3) ext=(1,1)，chExt=(1,1)，子左半 (0,0)-(0.5,1) 占坐标空间左半
    # flipH 绕组中心 x=2.5 镜像 → 子移入右半 [2.5,3]
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    g = slide.shapes.add_group_shape()
    c = g.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.5), Inches(1))
    c.text = "child"
    g.left, g.top = Inches(2), Inches(3)
    g.width, g.height = Inches(1), Inches(1)
    xfrm = g._element.xpath(".//a:xfrm")[0]
    xfrm.set("flipH", "1")
    # 强制 chExt=(1,1)：子只占坐标空间左半，flipH 才有可观察位移
    ch_ext = xfrm.find("{http://schemas.openxmlformats.org/drawingml/2006/main}chExt")
    ch_ext.set("cx", str(int(1.0 * 914400)))
    ch_ext.set("cy", str(int(1.0 * 914400)))
    ext = _make_extract(tmp_path, prs)
    rec = next(s for s in ext.slides[0].shapes if s.shape_id == c.shape_id)
    assert rec.left == pytest.approx(2.5)
    assert rec.top == pytest.approx(3.0)
    assert rec.width == pytest.approx(0.5)
    assert rec.height == pytest.approx(1.0)


def test_group_children_z_order_scoped(tmp_path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    g = slide.shapes.add_group_shape()
    a = g.shapes.add_textbox(Inches(0), Inches(0), Inches(1), Inches(1))
    b = g.shapes.add_textbox(Inches(0), Inches(0), Inches(1), Inches(1))
    ext = _make_extract(tmp_path, prs)
    ra = next(s for s in ext.slides[0].shapes if s.shape_id == a.shape_id)
    rb = next(s for s in ext.slides[0].shapes if s.shape_id == b.shape_id)
    assert ra.z_order == 0
    assert rb.z_order == 1


# ---------------------------------------------------------------- 隐藏


def test_hidden_shape_detection(tmp_path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    h = slide.shapes.add_shape(1, Inches(1), Inches(1), Inches(1), Inches(1))
    h._element.xpath("./p:nvSpPr/p:cNvPr")[0].set("hidden", "1")
    v = slide.shapes.add_shape(1, Inches(2), Inches(2), Inches(1), Inches(1))
    ext = _make_extract(tmp_path, prs)
    hr = next(s for s in ext.slides[0].shapes if s.shape_id == h.shape_id)
    vr = next(s for s in ext.slides[0].shapes if s.shape_id == v.shape_id)
    assert hr.is_hidden is True
    assert vr.is_hidden is False


# ---------------------------------------------------------------- 连接线 / 表


def test_connector_detection(tmp_path):
    ext = _make_extract(tmp_path, _build_pptx())
    conn = next(s for s in ext.slides[0].shapes if s.is_connector)
    assert conn.shape_type == "LINE"
    assert conn.is_connector is True
    assert conn.is_group is False


def test_connector_type_enum_line():
    assert MSO_SHAPE_TYPE.LINE is not None  # 探针保证判定入口存在


def test_table_detection(tmp_path):
    ext = _make_extract(tmp_path, _build_pptx())
    tbl = next(s for s in ext.slides[0].shapes if s.shape_type == "TABLE")
    assert tbl.has_table is True
    assert tbl.has_text_frame is False
    assert tbl.text == ""


# ---------------------------------------------------------------- 填充


def test_fill_kind_extracted(tmp_path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    solid = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(1), Inches(1))
    solid.fill.solid()  # 显式 a:solidFill
    hollow = slide.shapes.add_shape(1, Inches(1), Inches(0), Inches(1), Inches(1))
    hollow.fill.background()  # a:noFill → 透明
    grad = slide.shapes.add_shape(1, Inches(2), Inches(0), Inches(1), Inches(1))
    grad.fill.gradient()  # a:gradFill
    ext = _make_extract(tmp_path, prs)
    by_id = {r.shape_id: r for r in ext.slides[0].shapes}
    assert by_id[solid.shape_id].fill_kind == "solid"
    assert by_id[hollow.shape_id].fill_kind == "none"
    assert by_id[grad.shape_id].fill_kind == "gradient"
    # add_shape 默认无显式填充（主题继承）→ unknown（按不透明处理，防漏报真实遮挡）
    default = slide.shapes.add_shape(1, Inches(3), Inches(0), Inches(1), Inches(1))
    ext2 = _make_extract(tmp_path, prs)
    rec = next(r for r in ext2.slides[0].shapes if r.shape_id == default.shape_id)
    assert rec.fill_kind == "unknown"
    # 文本框默认显式 a:noFill（透明）→ none
    tb = slide.shapes.add_textbox(Inches(4), Inches(0), Inches(1), Inches(1))
    tb.text = "t"
    ext3 = _make_extract(tmp_path, prs)
    rec = next(r for r in ext3.slides[0].shapes if r.shape_id == tb.shape_id)
    assert rec.fill_kind == "none"


# ---------------------------------------------------------------- 占位符


def test_title_placeholder_type(tmp_path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title Slide
    title = slide.shapes.title
    title.text = "T"
    ext = _make_extract(tmp_path, prs)
    rec = next(s for s in ext.slides[0].shapes if s.shape_id == title.shape_id)
    assert rec.placeholder_type is not None
    assert rec.has_text_frame is True


# ---------------------------------------------------------------- 门禁


def test_import_extract_does_not_load_pptx(monkeypatch):
    orig_import = builtins.__import__

    def guarded(name, *args, **kwargs):
        if name == "pptx" or name.startswith("pptx."):
            raise AssertionError("import offipy.audit.extract 不应加载 python-pptx")
        return orig_import(name, *args, **kwargs)

    monkeypatch.setattr(builtins, "__import__", guarded)
    import offipy.audit.extract as m

    assert m is not None


def test_lazy_pptx_only_on_read(tmp_path):
    # extract_presentation 内部才加载 pptx；此处直接调它应成功
    ext = _make_extract(tmp_path, _build_pptx())
    assert len(ext.slides[0].shapes) >= 4
