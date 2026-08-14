"""A5 Task 3 — quote-mark renderer structural contract."""

from __future__ import annotations

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from offipy.assets.model import (
    AssetRect,
    AssetRef,
    AssetRenderContext,
    AssetRequest,
    NativeShapePayload,
)
from offipy.assets.primitives import get_native_renderer
from offipy.assets.providers.primitives import PrimitivesProvider

_THEME = {
    "accent": "#2251ff",
    "surface": "#F2F4F7",
    "ink": "#222222",
    "muted": "#667085",
    "bg": "#ffffff",
}

_GLYPH = "“"


def _ctx(**overrides) -> AssetRenderContext:
    kwargs = {
        "slide_index": 1,
        "rect": AssetRect(x=0, y=0, width=400, height=200),
        "theme_name": "mckinsey",
        "theme_vars": dict(_THEME),
        "placement": "replace",
    }
    kwargs.update(overrides)
    return AssetRenderContext(**kwargs)


def _blank_slide():
    prs = Presentation()
    return prs.slides.add_slide(prs.slide_layouts[6])


def _render(params: dict[str, str], **ctx_overrides):
    """Resolve through the provider (defaults filled), then render directly."""
    resolved = PrimitivesProvider().resolve(
        AssetRequest(
            ref=AssetRef("primitives", "primitive", "quote-mark"),
            params=tuple(params.items()),
        )
    )
    assert isinstance(resolved.payload, NativeShapePayload)
    slide = _blank_slide()
    get_native_renderer("quote-mark")(slide, dict(resolved.payload.params), _ctx(**ctx_overrides))
    return slide


def _runs(sp):
    return [r.text for p in sp.text_frame.paragraphs for r in p.runs]


def _texts(slide) -> list[str]:
    return [t for sp in slide.shapes if hasattr(sp, "text_frame") for t in _runs(sp)]


def _find_textbox(slide, text):
    for sp in slide.shapes:
        if hasattr(sp, "text_frame") and text in _runs(sp):
            return sp
    return None


def _autoshapes(slide):
    return [sp for sp in slide.shapes if sp.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE]


def _assert_within_rect(slide, rect: AssetRect, tol_px: float = 1.0) -> None:
    tol = round(tol_px * 6350)
    x0, y0 = int(rect.x * 6350), int(rect.y * 6350)
    x1, y1 = int((rect.x + rect.width) * 6350), int((rect.y + rect.height) * 6350)
    for sp in slide.shapes:
        assert sp.left >= x0 - tol, f"left overflow: {sp.left}"
        assert sp.top >= y0 - tol, f"top overflow: {sp.top}"
        assert sp.left + sp.width <= x1 + tol, f"right overflow: {sp.left + sp.width}"
        assert sp.top + sp.height <= y1 + tol, f"bottom overflow: {sp.top + sp.height}"


def test_quote_mark_text_round_trip() -> None:
    slide = _render({"text": "Hello world"})
    assert "Hello world" in _texts(slide)


def test_quote_mark_at_least_two_objects() -> None:
    slide = _render({"text": "Hi"})
    assert len(list(slide.shapes)) >= 2


def test_quote_mark_shapes_within_rect() -> None:
    slide = _render({"text": "A longer quote that wraps across lines."})
    _assert_within_rect(slide, _ctx().rect)


def test_quote_mark_default_theme_accent() -> None:
    slide = _render({"text": "Hi"})
    glyph = _find_textbox(slide, _GLYPH)
    assert glyph is not None
    assert str(glyph.text_frame.paragraphs[0].runs[0].font.color.rgb) == "2251FF"


def test_quote_mark_accent_override() -> None:
    slide = _render({"text": "Hi", "accent": "#FF0000"})
    glyph = _find_textbox(slide, _GLYPH)
    assert glyph is not None
    assert str(glyph.text_frame.paragraphs[0].runs[0].font.color.rgb) == "FF0000"


def test_quote_mark_transparent_fill_has_no_card() -> None:
    slide = _render({"text": "Hi"})  # default fill is transparent
    assert _autoshapes(slide) == []


def test_quote_mark_fill_creates_card() -> None:
    slide = _render({"text": "Hi", "fill": "surface"})
    cards = _autoshapes(slide)
    assert len(cards) == 1
    assert str(cards[0].fill.fore_color.rgb) == "F2F4F7"


def test_quote_mark_no_pictures() -> None:
    slide = _render({"text": "Hi"})
    assert not any(sp.shape_type == MSO_SHAPE_TYPE.PICTURE for sp in slide.shapes)
