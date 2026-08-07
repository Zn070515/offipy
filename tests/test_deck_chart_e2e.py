"""office-real 原生图表注入端到端测试（S3 Task 4）。

锁定完整管线：chart-dominant 布局 + 内置主题注入 → chromium 渲染 → convert →
postprocess_charts 把占位矩形替换成原生可编辑图表。验证用 python-pptx 读回产物
（不涉及 Office/COM）：图表形状存在、类型匹配、类别/序列数据与源一致、
序列/点颜色匹配主题 accent 或显式 data-chart-colors 覆盖。

marker 与 test_deck_e2e.py 一致：读回用 python-pptx，唯一环境依赖是 chromium
渲染管线（vendored 转换器 + 存活 server 8890）。无 Chromium / server 不可达 →
整模块 skip。
"""

import pytest

from offipy.client import _ping
from offipy.deck import CONVERT_PY

pytestmark = [
    pytest.mark.deck_render,
    pytest.mark.skipif(
        not CONVERT_PY.exists() or not _ping(),
        reason="需要 vendored 转换器 + 存活 server(8890)",
    ),
]

# 每个 chart-dominant 页都必须同时带 CSS class `chart-dominant` 与
# `data-layout="chart-dominant"`（布局 CSS 选择器是 `.chart-dominant .chart`）。
# <head> 放两个注入占位：主题 + 布局。
_BAR_PIE_HTML = """<!DOCTYPE html>
<html>
<head>
  <style data-theme="mckinsey"></style>
  <style data-layouts></style>
</head>
<body>
  <section class="slide chart-dominant" data-pptx-slide data-layout="chart-dominant">
    <div class="chart" data-chart="bar"
         data-chart-data='{"categories":["Q1","Q2","Q3","Q4"],"series":[{"name":"营收","values":[40,55,70,92]}]}'></div>
  </section>
  <section class="slide chart-dominant" data-pptx-slide data-layout="chart-dominant">
    <div class="chart" data-chart="pie"
         data-chart-data='{"categories":["直销","渠道","线上"],"series":[{"name":"营收占比","values":[45,35,20]}]}'></div>
  </section>
</body>
</html>"""

_VARIANT_HTML = """<!DOCTYPE html>
<html>
<head>
  <style data-theme="mckinsey"></style>
  <style data-layouts></style>
</head>
<body>
  <section class="slide chart-dominant dark" data-pptx-slide data-layout="chart-dominant">
    <div class="chart" data-chart="bar"
         data-chart-data='{"categories":["Q1","Q2"],"series":[{"name":"营收","values":[10,20]}]}'></div>
  </section>
  <section class="slide chart-dominant" data-pptx-slide data-layout="chart-dominant">
    <div class="chart" data-chart="bar"
         data-chart-data='{"categories":["Q1","Q2"],"series":[{"name":"营收","values":[10,20]}]}'></div>
  </section>
</body>
</html>"""

_COLORS_OVERRIDE_HTML = """<!DOCTYPE html>
<html>
<head>
  <style data-theme="mckinsey"></style>
  <style data-layouts></style>
</head>
<body>
  <section class="slide chart-dominant" data-pptx-slide data-layout="chart-dominant">
    <div class="chart" data-chart="bar"
         data-chart-colors='["#AABBCC","#DDEEFF"]'
         data-chart-data='{"categories":["Q1","Q2"],"series":[{"name":"A","values":[10,20]},{"name":"B","values":[30,40]}]}'></div>
  </section>
</body>
</html>"""

_LINE_HTML = """<!DOCTYPE html>
<html>
<head>
  <style data-theme="mckinsey"></style>
  <style data-layouts></style>
</head>
<body>
  <section class="slide chart-dominant" data-pptx-slide data-layout="chart-dominant">
    <div class="chart" data-chart="line"
         data-chart-data='{"categories":["J","F","M"],"series":[{"name":"trend","values":[1,3,2]}]}'></div>
  </section>
</body>
</html>"""

_CHART_DOM_HTML = """<!DOCTYPE html>
<html>
<head>
  <style data-theme="mckinsey"></style>
  <style data-layouts></style>
</head>
<body>
  <section class="slide chart-dominant" data-pptx-slide data-layout="chart-dominant">
    <div class="chart" data-chart="bar"
         data-chart-data='{"categories":["Q1","Q2"],"series":[{"name":"营收","values":[10,20]}]}'></div>
  </section>
</body>
</html>"""


def _render_deck(tmp_path, html_text: str, **kwargs) -> str:
    """写 HTML → 跑完整渲染管线 → 返回产出 .pptx 的绝对路径。"""
    from offipy.deck import render

    html = tmp_path / "deck.html"
    html.write_text(html_text, encoding="utf-8")
    out = tmp_path / "deck.pptx"
    return render(str(html), out=str(out), overwrite=True, **kwargs)


def test_e2e_bar_and_pie_chart_dominant_theme(tmp_path):
    """bar + pie 两页：图表存在、类型匹配、类别/数据一致、首色=主题 accent。"""
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.chart import XL_CHART_TYPE

    out = _render_deck(tmp_path, _BAR_PIE_HTML, theme="mckinsey", apply_layouts=True)

    prs = Presentation(out)
    assert len(prs.slides) == 2

    # ---- bar 页 ----
    bar_slide = prs.slides[0]
    bar_shapes = [s for s in bar_slide.shapes if s.has_chart]
    assert len(bar_shapes) == 1
    bar_chart = bar_shapes[0].chart
    assert bar_chart.chart_type == XL_CHART_TYPE.COLUMN_CLUSTERED
    assert list(bar_chart.plots[0].categories) == ["Q1", "Q2", "Q3", "Q4"]
    bar_series = bar_chart.plots[0].series[0]
    assert bar_series.name == "营收"
    assert list(bar_series.values) == [40.0, 55.0, 70.0, 92.0]
    assert bar_series.format.fill.fore_color.rgb == RGBColor(0x22, 0x51, 0xFF)

    # ---- pie 页 ----
    pie_slide = prs.slides[1]
    pie_shapes = [s for s in pie_slide.shapes if s.has_chart]
    assert len(pie_shapes) == 1
    pie_chart = pie_shapes[0].chart
    assert pie_chart.chart_type == XL_CHART_TYPE.PIE
    assert list(pie_chart.plots[0].categories) == ["直销", "渠道", "线上"]
    pie_series = pie_chart.plots[0].series[0]
    assert pie_series.name == "营收占比"
    assert list(pie_series.values) == [45.0, 35.0, 20.0]
    assert pie_series.points[0].format.fill.fore_color.rgb == RGBColor(0x22, 0x51, 0xFF)


def test_e2e_variant_dark_accent_is_page_specific(tmp_path):
    """同一 render 内：dark 变体页首色 #5B8CFF，普通页首色 #2251FF。"""
    from pptx import Presentation
    from pptx.dml.color import RGBColor

    out = _render_deck(tmp_path, _VARIANT_HTML, theme="mckinsey", apply_layouts=True)

    prs = Presentation(out)
    assert len(prs.slides) == 2

    def first_series_fill(slide_index: int):
        charts = [s for s in prs.slides[slide_index].shapes if s.has_chart]
        assert len(charts) == 1
        return charts[0].chart.plots[0].series[0].format.fill.fore_color.rgb

    assert first_series_fill(0) == RGBColor(0x5B, 0x8C, 0xFF)
    assert first_series_fill(1) == RGBColor(0x22, 0x51, 0xFF)


def test_e2e_explicit_colors_override_wins(tmp_path):
    """data-chart-colors 显式覆盖优先于主题 accent。"""
    from pptx import Presentation
    from pptx.dml.color import RGBColor

    out = _render_deck(tmp_path, _COLORS_OVERRIDE_HTML, theme="mckinsey", apply_layouts=True)

    prs = Presentation(out)
    charts = [s for s in prs.slides[0].shapes if s.has_chart]
    assert len(charts) == 1
    series = charts[0].chart.plots[0].series
    assert series[0].format.fill.fore_color.rgb == RGBColor(0xAA, 0xBB, 0xCC)
    assert series[1].format.fill.fore_color.rgb == RGBColor(0xDD, 0xEE, 0xFF)


def test_e2e_line_chart(tmp_path):
    """line 图表：类型 LINE_MARKERS，首序列 line.color 用主题 accent。"""
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.chart import XL_CHART_TYPE

    out = _render_deck(tmp_path, _LINE_HTML, theme="mckinsey", apply_layouts=True)

    prs = Presentation(out)
    charts = [s for s in prs.slides[0].shapes if s.has_chart]
    assert len(charts) == 1
    chart = charts[0].chart
    assert chart.chart_type == XL_CHART_TYPE.LINE_MARKERS
    assert list(chart.plots[0].categories) == ["J", "F", "M"]
    series = chart.plots[0].series[0]
    assert series.name == "trend"
    assert list(series.values) == [1.0, 3.0, 2.0]
    assert series.format.line.color.rgb == RGBColor(0x22, 0x51, 0xFF)


def test_e2e_chart_dominant_without_layouts_preflight_fails(tmp_path):
    """chart-dominant + apply_layouts=False → chromium 启动前 fail-fast。"""
    from offipy.deck import render
    from offipy.exceptions import InvalidArgumentError

    html = tmp_path / "deck.html"
    html.write_text(_CHART_DOM_HTML, encoding="utf-8")
    out = tmp_path / "deck.pptx"

    with pytest.raises(InvalidArgumentError):
        render(
            str(html),
            out=str(out),
            overwrite=True,
            theme="mckinsey",
            apply_layouts=False,
        )
