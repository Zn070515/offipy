"""v0.11 发布阻断项行为测试（ChatGPT_v13 审计修复）。

1. 四个 --no-…-ignore 开关真正门控规则：margin 角色豁免、overlap 背景参与、
   text-fit 角色跳过——关闭后同一 finding 恢复为普通 finding；
2. ignored_shapes / ignored_regions 对 overlap 的 primary/secondary 对称；
3. 损坏 PPTX（非 ZIP / 合法 ZIP 但 XML 损坏）→ ConversionError，内部 bug 原样抛。
"""

import zipfile

import pytest
from pptx import Presentation
from pptx.util import Inches

from offipy.audit.extract import extract_presentation
from offipy.audit.models import (
    RULE_MARGIN_BOTTOM,
    RULE_MARGIN_LEFT,
    RULE_MARGIN_TOP,
    AuditConfig,
)
from offipy.audit.rules import run_rules
from offipy.exceptions import ConversionError


def _run(prs, tmp_path, config=None):
    path = tmp_path / "x.pptx"
    prs.save(path)
    ext = extract_presentation(path)
    records = [r for slide in ext.slides for r in slide.shapes]
    findings, suppressed = run_rules(records, ext.slide_size, config or AuditConfig())
    return findings, suppressed, records


def _add_tb(slide, x, y, text, w=2, h=1):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tb.text = text
    return tb


# ---------------------------------------------------------------- 阻断项 1：ignore 开关门控


def test_page_number_margin_default_suppressed_then_recovered(tmp_path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_tb(slide, 9.5, 7.25, "12", w=0.3, h=0.2)  # 底部贴边纯数字 → page_number
    findings, suppressed, records = _run(prs, tmp_path)
    assert records[0].role == "page_number"
    assert findings == []  # 默认：页码 margin 全部进 suppressed
    assert any(s.reason == "page_number" for s in suppressed)

    findings, suppressed, _ = _run(prs, tmp_path, AuditConfig(ignore_page_numbers=False))
    assert any(f.rule_id == RULE_MARGIN_BOTTOM for f in findings)
    assert not any(s.reason == "page_number" for s in suppressed)


def test_background_margin_and_overlap_default_exempt_then_recovered(tmp_path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(7.5))  # 全页背景
    slide.shapes.add_shape(1, Inches(1), Inches(1), Inches(2), Inches(1))  # 背景内矩形（无文本）
    findings, suppressed, records = _run(prs, tmp_path)
    assert any(r.role == "background" for r in records)
    assert findings == []  # 默认：背景 margin 全 suppressed、不参与 overlap
    assert any(s.reason == "full_bleed" for s in suppressed)
    assert not any(s.reason == "intentional_containment" for s in suppressed)

    findings, suppressed, _ = _run(prs, tmp_path, AuditConfig(ignore_full_bleed_shapes=False))
    assert any(f.kind == "margin" for f in findings)
    assert any(f.kind == "overlap" for f in findings)
    assert not any(s.reason == "full_bleed" for s in suppressed)


def test_header_margin_default_suppressed_then_recovered(tmp_path):
    prs = Presentation()
    for _ in range(3):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _add_tb(slide, 4, 0.05, "ACME CORP", w=2, h=0.3)  # 顶部跨页重复 → header
    findings, suppressed, records = _run(prs, tmp_path)
    assert all(r.role == "header" for r in records)
    assert findings == []
    assert any(s.reason == "header_footer" for s in suppressed)

    findings, suppressed, _ = _run(prs, tmp_path, AuditConfig(ignore_headers_footers=False))
    assert any(f.rule_id == RULE_MARGIN_TOP for f in findings)
    assert not any(s.reason == "header_footer" for s in suppressed)


def test_decoration_margin_default_suppressed_then_recovered(tmp_path):
    prs = Presentation()
    for _ in range(3):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_shape(1, Inches(0.05), Inches(5), Inches(1), Inches(1))  # 贴左重复
    findings, suppressed, records = _run(prs, tmp_path)
    assert all(r.role == "decoration" for r in records)
    assert findings == []
    assert any(s.reason == "repeated_decoration" for s in suppressed)

    findings, suppressed, _ = _run(prs, tmp_path, AuditConfig(ignore_repeated_decorations=False))
    assert any(f.rule_id == RULE_MARGIN_LEFT for f in findings)
    assert not any(s.reason == "repeated_decoration" for s in suppressed)


# ---------------------------------------------------------------- 阻断项 4：primary/secondary 对称


def test_ignored_overlap_secondary_shape_suppressed(tmp_path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    a = slide.shapes.add_shape(1, Inches(2), Inches(2), Inches(4), Inches(2))
    b = slide.shapes.add_shape(1, Inches(3), Inches(2), Inches(4), Inches(2))
    # 忽略 primary → 抑制
    findings, suppressed, _ = _run(prs, tmp_path, AuditConfig(ignored_shapes={(1, a.shape_id)}))
    assert findings == []
    assert any(s.reason == "user_shape" for s in suppressed)
    # 忽略同一 overlap 的 secondary → 修复后同样抑制（不再依赖谁恰好是 primary）
    findings, suppressed, _ = _run(prs, tmp_path, AuditConfig(ignored_shapes={(1, b.shape_id)}))
    assert findings == []
    assert any(s.reason == "user_shape" for s in suppressed)


def test_ignored_region_secondary_shape_suppressed(tmp_path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    a = slide.shapes.add_shape(1, Inches(2), Inches(2), Inches(4), Inches(2))
    slide.shapes.add_shape(1, Inches(3), Inches(2), Inches(4), Inches(2))
    # region 只覆盖 secondary B 的中心 (5,3)，不覆盖 primary A 的中心 (4,3)
    findings, suppressed, _ = _run(
        prs, tmp_path, AuditConfig(ignored_regions=[(4.5, 2.5, 1.0, 1.0)])
    )
    assert a is not None
    assert findings == []
    assert any(s.reason == "user_region" for s in suppressed)


# ---------------------------------------------------------------- 阻断项 3：损坏文件包装


def test_non_zip_file_conversion_error(tmp_path):
    path = tmp_path / "x.pptx"
    path.write_bytes(b"this is not a zip archive")
    with pytest.raises(ConversionError):
        extract_presentation(path)


def _corrupt_presentation_xml(prs, tmp_path, name):
    good = tmp_path / "good.pptx"
    prs.save(good)
    bad = tmp_path / name
    with zipfile.ZipFile(good) as zin, zipfile.ZipFile(bad, "w") as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if item.filename == "ppt/presentation.xml":
                data = b"<p:presentation><broken></p:presentation>"  # tag 不匹配 → XMLSyntaxError
            zout.writestr(item, data)
    return bad


def test_valid_zip_corrupt_xml_conversion_error(tmp_path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_tb(slide, 1, 1, "hello")
    bad = _corrupt_presentation_xml(prs, tmp_path, "bad.pptx")
    with pytest.raises(ConversionError):
        extract_presentation(bad)


def test_internal_error_reraised(monkeypatch, tmp_path):
    def boom(*args, **kwargs):
        raise RuntimeError("internal bug")

    monkeypatch.setattr("pptx.Presentation", boom)
    path = tmp_path / "x.pptx"
    path.write_bytes(b"anything")
    with pytest.raises(RuntimeError, match="internal bug"):
        extract_presentation(path)
