"""A5 Tasks 10–13 — native primitives through the deck asset pipeline.

Task 10: canonical data-asset + data-primitive sugar + placements + theme
materialization + assets.json provenance. Task 11: editable-object structural
verification for every primitive (no raster, editable <a:t>, rect bounds,
contiguous placeholder slot, unique ids). Task 12/13: real-chromium gallery and
cross-provider mixed deck (deck_render gate, skipped without a browser).
"""

from __future__ import annotations

import json
import subprocess as _sp
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.oxml.ns import qn

from offipy import deck
from offipy.assets.declarations import preprocess_asset_declarations
from offipy.assets.model import AssetRef
from offipy.audit.pptx import audit_pptx
from offipy.envcheck import _check_browser
from offipy.exceptions import InvalidArgumentError

_PRIM_VARS = {
    "bg": "#ffffff",
    "surface": "#F2F4F7",
    "ink": "#222222",
    "muted": "#667085",
    "accent": "#2251ff",
}

_DARK_VARS = {
    "bg": "#051C2C",
    "surface": "#0C1C23",
    "ink": "#E5E7EB",
    "muted": "#94A3B8",
    "accent": "#5B8CFF",
}

_CUSTOM_VARS = {
    "bg": "#ffffff",
    "surface": "#ffffff",
    "ink": "#222222",
    "muted": "#667085",
    "accent": "#FF00AA",
}

# 通用 render rect：足够大，所有 primitive 都能按单位结构测试的几何落位。
_PRIM_RECT = {"x": 0, "y": 0, "w": 400, "h": 200}


@pytest.fixture(autouse=True)
def _no_browser(monkeypatch):
    monkeypatch.setattr(deck, "_preflight_browser", lambda: None)


def _prim_html(name: str, attrs: str) -> str:
    return (
        "<!doctype html>\n<html><head></head><body>\n"
        f"<section data-pptx-slide>"
        f'<div data-asset="asset://primitives/primitive/{name}" {attrs}></div></section>\n'
        "</body></html>"
    )


def _asset_meas(asset_id, x=0, y=0, w=400, h=200, theme_vars=None):
    return {
        "id": 1,
        "kind": "asset",
        "tag": "div",
        "assetId": asset_id,
        "rect": {"x": x, "y": y, "w": w, "h": h},
        "themeVars": theme_vars if theme_vars is not None else dict(_PRIM_VARS),
        "color": "rgb(255, 0, 0)",
    }


def _make_fake_convert(decls, *, meas_override=None):
    """fake subprocess.run：造带 OFFIPY_ASSET:: 占位符的 pptx + measurements.json。"""

    def fake_run(cmd, *a, **k):
        out = cmd[cmd.index("--out") + 1]
        prs = Presentation()
        for d in decls:
            while len(prs.slides) < d.slide_index:
                prs.slides.add_slide(prs.slide_layouts[6])
            slide = prs.slides[d.slide_index - 1]
            sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, 400, 200)
            sp.name = f"OFFIPY_ASSET::{d.declaration_id}"
        prs.save(out)
        max_slide = max((d.slide_index for d in decls), default=0)
        slides = [{"records": []} for _ in range(max_slide)]
        for d in decls:
            rec = _asset_meas(d.declaration_id)
            if meas_override:
                rec = meas_override(rec, d.declaration_id)
            slides[d.slide_index - 1]["records"].append(rec)
        audit = Path(out).with_name(f"{Path(out).stem}_audit")
        meas_dir = audit / "_cache"
        meas_dir.mkdir(parents=True)
        (meas_dir / "measurements.json").write_text(
            json.dumps({"slides": slides}, ensure_ascii=False), encoding="utf-8"
        )
        return _sp.CompletedProcess(args=[], returncode=0, stdout="", stderr="")

    return fake_run


def _make_marker_convert(decls):
    """在占位符两侧各放一个内容形状 A/B，验证 primitive 块落回占位符槽位。"""

    def fake_run(cmd, *a, **k):
        out = cmd[cmd.index("--out") + 1]
        prs = Presentation()
        for d in decls:
            while len(prs.slides) < d.slide_index:
                prs.slides.add_slide(prs.slide_layouts[6])
            slide = prs.slides[d.slide_index - 1]
            a = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, 10, 10)
            a.name = "A"
            ph = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, 400, 200)
            ph.name = f"OFFIPY_ASSET::{d.declaration_id}"
            b = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, 10, 10)
            b.name = "B"
        prs.save(out)
        max_slide = max((d.slide_index for d in decls), default=0)
        slides = [{"records": []} for _ in range(max_slide)]
        for d in decls:
            slides[d.slide_index - 1]["records"].append(_asset_meas(d.declaration_id))
        audit = Path(out).with_name(f"{Path(out).stem}_audit")
        meas_dir = audit / "_cache"
        meas_dir.mkdir(parents=True)
        (meas_dir / "measurements.json").write_text(
            json.dumps({"slides": slides}, ensure_ascii=False), encoding="utf-8"
        )
        return _sp.CompletedProcess(args=[], returncode=0, stdout="", stderr="")

    return fake_run


def _render_deck(tmp_path, html_text, *, theme=None):
    html = tmp_path / "d.html"
    html.write_text(html_text, encoding="utf-8")
    pptx = tmp_path / "d.pptx"
    pptx.write_bytes(b"placeholder")
    return deck.render(str(html), out=str(pptx), overwrite=True, theme=theme), pptx


def _load_manifest(tmp_path, stem="d"):
    p = tmp_path / f"{stem}_audit" / "assets.json"
    if not p.exists():
        return None
    return json.loads(p.read_text(encoding="utf-8"))


def _texts_of_shapes(shapes) -> list[str]:
    out: list[str] = []
    for s in shapes:
        if not hasattr(s, "text_frame"):
            continue
        for p in s.text_frame.paragraphs:
            for r in p.runs:
                out.append(r.text)
    return out


def _assert_within_rect(shapes, rect: dict, tol_px: float = 1.0) -> None:
    tol = int(round(tol_px * 6350))
    x0, y0 = int(rect["x"] * 6350), int(rect["y"] * 6350)
    x1, y1 = int((rect["x"] + rect["w"]) * 6350), int((rect["y"] + rect["h"]) * 6350)
    for s in shapes:
        assert s.left >= x0 - tol, f"left overflow: {s.left}"
        assert s.top >= y0 - tol, f"top overflow: {s.top}"
        assert s.left + s.width <= x1 + tol, f"right overflow: {s.left + s.width}"
        assert s.top + s.height <= y1 + tol, f"bottom overflow: {s.top + s.height}"


def _cNvPr_ids(slide) -> list[int]:
    ids = []
    for el in slide.shapes._spTree.iter():
        if el.tag.endswith("cNvPr"):
            ids.append(int(el.get("id") or 0))
    return ids


# ---------------------------------------------------------------------------
# Task 10 — canonical integration + data-primitive sugar + placements + theme
# ---------------------------------------------------------------------------


def test_data_primitive_sugar_preprocesses_to_canonical() -> None:
    html = (
        '<section data-pptx-slide><div data-primitive="metric-badge" '
        'data-asset-param-value="24%"></div></section>'
    )
    out, decls = preprocess_asset_declarations(html)
    assert len(decls) == 1
    assert decls[0].request.ref == AssetRef("primitives", "primitive", "metric-badge")
    assert 'data-asset="asset://primitives/primitive/metric-badge"' in out
    assert 'data-asset-placement="replace"' in out
    assert 'data-offipy-asset-id="asset-s01-001"' in out


def test_canonical_metric_badge_renders_native_shapes(tmp_path, monkeypatch) -> None:
    html_text = _prim_html(
        "metric-badge",
        'data-asset-param-value="24%" data-asset-param-label="YoY" data-asset-param-delta="+3.2%"',
    )
    _, decls = preprocess_asset_declarations(html_text)
    monkeypatch.setattr(deck.subprocess, "run", _make_fake_convert(decls))
    _render_deck(tmp_path, html_text)

    slide = Presentation(str(tmp_path / "d.pptx")).slides[0]
    assert not any(s.name.startswith("OFFIPY_ASSET::") for s in slide.shapes)
    assert not any(s.shape_type == MSO_SHAPE_TYPE.PICTURE for s in slide.shapes)
    assert len(list(slide.shapes)) == 4  # card + value + label + delta
    texts = _texts_of_shapes(slide.shapes)
    assert "24%" in texts and "YoY" in texts and "+3.2%" in texts


def test_primitive_decorative_placement_records_placement(tmp_path, monkeypatch) -> None:
    html_text = _prim_html(
        "label-pill",
        'data-asset-param-text="Hot" data-asset-placement="decorative"',
    )
    _, decls = preprocess_asset_declarations(html_text)
    monkeypatch.setattr(deck.subprocess, "run", _make_fake_convert(decls))
    _render_deck(tmp_path, html_text)

    manifest = _load_manifest(tmp_path)
    assert manifest["assets"][0]["placement"] == "decorative"
    assert not any(
        s.name.startswith("OFFIPY_ASSET::")
        for s in Presentation(str(tmp_path / "d.pptx")).slides[0].shapes
    )


def test_primitive_background_placement_rejected_preserves_output(tmp_path, monkeypatch) -> None:
    html_text = _prim_html(
        "label-pill",
        'data-asset-param-text="Hot" data-asset-placement="background"',
    )
    html = tmp_path / "d.html"
    html.write_text(html_text, encoding="utf-8")
    pptx = tmp_path / "d.pptx"
    pptx.write_bytes(b"placeholder")
    original = pptx.read_bytes()
    _, decls = preprocess_asset_declarations(html_text)
    monkeypatch.setattr(deck.subprocess, "run", _make_marker_convert(decls))

    with pytest.raises(InvalidArgumentError, match="background"):
        deck.render(str(html), out=str(pptx), overwrite=True)

    assert pptx.read_bytes() == original  # 失败不替换最终输出
    assert not list(tmp_path.glob("offipy-deck-*"))  # 无残留注入副本


@pytest.mark.parametrize(
    ("theme_vars", "expected"),
    [
        (dict(_PRIM_VARS), "2251FF"),  # mckinsey base accent
        (dict(_DARK_VARS), "5B8CFF"),  # dark variant
        (dict(_CUSTOM_VARS), "FF00AA"),  # custom
        ({**dict(_CUSTOM_VARS), "--accent": "#FF00AA"}, "FF00AA"),  # --key 形式
    ],
)
def test_measurement_theme_vars_drive_primitive_accent(
    tmp_path, monkeypatch, theme_vars, expected
) -> None:
    html_text = _prim_html(
        "timeline-node",
        'data-asset-param-label="Launch" data-asset-param-phase="current"',
    )
    _, decls = preprocess_asset_declarations(html_text)
    override = lambda rec, aid: {**rec, "themeVars": theme_vars}  # noqa: E731
    monkeypatch.setattr(deck.subprocess, "run", _make_fake_convert(decls, meas_override=override))
    _render_deck(tmp_path, html_text)

    slide = Presentation(str(tmp_path / "d.pptx")).slides[0]
    markers = [
        s
        for s in slide.shapes
        if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and s.auto_shape_type == MSO_SHAPE.OVAL
    ]
    assert len(markers) == 1
    assert str(markers[0].fill.fore_color.rgb) == expected


def test_explicit_hex_accent_overrides_theme_token(tmp_path, monkeypatch) -> None:
    html_text = _prim_html(
        "timeline-node",
        'data-asset-param-label="Launch" data-asset-param-phase="current"'
        ' data-asset-param-accent="#00FF00"',
    )
    _, decls = preprocess_asset_declarations(html_text)
    monkeypatch.setattr(deck.subprocess, "run", _make_fake_convert(decls))
    _render_deck(tmp_path, html_text)

    slide = Presentation(str(tmp_path / "d.pptx")).slides[0]
    marker = next(
        s
        for s in slide.shapes
        if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and s.auto_shape_type == MSO_SHAPE.OVAL
    )
    assert str(marker.fill.fore_color.rgb) == "00FF00"


def test_explicit_hex_fill_materializes_quote_card(tmp_path, monkeypatch) -> None:
    html_text = _prim_html(
        "quote-mark", 'data-asset-param-text="Go" data-asset-param-fill="#123456"'
    )
    _, decls = preprocess_asset_declarations(html_text)
    monkeypatch.setattr(deck.subprocess, "run", _make_fake_convert(decls))
    _render_deck(tmp_path, html_text)

    slide = Presentation(str(tmp_path / "d.pptx")).slides[0]
    cards = [
        s
        for s in slide.shapes
        if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
        and s.auto_shape_type == MSO_SHAPE.ROUNDED_RECTANGLE
    ]
    assert len(cards) == 1
    assert str(cards[0].fill.fore_color.rgb) == "123456"


def test_primitives_manifest_provenance(tmp_path, monkeypatch) -> None:
    html_text = (
        "<!doctype html>\n<html><head></head><body>\n"
        '<section data-pptx-slide><div data-asset="asset://primitives/primitive/label-pill"'
        ' data-asset-param-text="A"></div></section>\n'
        '<section data-pptx-slide><div data-asset="asset://primitives/primitive/section-number"'
        ' data-asset-param-number="1"></div></section>\n'
        '<section data-pptx-slide><div data-asset="asset://primitives/primitive/process-arrow"'
        ' data-asset-param-steps="Plan,Build"></div></section>\n'
        "</body></html>"
    )
    _, decls = preprocess_asset_declarations(html_text)
    monkeypatch.setattr(deck.subprocess, "run", _make_fake_convert(decls))
    _render_deck(tmp_path, html_text)

    manifest = _load_manifest(tmp_path)
    assert manifest["schema"] == 1
    assets = manifest["assets"]
    assert [a["declaration_id"] for a in assets] == [
        "asset-s01-001",
        "asset-s02-001",
        "asset-s03-001",
    ]
    assert [a["slide_index"] for a in assets] == [1, 2, 3]
    for a in assets:
        p = a["provider"]
        assert p["id"] == "primitives"
        assert p["license"] == "MIT"
        assert p["first_party"] is True
        assert p["redistributable"] is True
        assert p["source_url"] == "https://github.com/Zn070515/offipy"


# ---------------------------------------------------------------------------
# Task 11 — editable-object structural verification for every primitive
# ---------------------------------------------------------------------------

_PRIM_CASES = [
    {
        "name": "quote-mark",
        "attrs": 'data-asset-param-text="Work smarter"',
        "count": 2,
        "texts": ["Work smarter"],
    },
    {
        "name": "section-number",
        "attrs": 'data-asset-param-number="42"',
        "count": 2,
        "texts": ["42"],
    },
    {
        "name": "label-pill",
        "attrs": 'data-asset-param-text="Hot"',
        "count": 2,
        "texts": ["Hot"],
    },
    {
        "name": "metric-badge",
        "attrs": 'data-asset-param-value="24%" data-asset-param-label="YoY"'
        ' data-asset-param-delta="+3.2%"',
        "count": 4,
        "texts": ["24%", "YoY", "+3.2%"],
    },
    {
        "name": "timeline-node",
        "attrs": 'data-asset-param-label="Launch" data-asset-param-phase="current"',
        "count": 2,
        "texts": ["Launch"],
    },
    {
        "name": "process-arrow",
        "attrs": 'data-asset-param-steps="Plan,Build,Launch"',
        "count": 3,
        "texts": ["Plan", "Build", "Launch"],
    },
    {
        "name": "device-frame",
        "attrs": 'data-asset-param-device="phone"',
        "count": 4,
        "texts": [],
    },
    {
        "name": "browser-mockup",
        "attrs": 'data-asset-param-title="Q3" data-asset-param-url="https://offipy.dev"',
        "count": 10,
        "texts": ["Q3", "https://offipy.dev"],
    },
]


@pytest.mark.parametrize("case", _PRIM_CASES, ids=lambda c: c["name"])
def test_primitive_structural_contract_in_placeholder_slot(tmp_path, monkeypatch, case) -> None:
    """每个 primitive 落位后：无位图、文本可编辑、rect 内、占位符槽位连续、id 唯一。"""
    html_text = _prim_html(case["name"], case["attrs"])
    _, decls = preprocess_asset_declarations(html_text)
    monkeypatch.setattr(deck.subprocess, "run", _make_marker_convert(decls))
    _render_deck(tmp_path, html_text)

    slide = Presentation(str(tmp_path / "d.pptx")).slides[0]
    names = [s.name for s in slide.shapes]
    assert not any(n.startswith("OFFIPY_ASSET::") for n in names)  # 占位符全消
    assert names[0] == "A" and names[-1] == "B"  # 环绕内容形状保留
    prim_shapes = list(slide.shapes)[1:-1]
    assert len(prim_shapes) == case["count"]  # 连续块 = 环绕形状之间的全部形状
    assert not any(s.shape_type == MSO_SHAPE_TYPE.PICTURE for s in prim_shapes)
    if case["texts"]:
        runs = _texts_of_shapes(prim_shapes)
        for t in case["texts"]:
            assert t in runs, f"text {t!r} not editable in <a:t>"
    _assert_within_rect(prim_shapes, _PRIM_RECT)
    ids = _cNvPr_ids(slide)
    assert len(ids) == len(set(ids)), "duplicate cNvPr id generated"
    assert len({n for n in names if n not in ("A", "B")}) == len(
        [n for n in names if n not in ("A", "B")]
    ), "duplicate shape name generated"


# ---------------------------------------------------------------------------
# Task 12 — primitive gallery smoke deck (real chromium)
# ---------------------------------------------------------------------------

_GALLERY_HTML = """<!doctype html>
<html><head><style data-theme="mckinsey">
  :root { --bg:#ffffff; --surface:#f3f4f6; --ink:#222222; --muted:#667085;
          --accent:#2251ff; }
</style></head><body>
<section data-pptx-slide>
  <h1>quote</h1>
  <div data-asset="asset://primitives/primitive/quote-mark"
       data-asset-param-text="Work smarter, not harder"
       style="position:absolute;left:100px;top:140px;width:500px;height:300px"></div>
</section>
<section data-pptx-slide>
  <h1>section</h1>
  <div data-asset="asset://primitives/primitive/section-number"
       data-asset-param-number="7" data-asset-param-label="Chapter"
       style="position:absolute;left:100px;top:140px;width:500px;height:300px"></div>
</section>
<section data-pptx-slide>
  <h1>pill</h1>
  <div data-asset="asset://primitives/primitive/label-pill"
       data-asset-param-text="Accent override" data-asset-param-accent="#FF00AA"
       style="position:absolute;left:100px;top:140px;width:500px;height:300px"></div>
</section>
<section data-pptx-slide>
  <h1>metric</h1>
  <div data-asset="asset://primitives/primitive/metric-badge"
       data-asset-param-value="24%" data-asset-param-label="YoY"
       data-asset-param-delta="+3.2%"
       style="position:absolute;left:100px;top:140px;width:500px;height:300px"></div>
</section>
<section data-pptx-slide>
  <h1>timeline</h1>
  <div data-asset="asset://primitives/primitive/timeline-node"
       data-asset-param-label="Launch" data-asset-param-phase="future"
       style="position:absolute;left:100px;top:140px;width:500px;height:300px"></div>
</section>
<section data-pptx-slide>
  <h1>arrow</h1>
  <div data-asset="asset://primitives/primitive/process-arrow"
       data-asset-param-steps="Plan,Build,Launch,Grow"
       style="position:absolute;left:100px;top:140px;width:700px;height:300px"></div>
</section>
<section data-pptx-slide>
  <h1>phone</h1>
  <div data-asset="asset://primitives/primitive/device-frame" data-asset-param-device="phone"
       style="position:absolute;left:100px;top:140px;width:500px;height:300px"></div>
</section>
<section data-pptx-slide>
  <h1>tablet</h1>
  <div data-asset="asset://primitives/primitive/device-frame" data-asset-param-device="tablet"
       style="position:absolute;left:100px;top:140px;width:500px;height:300px"></div>
</section>
<section data-pptx-slide>
  <h1>desktop</h1>
  <div data-asset="asset://primitives/primitive/device-frame" data-asset-param-device="desktop"
       style="position:absolute;left:100px;top:140px;width:500px;height:300px"></div>
</section>
<section data-pptx-slide>
  <h1>browser</h1>
  <div data-asset="asset://primitives/primitive/browser-mockup" data-asset-param-title="Q3 Report" data-asset-param-url="https://offipy.dev"
       style="position:absolute;left:100px;top:140px;width:500px;height:300px"></div>
</section>
</body></html>"""


@pytest.mark.deck_render
@pytest.mark.skipif(
    not deck.CONVERT_PY.exists() or not _check_browser().ok,
    reason="vendored 转换器缺失或 chromium 不可用",
)
def test_real_render_primitive_gallery(tmp_path):
    """真实浏览器渲染：8 个 primitive 全渲染、占位符全消、清单 10 条、无位图。"""
    html = tmp_path / "gallery.html"
    html.write_text(_GALLERY_HTML, encoding="utf-8")
    out = tmp_path / "gallery.pptx"
    deck.render(str(html), out=str(out), overwrite=True, timeout=300)

    manifest = _load_manifest(tmp_path, stem="gallery")
    assert manifest is not None
    assert len(manifest["assets"]) == 10
    assert all(a["provider"]["id"] == "primitives" for a in manifest["assets"])

    prs = Presentation(str(out))
    assert len(prs.slides) == 10
    for slide in prs.slides:
        assert not any(s.name.startswith("OFFIPY_ASSET::") for s in slide.shapes)
        assert not any(s.shape_type == MSO_SHAPE_TYPE.PICTURE for s in slide.shapes)
        assert len(list(slide.shapes)) >= 2  # 每页至少内容 + primitive 对象


# ---------------------------------------------------------------------------
# Task 13 — full Asset System cross-provider integration (real chromium)
# ---------------------------------------------------------------------------

_CROSS_HTML = """<!doctype html>
<html><head><style data-theme="mckinsey">
  :root { --bg:#ffffff; --surface:#f3f4f6; --ink:#222222; --muted:#667085;
          --accent:#2251ff; }
</style></head><body>
<section data-pptx-slide>
  <h1>Assets</h1>
  <svg data-icon="ph:check" viewBox="0 0 256 256"
       style="position:absolute;left:80px;top:140px;width:120px;height:120px"></svg>
  <div data-asset="asset://lu/icon/settings"
       style="position:absolute;left:240px;top:140px;width:120px;height:120px"></div>
  <div data-asset="asset://procedural/pattern/topography?seed=42" data-asset-placement="background"
       style="position:absolute;left:0;top:0;width:100%;height:100%"></div>
  <div data-asset="asset://procedural/pattern/rings?count=4"
       style="position:absolute;left:420px;top:140px;width:200px;height:200px"></div>
  <div data-asset="asset://primitives/primitive/metric-badge"
       data-asset-param-value="24%" data-asset-param-label="YoY"
       style="position:absolute;left:680px;top:140px;width:300px;height:180px"></div>
  <div data-asset="asset://primitives/primitive/process-arrow" data-asset-param-steps="Plan,Build"
       style="position:absolute;left:80px;top:420px;width:500px;height:140px"></div>
</section>
<section class="slide chart-dominant" data-pptx-slide data-layout="chart-dominant">
  <h1>Chart</h1>
  <div class="chart" data-chart="bar"
       data-chart-data='{"categories":["Q1","Q2","Q3"],"series":[{"name":"营收","values":[40,55,70]}]}'></div>
</section>
</body></html>"""


def _shape_sequence(slide) -> list[tuple[str, str]]:
    """spTree 顺序下每个形状的 (类别, 名字)。类别: pic/freeform/prst/textbox。"""
    out: list[tuple[str, str]] = []
    for el in slide.shapes._spTree:
        if el.tag == qn("p:pic"):
            out.append(("pic", "SVG Asset"))
        elif el.tag == qn("p:sp"):
            spPr = el.find(qn("p:spPr"))
            if spPr is not None and spPr.find(qn("a:custGeom")) is not None:
                out.append(("freeform", ""))
            else:
                nv = el.find(f"{qn('p:nvSpPr')}/{qn('p:cNvPr')}")
                out.append(("prst", nv.get("name", "") if nv is not None else ""))
    return out


@pytest.mark.deck_render
@pytest.mark.skipif(
    not deck.CONVERT_PY.exists() or not _check_browser().ok,
    reason="vendored 转换器缺失或 chromium 不可用",
)
def test_real_render_cross_provider_mixed_deck(tmp_path):
    """charts + icons + procedural + primitives 共存；图表保持原生、清单只列非图表资产。"""
    html = tmp_path / "cross.html"
    html.write_text(_CROSS_HTML, encoding="utf-8")
    out = tmp_path / "cross.pptx"
    deck.render(str(html), out=str(out), overwrite=True, timeout=300, apply_layouts=True)

    prs = Presentation(str(out))
    slide1 = prs.slides[0]
    assert not any(s.name.startswith("OFFIPY_ASSET::") for s in slide1.shapes)

    seq = _shape_sequence(slide1)
    kinds = [k for k, _ in seq]
    # 背景 procedural 沉底：第一个形状就是 pic（grpSpPr 之后、内容之前）
    assert kinds[0] == "pic"
    # 图标 = freeform；procedural 装饰 = pic；primitive = prst 原生
    assert "freeform" in kinds
    assert kinds.count("pic") == 2  # background + decorative
    assert any(k == "prst" for k in kinds)
    # 背景 pic 在所有 freeform / native 之前
    assert seq[0][0] == "pic"
    first_freeform = kinds.index("freeform")
    first_prst = kinds.index("prst")
    assert first_freeform > 0 and first_prst > 0

    # 图表页保持原生 chart
    assert any(s.has_chart for s in prs.slides[1].shapes)

    # 清单只列非图表资产（图表不是 Asset System v1）
    manifest = _load_manifest(tmp_path, stem="cross")
    assert manifest is not None
    assert len(manifest["assets"]) == 6
    requests = [a["request"] for a in manifest["assets"]]
    assert any(r.startswith("asset://ph/") for r in requests)
    assert any(r.startswith("asset://lu/") for r in requests)
    assert any(r.startswith("asset://procedural/") for r in requests)
    assert any(r.startswith("asset://primitives/") for r in requests)
    assert not any("chart" in r for r in requests)

    # 严格几何审计在 assets 后处理之后仍能跑（kind=asset 记录不崩溃）
    report = audit_pptx(str(out))
    assert report is not None
