"""A5 Task 6 — timeline-node renderer structural contract."""

from __future__ import annotations

from pptx import Presentation
from pptx.enum.dml import MSO_FILL_TYPE
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE

from offipy.assets.model import (
    AssetRect,
    AssetRef,
    AssetRenderContext,
    AssetRequest,
    NativeShapePayload,
)
from offipy.assets.primitives import get_native_renderer
from offipy.assets.providers.primitives import PrimitivesProvider

_THEME = {
    "accent": "#2251ff",
    "surface": "#F2F4F7",
    "ink": "#222222",
    "muted": "#667085",
    "bg": "#ffffff",
}


def _ctx(**overrides) -> AssetRenderContext:
    kwargs = dict(
        slide_index=1,
        rect=AssetRect(x=0, y=0, width=400, height=200),
        theme_name="mckinsey",
        theme_vars=dict(_THEME),
        placement="replace",
    )
    kwargs.update(overrides)
    return AssetRenderContext(**kwargs)


def _blank_slide():
    prs = Presentation()
    return prs.slides.add_slide(prs.slide_layouts[6])


def _render(params: dict[str, str], **ctx_overrides):
    resolved = PrimitivesProvider().resolve(
        AssetRequest(
            ref=AssetRef("primitives", "primitive", "timeline-node"),
            params=tuple(params.items()),
        )
    )
    assert isinstance(resolved.payload, NativeShapePayload)
    slide = _blank_slide()
    get_native_renderer("timeline-node")(
        slide, dict(resolved.payload.params), _ctx(**ctx_overrides)
    )
    return slide


def _runs(sp):
    return [r.text for p in sp.text_frame.paragraphs for r in p.runs]


def _texts(slide) -> list[str]:
    return [t for sp in slide.shapes if hasattr(sp, "text_frame") for t in _runs(sp)]


def _marker(slide):
    for sp in slide.shapes:
        if sp.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and sp.auto_shape_type == MSO_SHAPE.OVAL:
            return sp
    raise AssertionError("no OVAL marker found")


def _assert_within_rect(slide, rect: AssetRect, tol_px: float = 1.0) -> None:
    tol = int(round(tol_px * 6350))
    x0, y0 = int(rect.x * 6350), int(rect.y * 6350)
    x1, y1 = int((rect.x + rect.width) * 6350), int((rect.y + rect.height) * 6350)
    for sp in slide.shapes:
        assert sp.left >= x0 - tol, f"left overflow: {sp.left}"
        assert sp.top >= y0 - tol, f"top overflow: {sp.top}"
        assert sp.left + sp.width <= x1 + tol, f"right overflow: {sp.left + sp.width}"
        assert sp.top + sp.height <= y1 + tol, f"bottom overflow: {sp.top + sp.height}"


def test_timeline_node_default_phase_is_current() -> None:
    slide = _render({})
    assert str(_marker(slide).fill.fore_color.rgb) == "2251FF"


def test_timeline_node_past_uses_muted_fill() -> None:
    slide = _render({"phase": "past"})
    assert str(_marker(slide).fill.fore_color.rgb) == "667085"


def test_timeline_node_current_uses_accent() -> None:
    slide = _render({"phase": "current"})
    assert str(_marker(slide).fill.fore_color.rgb) == "2251FF"


def test_timeline_node_current_fill_param_applies() -> None:
    # #59：fill 公共参数应驱动 current marker 填充
    slide = _render({"phase": "current", "fill": "#E8635A"})
    assert str(_marker(slide).fill.fore_color.rgb) == "E8635A"


def test_timeline_node_current_accent_override_follows_fill() -> None:
    # fill 默认 transparent → 回退 accent；accent param 应生效
    slide = _render({"phase": "current", "accent": "#FF0000"})
    assert str(_marker(slide).fill.fore_color.rgb) == "FF0000"


def test_timeline_node_future_is_outline() -> None:
    slide = _render({"phase": "future"})
    marker = _marker(slide)
    assert marker.fill.type == MSO_FILL_TYPE.BACKGROUND
    assert str(marker.line.color.rgb) == "667085"


def test_timeline_node_three_phases_differ() -> None:
    past = _render({"phase": "past"})
    current = _render({"phase": "current"})
    future = _render({"phase": "future"})
    fills = {
        str(_marker(past).fill.fore_color.rgb),
        str(_marker(current).fill.fore_color.rgb),
    }
    assert len(fills) == 2  # past and current both solid, differ in color
    assert _marker(future).fill.type != MSO_FILL_TYPE.SOLID


def test_timeline_node_label_editable() -> None:
    slide = _render({"phase": "current", "label": "Launch"})
    assert "Launch" in _texts(slide)


def test_timeline_node_no_label_renders_marker_only() -> None:
    slide = _render({"phase": "past"})
    assert len(list(slide.shapes)) == 1


def test_timeline_node_shapes_within_rect() -> None:
    slide = _render({"phase": "future", "label": "Launch"})
    _assert_within_rect(slide, _ctx().rect)


def test_timeline_node_no_pictures() -> None:
    slide = _render({"phase": "current", "label": "L"})
    assert not any(sp.shape_type == MSO_SHAPE_TYPE.PICTURE for sp in slide.shapes)
