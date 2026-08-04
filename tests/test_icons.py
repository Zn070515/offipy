# tests/test_icons.py
"""图标声明解析 / SVG path 解析 / 注入测试（纯 Python，不依赖 Office）。"""

import io
import json
import math
import struct
import zlib

import pytest
from pptx import Presentation

from offipy.icons import (
    IconDecl,
    _geom_points,
    _parse_path,
    _parse_points_list,
    load_icon_svg,
    parse_icon_declarations,
)

# ---------- path parser ----------


def test_parse_path_m_l_z():
    sp = _parse_path("M10 10 L20 10 L20 20 Z")
    assert len(sp) == 1
    assert sp[0].close is True
    assert sp[0].points == [(10.0, 10.0), (20.0, 10.0), (20.0, 20.0)]


def test_parse_path_relative():
    sp = _parse_path("M10 10 l5 0 l0 5 z")
    assert len(sp) == 1
    assert sp[0].close is True
    assert sp[0].points == [(10.0, 10.0), (15.0, 10.0), (15.0, 15.0)]


def test_parse_path_hv_commands():
    sp = _parse_path("M0 0 H10 V10 h-5 z")
    assert sp[0].points == [(0.0, 0.0), (10.0, 0.0), (10.0, 10.0), (5.0, 10.0)]


def test_parse_path_cubic_flattens():
    sp = _parse_path("M0 0 C0 10 10 10 10 0 Z")
    assert sp[0].close is True
    assert len(sp[0].points) >= 5  # 压平出多段
    assert sp[0].points[-1] == (10.0, 0.0)
    assert sp[0].points[0] == (0.0, 0.0)


def test_parse_path_quad_flattens():
    sp = _parse_path("M0 0 Q5 10 10 0 Z")
    assert len(sp[0].points) >= 5
    assert sp[0].points[-1] == (10.0, 0.0)


def test_parse_path_arc_flattens():
    sp = _parse_path("M0 0 A5 5 0 0 1 10 0 Z")
    assert len(sp[0].points) >= 5
    assert sp[0].points[-1] == (10.0, 0.0)


def test_parse_path_multi_subpath():
    sp = _parse_path("M0 0 L5 5 Z M10 10 L15 15")
    assert len(sp) == 2
    assert sp[0].close is True
    assert sp[1].close is False
    assert sp[1].points == [(10.0, 10.0), (15.0, 15.0)]


def test_parse_path_smooth_commands():
    # S 反射上一个 C 控制点；T 反射上一个 Q 控制点（不抛错即可）
    sp = _parse_path("M0 0 C0 5 5 5 5 0 S10 -5 10 0 Z")
    assert len(sp) == 1
    sp2 = _parse_path("M0 0 Q2 5 5 0 T10 0 Z")
    assert len(sp2) == 1


def test_parse_path_empty_raises():
    with pytest.raises(ValueError):
        _parse_path("")


def test_parse_path_bad_first_token_raises():
    with pytest.raises(ValueError):
        _parse_path("10 10 L20 20")


# ---------- geometry 元素 ----------


def test_geom_line():
    pts, close = _geom_points("line", {"x1": "0", "y1": "0", "x2": "10", "y2": "5"})
    assert pts == [(0.0, 0.0), (10.0, 5.0)]
    assert close is False


def test_geom_polyline_and_polygon():
    pts, close = _geom_points("polyline", {"points": "6 9 12 15 18 9"})
    assert pts == [(6.0, 9.0), (12.0, 15.0), (18.0, 9.0)]
    assert close is False
    _, close = _geom_points("polygon", {"points": "0 0 4 0 2 3"})
    assert close is True


def test_geom_rect():
    pts, close = _geom_points("rect", {"x": "1", "y": "2", "width": "4", "height": "3"})
    assert close is True
    assert pts == [(1.0, 2.0), (5.0, 2.0), (5.0, 5.0), (1.0, 5.0)]


def test_geom_circle():
    from offipy.icons import _CIRCLE_SAMPLES

    pts, close = _geom_points("circle", {"cx": "5", "cy": "5", "r": "5"})
    assert close is True
    assert len(pts) == _CIRCLE_SAMPLES
    assert abs(sum(x for x, _ in pts) / len(pts) - 5.0) < 1e-6  # 中心≈cx
    assert abs(sum(y for _, y in pts) / len(pts) - 5.0) < 1e-6


def test_parse_points_list():
    assert _parse_points_list("6 9 12 15 18 9") == [(6.0, 9.0), (12.0, 15.0), (18.0, 9.0)]


# ---------- 资产访问 ----------

PH_SVG = (
    '<svg xmlns="http://www.w3.org/2000/svg" width="256" height="256" '
    'viewBox="0 0 256 256" fill="currentColor">'
    '<path d="M213.66 82.34a8 8 0 0 0-11.32 0L128 156.69l-74.34-74.35'
    'a8 8 0 0 0-11.32 11.32l80 80a8 8 0 0 0 11.32 0l80-80a8 8 0 0 0 0-11.32Z"/>'
    "</svg>"
)

LU_SVG = (
    '<svg xmlns="http://www.w3.org/2000/svg" width="24" height="24" '
    'viewBox="0 0 24 24" fill="none" stroke="currentColor" stroke-width="2" '
    'stroke-linecap="round" stroke-linejoin="round">'
    '<line x1="12" x2="12" y1="2" y2="22"/>'
    '<polyline points="6 9 12 15 18 9"/>'
    "</svg>"
)


@pytest.fixture
def fake_assets(tmp_path, monkeypatch):
    from offipy import icons

    assets = tmp_path / "assets" / "icons"
    (assets / "phosphor").mkdir(parents=True)
    (assets / "lucide").mkdir(parents=True)
    # ph 资产按真实库约定 <name>-fill.svg 命名
    (assets / "phosphor" / "check-fill.svg").write_text(PH_SVG, encoding="utf-8")
    (assets / "lucide" / "zap.svg").write_text(LU_SVG, encoding="utf-8")
    monkeypatch.setattr(icons, "ASSETS_DIR", assets)
    return assets


def test_load_icon_svg(fake_assets):
    text = load_icon_svg("ph:check")
    assert 'fill="currentColor"' in text
    text = load_icon_svg("lu:zap")
    assert 'stroke="currentColor"' in text


def test_load_icon_svg_missing_raises(fake_assets):
    with pytest.raises(ValueError, match="ph:missing"):
        load_icon_svg("ph:missing")


def test_load_icon_svg_bad_set(fake_assets):
    with pytest.raises(ValueError, match="xx"):
        load_icon_svg("xx:check")


def test_load_icon_svg_ph_resolves_fill_suffix():
    # 真实资产库（不 monkeypatch ASSETS_DIR）：ph 资产按 <name>-fill.svg 命名，
    # ph:check 必须能解析到 check-fill.svg；ph:check-fill 已带后缀不重复加。
    from offipy import icons

    content = icons.load_icon_svg("ph:check")
    assert content.startswith("<svg")  # 解析成功，读到真实资产
    assert content == icons.load_icon_svg("ph:check-fill")  # 两种写法解析到同一资产


def test_svg_to_subpaths_phosphor(fake_assets):
    from offipy.icons import _svg_to_subpaths

    subpaths, sw = _svg_to_subpaths(load_icon_svg("ph:check"))
    assert len(subpaths) == 1
    assert subpaths[0].close is True
    assert len(subpaths[0].points) >= 3


def test_svg_to_subpaths_lucide(fake_assets):
    from offipy.icons import _svg_to_subpaths

    subpaths, sw = _svg_to_subpaths(load_icon_svg("lu:zap"))
    assert len(subpaths) == 2  # line + polyline
    assert subpaths[0].close is False
    assert sw == 2.0


# ---------- 健壮性（M3-T2 审查补测） ----------


def test_parse_path_truncated_raises():
    with pytest.raises(ValueError):
        _parse_path("M10 10 C0 5 5")


def test_parse_path_degenerate_arc_lineto():
    # 零半径弧按 SVG F.6.2 应 lineto 到终点，不丢终点
    sp = _parse_path("M0 0 A0 0 0 0 1 10 0")
    assert sp[0].points[-1] == (10.0, 0.0)


def test_parse_path_first_not_m_raises():
    with pytest.raises(ValueError):
        _parse_path("L10 10")


def test_geom_ellipse():
    from offipy.icons import _CIRCLE_SAMPLES

    pts, close = _geom_points("ellipse", {"cx": "0", "cy": "0", "rx": "4", "ry": "2"})
    assert close is True
    assert len(pts) == _CIRCLE_SAMPLES
    # 长半轴 4，所有点到中心距离 ∈ [2, 4]
    assert all(2.0 - 1e-9 <= math.hypot(x, y) <= 4.0 + 1e-9 for x, y in pts)


def test_load_icon_svg_invalid_name_raises(fake_assets):
    with pytest.raises(ValueError):
        load_icon_svg("ph:../x")


# ---------- HTML 声明解析（M3-T3） ----------

ATTR_HTML = (
    "<!DOCTYPE html>\n<html><body>\n"
    '<section class="slide" data-pptx-slide>\n'
    '  <svg class="icon" data-icon="ph:check" viewBox="0 0 256 256" width="48" height="48"></svg>\n'
    "</section>\n"
    "</body></html>"
)

MULTI_HTML = (
    "<html><body>\n"
    '<section class="slide" data-pptx-slide>\n'
    '  <svg class="icon" data-icon="ph:check" viewBox="0 0 256 256"></svg>\n'
    '  <svg class="icon" data-icon="lu:zap" viewBox="0 0 24 24"></svg>\n'
    "</section>\n"
    "</body></html>"
)


def test_parse_icon_attribute_format():
    decls = parse_icon_declarations(ATTR_HTML)
    assert len(decls) == 1
    assert decls[0].slide_index == 1
    assert decls[0].data_icon == "ph:check"
    assert decls[0].view_box == (0.0, 0.0, 256.0, 256.0)


def test_parse_icon_multiple_on_page():
    decls = parse_icon_declarations(MULTI_HTML)
    assert len(decls) == 2
    assert [d.data_icon for d in decls] == ["ph:check", "lu:zap"]
    assert decls[1].view_box == (0.0, 0.0, 24.0, 24.0)


def test_parse_icon_no_icon_returns_empty():
    assert parse_icon_declarations("<section data-pptx-slide><h2>x</h2></section>") == []


def test_parse_icon_outside_slide_raises():
    with pytest.raises(ValueError):
        parse_icon_declarations(
            '<svg class="icon" data-icon="ph:check" viewBox="0 0 256 256"></svg>'
            "<section data-pptx-slide><h2>x</h2></section>"
        )


def test_parse_icon_bad_viewbox_raises():
    with pytest.raises(ValueError):
        parse_icon_declarations(
            "<section data-pptx-slide>"
            '<svg class="icon" data-icon="ph:check" viewBox="0 0 256"></svg>'
            "</section>"
        )


def test_parse_icon_missing_viewbox_lu():
    # lu 缺省 viewBox 是 24 坐标系；按 256 兜底会静默缩小到约 1/10
    decls = parse_icon_declarations(
        '<section data-pptx-slide><svg class="icon" data-icon="lu:zap"></svg></section>'
    )
    assert len(decls) == 1
    assert decls[0].view_box == (0.0, 0.0, 24.0, 24.0)


def _blank_pptx(tmp_path):
    from pptx.util import Emu

    prs = Presentation()
    prs.slide_width = Emu(12192000)
    prs.slide_height = Emu(6858000)
    prs.slides.add_slide(prs.slide_layouts[6])
    return prs


def _png_bytes() -> bytes:
    def chunk(typ, data):
        c = typ + data
        return struct.pack(">I", len(data)) + c + struct.pack(">I", zlib.crc32(c) & 0xFFFFFFFF)

    ihdr = struct.pack(">IIBBBBB", 1, 1, 8, 6, 0, 0, 0)
    idat = zlib.compress(b"\x00\x00\x00\x00")
    return b"\x89PNG\r\n\x1a\n" + chunk(b"IHDR", ihdr) + chunk(b"IDAT", idat) + chunk(b"IEND", b"")


def test_load_icon_boxes(tmp_path):
    from offipy.icons import load_icon_boxes

    meas = tmp_path / "measurements.json"
    meas.write_text(
        json.dumps(
            {
                "slides": [
                    {
                        "slide": {},
                        "records": [
                            {
                                "id": 1,
                                "kind": "text",
                                "tag": "h2",
                                "className": "title",
                                "rect": {"x": 96, "y": 96, "w": 800, "h": 80},
                            },
                            {
                                "id": 2,
                                "kind": "svg",
                                "tag": "svg",
                                "rect": {"x": 96, "y": 225, "w": 48, "h": 48},
                                "outerHTML": (
                                    '<svg data-icon="ph:check" viewBox="0 0 256 256"></svg>'
                                ),
                                "color": "rgb(34, 81, 255)",
                            },
                            {
                                "id": 3,
                                "kind": "svg",
                                "tag": "svg",
                                "rect": {"x": 160, "y": 225, "w": 48, "h": 48},
                                "outerHTML": '<svg data-icon="lu:zap" viewBox="0 0 24 24"></svg>',
                                "color": "rgb(51, 65, 85)",
                            },
                        ],
                    }
                ]
            }
        ),
        encoding="utf-8",
    )
    boxes = load_icon_boxes(str(meas))
    assert len(boxes[1]) == 2
    assert boxes[1][0]["rect"] == {"x": 96, "y": 225, "w": 48, "h": 48}
    assert boxes[1][0]["color"] == "rgb(34, 81, 255)"


def test_inject_icon_ph(tmp_path, fake_assets):
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.util import Emu

    from offipy.icons import inject_icons

    prs = _blank_pptx(tmp_path)
    slide = prs.slides[0]
    # 模拟 convert 生成的 SVG 占位（PNG picture）
    slide.shapes.add_picture(
        io.BytesIO(_png_bytes()), Emu(96 * 6350), Emu(225 * 6350), Emu(48 * 6350), Emu(48 * 6350)
    )
    out = tmp_path / "out.pptx"
    prs.save(str(out))

    decl = IconDecl(slide_index=1, data_icon="ph:check", view_box=(0.0, 0.0, 256.0, 256.0))
    matched = {
        1: [
            (
                decl,
                {
                    "rect": {"x": 96, "y": 225, "w": 48, "h": 48},
                    "color": "rgb(34, 81, 255)",
                    "outerHTML": '<svg data-icon="ph:check"></svg>',
                },
            )
        ]
    }
    inject_icons(str(out), matched)

    prs2 = Presentation(str(out))
    slide2 = prs2.slides[0]
    freeforms = [s for s in slide2.shapes if s.shape_type == MSO_SHAPE_TYPE.FREEFORM]
    assert len(freeforms) >= 1
    # 占位 picture 被移除
    from pptx.enum.shapes import MSO_SHAPE_TYPE as ST

    assert not any(s.shape_type == ST.PICTURE for s in slide2.shapes)
    # fill 模式：非透明填充
    assert freeforms[0].fill.type is not None


def test_inject_icon_lu(tmp_path, fake_assets):
    from pptx.enum.dml import MSO_FILL_TYPE
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.util import Emu

    from offipy.icons import inject_icons

    prs = _blank_pptx(tmp_path)
    slide = prs.slides[0]
    slide.shapes.add_picture(
        io.BytesIO(_png_bytes()), Emu(96 * 6350), Emu(225 * 6350), Emu(48 * 6350), Emu(48 * 6350)
    )
    out = tmp_path / "out.pptx"
    prs.save(str(out))

    decl = IconDecl(slide_index=1, data_icon="lu:zap", view_box=(0.0, 0.0, 24.0, 24.0))
    matched = {
        1: [
            (
                decl,
                {
                    "rect": {"x": 96, "y": 225, "w": 48, "h": 48},
                    "color": "rgb(51, 65, 85)",
                    "outerHTML": '<svg data-icon="lu:zap"></svg>',
                },
            )
        ]
    }
    inject_icons(str(out), matched)

    prs2 = Presentation(str(out))
    slide2 = prs2.slides[0]
    freeforms = [s for s in slide2.shapes if s.shape_type == MSO_SHAPE_TYPE.FREEFORM]
    assert len(freeforms) >= 2  # line + polyline 各一个
    # stroke 模式：显式 noFill（否则 PowerPoint 会按隐式闭合用主题蓝涂实心）
    assert freeforms[0].fill.type == MSO_FILL_TYPE.BACKGROUND
    assert freeforms[0].line.width > 0


def test_inject_icon_lu_round_stroke(tmp_path, fake_assets):
    """stroke 集对齐 Lucide 设计：freeform 折线设 round cap + round join（默认 flat 尖角）。"""
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.oxml.ns import qn
    from pptx.util import Emu

    from offipy.icons import inject_icons

    prs = _blank_pptx(tmp_path)
    slide = prs.slides[0]
    slide.shapes.add_picture(
        io.BytesIO(_png_bytes()), Emu(96 * 6350), Emu(225 * 6350), Emu(48 * 6350), Emu(48 * 6350)
    )
    out = tmp_path / "out.pptx"
    prs.save(str(out))

    decl = IconDecl(slide_index=1, data_icon="lu:zap", view_box=(0.0, 0.0, 24.0, 24.0))
    matched = {
        1: [
            (
                decl,
                {
                    "rect": {"x": 96, "y": 225, "w": 48, "h": 48},
                    "color": "rgb(51, 65, 85)",
                    "outerHTML": '<svg data-icon="lu:zap"></svg>',
                },
            )
        ]
    }
    inject_icons(str(out), matched)

    prs2 = Presentation(str(out))
    slide2 = prs2.slides[0]
    freeforms = [s for s in slide2.shapes if s.shape_type == MSO_SHAPE_TYPE.FREEFORM]
    for shp in freeforms:
        ln = shp.line._get_or_add_ln()
        assert ln.get("cap") == "rnd"
        assert ln.find(qn("a:round")) is not None


def test_inject_icons_multiple_on_page(tmp_path, fake_assets):
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.util import Emu

    from offipy.icons import inject_icons

    prs = _blank_pptx(tmp_path)
    slide = prs.slides[0]
    # 一行两个图标占位（ph:check + lu:zap）
    slide.shapes.add_picture(
        io.BytesIO(_png_bytes()), Emu(96 * 6350), Emu(225 * 6350), Emu(48 * 6350), Emu(48 * 6350)
    )
    slide.shapes.add_picture(
        io.BytesIO(_png_bytes()), Emu(160 * 6350), Emu(225 * 6350), Emu(48 * 6350), Emu(48 * 6350)
    )
    out = tmp_path / "out.pptx"
    prs.save(str(out))

    decl_ph = IconDecl(slide_index=1, data_icon="ph:check", view_box=(0.0, 0.0, 256.0, 256.0))
    decl_lu = IconDecl(slide_index=1, data_icon="lu:zap", view_box=(0.0, 0.0, 24.0, 24.0))
    matched = {
        1: [
            (
                decl_ph,
                {
                    "rect": {"x": 96, "y": 225, "w": 48, "h": 48},
                    "color": "rgb(34, 81, 255)",
                    "outerHTML": '<svg data-icon="ph:check"></svg>',
                },
            ),
            (
                decl_lu,
                {
                    "rect": {"x": 160, "y": 225, "w": 48, "h": 48},
                    "color": "rgb(51, 65, 85)",
                    "outerHTML": '<svg data-icon="lu:zap"></svg>',
                },
            ),
        ]
    }
    inject_icons(str(out), matched)

    prs2 = Presentation(str(out))
    slide2 = prs2.slides[0]
    freeforms = [s for s in slide2.shapes if s.shape_type == MSO_SHAPE_TYPE.FREEFORM]
    assert len(freeforms) >= 3  # ph:check(1) + lu:zap(line+polyline 2)
    # 两个占位 picture 都被移除
    assert not any(s.shape_type == MSO_SHAPE_TYPE.PICTURE for s in slide2.shapes)


def test_inject_icon_lu_mixed_fill(tmp_path, fake_assets):
    from pptx.enum.dml import MSO_FILL_TYPE
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.util import Emu

    from offipy.icons import inject_icons

    # Lucide 混 fill 图标：一个 fill="currentColor" 闭合路径（该画实心）+ 一个纯 stroke 路径
    (fake_assets / "lucide" / "fill.svg").write_text(
        '<svg xmlns="http://www.w3.org/2000/svg" width="24" height="24" '
        'viewBox="0 0 24 24" fill="none" stroke="currentColor" stroke-width="2" '
        'stroke-linecap="round" stroke-linejoin="round">'
        '<path d="M12 2 L22 22 L2 22 Z" fill="currentColor"/>'
        '<path d="M4 4 L8 4 L8 8 L4 8 Z"/>'
        "</svg>",
        encoding="utf-8",
    )

    prs = _blank_pptx(tmp_path)
    slide = prs.slides[0]
    slide.shapes.add_picture(
        io.BytesIO(_png_bytes()), Emu(96 * 6350), Emu(225 * 6350), Emu(48 * 6350), Emu(48 * 6350)
    )
    out = tmp_path / "out.pptx"
    prs.save(str(out))

    decl = IconDecl(slide_index=1, data_icon="lu:fill", view_box=(0.0, 0.0, 24.0, 24.0))
    matched = {
        1: [
            (
                decl,
                {
                    "rect": {"x": 96, "y": 225, "w": 48, "h": 48},
                    "color": "rgb(51, 65, 85)",
                    "outerHTML": '<svg data-icon="lu:fill"></svg>',
                },
            )
        ]
    }
    inject_icons(str(out), matched)

    prs2 = Presentation(str(out))
    slide2 = prs2.slides[0]
    freeforms = [s for s in slide2.shapes if s.shape_type == MSO_SHAPE_TYPE.FREEFORM]
    assert len(freeforms) == 2
    solid = [s for s in freeforms if s.fill.type == MSO_FILL_TYPE.SOLID]
    bg = [s for s in freeforms if s.fill.type == MSO_FILL_TYPE.BACKGROUND]
    assert len(solid) == 1  # fill="currentColor" 闭合路径 → 实心填充
    assert len(bg) == 1  # 纯 stroke 路径 → noFill + 描边


def test_inject_icon_color_fallback_theme_accent(tmp_path, fake_assets):
    """svg record color 缺失 → 用传入的 fallback（主题 --accent），而非写死主蓝。"""
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.util import Emu

    from offipy.icons import inject_icons

    prs = _blank_pptx(tmp_path)
    slide = prs.slides[0]
    slide.shapes.add_picture(
        io.BytesIO(_png_bytes()), Emu(96 * 6350), Emu(225 * 6350), Emu(48 * 6350), Emu(48 * 6350)
    )
    out = tmp_path / "out.pptx"
    prs.save(str(out))

    decl = IconDecl(slide_index=1, data_icon="ph:check", view_box=(0.0, 0.0, 256.0, 256.0))
    matched = {
        1: [
            (
                decl,
                {
                    "rect": {"x": 96, "y": 225, "w": 48, "h": 48},
                    "color": None,  # 容器没测到颜色 → 走 fallback
                    "outerHTML": '<svg data-icon="ph:check"></svg>',
                },
            )
        ]
    }
    inject_icons(str(out), matched, fallback=RGBColor(0x1F, 0x3A, 0x5F))  # academic --accent

    prs2 = Presentation(str(out))
    freeforms = [s for s in prs2.slides[0].shapes if s.shape_type == MSO_SHAPE_TYPE.FREEFORM]
    assert freeforms[0].fill.fore_color.rgb == RGBColor(0x1F, 0x3A, 0x5F)


def test_postprocess_color_fallback_uses_theme_accent(tmp_path, fake_assets):
    """容器未设 color → 图标缺省色从 HTML <style data-theme> 的主题 --accent 取。"""
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.util import Emu

    from offipy.icons import postprocess_icons

    html = tmp_path / "d.html"
    html.write_text(
        '<html><head><style data-theme="academic"></style></head><body>'
        '<section class="slide icons-row" data-pptx-slide data-layout="icons-row">'
        '<svg class="icon" data-icon="ph:check" viewBox="0 0 256 256"></svg>'
        "</section></body></html>",
        encoding="utf-8",
    )
    prs = _blank_pptx(tmp_path)
    prs.slides[0].shapes.add_picture(
        io.BytesIO(_png_bytes()),
        Emu(96 * 6350),
        Emu(225 * 6350),
        Emu(48 * 6350),
        Emu(48 * 6350),
    )
    pptx = tmp_path / "d.pptx"
    prs.save(str(pptx))
    meas_dir = tmp_path / "d_audit" / "_cache"
    meas_dir.mkdir(parents=True)
    (meas_dir / "measurements.json").write_text(
        json.dumps(
            {
                "slides": [
                    {
                        "slide": {},
                        "records": [
                            {
                                "id": 1,
                                "kind": "svg",
                                "tag": "svg",
                                "rect": {"x": 96, "y": 225, "w": 48, "h": 48},
                                "outerHTML": (
                                    '<svg class="icon" data-icon="ph:check" '
                                    'viewBox="0 0 256 256"></svg>'
                                ),
                                # 无 color 字段 → 走主题 fallback
                            }
                        ],
                    }
                ]
            }
        ),
        encoding="utf-8",
    )

    postprocess_icons(str(html), str(pptx))

    prs2 = Presentation(str(pptx))
    freeforms = [s for s in prs2.slides[0].shapes if s.shape_type == MSO_SHAPE_TYPE.FREEFORM]
    assert freeforms
    assert freeforms[0].fill.fore_color.rgb == RGBColor(0x1F, 0x3A, 0x5F)  # academic --accent


def test_postprocess_skips_without_icons(tmp_path, monkeypatch):
    from offipy import icons
    from offipy.icons import postprocess_icons

    called = {}
    monkeypatch.setattr(icons, "load_icon_boxes", lambda *a, **k: called.setdefault("load", True))
    html = tmp_path / "d.html"
    html.write_text("<section data-pptx-slide><h2>x</h2></section>", encoding="utf-8")
    pptx = tmp_path / "d.pptx"
    pptx.write_bytes(b"not a real pptx")
    postprocess_icons(str(html), str(pptx))
    assert "load" not in called


def test_postprocess_missing_measurements_raises(tmp_path):
    from offipy.icons import postprocess_icons

    html = tmp_path / "d.html"
    html.write_text(ATTR_HTML, encoding="utf-8")
    pptx = tmp_path / "d.pptx"
    pptx.write_bytes(b"x")
    with pytest.raises(RuntimeError, match="measurements"):
        postprocess_icons(str(html), str(pptx))


def test_postprocess_missing_box_raises(tmp_path, fake_assets):
    from offipy.icons import postprocess_icons

    html = tmp_path / "d.html"
    html.write_text(ATTR_HTML, encoding="utf-8")
    pptx = tmp_path / "d.pptx"
    pptx.write_bytes(b"x")
    meas_dir = tmp_path / "d_audit" / "_cache"
    meas_dir.mkdir(parents=True)
    (meas_dir / "measurements.json").write_text(
        json.dumps({"slides": [{"slide": {}, "records": []}]}), encoding="utf-8"
    )
    with pytest.raises(RuntimeError, match="ph:check"):
        postprocess_icons(str(html), str(pptx))
