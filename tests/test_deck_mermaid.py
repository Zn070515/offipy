"""deck 集成：HTML <pre class="mermaid"> 块 → measurements bbox → 可编辑形状注入。"""

from __future__ import annotations

import json

import pytest
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu, Inches

from offipy.diagrams import parse_mermaid_declarations, postprocess_mermaid

HTML = """\
<!doctype html>
<html><head><meta charset="utf-8"></head><body>
<section data-pptx-slide>
  <h1>系统流程</h1>
  <div class="mermaid-box">
    <pre class="mermaid">graph TD
    A[开始] --> B[处理]
    B --> C[输出]</pre>
  </div>
</section>
</body></html>
"""


def _measurements(slide_index=1):
    return {
        "slides": [
            {
                "index": slide_index,
                "records": [
                    {"className": "mermaid-box", "rect": {"x": 0, "y": 0, "w": 0, "h": 0}},
                    {"className": "mermaid", "rect": {"x": 60, "y": 120, "w": 400, "h": 240}},
                ],
            }
        ]
    }


def _write(tmp_path):
    html = tmp_path / "deck.html"
    html.write_text(HTML, encoding="utf-8")
    meas = tmp_path / "deck_audit" / "_cache" / "measurements.json"
    meas.parent.mkdir(parents=True)
    meas.write_text(json.dumps(_measurements()), encoding="utf-8")
    return str(html)


def test_parse_mermaid_declarations():
    decls = parse_mermaid_declarations(HTML)
    assert len(decls) == 1
    assert decls[0]["slide"] == 1
    assert "A[开始]" in decls[0]["source"]


def test_parse_mermaid_declarations_skips_non_mermaid_pre():
    html = HTML.replace('class="mermaid"', 'class="code"')
    assert parse_mermaid_declarations(html) == []


def test_parse_mermaid_declarations_multi_class_tokens():
    # class 值是分词匹配：多 class（mermaid chart）也必须命中。postprocess_mermaid
    # 不设子串守卫（会漏掉此类变体），正确性全押在 parser 的 token 语义上。
    html = HTML.replace('<pre class="mermaid">', '<pre class="mermaid chart">')
    decls = parse_mermaid_declarations(html)
    assert len(decls) == 1
    assert "A[开始]" in decls[0]["source"]


def test_parse_mermaid_declarations_rejects_outside_section():
    html = '<html><body><pre class="mermaid">graph TD\nA --> B</pre></body></html>'
    with pytest.raises(ValueError, match="slide 之外"):
        parse_mermaid_declarations(html)


def test_postprocess_mermaid_injects_shapes(tmp_path):
    html = _write(tmp_path)
    # 模拟 convert 产物：<pre class="mermaid"> 渲染成占位文本框（bbox px 60,120,400,240
    # → inch 0.42,0.83,2.78,1.67，中心落在 bbox 内）。不能拿 mermaid_to_pptx 产物当
    # 占位——它已含可编辑形状，会让"postprocess 没注入也通过"的假阳性逃过断言。
    pptx_path = str(tmp_path / "deck.pptx")
    prs = Presentation()
    prs.slide_width = Emu(12192000)
    prs.slide_height = Emu(6858000)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(0.42), Inches(0.83), Inches(2.78), Inches(1.67))
    tb.text = "graph TD\n    A[开始] --> B[处理]"
    prs.save(pptx_path)

    postprocess_mermaid(html, pptx_path)
    prs = Presentation(pptx_path)
    auto = [
        sh for sl in prs.slides for sh in sl.shapes if sh.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
    ]
    texts = {sh.text_frame.text for sh in auto}
    assert "开始" in texts  # 注入的可编辑形状出现
    # 占位文本被移除（不再是字面 mermaid 源码）——证明 postprocess 真的替换了
    assert not any(
        sh.has_text_frame and "graph TD" in sh.text_frame.text for sh in prs.slides[0].shapes
    )
