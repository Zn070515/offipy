"""deck 图标迁移 legacy/new 行为等价套件（A3 Task 11，计划 §12）。

同一份注入 HTML 同时携带 legacy `data-icon`（postprocess_icons 经 outerHTML
`_match_svg` 读取）与确定性 `data-offipy-asset-id`（postprocess_assets 读取）。
本文件对同源 HTML 双路渲染：

- legacy 路：pptx 在图标 rect 放 PICTURE 占位（_remove_placeholder 按中心点
  删除），跑 `postprocess_icons`；
- new 路：pptx 在图标 rect 放 `OFFIPY_ASSET::{id}` 矩形占位（按精确名绑定），
  跑 `postprocess_assets`。

两路共用同一份 measurements.json（同时含 kind='svg' 与 kind='asset' 记录），
并且共享同一个 `_build_icon_shapes` → freeform 几何/样式数学上完全一致。断言
（§12.1 容忍度：EMU 精确、颜色精确、形状数精确）：

- 每页 freeform 数量、逐条几何（path w/h、压平后所有点、close）、线宽/cap/
  round、fill/noFill、RGB 完全一致；
- 图标在末尾/独占的 fixture：整页形状序列（含环绕内容形状顺序）完全一致；
- 图标夹在两个内容形状之间的 fixture：new 路把 freeform 块插回占位符槽位
  （内容形状之间）——这是相对 legacy 末尾追加的刻意修复，单独断言。

monkeypatch 浏览器/转换子进程——纯 python-pptx + lxml，任何环境可跑。
"""

from __future__ import annotations

import json
from dataclasses import dataclass
from io import BytesIO
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.util import Emu

from offipy.assets.declarations import preprocess_asset_declarations
from offipy.assets.render import PLACEHOLDER_PREFIX, postprocess_assets
from offipy.icons import postprocess_icons

_PX_TO_EMU = 6350

# 1×1 透明 PNG：add_picture 占位用，_remove_placeholder 按中心点删除。
_TINY_PNG = bytes.fromhex(
    "89504e470d0a1a0a0000000d49484452000000010000000108060000001f15c489"
    "0000000b4944415478da63f80f00050001010001020000000049454e44ae426082"
)

_MK_THEME_VARS = {"accent": "#2251FF", "bg": "#ffffff", "surface": "#f3f4f6"}
_DARK_THEME_VARS = {"accent": "#38BDF8", "bg": "#0B1B33", "surface": "#0B1B33"}

_TITLE = {"x": 200, "y": 100, "w": 600, "h": 60}
_ICON_A = {"x": 800, "y": 400, "w": 200, "h": 200}
_ICON_B = {"x": 500, "y": 700, "w": 120, "h": 120}
_BOX_A = {"x": 200, "y": 300, "w": 400, "h": 200}
_BOX_B = {"x": 200, "y": 700, "w": 400, "h": 150}


@dataclass
class _SlideSpec:
    """一页的装配顺序：内容形状与图标占位按 list 顺序落 spTree。"""

    shapes: list[dict]


def _content(kind: str, name: str, rect: dict) -> dict:
    return {"kind": kind, "name": name, "rect": rect}


def _icon(data_icon: str, rect: dict, html_tag: str = "svg") -> dict:
    return {"kind": "icon", "data_icon": data_icon, "html_tag": html_tag, "rect": rect}


def _decl_id(slide: int, ordinal: int) -> str:
    return f"asset-s{slide:02d}-{ordinal:03d}"


def _icon_decl_tag(data_icon: str) -> str:
    set_ = data_icon.split(":", 1)[0]
    vb = "0 0 256 256" if set_ == "ph" else "0 0 24 24"
    return f'<svg data-icon="{data_icon}" viewBox="{vb}"></svg>'


def _source_html(specs: list[_SlideSpec], theme_style: str | None = None) -> str:
    style = f'<style data-theme="{theme_style}"></style>' if theme_style else ""
    sections = "\n".join(
        f"<section data-pptx-slide>"
        f"{''.join(_icon_decl_tag(s['data_icon']) for s in spec.shapes if s['kind'] == 'icon')}"
        f"</section>"
        for spec in specs
    )
    return f"<!doctype html>\n<html><head>{style}</head><body>{sections}</body></html>"


def _to_emu(v: float) -> int:
    return int(round(v * _PX_TO_EMU))


def _add_content_shape(slide, s: dict) -> None:
    x, y, w, h = (
        _to_emu(v) for v in (s["rect"]["x"], s["rect"]["y"], s["rect"]["w"], s["rect"]["h"])
    )
    if s["kind"] == "textbox":
        tb = slide.shapes.add_textbox(Emu(x), Emu(y), Emu(w), Emu(h))
        tb.text_frame.text = s["name"]
        tb.name = s["name"]
        return
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(x), Emu(y), Emu(w), Emu(h))
    sp.name = s["name"]


def _build_deck(path: Path, specs: list[_SlideSpec], *, legacy: bool) -> None:
    """legacy：图标 rect 放 PICTURE 占位；new：放 OFFIPY_ASSET::{id} 矩形。"""
    prs = Presentation()
    blank = prs.slide_layouts[6]
    for si, spec in enumerate(specs, start=1):
        slide = prs.slides.add_slide(blank)
        icon_ord = 0
        for s in spec.shapes:
            if s["kind"] != "icon":
                _add_content_shape(slide, s)
                continue
            icon_ord += 1
            x, y, w, h = (
                _to_emu(s["rect"]["x"]),
                _to_emu(s["rect"]["y"]),
                _to_emu(s["rect"]["w"]),
                _to_emu(s["rect"]["h"]),
            )
            if legacy:
                pic = slide.shapes.add_picture(BytesIO(_TINY_PNG), Emu(x), Emu(y), Emu(w), Emu(h))
                pic.name = "icon-placeholder"
            else:
                sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(x), Emu(y), Emu(w), Emu(h))
                sp.name = f"{PLACEHOLDER_PREFIX}{_decl_id(si, icon_ord)}"
    prs.save(path)


def _svg_outer_html(data_icon: str) -> str:
    set_ = data_icon.split(":", 1)[0]
    vb = "0 0 256 256" if set_ == "ph" else "0 0 24 24"
    return (
        f'<svg class="icon" data-icon="{data_icon}" viewBox="{vb}" width="120" height="120"></svg>'
    )


def _write_measurements(path: Path, specs: list[_SlideSpec], color, theme_vars) -> None:
    """一份 measurements.json 同时带 kind='svg' 与 kind='asset' 记录（两路共用）。"""
    slides = [{"records": []} for _ in specs]
    for si, spec in enumerate(specs, start=1):
        icon_ord = 0
        for s in spec.shapes:
            if s["kind"] != "icon":
                continue
            icon_ord += 1
            rect = s["rect"]
            slides[si - 1]["records"].extend(
                (
                    {
                        "kind": "svg",
                        "rect": rect,
                        "color": color,
                        "outerHTML": _svg_outer_html(s["data_icon"]),
                    },
                    {
                        "kind": "asset",
                        "assetId": _decl_id(si, icon_ord),
                        "rect": rect,
                        "color": color,
                        "themeVars": theme_vars,
                        "tag": s.get("html_tag", "svg"),
                    },
                )
            )
    meas = path.with_name(f"{path.stem}_audit") / "_cache" / "measurements.json"
    meas.parent.mkdir(parents=True, exist_ok=True)
    meas.write_text(json.dumps({"slides": slides}, ensure_ascii=False), encoding="utf-8")


def _run_dual(
    tmp_path: Path, specs: list[_SlideSpec], *, color, theme_vars, theme_style=None, label="case"
):
    """同源注入 HTML 双路渲染，返回 (legacy_prs, new_prs, injected, decls)。

    注入副本同时携带 legacy data-icon 与 data-offipy-asset-id，且 legacy 语法
    原样保留（迁移不破坏旧 source 语法，§14 hard merge checklist）。
    """
    work = tmp_path / label
    work.mkdir()
    source = _source_html(specs, theme_style)
    injected, decls = preprocess_asset_declarations(source)
    n_icons = sum(1 for spec in specs for s in spec.shapes if s["kind"] == "icon")
    assert len(decls) == n_icons
    assert all(
        f'data-icon="{s["data_icon"]}"' in injected
        for spec in specs
        for s in spec.shapes
        if s["kind"] == "icon"
    )

    html_file = work / "d.audited.html"
    html_file.write_text(injected, encoding="utf-8")
    legacy = work / "legacy.pptx"
    new = work / "new.pptx"
    _build_deck(legacy, specs, legacy=True)
    _build_deck(new, specs, legacy=False)
    _write_measurements(legacy, specs, color, theme_vars)
    _write_measurements(new, specs, color, theme_vars)

    postprocess_icons(str(html_file), str(legacy))
    report = postprocess_assets(str(html_file), str(new))
    assert len(report.records) == len(decls)
    return Presentation(str(legacy)), Presentation(str(new)), injected, decls


# ---------------------------------------------------------------------------
# freeform 描述提取（§12.1：EMU 精确、颜色精确、形状数精确）
# ---------------------------------------------------------------------------


def _fill_desc(spPr) -> str:
    solid = spPr.find(qn("a:solidFill"))
    if solid is not None:
        clr = solid.find(qn("a:srgbClr"))
        return f"solid:{clr.get('val')}" if clr is not None else "solid"
    if spPr.find(qn("a:noFill")) is not None:
        return "noFill"
    return "none"


def _ln_desc(spPr):
    ln = spPr.find(qn("a:ln"))
    if ln is None:
        return None
    solid = ln.find(qn("a:solidFill"))
    clr = solid.find(qn("a:srgbClr")) if solid is not None else None
    fill = (
        f"solid:{clr.get('val')}"
        if clr is not None
        else ("noFill" if ln.find(qn("a:noFill")) is not None else "none")
    )
    return {
        "w": ln.get("w"),
        "cap": ln.get("cap"),
        "round": ln.find(qn("a:round")) is not None,
        "fill": fill,
    }


def _freeform_descriptor(sp: object) -> dict:
    spPr = sp.find(qn("p:spPr"))
    path = spPr.find(qn("a:custGeom")).find(qn("a:pathLst")).find(qn("a:path"))
    points: list[tuple[str, str]] = []
    close = False
    for child in path:
        if child.tag == qn("a:close"):
            close = True
        elif child.tag in (qn("a:moveTo"), qn("a:lnTo")):
            pt = child.find(qn("a:pt"))
            points.append((pt.get("x"), pt.get("y")))
    return {
        "path_w": path.get("w"),
        "path_h": path.get("h"),
        "points": tuple(points),
        "close": close,
        "fill": _fill_desc(spPr),
        "ln": _ln_desc(spPr),
    }


def _freeform_descriptors(slide) -> list[dict]:
    descs = []
    for el in slide.shapes._spTree:
        if el.tag != qn("p:sp"):
            continue
        spPr = el.find(qn("p:spPr"))
        if spPr is None or spPr.find(qn("a:custGeom")) is None:
            continue
        descs.append(_freeform_descriptor(el))
    return descs


def _shape_sequence(slide) -> list[str]:
    """文档顺序的形状序列：pic / freeform / prst:<name>（内容形状显式命名）。"""
    seq: list[str] = []
    for el in slide.shapes._spTree:
        if el.tag == qn("p:pic"):
            seq.append("pic")
        elif el.tag == qn("p:sp"):
            spPr = el.find(qn("p:spPr"))
            if spPr is not None and spPr.find(qn("a:custGeom")) is not None:
                seq.append("freeform")
            else:
                nv = el.find(f"{qn('p:nvSpPr')}/{qn('p:cNvPr')}")
                name = nv.get("name") if nv is not None else ""
                seq.append(f"prst:{name}")
    return seq


def _freeform_colors(slide) -> set[str]:
    colors: set[str] = set()
    for d in _freeform_descriptors(slide):
        if d["fill"].startswith("solid:"):
            colors.add(d["fill"][len("solid:") :])
        if d["ln"] is not None and d["ln"]["fill"].startswith("solid:"):
            colors.add(d["ln"]["fill"][len("solid:") :])
    return colors


def _assert_freeform_equivalence(legacy_prs, new_prs, label: str) -> None:
    assert len(legacy_prs.slides) == len(new_prs.slides)
    for i, (ls, ns) in enumerate(zip(legacy_prs.slides, new_prs.slides, strict=True), start=1):
        assert _freeform_descriptors(ls) == _freeform_descriptors(ns), (
            f"{label} slide {i}: legacy/new freeform geometry or style differs"
        )


def _assert_same_shape_sequence(legacy_prs, new_prs, label: str) -> None:
    for i, (ls, ns) in enumerate(zip(legacy_prs.slides, new_prs.slides, strict=True), start=1):
        assert _shape_sequence(ls) == _shape_sequence(ns), (
            f"{label} slide {i}: surrounding shape order differs"
        )


# ---------------------------------------------------------------------------
# §12.2 必测 fixture
# ---------------------------------------------------------------------------


def test_phosphor_single_path_fill(tmp_path):
    """ph 单 path fill 图标：两路 freeform 几何/样式一致，环绕顺序一致。"""
    specs = [
        _SlideSpec(shapes=[_content("textbox", "Title", _TITLE), _icon("ph:airplane", _ICON_A)])
    ]
    leg, new, _, _ = _run_dual(
        tmp_path, specs, color=None, theme_vars=_MK_THEME_VARS, theme_style="mckinsey"
    )
    _assert_freeform_equivalence(leg, new, "ph single-path")
    _assert_same_shape_sequence(leg, new, "ph single-path")
    assert _freeform_colors(new.slides[0]) == {"2251FF"}  # theme accent fallback


def test_lucide_stroke(tmp_path):
    """lu stroke 图标：线宽/cap=round/fill noFill 两路一致。"""
    specs = [
        _SlideSpec(shapes=[_content("textbox", "Title", _TITLE), _icon("lu:settings", _ICON_A)])
    ]
    leg, new, _, _ = _run_dual(
        tmp_path, specs, color=None, theme_vars=_MK_THEME_VARS, theme_style="mckinsey"
    )
    _assert_freeform_equivalence(leg, new, "lu stroke")
    _assert_same_shape_sequence(leg, new, "lu stroke")
    for d in _freeform_descriptors(new.slides[0]):
        assert d["fill"] == "noFill"
        assert d["ln"] is not None and d["ln"]["cap"] == "rnd" and d["ln"]["round"] is True


def test_lucide_mixed_filled_subpath(tmp_path):
    """lu 混 fill 子路径（key-round 的实心圈）：同一图标内 fill+stroke 并存，两路一致。"""
    specs = [
        _SlideSpec(shapes=[_content("textbox", "Title", _TITLE), _icon("lu:key-round", _ICON_A)])
    ]
    leg, new, _, _ = _run_dual(
        tmp_path, specs, color=None, theme_vars=_MK_THEME_VARS, theme_style="mckinsey"
    )
    _assert_freeform_equivalence(leg, new, "lu mixed")
    _assert_same_shape_sequence(leg, new, "lu mixed")
    fills = [d["fill"] for d in _freeform_descriptors(new.slides[0])]
    assert "noFill" in fills and any(f.startswith("solid:") for f in fills)


def test_lucide_arc_path_complex(tmp_path):
    """arc/path 复杂图标（airplay：开放弧 + 三角闭合）：压平后点两路逐点一致。"""
    specs = [
        _SlideSpec(shapes=[_content("textbox", "Title", _TITLE), _icon("lu:airplay", _ICON_A)])
    ]
    leg, new, _, _ = _run_dual(
        tmp_path, specs, color=None, theme_vars=_MK_THEME_VARS, theme_style="mckinsey"
    )
    _assert_freeform_equivalence(leg, new, "lu arc")
    _assert_same_shape_sequence(leg, new, "lu arc")


def test_two_identical_icons_same_slide(tmp_path):
    """同页两个相同图标：声明序 asset-s01-001/002，两路 freeform 数/几何一致。"""
    specs = [
        _SlideSpec(
            shapes=[
                _content("textbox", "Title", _TITLE),
                _icon("ph:check", _ICON_A),
                _icon("ph:check", _ICON_B),
            ]
        )
    ]
    leg, new, _, decls = _run_dual(
        tmp_path, specs, color=None, theme_vars=_MK_THEME_VARS, theme_style="mckinsey"
    )
    assert [d.declaration_id for d in decls] == ["asset-s01-001", "asset-s01-002"]
    _assert_freeform_equivalence(leg, new, "two identical")
    _assert_same_shape_sequence(leg, new, "two identical")
    # ph:check 2 子路径 ×2 图标 = 4 freeforms
    assert len(_freeform_descriptors(new.slides[0])) == 4


def test_icons_on_two_slides(tmp_path):
    """两页各一图标：页码/声明序跨页正确，两路逐页一致。"""
    specs = [
        _SlideSpec(shapes=[_content("textbox", "Title", _TITLE), _icon("ph:airplane", _ICON_A)]),
        _SlideSpec(shapes=[_content("textbox", "Title", _TITLE), _icon("lu:settings", _ICON_B)]),
    ]
    leg, new, _, decls = _run_dual(
        tmp_path, specs, color=None, theme_vars=_MK_THEME_VARS, theme_style="mckinsey"
    )
    assert [d.declaration_id for d in decls] == ["asset-s01-001", "asset-s02-001"]
    assert [d.slide_index for d in decls] == [1, 2]
    _assert_freeform_equivalence(leg, new, "two slides")
    _assert_same_shape_sequence(leg, new, "two slides")


def test_custom_css_color(tmp_path):
    """computed CSS color（rgb(255,0,0)）到达两路，结果同为 #FF0000。"""
    specs = [
        _SlideSpec(shapes=[_content("textbox", "Title", _TITLE), _icon("ph:airplane", _ICON_A)])
    ]
    leg, new, _, _ = _run_dual(
        tmp_path, specs, color="rgb(255, 0, 0)", theme_vars=_MK_THEME_VARS, theme_style=None
    )
    _assert_freeform_equivalence(leg, new, "custom color")
    _assert_same_shape_sequence(leg, new, "custom color")
    assert _freeform_colors(new.slides[0]) == {"FF0000"}


def test_theme_accent_fallback(tmp_path):
    """color=None → 主题 accent 兜底：_accent_rgb(context) == _theme_accent_fallback(html)。"""
    specs = [
        _SlideSpec(shapes=[_content("textbox", "Title", _TITLE), _icon("lu:settings", _ICON_A)])
    ]
    leg, new, _, _ = _run_dual(
        tmp_path, specs, color=None, theme_vars=_MK_THEME_VARS, theme_style="mckinsey"
    )
    _assert_freeform_equivalence(leg, new, "accent fallback")
    # 两路都落在 THEMES['mckinsey'].base_vars['--accent'] = #2251FF
    assert _freeform_colors(leg.slides[0]) == {"2251FF"}
    assert _freeform_colors(new.slides[0]) == {"2251FF"}


def test_dark_light_variant(tmp_path):
    """明/暗主题不同 accent：两路在各自主题下一致，且主题间颜色不同（主题敏感）。"""
    light_specs = [
        _SlideSpec(shapes=[_content("textbox", "Title", _TITLE), _icon("ph:airplane", _ICON_A)])
    ]
    leg_l, new_l, _, _ = _run_dual(
        tmp_path,
        light_specs,
        color=None,
        theme_vars=_MK_THEME_VARS,
        theme_style="mckinsey",
        label="light",
    )
    dark_specs = [
        _SlideSpec(shapes=[_content("textbox", "Title", _TITLE), _icon("ph:airplane", _ICON_A)])
    ]
    leg_d, new_d, _, _ = _run_dual(
        tmp_path,
        dark_specs,
        color=None,
        theme_vars=_DARK_THEME_VARS,
        theme_style="dark-tech",
        label="dark",
    )
    _assert_freeform_equivalence(leg_l, new_l, "light")
    _assert_freeform_equivalence(leg_d, new_d, "dark")
    light = _freeform_colors(new_l.slides[0])
    dark = _freeform_colors(new_d.slides[0])
    assert light == {"2251FF"}  # mckinsey base --accent
    assert dark == {"38BDF8"}  # dark-tech base --accent
    assert light != dark


def test_icon_between_content_shapes_restores_slot(tmp_path):
    """图标夹在内容形状之间：new 路把 freeform 块插回占位符槽位（内容之间）。

    相对 legacy 末尾追加的刻意修复：几何两路一致，但顺序不同——legacy
    [A, B, freeforms]，new [A, freeforms, B]（ph:check 2 子路径 → 连续 2 个）。
    """
    specs = [
        _SlideSpec(
            shapes=[
                _content("rect", "ContentA", _BOX_A),
                _icon("ph:check", _ICON_A),
                _content("rect", "ContentB", _BOX_B),
            ]
        )
    ]
    leg, new, _, _ = _run_dual(
        tmp_path, specs, color=None, theme_vars=_MK_THEME_VARS, theme_style="mckinsey"
    )
    _assert_freeform_equivalence(leg, new, "between content")
    assert _shape_sequence(leg.slides[0]) == [
        "prst:ContentA",
        "prst:ContentB",
        "freeform",
        "freeform",
    ]
    assert _shape_sequence(new.slides[0]) == [
        "prst:ContentA",
        "freeform",
        "freeform",
        "prst:ContentB",
    ]
