"""A3 Task 7 — render_asset 分发 + place_rendered_elements 精确 z-order 落位。

render_asset 按 payload 类型分发：freeform_svg → 既有图标 freeform 渲染器；
RasterPayload → add_picture；SvgPayload(svg)/SvgTemplatePayload → OOXML SVG picture
（A1 选定的 P2 svgBlip 路线）；NativeShapePayload → A5 前未注册，明确报错。

place_rendered_elements 把渲染产物精确放入占位符 XML 槽位：replace/decorative
按渲染顺序插回占位符所在槽；background 移到 grpSpPr 之后、所有内容形状之前。
多路径图标必须是连续块、不反转顺序、不散落到别处。
"""

import struct
import zlib

import pytest

from offipy.assets.model import (
    AssetMeta,
    AssetProviderMeta,
    AssetRect,
    AssetRef,
    AssetRenderContext,
    AssetRequest,
    NativeShapePayload,
    RasterPayload,
    ResolvedAsset,
    SvgPayload,
    SvgTemplatePayload,
)
from offipy.assets.providers.icons import IconProvider
from offipy.assets.render import (
    _as_element,
    _svg_page_screenshot,
    place_rendered_elements,
    render_asset,
)
from offipy.exceptions import InvalidArgumentError

_ASVG_NS = "http://schemas.microsoft.com/office/drawing/2016/SVG/main"
_R = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed"
_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
_P = "http://schemas.openxmlformats.org/presentationml/2006/main"


def _tiny_png() -> bytes:
    """1×1 RGBA 透明 PNG（python-pptx add_picture 可解析）。"""
    sig = b"\x89PNG\r\n\x1a\n"

    def chunk(tag, data):
        body = tag + data
        return struct.pack(">I", len(data)) + body + struct.pack(">I", zlib.crc32(body))

    ihdr = struct.pack(">IIBBBBB", 1, 1, 8, 6, 0, 0, 0)
    idat = zlib.compress(b"\x00\xff\x00\x00\xff")
    return sig + chunk(b"IHDR", ihdr) + chunk(b"IDAT", idat) + chunk(b"IEND", b"")


def _blank_slide():
    from pptx import Presentation

    prs = Presentation()
    return prs.slides.add_slide(prs.slide_layouts[6])


def _context(rect=None, placement="replace", theme_vars=None):
    return AssetRenderContext(
        slide_index=1,
        rect=rect or AssetRect(20, 30, 120, 60),
        theme_name="test",
        theme_vars=theme_vars or {"accent": "#2251ff"},
        placement=placement,
    )


def _resolved(payload, provider="ph", name="check"):
    ref = AssetRef(provider, "icon", name)
    meta = AssetMeta(ref=ref, title=name, tags=(name,))
    pmeta = AssetProviderMeta(
        provider_id=provider,
        license="test",
        source_url=None,
        source_commit=None,
        attribution=None,
        redistributable=True,
    )
    return ResolvedAsset(request=AssetRequest(ref), meta=meta, provider_meta=pmeta, payload=payload)


def _rect_svg():
    return (
        '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 24 24">'
        '<path d="M0 0h24v24H0z"/></svg>'
    )


def _render_single_p2(slide):
    """确定性单元素渲染产物：SvgPayload(svg) → 一个 P2 svgBlip picture。"""
    resolved = _resolved(SvgPayload(svg=_rect_svg(), render_mode="svg", view_box=(0, 0, 24, 24)))
    return render_asset(slide, resolved, _context())


def _slide_with_three():
    """spTree 子顺序: nvGrpSpPr, grpSpPr, A, placeholder, B。"""
    from pptx.enum.shapes import MSO_SHAPE

    slide = _blank_slide()
    a = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, 100, 100)
    a.name = "A"
    ph = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 100, 0, 100, 100)
    ph.name = "OFFIPY_ASSET::asset-s01-001"
    b = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 200, 0, 100, 100)
    b.name = "B"
    return slide, ph, a, b


def _children(slide):
    sp_tree = slide.shapes._spTree
    return [sp_tree.find(f"{{{_P}}}nvGrpSpPr"), sp_tree.find(f"{{{_P}}}grpSpPr")] + list(sp_tree)[
        2:
    ]


# ---------------------------------------------------------------------------
# render_asset 分发
# ---------------------------------------------------------------------------


class TestRenderAssetDispatch:
    def test_freeform_svg_renders_icon_shapes(self):
        resolved = IconProvider("ph").resolve(AssetRequest(AssetRef("ph", "icon", "check")))
        slide = _blank_slide()
        rendered = render_asset(slide, resolved, _context())
        assert len(rendered) >= 1
        for shape in rendered:
            assert _as_element(shape).getparent() is not None  # 已在 spTree

    def test_raster_payload_adds_picture(self):
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        payload = RasterPayload(
            data=_tiny_png(), media_type="image/png", pixel_width=1, pixel_height=1
        )
        rendered = render_asset(_blank_slide(), _resolved(payload), _context())
        assert len(rendered) == 1
        pic = rendered[0]
        assert pic.shape_type == MSO_SHAPE_TYPE.PICTURE
        assert (pic.left, pic.top, pic.width, pic.height) == (127000, 190500, 762000, 381000)

    def test_svg_payload_renders_p2_svg_picture(self):
        svg = _rect_svg()
        resolved = _resolved(SvgPayload(svg=svg, render_mode="svg", view_box=(0, 0, 24, 24)))
        slide = _blank_slide()
        rendered = render_asset(slide, resolved, _context())
        assert len(rendered) == 1
        pic = rendered[0]
        assert pic.tag == f"{{{_P}}}pic"
        blip = pic.find(f".//{{{_ASVG_NS}}}svgBlip")
        assert blip is not None
        r_id = blip.get(_R)
        assert r_id
        part = slide.part.related_part(r_id)
        assert part.partname.endswith(".svg")
        assert part.blob.decode("utf-8") == svg
        off = pic.find(f".//{{{_A}}}off")
        ext = pic.find(f".//{{{_A}}}xfrm/{{{_A}}}ext")
        assert (off.get("x"), off.get("y")) == ("127000", "190500")
        assert (ext.get("cx"), ext.get("cy")) == ("762000", "381000")

    def test_svg_template_materializes_then_renders_p2(self):
        template = (
            '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 24 24">'
            '<path d="M0 0h24v24H0z" fill="__ACCENT__"/></svg>'
        )
        payload = SvgTemplatePayload(
            template=template,
            render_mode="svg",
            view_box=(0, 0, 24, 24),
            color_slots=(("__ACCENT__", "accent"),),
        )
        slide = _blank_slide()
        rendered = render_asset(slide, _resolved(payload), _context())
        pic = rendered[0]
        r_id = pic.find(f".//{{{_ASVG_NS}}}svgBlip").get(_R)
        assert "#2251FF" in slide.part.related_part(r_id).blob.decode("utf-8")

    def test_native_shape_payload_not_registered(self):
        payload = NativeShapePayload(primitive="rect", params=())
        with pytest.raises(InvalidArgumentError, match="unknown native primitive"):
            render_asset(_blank_slide(), _resolved(payload), _context())

    def test_svg_payload_png_fallback_embeds_png_and_svg(self):
        # #58：svg_to_png 提供时主 <a:blip> 挂 PNG embed，svgBlip 仍指 SVG——
        # 支持 SVG 的宿主用矢量，不支持的查看器回退 PNG。
        svg = _rect_svg()
        resolved = _resolved(SvgPayload(svg=svg, render_mode="svg", view_box=(0, 0, 24, 24)))
        slide = _blank_slide()
        rendered = render_asset(slide, resolved, _context(), svg_to_png=lambda s: _tiny_png())
        pic = rendered[0]
        blip = pic.find(f".//{{{_A}}}blip")
        png_rid = blip.get(_R)
        assert png_rid
        png_part = slide.part.related_part(png_rid)
        assert png_part.partname.endswith(".png")
        assert png_part.blob == _tiny_png()
        svg_rid = pic.find(f".//{{{_ASVG_NS}}}svgBlip").get(_R)
        assert slide.part.related_part(svg_rid).blob.decode("utf-8") == svg

    def test_svg_payload_without_svg_to_png_keeps_pure_svg(self):
        svg = _rect_svg()
        resolved = _resolved(SvgPayload(svg=svg, render_mode="svg", view_box=(0, 0, 24, 24)))
        slide = _blank_slide()
        rendered = render_asset(slide, resolved, _context())
        pic = rendered[0]
        assert pic.find(f".//{{{_A}}}blip").get(_R) is None  # 纯 SVG：主 blip 无 embed
        svg_rid = pic.find(f".//{{{_ASVG_NS}}}svgBlip").get(_R)
        assert slide.part.related_part(svg_rid).blob.decode("utf-8") == svg


# ---------------------------------------------------------------------------
# place_rendered_elements z-order
# ---------------------------------------------------------------------------


class TestPlaceRenderedElements:
    def test_replace_keeps_placeholder_slot(self):
        slide, ph, a, b = _slide_with_three()
        rendered = _render_single_p2(slide)
        place_rendered_elements(slide, ph, rendered, "replace")
        assert ph._element.getparent() is None
        assert _children(slide) == [
            slide.shapes._spTree.find(f"{{{_P}}}nvGrpSpPr"),
            slide.shapes._spTree.find(f"{{{_P}}}grpSpPr"),
            a._element,
            _as_element(rendered[0]),
            b._element,
        ]

    def test_decorative_same_slot_as_replace(self):
        slide, ph, a, b = _slide_with_three()
        rendered = _render_single_p2(slide)
        place_rendered_elements(slide, ph, rendered, "decorative")
        assert ph._element.getparent() is None
        assert _children(slide) == [
            slide.shapes._spTree.find(f"{{{_P}}}nvGrpSpPr"),
            slide.shapes._spTree.find(f"{{{_P}}}grpSpPr"),
            a._element,
            _as_element(rendered[0]),
            b._element,
        ]

    def test_background_below_all_content(self):
        slide, ph, a, b = _slide_with_three()
        rendered = _render_single_p2(slide)
        place_rendered_elements(slide, ph, rendered, "background")
        assert ph._element.getparent() is None
        assert _children(slide) == [
            slide.shapes._spTree.find(f"{{{_P}}}nvGrpSpPr"),
            slide.shapes._spTree.find(f"{{{_P}}}grpSpPr"),
            _as_element(rendered[0]),
            a._element,
            b._element,
        ]

    def test_multi_path_icon_contiguous_no_reversal(self):
        # 回归：多子路径图标全部占住旧占位符槽，作为连续块、不反转顺序
        slide, ph, a, b = _slide_with_three()
        resolved = IconProvider("lu").resolve(AssetRequest(AssetRef("lu", "icon", "settings")))
        rendered = render_asset(slide, resolved, _context())
        assert len(rendered) == 2  # settings.svg 多子路径
        place_rendered_elements(slide, ph, rendered, "replace")
        assert ph._element.getparent() is None
        assert _children(slide) == [
            slide.shapes._spTree.find(f"{{{_P}}}nvGrpSpPr"),
            slide.shapes._spTree.find(f"{{{_P}}}grpSpPr"),
            a._element,
            _as_element(rendered[0]),
            _as_element(rendered[1]),
            b._element,
        ]


class _FakeSvgPage:
    """_svg_page_screenshot 的最小 page 替身：记录 viewport / clip。"""

    def __init__(self):
        self.viewport = None
        self.clip = None

    def set_viewport_size(self, viewport):
        self.viewport = viewport

    def set_content(self, _content):
        pass

    def screenshot(self, **kwargs):
        self.clip = kwargs.get("clip")
        return b"PNG-bytes"


def test_svg_page_screenshot_handles_comma_viewbox():
    # #64：逗号分隔 viewBox（合法 SVG）此前 float("0,0") 抛 ValueError → fallback 静默丢失
    svg = '<svg viewBox="0,0 100,100"><rect width="100" height="100"/></svg>'
    page = _FakeSvgPage()
    assert _svg_page_screenshot(page, svg) == b"PNG-bytes"
    assert page.viewport == {"width": 100, "height": 100}
    assert page.clip == {"x": 0, "y": 0, "width": 100, "height": 100}


def test_svg_page_screenshot_handles_mixed_viewbox_separators():
    svg = '<svg viewBox="10,20 30 40"><rect width="100" height="100"/></svg>'
    page = _FakeSvgPage()
    assert _svg_page_screenshot(page, svg) == b"PNG-bytes"
    assert page.clip == {"x": 10, "y": 20, "width": 30, "height": 40}
